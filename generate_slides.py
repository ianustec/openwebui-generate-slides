"""
title: Generate Slides
author: IANUSTEC
author_url: https://ianustec.com
funding_url: https://github.com/ianustec
description: Generate high-quality native PowerPoint (.pptx) presentations from a JSON spec - layered graphics, native charts, icons, rich layouts
requirements: python-pptx, pillow
required_open_webui_version: 0.4.0
version: 1.0.1
license: MIT
"""

# ============================================================================
# Native PPTX Slides Engine — open-source OpenWebUI tool
# ----------------------------------------------------------------------------
# Server-side OpenWebUI tool `generate_slides`. The model emits a structured
# JSON spec ({title, subtitle, author, theme, slides:[{layout, ...}]}) and this
# engine renders a NATIVE .pptx with python-pptx: coherent visual system,
# layered decorative shapes, native charts, icons-in-circles, rich layouts.
#
# The deck is saved via the OpenWebUI Files API (with a /cache/files fallback)
# and a clickable download link is emitted in chat.
#
# License: MIT — Copyright (c) IANUSTEC.
# ============================================================================

import os
import re
import json
import uuid
import base64
import unicodedata
from datetime import datetime, timezone
from io import BytesIO
from typing import Any, Optional

from pydantic import BaseModel, Field

# --- python-pptx (rendering engine) -----------------------------------------
try:
    from pptx import Presentation  # type: ignore
    from pptx.util import Inches, Pt  # type: ignore
    from pptx.dml.color import RGBColor  # type: ignore
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR  # type: ignore
    from pptx.enum.shapes import MSO_SHAPE  # type: ignore
    from pptx.oxml.ns import qn  # type: ignore
    from pptx.chart.data import CategoryChartData  # type: ignore
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION  # type: ignore

    _HAS_PPTX = True
except Exception:  # pragma: no cover
    _HAS_PPTX = False

# --- Pillow (icons + image resize) ------------------------------------------
try:
    from PIL import Image  # type: ignore

    _HAS_PIL = True
except Exception:
    _HAS_PIL = False

# --- httpx (image fetch) ----------------------------------------------------
try:
    import httpx  # type: ignore

    _HAS_HTTPX = True
except Exception:
    _HAS_HTTPX = False

# --- OpenWebUI Files API (native download chip) -----------------------------
try:
    from fastapi import UploadFile  # type: ignore
    from starlette.datastructures import Headers  # type: ignore
    from open_webui.routers.files import upload_file_handler  # type: ignore
    from open_webui.models.users import Users  # type: ignore

    _HAS_OWUI_FILES = True
except Exception:
    _HAS_OWUI_FILES = False

# --- OpenWebUI AI image generation (optional) -------------------------------
try:
    from open_webui.routers.images import image_generations as _owui_image_generations  # type: ignore
    from open_webui.routers.images import GenerateImageForm as _OwuiImageForm  # type: ignore

    _HAS_OWUI_IMAGES = True
except Exception:
    _HAS_OWUI_IMAGES = False


# ============================================================================
# Geometry (LAYOUT_WIDE 16:9)
# ============================================================================
EMU_IN = 914400
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
MARGIN = 0.62          # outer margin (in)
CONTENT_TOP = 1.9      # y where content starts on standard slides (in)
FOOTER_Y = 7.02        # footer band y (in)
BODY_FONT = "Calibri"
HEAD_FONT = "Georgia"  # serif display for a premium, editorial feel


# ============================================================================
# Colour utilities
# ============================================================================

def _hex(value: str) -> str:
    """Return a clean 6-char uppercase hex (no #)."""
    if not isinstance(value, str):
        return "1E2761"
    h = value.strip().lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    if len(h) != 6:
        return "1E2761"
    try:
        int(h, 16)
        return h.upper()
    except ValueError:
        return "1E2761"


def _to_rgb(value: str) -> tuple[int, int, int]:
    h = _hex(value)
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)


def _from_rgb(r: int, g: int, b: int) -> str:
    def clamp(x):
        return max(0, min(255, int(round(x))))
    return f"{clamp(r):02X}{clamp(g):02X}{clamp(b):02X}"


def _mix(a: str, b: str, t: float) -> str:
    ra, ga, ba = _to_rgb(a)
    rb, gb, bb = _to_rgb(b)
    t = max(0.0, min(1.0, t))
    return _from_rgb(ra + (rb - ra) * t, ga + (gb - ga) * t, ba + (bb - ba) * t)


def _darken(c: str, t: float) -> str:
    return _mix(c, "000000", t)


def _lighten(c: str, t: float) -> str:
    return _mix(c, "FFFFFF", t)


def _luminance(c: str) -> float:
    def ch(x: float) -> float:
        x /= 255.0
        return x / 12.92 if x <= 0.03928 else ((x + 0.055) / 1.055) ** 2.4

    r, g, b = _to_rgb(c)
    return 0.2126 * ch(r) + 0.7152 * ch(g) + 0.0722 * ch(b)


def _on(bg: str) -> str:
    """Readable ink (near-black or white) for text on ``bg``."""
    return "1A1A1A" if _luminance(bg) > 0.42 else "FFFFFF"


def _rgb(value: str) -> "RGBColor":
    return RGBColor(*_to_rgb(value))


# ============================================================================
# Named palettes (curated theme presets)
# ============================================================================
_PALETTES: dict[str, dict[str, str]] = {
    "midnight":   {"primary": "1E2761", "dark": "141A45", "accent": "C99A3B", "light_accent": "CADCFC"},
    "forest":     {"primary": "2C5F2D", "dark": "1A3D1B", "accent": "97BC62", "light_accent": "B7D48A"},
    "coral":      {"primary": "2F3C7E", "dark": "1E2A5C", "accent": "F96167", "light_accent": "F9E795"},
    "terracotta": {"primary": "B85042", "dark": "8A3328", "accent": "C9A227", "light_accent": "A7BEAE"},
    "ocean":      {"primary": "21295C", "dark": "0A0F33", "accent": "1C7293", "light_accent": "7FBADC"},
    "charcoal":   {"primary": "2B3640", "dark": "1A2128", "accent": "C08A2D", "light_accent": "9FB0BB"},
    "teal":       {"primary": "013F46", "dark": "012A2F", "accent": "02C39A", "light_accent": "7FE3CD"},
    "berry":      {"primary": "6D2E46", "dark": "3F1A28", "accent": "C99A3B", "light_accent": "D6A5A6"},
    "sage":       {"primary": "2E4D58", "dark": "1F343C", "accent": "84B59F", "light_accent": "A4C7BD"},
    "cherry":     {"primary": "2F3C7E", "dark": "1A2356", "accent": "C0392B", "light_accent": "E8636E"},
    "slate":      {"primary": "27303F", "dark": "171C26", "accent": "3B82F6", "light_accent": "93C5FD"},
    # legacy aliases
    "blue":       {"primary": "1E2761", "dark": "141A45", "accent": "C99A3B", "light_accent": "CADCFC"},
    "green":      {"primary": "2C5F2D", "dark": "1A3D1B", "accent": "97BC62", "light_accent": "B7D48A"},
    "dark":       {"primary": "1F2937", "dark": "111827", "accent": "60A5FA", "light_accent": "93C5FD"},
}
_DEFAULT_PALETTE = "midnight"

_THEME_KEYWORDS: dict[str, tuple[str, ...]] = {
    "forest": ("sostenib", "ambient", "green", "eco", "natur", "agricol", "environment"),
    "ocean": ("data", "analy", "research", "scienz", "ricerca", "dati", "tech"),
    "coral": ("creativ", "design", "marketing", "brand"),
    "terracotta": ("food", "hospitality", "travel", "ristorante", "cucina", "viaggio"),
    "teal": ("health", "medic", "wellness", "fintech", "salute", "benessere"),
    "berry": ("luxury", "fashion", "retail", "lusso", "moda"),
    "sage": ("yoga", "mindfulness", "spa"),
    "cherry": ("urgent", "alert", "energy", "urgente", "energia"),
    "charcoal": ("legal", "consult", "professional", "corporate", "legale", "avvocat"),
    "midnight": ("finance", "banking", "executive", "enterprise", "business", "strategy",
                 "finanz", "strategia", "startup", "invest", "plan"),
}


def _pick_palette_name(requested: Optional[str], title: str, slides: list) -> str:
    if requested:
        r = str(requested).strip().lower()
        if r in _PALETTES:
            return r
        if r not in ("auto", ""):
            return _DEFAULT_PALETTE
    hay = " ".join([title or ""] + [str(s.get("title", "")) for s in slides]).lower()
    for name, kws in _THEME_KEYWORDS.items():
        if any(k in hay for k in kws):
            return name
    return _DEFAULT_PALETTE


def _resolve_theme(spec: dict, slides: list) -> dict:
    """Build the full design-token palette for the deck."""
    name = _pick_palette_name(spec.get("theme"), spec.get("title", ""), slides)
    base = dict(_PALETTES.get(name, _PALETTES[_DEFAULT_PALETTE]))

    # Explicit overrides from the spec (palette{} or accent).
    pal = spec.get("palette") if isinstance(spec.get("palette"), dict) else {}
    for k in ("primary", "dark", "accent", "light_accent"):
        if pal.get(k):
            base[k] = _hex(pal[k])
    if spec.get("accent"):
        base["accent"] = _hex(spec["accent"])
    if spec.get("primary"):
        base["primary"] = _hex(spec["primary"])

    primary = _hex(base["primary"])
    dark = _hex(base.get("dark") or _darken(primary, 0.35))
    accent = _hex(base["accent"])

    theme = {
        "name": name,
        # dark surfaces (cover / section / closing)
        "bg_dark": dark,
        "bg_dark2": _lighten(dark, 0.06),
        "panel_dark": _lighten(dark, 0.10),
        "on_dark": "FFFFFF",
        "on_dark_soft": _mix(dark, "FFFFFF", 0.72),
        "on_dark_faint": _mix(dark, "FFFFFF", 0.50),
        # light surfaces (content)
        "bg_light": "FFFFFF",
        "bg_soft": _lighten(primary, 0.955),
        "ink": _darken(primary, 0.15) if _luminance(primary) > 0.12 else "1A1A2E",
        "muted": _mix(primary, "6B7280", 0.5),
        "faint": "9AA0AB",
        # brand
        "primary": primary,
        "accent": accent,
        "accent_soft": _lighten(accent, 0.82),
        "light_accent": _hex(base.get("light_accent") or _lighten(accent, 0.5)),
        # cards
        "card_bg": "FFFFFF",
        "card_border": _lighten(primary, 0.86),
        "card_soft": _lighten(primary, 0.955),
        "hairline": "E4E7EC",
        # fonts
        "head_font": _hex_font(spec.get("heading_font")) or HEAD_FONT,
        "body_font": _hex_font(spec.get("body_font")) or BODY_FONT,
    }
    return theme


def _hex_font(v):
    if isinstance(v, str) and v.strip():
        return v.strip()
    return None


# ============================================================================
# Layout vocabulary (canonical layout names + aliases)
# ============================================================================
_LAYOUT_ALIASES: dict[str, str] = {
    "title": "cover", "copertina": "cover",
    "content": "content",
    "image": "image_full_caption", "full_image": "image_full_caption",
    "hero_image": "image_full_caption",
    "stat": "kpi_row", "kpi": "kpi_row", "kpis": "kpi_row", "stats": "kpi_row",
    "metrics": "kpi_row", "big_number": "kpi_row", "bignumber": "kpi_row",
    "timeline": "timeline_horizontal", "tempistica": "timeline_horizontal",
    "process": "process_flow", "flow": "process_flow",
    "comparison": "comparison_two", "compare": "comparison_two", "vs": "comparison_two",
    "confronto": "comparison_two", "tabella": "table",
    "conclusion": "closing", "chiusura": "closing", "fine": "closing",
    "icon-list": "icon_list_vertical", "icon_list": "icon_list_vertical",
    "iconlist": "icon_list_vertical", "vertical_list": "icon_list_vertical",
    "icon-grid": "icon_grid_3", "icon_grid": "icon_grid_3",
    "image-grid": "image_grid", "gallery": "image_grid",
    "two_column": "two_column_text", "two_col": "two_column_text",
    "due_colonne": "two_column_text",
    "image_right": "text_image_right", "image_left": "image_left_text_right",
    "bullet_list": "title_bullets", "bulletlist": "title_bullets",
    "bullets": "title_bullets", "bullet-list": "title_bullets",
    "list": "title_bullets", "elenco": "title_bullets", "punti": "title_bullets",
    "highlights": "title_bullets", "key_points": "title_bullets",
    "keypoints": "title_bullets", "ranking": "title_bullets",
    "sezione": "section", "titolo": "title_only",
    "citazione": "quote", "grafico": "chart", "vuoto": "blank",
    "alert": "alert", "callout": "alert", "note": "alert", "warning": "alert",
    "info": "alert", "tip": "alert", "danger": "alert",
    "scheda": "title_body", "paragrafo": "title_body",
    # diagrams
    "diagram": "diagram", "diagramma": "diagram", "schema": "diagram",
    "imbuto": "funnel", "piramide": "pyramid", "ciclo": "cycle",
    "occhio_di_bue": "bullseye", "pilastri": "pillars", "pillar": "pillars",
    "strada": "roadmap", "road": "roadmap", "matrix": "quadrant",
    "matrice": "quadrant", "quadrante": "quadrant", "venn_diagram": "venn",
}

_DIAGRAM_LAYOUTS = ("funnel", "pyramid", "cycle", "venn", "quadrant",
                    "bullseye", "pillars", "iceberg", "roadmap")

_LIST_FIELD_ALIASES = (
    "bullets", "bullet_points", "bullet_list", "points", "items", "key_points",
    "highlights", "features", "benefits", "challenges", "problems", "risks",
    "advantages", "pros", "cons", "steps", "actions", "tasks", "takeaways",
    "conclusions", "recommendations", "details", "notes", "facts", "objectives",
    "goals", "results", "reasons", "examples", "entries", "lines",
)
_TEXT_FIELD_ALIASES = (
    "body", "content", "text", "description", "summary", "paragraph",
    "detail", "note", "explanation", "overview", "lead", "intro",
)


# ============================================================================
# Text / helpers
# ============================================================================
_MD_STRONG = re.compile(r"\*\*(.+?)\*\*|__(.+?)__")
_MD_EM = re.compile(r"(?<!\*)\*(?!\*)(.+?)(?<!\*)\*(?!\*)|(?<!_)_(?!_)(.+?)(?<!_)_(?!_)")
_MARKER_RE = re.compile(r"^\s*(?:[\u2022\u25E6\u2013\u2014\-\*\u2611\u2610\u2713\u2714\u2192>]\s+)")


def _smart(text: str) -> str:
    if not text:
        return ""
    return str(text).replace("--", "\u2013")


def _strip_md(text: str) -> str:
    """Remove **bold**/*em*/`code` markers, returning plain text."""
    if not text:
        return ""
    t = str(text)
    t = re.sub(r"\*\*(.+?)\*\*", r"\1", t)
    t = re.sub(r"__(.+?)__", r"\1", t)
    t = re.sub(r"`(.+?)`", r"\1", t)
    t = re.sub(r"(?<!\*)\*(?!\*)(.+?)\*(?!\*)", r"\1", t)
    return t.strip()


def _clean_bullet(text: str) -> str:
    t = _strip_md(str(text or ""))
    t = _MARKER_RE.sub("", t)
    return t.strip()


def _md_runs(text: str) -> list[tuple[str, bool, bool]]:
    """Split text into (chunk, bold, italic) runs from lightweight markdown."""
    if text is None:
        return [("", False, False)]
    s = str(text)
    runs: list[tuple[str, bool, bool]] = []
    # handle **bold** first via tokenization
    pattern = re.compile(r"\*\*(.+?)\*\*|`(.+?)`")
    pos = 0
    for m in pattern.finditer(s):
        if m.start() > pos:
            runs.append((s[pos:m.start()], False, False))
        if m.group(1) is not None:
            runs.append((m.group(1), True, False))
        else:
            runs.append((m.group(2), False, True))
        pos = m.end()
    if pos < len(s):
        runs.append((s[pos:], False, False))
    return runs or [("", False, False)]


def _slugify(text: str, *, max_len: int = 60) -> str:
    t = unicodedata.normalize("NFKD", text or "")
    t = "".join(c for c in t if not unicodedata.combining(c))
    t = t.lower()
    t = re.sub(r"[^a-z0-9]+", "-", t)
    t = re.sub(r"-{2,}", "-", t).strip("-")
    return (t or "presentation")[:max_len]


def _as_list(val) -> list:
    if val is None:
        return []
    if isinstance(val, list):
        return val
    if isinstance(val, str):
        parts = [p.strip() for p in val.split("\n") if p.strip()]
        return parts or [val]
    return [val]


def _first(d: dict, *keys, default=None):
    for k in keys:
        v = d.get(k)
        if v not in (None, "", [], {}):
            return v
    return default


def _harvest_bullets(slide: dict) -> list[str]:
    """Collect bullet strings from the many alias fields the LLM may use."""
    for key in _LIST_FIELD_ALIASES:
        val = slide.get(key)
        if isinstance(val, list) and val:
            out = []
            for item in val:
                if isinstance(item, dict):
                    t = _first(item, "title", "label", "heading", "name", "text", default="")
                    d = _first(item, "description", "desc", "detail", "body", "text", default="")
                    line = t
                    if d and d != t:
                        line = f"**{t}** — {d}" if t else d
                    if line:
                        out.append(str(line))
                elif item not in (None, ""):
                    out.append(str(item))
            if out:
                return out
        if isinstance(val, str) and val.strip():
            lines = [ln.strip() for ln in val.split("\n") if ln.strip()]
            if lines:
                return lines
    # text fields → split into paragraphs/bullets
    for key in _TEXT_FIELD_ALIASES:
        val = slide.get(key)
        if isinstance(val, str) and val.strip():
            lines = [ln.strip() for ln in val.split("\n") if ln.strip()]
            if len(lines) > 1:
                return lines
    return []


def _resolve_layout(slide: dict) -> str:
    raw = str(_first(slide, "layout", "type", default="content")).strip().lower()
    return _LAYOUT_ALIASES.get(raw, raw)


# ============================================================================
# Low-level PPTX drawing helpers
# ============================================================================

def _set_bg(slide, color: str) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(color)


def _no_line(shape) -> None:
    shape.line.fill.background()


def _set_shadow(shape, *, blur=8, dist=3, alpha=78000, direction=5400000) -> None:
    """Apply a soft outer shadow to a shape via OOXML."""
    try:
        spPr = shape._element.spPr
        # remove existing effectLst
        for el in spPr.findall(qn("a:effectLst")):
            spPr.remove(el)
        effect = spPr.makeelement(qn("a:effectLst"), {})
        shadow = effect.makeelement(qn("a:outerShdw"), {
            "blurRad": str(int(blur * 12700)),
            "dist": str(int(dist * 12700)),
            "dir": str(int(direction)),
            "rotWithShape": "0",
        })
        clr = shadow.makeelement(qn("a:srgbClr"), {"val": "1A1A2E"})
        a = clr.makeelement(qn("a:alpha"), {"val": str(int(alpha))})
        clr.append(a)
        shadow.append(clr)
        effect.append(shadow)
        spPr.append(effect)
    except Exception:
        pass


def _fill_alpha(shape, color: str, alpha_pct: float) -> None:
    """Solid fill with transparency (alpha_pct 0..100 = opacity)."""
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)
    try:
        sf = shape.fill.fore_color._xFill.find(qn("a:srgbClr"))
        if sf is not None:
            for el in sf.findall(qn("a:alpha")):
                sf.remove(el)
            a = sf.makeelement(qn("a:alpha"), {"val": str(int(alpha_pct * 1000))})
            sf.append(a)
    except Exception:
        pass


def _card(slide, x, y, w, h, *, fill="FFFFFF", radius=0.09, line=None,
          line_w=1.0, shadow=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(fill)
    if line:
        shp.line.color.rgb = _rgb(line)
        shp.line.width = Pt(line_w)
    else:
        _no_line(shp)
    if shadow:
        _set_shadow(shp, blur=10, dist=3, alpha=62000)
    shp.shadow.inherit = False
    return shp


def _rect(slide, x, y, w, h, *, fill, alpha=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    if alpha is not None:
        _fill_alpha(shp, fill, alpha)
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = _rgb(fill)
    _no_line(shp)
    shp.shadow.inherit = False
    return shp


def _oval(slide, x, y, d, *, fill, alpha=None, line=None, line_w=1.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(x), Inches(y), Inches(d), Inches(d))
    if alpha is not None:
        _fill_alpha(shp, fill, alpha)
    elif fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = _rgb(fill)
    if line:
        shp.line.color.rgb = _rgb(line)
        shp.line.width = Pt(line_w)
    else:
        _no_line(shp)
    shp.shadow.inherit = False
    return shp


def _set_tracking(run, pts: float) -> None:
    """Letter-spacing in points (OOXML spc is in 1/100 pt)."""
    try:
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(int(pts * 100)))
    except Exception:
        pass


def _textbox(slide, x, y, w, h, *, anchor="top", wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = {
        "top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM,
    }.get(anchor, MSO_ANCHOR.TOP)
    return tb, tf


def _style_para(p, *, size=None, color=None, bold=None, italic=None, font=None,
                align=None, space_after=None, space_before=None, line=None,
                tracking=None):
    if align is not None:
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                       "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}.get(align, PP_ALIGN.LEFT)
    if space_after is not None:
        p.space_after = Pt(space_after)
    if space_before is not None:
        p.space_before = Pt(space_before)
    if line is not None:
        p.line_spacing = line
    for r in p.runs:
        if size is not None:
            r.font.size = Pt(size)
        if color is not None:
            r.font.color.rgb = _rgb(color)
        if bold is not None:
            r.font.bold = bold
        if italic is not None:
            r.font.italic = italic
        if font is not None:
            r.font.name = font
        if tracking is not None:
            _set_tracking(r, tracking)


def _add_para(tf, text, *, first=False, **style):
    p = tf.paragraphs[0] if (first and not tf.paragraphs[0].runs) else tf.add_paragraph()
    if isinstance(text, list):
        for chunk, b, it in text:
            if chunk == "":
                continue
            r = p.add_run()
            r.text = _smart(chunk)
    else:
        r = p.add_run()
        r.text = _smart(text)
    _style_para(p, **style)
    # apply bold/italic per-run for markdown runs
    if isinstance(text, list):
        runs = p.runs
        idx = 0
        for chunk, b, it in text:
            if chunk == "":
                continue
            if idx < len(runs):
                if b:
                    runs[idx].font.bold = True
                if it:
                    runs[idx].font.italic = True
                idx += 1
    return p


def _add_bullet(tf, text, *, first=False, size=13.5, color="333333",
                accent="1E2761", font=BODY_FONT, space_after=7, indent=0.24,
                line=1.12):
    """Add a real PPTX bullet paragraph (buChar dot in accent colour)."""
    p = tf.paragraphs[0] if (first and not tf.paragraphs[0].runs) else tf.add_paragraph()
    for chunk, b, it in _md_runs(str(text)):
        if chunk == "":
            continue
        r = p.add_run()
        r.text = _smart(chunk)
        r.font.size = Pt(size)
        r.font.name = font
        r.font.color.rgb = _rgb(color)
        if b:
            r.font.bold = True
        if it:
            r.font.italic = True
    p.space_after = Pt(space_after)
    p.line_spacing = line
    # bullet formatting via pPr
    pPr = p._pPr if p._pPr is not None else p.get_or_add_pPr()
    pPr.set("marL", str(int(indent * EMU_IN)))
    pPr.set("indent", str(-int(0.18 * EMU_IN)))
    # remove existing bullet defs
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum", "a:buClr", "a:buSzPct", "a:buFont"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    buClr = pPr.makeelement(qn("a:buClr"), {})
    srgb = buClr.makeelement(qn("a:srgbClr"), {"val": _hex(accent)})
    buClr.append(srgb)
    buSz = pPr.makeelement(qn("a:buSzPct"), {"val": "90000"})
    buFont = pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"})
    buChar = pPr.makeelement(qn("a:buChar"), {"char": "\u25CF"})
    pPr.append(buClr)
    pPr.append(buSz)
    pPr.append(buFont)
    pPr.append(buChar)
    return p


# ============================================================================
# Page chrome (eyebrow, title, footer) — shared by content slides
# ============================================================================

def _eyebrow(slide, theme, text, *, x=MARGIN, y=0.62, w=None, dark=False):
    if not text:
        return
    w = w or (SLIDE_W_IN - 2 * MARGIN)
    tb, tf = _textbox(slide, x, y, w, 0.32)
    color = theme["accent"] if not dark else theme["accent"]
    _add_para(tf, _strip_md(str(text)).upper(), first=True, size=11,
              color=color, bold=True, font=theme["body_font"],
              align="left", tracking=2.2, space_after=0)


def _title(slide, theme, text, *, x=MARGIN, y=0.98, w=None, size=30, dark=False,
           color=None):
    w = w or (SLIDE_W_IN - 2 * MARGIN)
    tb, tf = _textbox(slide, x, y, w, 0.95, anchor="top")
    c = color or (theme["on_dark"] if dark else theme["ink"])
    _add_para(tf, _strip_md(str(text)), first=True, size=size, color=c,
              bold=True, font=theme["head_font"], align="left",
              line=1.02, space_after=0)
    return tb


def _footer(slide, theme, *, page=None, label="", dark=False):
    color = theme["on_dark_faint"] if dark else theme["faint"]
    # left label
    tb, tf = _textbox(slide, MARGIN, FOOTER_Y, SLIDE_W_IN - 2 * MARGIN - 0.6, 0.3,
                      anchor="middle")
    _add_para(tf, _strip_md(label).upper() if label else "", first=True, size=8,
              color=color, bold=False, font=theme["body_font"], align="left",
              tracking=1.2, space_after=0)
    if page is not None:
        tb2, tf2 = _textbox(slide, SLIDE_W_IN - MARGIN - 0.8, FOOTER_Y, 0.8, 0.3,
                            anchor="middle")
        _add_para(tf2, str(page), first=True, size=9, color=color,
                  font=theme["body_font"], align="right", space_after=0)


# ============================================================================
# Icon system (icon-in-circle motif)
# ----------------------------------------------------------------------------
# Baked monochrome PNGs (white glyph on transparent) live in _ICON_PNG (filled
# by dev-time baking from Lucide SVGs). At render time an icon is drawn as a
# coloured circle with the white glyph centred. When a glyph is unavailable we
# fall back to a clean numbered/lettered circle so slides never break.
# ============================================================================
_ICON_PNG: dict[str, str] = {
    "activity": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAJRUlEQVR4nO2dXawdVRXHf7vFQivQ4i3FapUPQ8QPbIoaaJSKpaWEB43RJiRqNCb6SGLUmPjkg/piSHwUH3gimhCtIQL2CxqgUIMFoyUGUUBs2tvbLyKCpS23fx9mrr33cuaeWTN7zuwzs37JSU7u3XvtNbPXWWvvPXvWBsdxHMdxHMdJHUmXSfqepL2SJvPPE5K+K2l52/o5DSLpdklTKmZK0m1t6+k0gKQ7JJ1doPNnOONG0DEkrZB0tETnz3DEw0E5FrWtQEm+BVxuKH8F8M2GdOkU42IAX6hQ54vRteggoW0FhiFpApgCFhurTgNXhBBOxNeqO4yDB7gNe+eT19kUWZfOMS4GUJXbo2nRUZIOAZICcAhYXVHEJPDeEILiadUtUvcAa6ne+eR110bSpZOkbgAxXPiWCDI6S+oGEKPz3AAWINkxgKRLgOPAkpqizgIrQwiv1deqe6TsAW6lfucDvAP4bAQ5nSRlA4jpuj0MFJByCHgRuCaSuH+GEK6OJKtTJOkBJF1HvM4HuErSByPK6wxJGgDNrOD5quAAUjWAJmK2jwMGkNwYQNJFwAlgWWTRp4CJEMKpyHLHmhQ9wC3E73yApcCGBuSONSkagMVVn8s/ZfFxwDxSNABLJ+0HnjWU93HAPJIyAElXAdcZqmzPP2X5kKQrTUp1nKQMALuL3pF/LLgXmEVqBmDpnFeBp4F9wMmG2ug8yRiAJOtDm10hhLdCCNPAo4Z6m/K2HBIyAOBTgOVljh0F34dxKbDeUL7TpGQAVte8a9b37YBl35+HgdSQ9CfDq19/HlD/gKH+M21cY4ok4QEkvRvb5s1BUz/LdHCdpDqbTTtDEgZANv2zPJcYFPMt44AAbDaU7yypGIAlJr8BPDng748DrzfUptMUkhZLOm6I3w8sIOt3BjnHJKXyA2iNFG7AJ4EJQ/mFXL0lDKwEPmEo30lSMIAqy79FPNRw205sJO0zuO2/lZD3d4O8QWOJXtGqB5B0GVkIKEuZqZ5lOnijpHcZyneOtkPAFmzv/peJ8ZZxQO9zCKRgAGU5DTxWotyevGwTOjixkBQkHTLE650G2bsNcg8ry0PQS9r0AGuB9xjKW1y7pexq4KOG8p2iTQOwTsEsgztLWfDp4OiRtMfgpg9WkP8vg/zdTVzjONCKB5B0MbZNGb+v0EzpMQNwc65T72grBGwCLjSUt278tNZZQk9zCLRlAJbUb28BVVz0rrxuWXwcMCokvWSIz4/XaGevoZ0XY17juDByD6DsPX1LsoYq7n8Gy2zgGknX1mhrLGkjBHzNWN46patT18NAk0haJ+m/Brc8pRqbNiQtku2cAevj5LFnJB5A2bLvVrLB3FJD1Z0hBMvbv3PI61qmg7coy0/QGy4YVkBZvp7Pk8XtKmvmFwM3US3nT+H2LwPbgS+XLLsMuJm57xyMDEmLgRuAdcD7sR2SMUcU8DLwQAjh+arKXCjpF5LOGVxoTI7F+DVKWiVp2tDu3XXbrKDjTZLulXQy8j08p6wPLWsu/3fZ2yIrY+WHEW/wfkO7z8Vqt4ReH5O0q5nbN4dtsjzxlLR1BEotxBFJl0a80T8ytv++WG0X6LNI0g+UnXA2KrYO0qVoEPiNBq+/DHdFzu1rXUtobDooaSnwa+DHZGlsR8XAPi0ygDZz7P88hHB/ZJn7gH8byjeyS0jSEuC3VDsEqy4D+7TtLWHzeRC4K7bQEIL1ecImSUNnSBW4h8S2oBUZwNvevh0B9wNfCiGcbUi+JQwsJ3IOAUlfAb4eU6aRgX1aZAD3NqjIfE4D3wfuDCFYNnNasS4LR/ulKtt6/rNY8ipSvk81umngg8oWmkaCpOcMuu2P2O5Pmrl9pbFNA3Olm1oIOirpHknrYt1gwzXdbdBzWtKqCG0uk/Rq5HtYlqELQUOtQvWXgqeBY8ArwDPAgTrr+3WQtBnbs4GvhhDuq9nmncCvKlR9AfgD8GaVZqm7FNxFJF0k6Q3DL6hW5+dt3mf81R6V9LkY1+sMQNJDxs4Y5ePo1yVdH/N6h5HaOsAosMwGLid7OleVG7A90ftpCOFAjfbMuAEMp86ysLXuL2u05ZRF0j8MbvmJGu1YNqW+FPMay9JHDwC2VcH1qpBDQNIK4EZDlYetbcTADWA4i4GNFdrYTIkdV7Oos/u5Mn01gEdoPoeApc4ZsrwGzqiQ9IghPh+UcSlVY/Jyal89ANhc7hrgI2ULK5vLW3YVteL+od8G0OR0sMncB1HpswEcAA4byltiuqXsJDCyjajz6a0BhBCEzfVuUIkcApLeCXzaIPfhXJdW6K0B5FhzCHymRLmNNJ/7IBp9N4CdZI+ry1LGtVvc/zTZlNRpC0lPGaZrL5SQZ1lm3juKa1yIvnsAsLngayV9oOifyvILFP6/ZtuN4AYQdzoYM/P5SHADgD+SbVkry0Ix3hL/j5OdfdwqvTeAfH+iZSl2owZsssz/VmaWMEOt3Aex6L0B5FhccdE8fwNZLoQm2mwMN4CMHdQ/eNLi/kVLSSicAiQ9a5i+/WVAfctLJ8kcXOke4DyW2cD1mpVDQNIa4MMNtdUobgDnscbk2dlOYxx82QpuAOd5kuo5BCzx/zWyfAVJ4AaQk+cQeNRQZbOkC5Rl9rLsGdzd4CvwZtwA5mJxzTO7ftcDll3Dybh/sO1a7QPWrdlbsL8w6waQKiGEg5KeB8rmLLAO/v4aQnjFrllzuAG8ne2UN4CPG2Un9esHHwMMwtJJi7Ddw2Tm/zP09ry8IpSlpz1Bljc4JqeAiRDCqchya+EeYB4hhDeByqeULMCe1Dof3ACKaCJWJxf/wQ2giCZidXLxH3wMUIiyQ6SqnHEwiJdDCLFkRcU9QDGWbGLDqHLw5UhwAygmZsxOMv6Dh4BCJF1CtnFzSU1RZ4CVIYT/1NcqPu4BCsg77KkIovam2vngBjCMGCP3ZN0/uAEMI0bnJTn9m8HHAAugLC3MIWB1RRGHgTVtvv49DPcAC1Ahh8B8dqTc+eAGUIZRHV7dCh4ChiBpApgiyxdoYRpYFUI4GV+reLgHGEII4QTZC6RWnk6988ENoCzbKtT5TXQtnHaQtELZUfZlmVTEk0+dBJC0ReWOej0t6da29XUaIDeCI0N++d75XUbScknfkfRY3uGT+fdvu9t3HMdxHMdxxoL/ARXtedStfw8hAAAAAElFTkSuQmCC",
    "alert-triangle": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAITklEQVR4nO2dXaxdRRmG3w8qDSgUDWhUjCilIRCCkijBppqY6A0q9UJAVLgxUQwQhcZ44R+BaERqRVMR9UJDREETCTUaUDGooA0BhQLR+ldisGqtrVUSWnr6eLH2IcfTtc9a39p7Zs3a53uSc7P3/HxzZvbs9ayZWVsKgiAIgiAIgiAIgiAIgmD2sb4DyAnwAknrJJ0p6eWSVo3e2i3pCUkPS/q5mf23nwiDqQOsAC4C7gEO0szTwJ3AecCy+oDMHMD5wPYWnT6OB4HX992OwAlwLPCtCTp+IYeAG4Gj+m5X0ALgJODRKXX+Qu4FVjVHEPQG8GLg9wk6f55fAc/ru51BDcBRwNaEnT/P94iLw/IArs/Q+fNc3nd7p8VMjGTgTEkPSVqRqcr/SDrNzP6aqb5kHNF3AFPiBuXrfEk6VtKnMtYXjIPK9fvgEPDavts/KYP+CqBy80clnerMOifpp5IeUTULniPp3A4h3CdpnZnRIW8wKcDVHT65DwNn1JT1OmBHh/Iu6qPtyx7gRGCPs7P+RLUgNK7MVwD/cpa5Azg6Z9sDScBNzo4COL9FuZd1KPejOdocjABOB55xdtJPWpZ9JPCIs+yngJelbncwArjL2UEHqe4VtC3/jc7yAb6esMnBPHTTvi91qOdOZx0zoYVFQ3W//3fOjvk38KIOdZ1CtTnEw/3EOkE6gA3ODgG4aoL6PtehvtDCFNBN+/4ArJygzuOBXc46/wIcM822B5KALzs7AuBtU6j3Ax3q/dg02hyMIKH2tag7tLBvSKx9LeoPLewLMmlfizhCC3NDRu1rEUtoYW7IrH0t4gktzAU9aF+LmEILc0FP2tcirtDC1NCj9rWILbQwNfSsfS3iCy1MBYVoX4s4Z0ILi1IUqk2e2yStcWTbJ2mNmf3dUc9xktZLOnv00kOS7jCzfY4yTpH0mCTPRecvJa2NTaRjIIP2ARcAu2vK2Q1c4CwrtHBakEH7gAuppuJxHAIudJQXWjgtSKx9wHHUf/IXs5vqK6JtuaGFk0IG7QMudZR9iaPcQWthKWcDN8l3tm9O0geddZyVIq2ZdYnlGEnXOvMkofcBQLVX/83ObF8xs23OPM91pHU9BMLM7pG0xReOLqEALex1AFBp3/XObPskXZMgnEn5kKT9jvQm6fP0vFrY9wxwpXzOL0nXeJw/F2b2R0neG1LnSmptHTMFmVf7gJsd9dzcsY7BaWGfM8C1ko535rnKzDzTbFbMbK+kTziznSTp6gThlAs9rPaRYQYY1TMoLexrBsihfb0wNC3MPgDIp329MSQtzDoAmC3ta2IQWph7BpgZ7WsitHAR9LzJk0wXgYvqLF4Lc84AM6d9TYQWjqCATZ70MAOM6i1aC3PNADOrfU1MoIXXJQjnMJIPAGC9Zlz7muiohe8hgxYmHQBU2vcZZ7ahal8TRWph6hlg2WhfEyMt3OzMNlwtpLCzffR0EbgohuK0MOUMsOy0r4llo4UUoH01MfU+A4ziKEoLU80Ay1b7mihNC6c+AChX+55JlNZNSVo41QFA2dr350Rpu1KkFk4EhT3SZVFsa4C5FvHMAV517RrTxg7/rzLPFlKY9o2JcXOLmLyuPkk8xWlhZyj0kS6LYlwJ3L5EPLeTcUCOYhr+2UIK1L6GeNcDW4AnR39bqC5e+4ilKC3sBIU/0qV06PbImW/0HbekZz9NXrI/0qV0GOIjZyjoSZ5DhyE+iZSCtW+IMCQtZADaNzQYkhYyAO0bIgxBCxmY9g0JhqCFhPYlhZK1kNC+LFCiFhLalw1K1EJmRPv4/1vBO4EfUj1EsqilVkrSQmZA+2heDPo+Ba20UZIWMgPaR7vl4Fv6jnMhlKCFzID20X5DCMCr+453HhJrYdstYbOwyfOtat/eXpaG60i9ibTxH0I1jXs3eX61wLN9JydKm5yOm0jfDbymKdGSA4Bqk+cNzor3Svq4M08O9iZKm4sNkg440h8haRMNZtM0A1wh6VRHpZJ0nZntcubJwb2J0mbBzLbLf7ZwrbqeLQROAPY6Lz62j2aN4gAMeKBFGx4HPNc72aCbFu6gi4oDn3ZWBIVp32KA1VQ3f8bxDwpfswAu69Av7/dWshL/SPtxojZPFeClwK3AgQWxHwS+C5zcd3xNACuAbc6+8V2QA29xVjC41T5gFbAWWAc8v+94PABvcvYPrv4Bvugs/KaE7Q1qwL9auMFT+FZHwXuAExO2NaiB6s7mgYa+WchtdeWM08DVjlg2F6p9M81IC7/tyFJ73vGwAUClQJ7vxLscaYPpcrcjbe0sXTcDrJDvJ2VLvGu2XNjjSPucuhfrBsB+SQcdBb/SkTaYLp6v6qfqXjxsAIx+2Hino+CLHWmD6fJOR9on614cdxH4uKPgd1D4HcBZBLhc0jmOLI/VvThuAPzCUbBJuo1q50qR99BnCeBo4JOSbnRmva/uxdqLPeBVkn7trECS/ibpZ4oLw1S8UNIb5LM0qXro1UvM7J+L3xh7tQ/8Rr7f2w3K5Q4ze3vdG0vtB9iYKJggP5vGvbHUDHCkpAcVs8DQ+YGZnTfuzSVv+FDtKbtfvg2hQTnsk3SWme0Yl2DJLWFm9oCkj0w5qCAPSHrvUp0vtdgVbGYb5d+LFvTPh83sO02J2u6Tv0LSZyeLJ8jEnKQrzazVbm7XYUjgXapmg1UdAgvSs1PSpWb2o7YZXA+LNrNvSjpD0i2SDvliCxKyX9IXJJ3u6XzJOQMsBFgt6X2q9p2X89TK5cVvJd0q6Wtm5lnAe5apnIcHTpN0tqpDJCdIKvJswAzwtKRdkrZL2mpmT/QcTxAEQRAEQRAEQRAEQRAEg+B/dEf9zrev5UQAAAAASUVORK5CYII=",
    "award": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAALc0lEQVR4nO2daaxdVRXH/wvLWIZSCxQiQ4W2vAJVJIACUSExEhASJOIHSkGRmEikDnw38QPRxKiAmkiQwcYPIqESSQWlUAVkaoJ0eKUFSmwFKrSUodPr9PPDPo/3enmv764z7L3vffeXNE1e9jnnv9dZd+999rCW1KNHjx49evTo0WPcYakFNAmwn6RZkvokzZR0qqQTJE2UNEnSoZKQtEXSpuL/dZJekrRK0kpJ/Wa2J7r4SHSdAwDTJF0i6SJJX5D08Yq33ChpsaTHJS00s9cq3q9H3QBHAHOBvwN7aJYlwDxgSup6j3uAY4CfAFsafukjsR34LfCJ1HYYdwDHAXcAAwlefCuDjnBcart0PcAE4PvA+0lf+ci8B3wPmJDaTl0JcCbwYtp33Bb/Bj6d2l5dBWGAtzXxi/WwHZiX2m4dDzARuC/xy6zCH4GJqe24L7KdBwCOlPSQpPNSa6nI85IuNbO3UwsZiSwdgDCifljSGam11MRKSReb2drUQlrJzgGKl/+UpJNqvvVqScsUpnjXSnpHYepXClPDkxWmiWdKmi1pes3Pf03S+Wb2Zs337R6AScDSmvrfAWABcDUlvtEJcw3XAH+mvvmGF4EjmrBdxwMcBDxRg5FfB34IVF0DGK5tCnAz8EYN+hYDB9alrWsAfl3RsBuA7zRpXOBA4EZgY0WttzalsSMBvlbBmHuAO6nxF9+G3inA7ypq/mosvVkDTAPeLWnITcAVCbVfWVH7iam0ZwPwUEkDvkRY+0+t/2RgVck6PJhaf1KAK0oa7nngqNT6BwEmA/8qWZfLUutPAnAI8J8SBnuBDD+lCJtSyixWrQEOSq0/OoRdNV5eAaam1j4awLHFC/VyY2rtUQH2x//r3wacmVr7WACz8a9crgUOSK09GsANTgMBfDu17nYhzEd4+WZq3dHAP937aGrNHgADHnfW8YXUuqNA2NnjYQDoS63bCzAL2OGsa/TVz/1iP1DSHGf5O8xsZSNKGsTM+iXd6bzsmia0ZAOhafQsqOwATkituyzASc5WILv9ArUCnO4wBsD81JqrAvzBWeeo3V3sLuAiZ/l7mhARmXud5b02qkRsB7jQUXadwnm8TmeRpDcc5bvaAc51lH24G07lmtluhf2N7eKxUWWiOQBwuKRjHZc81pSWBHhasuOAQxtT0kLMFmCms/xTjahIwxOOsia/rUqTqwNslvTfpoQkYK2GdiC3Q1c6gGf9fpWZ0ZiSyBR1We245OimtLQS0wEOc5Ttxr3z6x1lu3IM4HGAzY2pSMcHjrIeW1UipgN4DklubUxFOjxO3ZUOMOAo242bIzzbvrY3pqKFmA6Q5S8gIp46vd+YihZiOoCnD5zcmIp0HOkoG20MFNMB3nWUrftkbg7McJTd1JiKFmI6wCuOslOBSY0piUxRF8+3vcdWlYjpAKuc5c9qREUavAs8XluVJpoDmNnr8g1uPEvHueOpy3tm5pk0qkTs5eBljrJfakxFfDx1WdqYihGI7QCeZdFzgGiLIk0BTJf0GcclUZfBYzuAt3JXN6IiLtc6y3fTPoi9IYSB8RybehM4OLXushAOwK531HczkcPHRG0BzGy7pIWOS6ZK6uQjUzdIOsZR/i9m5pky7zyAyx2/CIB1RNwiVRfAYfiDSl2aWnfjEE4Gv+U0zE9T6/YC/NxZx/8B+6fWHQXgZ07j7KADjoYPApwF7HTWseOcvDTAVPxn6F8m7CzOGkLT740ZtI3xlnAC+I3TSAB/ImQCyxJgP0J0Ui+3p9YeHcLBye0ljHVbau2jAfyqRH22Acen1p4E4MclDAZwC5BNoGvCqedbStblR6n1JwM4GHi1pOHuIYP8PMDHCEmjyvAK4zFC2HCAS0oaD2A1Ya49lfbpxUsswx7gy6m0ZwVwWwUn2AV8N4Hmm4pnl+UXsTVnCyEK95IKxgRYSIRoIoTB618ran2O8RQWrh0IMXc3VDTsAKE//mRD+upIVvk2GcQ4zhLgc9STBnYX8Agh40fpySPgcELKur9RrbkfZDMQ9fz/WGTzKTUIYUFkgaS65sV3KWTuWqyw22aVpLVmtnHYM01hK/qJCrt3Z0v6oqSzJdX1pbFT0uVm5gkW0TjZOYAkAVdJmq/mTwgNTxrVJAOS5pjZ/Q0/x02WDiBJwIUKoVU6fbC0U9JlZvZIaiEjka0DSBLwtKTPptZRkafNLNvkl9kurBR0Q0qVk1IL2BfZtgDAFElZplstwdG5po7NuQWYnVpAjZyeWsBo5OwA3ZI3WMq4Lj0HiEO2dcnZATxdwKOSHpG0uyEtw9ldPMuTxCLb7ixLByBs+5rluOQBM7tY0vGSblY4glZnmJVtCjF/fyDp+OJZCxzXn0amW9my/AogrPF74updYGZ7RRYlbLQ4T9LnJfUpBF+cIWmsk0bbFKaLV0vql/RPhW/5vRwKuEC+CKDTzSzauf92Sb6jZhQ8TSaSVrT+sXhhj2nYWbtizv8ohTh8R2ooHt9mhagcH0ja0GaQymXFs9v9Ec1WxMAP7ZKrA3gGTWvNrK3wM8WLfav4Vwkzew9YJ6ndPQhnSHqg6nPrJst+ST4H8MQcqBvPWf4svwS6wQGiBlRoweN8WX4JZOcAwCGSTnZckrIF8Dz7ZKDpZWc32TmAwrSpR1endAHeT9so5OgAnuZ/QL7PxbpZJV8I3Oy6gU53gJVmtrMxJWNgZrskveS4JLuBYI4O4PmVpGz+B+noL4EcHcCzdJqDA3T0l0BWDgAcK19qmZSfgIN4HGBKUcdsyMoB5P+FdFoLIGXWDeTmAB7jvGNmnoycjVCEwN04ZsEhsuoGOtkBcmj+B/G0Ar0WYB94fh09B6iBbFYDCcEeTnVcsrzCs/aXdIqGUtmul/RyhTkFjwP0AROKOYQegwB9zoOW7gMjhGAUC4pDmq1sAR4EvoIz/AxwrlN7dlPCyQGuchhwN47oocBM4EnH/Z8BTnPc/9BCU7t8vZyV6ienMYCn/19jZm0lVgKulPSspPMd9z9X0hJgXjuFCy2vOe6fzTigUx1gzD6XEJn8Vkn3SzqihJ6DJP0SmN9ma9ORU8I5OUBtn4CERBPPSLqpkqLAHIXWYCwH7cgp4SwcADhMvoOgoxobmCNpiaRPVdU1jJmSnh2jS/A4wIl0UVa0ygDnOQZQMEIqGULMwTud9ynDiF0CMMN5H8+YpDGyaAHkaxK3qWV7NUNN/vV1ihqF0bqEV+VLep3FOCAXB/AYY4WZfXgEjKEmP2a/+pEuodDU77hHzwGG4Z4CZmiUP19DBzxiMtJXQsd9CeTiAJ5NIMuBPoXIX3WM8qsyR9JzwOlyrgmQQcDr5A5ACJXuGRHPVnj5VYMubJX0LUlzNRQtrCx9CpNNnqwmkxQOs45vqBYsuiz9xS92UEMfsCyBjuRJopK3AIo/KTJf0tlm9uFqopmtVJj+vSuyluQTQjk4QKzB0FZJ15vZXDP7SJNvZlvN7HqFTJ9Vu4R2yWIgmBRgaYSmdq8mvw1NsbqEHPY0pgM4gOrRt8fiXkqcySOkfb2rYW07GM9h44HTGjTuFuAbNWicy8gbSOqi7X0HTZB6DOA5A+ChX9I5ZnZ31RuZ2e8lnaMRopDUxJSG7tsWqR3AM3feLvcqvPzaXpiZ9Ss4wT113XMYbW1s6UoIWTbrSMQAoZm+LoLma6mvS9hJBybGrhVgUQ2GXE7EjZbArOKZVVkUS3O2EHbhVuFuQlSR2LoPKZ5dheQzgVlAyAnsZTMwNwPt11KuS7gvtfZsIGyr9mzbHlwRzAJCl7DCof9Jxnvf3wphff829j0o3AXcToImfyyAiYSM6Ps6HzCoP5t0scnXo1sBTpF0nULWrmkKn6prFLJ+3Z1juNXhELanXaeg/wQFG6+R9A91gP4ePXr06NGjR48e44D/A0OttyhWKHsKAAAAAElFTkSuQmCC",
    "bar-chart": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAFBUlEQVR4nO3du49VVRTH8e/mpQkik7EyKAISJTYQGgsfUbS0UAqNQU3GwsRor5V/gCGxMeiEhA5fBdEGK2M0mpAIBgsKCilQNCiRkcI4CPOzmD0JDHeGs+95Muv3SWiGtc9acH/cuXfPPRswMzMzMzMzMzOzlS+Ns0jSBLAX2JW/dBI4klKaaWowGyhJ+yRd1I0uStrX93zWovzgz4148BfMOQQrlKSJJf7lj3ommOh7XqtmVUHtXqDKA7vw+sBuASUB2HXzkrFqrUclAbitpVrrUUkAbAVyAIJzAIJzAIJzADJJOyR9KOm0pAuSjkl6S9L6vmcbBEnTFTaBFkz3PW8JSa9Kml3iz3JG0oN9z9iW8M8Akp4CDgLrlijZChyVtKG7qboTPgDAu9z872Eb8EYHs3QudAAk3Qfsrlj+bJuz9CV0AIAtBbVb2xqiT9EDsNT3/VHWtjZFj6IHIDwHIDgHIDgHIDgHIDgHIDgHIDgHIDgHIDgHIDgHIDgHIDgHIDgHILg1fQ8QkaQ7mf+AycKHUX4EPk8pXep6FgegY5KeBz4AJhf91l+SXk8pfdblPA5AhyS9AHzM6JNZJoFPJKWU0qddzeTXAB3JT/sHWP5YngQcyLWdcAC68xw3Pu2PMkmHH0B1ALqzs6XaWhyA7pTcYnZHa1Ms4gAE5wAE5wAE5wAE5wAE5wAE5wAE5wAE5wAE5wAE5wAE5wAE5wAE5wAE5wAE5wAE5wAE5wAE5wAE5wAEN+gASFqfz+w/ls/wP53P9N/R92wrxWDvDMpn9H/J9Wf03gU8AEzl26gO9TLcCjLIZ4B8Nv9Rlj6geR1wMJ/1bzUMMgDMn82/7SY1q5g/699qGGoAqt4atTuf+W9jGmoASs7m39LWEBEMNQAlZ/OXnPlviww1ANYRByA4ByA4ByA4ByA4ByA4ByA4ByA4ByA4ByA4ByA4ByA4ByA4ByC4oQZABbVzNfqUrC2Zqe76ur0qG2oA/myptq8+XfeqbKgB+KZi3QXgVI0+p/I1qqg6UxPr6/aqbKgBeA+YrVC3P6V0ddwmee3+CqWzeaY6vgKOV6g7nmuHRdK0qptuoN+UpKvL9PhC0uoG+qzO11rKVUlTdfvkXtslnVum1zlJ25vo1biuA5B77pF0YtG1z0t6W1JjN7VIWpOveX5RrxOS9jTVJ/faJOkjSZev6XM5f21Tk70a1UcArum9OYdhZ5MP/Ig+a3KPPZI2t9Un99oo6ZH8a2ObvZYz2FvDrpVSOguc7aDPFeCntvvkXn8D33fRazlDfRFoHXEAgnMAgisJQMn2pIN1iyh5oP4pqL2/dBDrR0kAfiuofbTtt1HWjJIAlOy5rwWm23zPbh2TNCHpSsFmkCQdlXRv37Pb0pb7j4xvIOlr4InCHv8B3wE/U+9n9zbaLHASOJJSmildXBqAKcAHMw3TDPBmSulwyaLSANwOnAHuLllnnRHwckkIit6vp5T+Bd4pnco6k4D3JU1UXTDOhs0h4Nsx1lk3JoC9VYuLA5BSmgNeAv4oXWud2VW1cKwt25TSL8AzwKVx1ttwjL1nn1L6AXgaON/cONaQk1ULi94FjCLpHuAw8Hjda1kjZoCtVfcEav/ULqX0K/Ak8Brwe93rWS1ifi+g8oZQ7WeA67rP7xO8CLwCPAbU/tSuVdb+RlCJ/F70YeAhYBOwoa1ewdXaCjYzMzMzMzMzM7MI/gfYDV5J/ydD+AAAAABJRU5ErkJggg==",
    "bar-chart-3": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAFBUlEQVR4nO3du49VVRTH8e/mpQkik7EyKAISJTYQGgsfUbS0UAqNQU3GwsRor5V/gCGxMeiEhA5fBdEGK2M0mpAIBgsKCilQNCiRkcI4CPOzmD0JDHeGs+95Muv3SWiGtc9acH/cuXfPPRswMzMzMzMzMzOzlS+Ns0jSBLAX2JW/dBI4klKaaWowGyhJ+yRd1I0uStrX93zWovzgz4148BfMOQQrlKSJJf7lj3ommOh7XqtmVUHtXqDKA7vw+sBuASUB2HXzkrFqrUclAbitpVrrUUkAbAVyAIJzAIJzAIJzADJJOyR9KOm0pAuSjkl6S9L6vmcbBEnTFTaBFkz3PW8JSa9Kml3iz3JG0oN9z9iW8M8Akp4CDgLrlijZChyVtKG7qboTPgDAu9z872Eb8EYHs3QudAAk3Qfsrlj+bJuz9CV0AIAtBbVb2xqiT9EDsNT3/VHWtjZFj6IHIDwHIDgHIDgHIDgHIDgHIDgHIDgHIDgHIDgHIDgHIDgHIDgHIDgHILg1fQ8QkaQ7mf+AycKHUX4EPk8pXep6FgegY5KeBz4AJhf91l+SXk8pfdblPA5AhyS9AHzM6JNZJoFPJKWU0qddzeTXAB3JT/sHWP5YngQcyLWdcAC68xw3Pu2PMkmHH0B1ALqzs6XaWhyA7pTcYnZHa1Ms4gAE5wAE5wAE5wAE5wAE5wAE5wAE5wAE5wAE5wAE5wAE5wAE5wAE5wAE5wAE5wAE5wAE5wAE5wAE5wAEN+gASFqfz+w/ls/wP53P9N/R92wrxWDvDMpn9H/J9Wf03gU8AEzl26gO9TLcCjLIZ4B8Nv9Rlj6geR1wMJ/1bzUMMgDMn82/7SY1q5g/699qGGoAqt4atTuf+W9jGmoASs7m39LWEBEMNQAlZ/OXnPlviww1ANYRByA4ByA4ByA4ByA4ByA4ByA4ByA4ByA4ByA4ByA4ByA4ByA4ByC4oQZABbVzNfqUrC2Zqe76ur0qG2oA/myptq8+XfeqbKgB+KZi3QXgVI0+p/I1qqg6UxPr6/aqbKgBeA+YrVC3P6V0ddwmee3+CqWzeaY6vgKOV6g7nmuHRdK0qptuoN+UpKvL9PhC0uoG+qzO11rKVUlTdfvkXtslnVum1zlJ25vo1biuA5B77pF0YtG1z0t6W1JjN7VIWpOveX5RrxOS9jTVJ/faJOkjSZev6XM5f21Tk70a1UcArum9OYdhZ5MP/Ig+a3KPPZI2t9Un99oo6ZH8a2ObvZYz2FvDrpVSOguc7aDPFeCntvvkXn8D33fRazlDfRFoHXEAgnMAgisJQMn2pIN1iyh5oP4pqL2/dBDrR0kAfiuofbTtt1HWjJIAlOy5rwWm23zPbh2TNCHpSsFmkCQdlXRv37Pb0pb7j4xvIOlr4InCHv8B3wE/U+9n9zbaLHASOJJSmildXBqAKcAHMw3TDPBmSulwyaLSANwOnAHuLllnnRHwckkIit6vp5T+Bd4pnco6k4D3JU1UXTDOhs0h4Nsx1lk3JoC9VYuLA5BSmgNeAv4oXWud2VW1cKwt25TSL8AzwKVx1ttwjL1nn1L6AXgaON/cONaQk1ULi94FjCLpHuAw8Hjda1kjZoCtVfcEav/ULqX0K/Ak8Brwe93rWS1ifi+g8oZQ7WeA67rP7xO8CLwCPAbU/tSuVdb+RlCJ/F70YeAhYBOwoa1ewdXaCjYzMzMzMzMzM7MI/gfYDV5J/ydD+AAAAABJRU5ErkJggg==",
    "book": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAEoElEQVR4nO3cPYhcZRSH8eesIa5sFI0QTdQmbmO7mBRRJEYSibCN2thExKjxAyyjKYLBxsKPJn6AWhix0AgpVCImIUrEiJV2amHjBysjC7ILIevqsZgJblLs7p333HfGOf9ffec9l73P3rlz7zAgIiIiIulY7YHuPgZsBXYAU8AksAEYr70vQ24W+AE4DbxvZj+3MaRaAO5+E/AksAfYWGvuiFgE3gOeMbOZyIVbD8Dd1wOHgEeBtW3PG3F/AA+Y2cmoBVsNwN2ngTeB69qck8wCcJ+ZfRyxWCsBuLsBzwMH2pqR3Dywxcy+L10o/OD0LvLeAB6JXlsu8hVwu5l5ySJjQTuz1Mvo4NewDdhdukhoAO7+MPB05JqyrL2lC4S9Bbj7JPAtMBG1pqxoHrjGzBb7XSDyDPAqOvi1rQNuLlkgJAB3vwvYFbGWNHZ9yYvXBO3E/j5fNwMcBc4CnaB9GQXH6P53r8Zgb665+2Z3/8ebOe/uB9z98oHu/JBy99kGf8udJbMizgD30+xi8hwwbWanAmZLoYhrgKbv/ft08IdHUQC9W75bGrzkjJkdKZkpsUrPADcAVzXY/nDhPAlWGsCmhtuHPcaUGKUBrPajCsCcmc0WzpNgpQFc1mDbvm9XSnvaeBoo/yMKIDkFkJwCSE4BJKcAklMAySmA5BRAcgogOQWQnAJITgEkpwCSUwDJKYDkFEByCiA5BZCcAkhOASSnAJJTAMkpgOQUQHIKIDkFkJwCSE4BJKcAklMAySmA5BRAcgogOQWQnAJITgEkpwCSUwDJKYDkFEByCiA5BZCcAkhOASSnAJJTAMkpgOQUQHIKIDkFkJwCSE4BJKcAklMAySmA5BRAcgogOQWQnAJIbk3tge4+BmwFdgBTwCSwARivvS9D7Opag2oGsNbdXwD2ABsrzpVl1AxgAthfcZ6sgq4BklMAySmA5BRAcgogOQWQnAJITgEkV/1W8CVmgKPAWaAz4H0ZJseAdTUGDSqABeAQ8JKZnR/QPgwtd/+r1qxBBHAOmDazUwOYLZcYxDXAPh384VE7gDNmdqTyTFlG7QAOV54nK6gdwMnK82QFNQOYM7PZivNkFUoD+DtkL2RgSgOYb7Dtle5+beE8CVYawK8Nt99ZOE+ClQbwG/Bng+2fKpwnwYoCMDMHvmnwktvc/aGSmRIr4lPAZw23f83d7w6YKwEiAvgQ8AbbjwMfuftBd78iYL4UsIhF3P1ToJ//6g7dgL6mez0hXU0eB+8ysxP9DooKYDtwOmItaWy7mX3R74tD7gSa2efA8Yi1pLGZkheHnAEA3H0z8B2VvskiAMwB681ssd8Fwp4FmNlP6HN+bSdKDj4EPwwys3eAFyPXlGW9VbpA2FvABe5uwOvAY9Fry0W+BO7o3YzrW/jj4N4OPQ48R7P7A7J688De0oMPLZwBlnL33cDb6AchIi0A95rZJxGLtfqFEDM7DtwCvALo69/lOsA9UQcfWj4DLOXum4AngAeBG2vNHRGLwLvAs2b2e+TC1QK4oPcjUVPAncCt/PcjURO192XIdYAf6d5h/cDMfhnw/oiIiIjIqPgXLHIsRkaF8NUAAAAASUVORK5CYII=",
    "book-open": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAGMElEQVR4nO3dTegUdRzH8ff3b1aiaVaKZWZgYWGRWlpUlBIe1ENlIkbRyfRY2MMhj2VERAc9hYaEQp4Ug9QwTEOjKB8OPoaZIiWYmZGPqPvpMPMHkZ2/O7vz29md3/cFe5n5zW++s/OZp98OLDjnnHPOOeecc84556rP8i4gaQgwBRgPjAB6Cq7pek4DJ4CjwD7goJldaXMNAEjqB4wFxgGjgWHArW0u4wpwHNgNbDGz//Is3HAAJI0FFgFzgJvyrCSwM8B2YAOwxsyOhVyZpHuAWcB04ClgYMj15XQBWA18YGa/FdKjJJP0rqSL6nxXJH0r6SUlR2dR30E/SbMlbZZUK3H7GnVe0kJJuc/w1254j6TPy92Wph2QNFctfAlKwv+ypF9L3ZLmLZPU/CVa0sdlb0EBtkl6qIltf1jSDyXXXoTFfW1n5tEh6TlgU19tushFkvuXT81MfTVUcsZ4B3gfuLENtYVWA6aa2ff1ZtbduemXsJPkTr9K1gKvmdmZejMlDQJWAc+3tarwfjGzSfVmZAXgGWBr0JLKswOYbmZ/XT1R0nCSJ4mJpVQV3tNmtv3aiVk3CC8GLqZMjwLfSRrWOyHd+Zup7s6HjH2aFYDHAxbSCcYB6yUNTE/7G9JpVfZEvYk3ZDQeHbCQTvEYsJLkIKjykd+r7j7Nugc4BQxtsOOvgT+aLCqvfsBw4AHgPjrnCUXAIeAAyTB1u4amRwIzG2z7j5nd1lBLSadyPGdOa7r8Fki6Q9KrkjaqnNG5K+m6X5F0e0nfwbQc9Z7K03HHB+CaesdJWpej5latkfRgB2x3ywFo9y95QZjZXjN7nuRHmr8Druok8IKZzTKz/QHX0zaVCEAvM1sLTCK5FhdtPzDJzNYF6Ls0lQoAgJn9DjwLHCyw2/3AFDM7UmCfHaFyAQAwsxPADKDxG59sJ4EZaZ+VU8kAAJjZYeD1ArqaV8Ujv1dlAwBgZmuAr1roYm3VrvnXqnQAUu+RDNTkVSP5CbnSKh8AM9tL8l5DXpuq8qjXl8oHILWqiWVWFl5FB4olAN+Q7zJQAzYGqqWjRBGA9BHucI5FDplZyBHFjhFFAFJ5dmgR4wddIaYA3Jyj7YBgVXSYmAKQ9fJLPf2DVdFhYgpA3TeBC2jb1WIKwJFAbbtaTAE4HahtV4spAK4OD0DkPACR8wBEzgMQOQ9A5DwAkfMARM4DEDkPQOQ8AJHzAETOAxA5D0DkPACR8wBEzgMQOQ9A5DwAkfMARM4DEDkPQOQ8AJHzAETOAxA5D0DkPACR8wBEzgMQOQ9A5DwAkfMARM4DEDkPQOQ8AJHzAETOAxA5D0DkPACR8wBEzgMQOQ9A5DwAkfMARM4DEDkPQOQ8AJHzAETOAxC5PH+klOUNSbML6KcVZ4DjwCFgl5kdDbESSfcCE4AxwAjglhDryWFkqx1kBeA8MLTBPma2WkTRJB0D1gOrga1m1syfRyPJgCnAXGAGcHdRNZbgbL2JWQE4AtwVrJTwRgEL0s9BSZ+Q73LXI2k+8DZwf4D6ynCk3sSsAPwEPBmslPYaCywD/s2xzBxgXphySvNjvYlZR8WagIWUZUiOtoODVVGetfUmWlZrSTuAicHKce30s5lNrjejr+viW+T7y3XXmWrAwqyZmQEwsy3ARwEKcu31oZlty5qZeQkAkNQDfEb1bohisQxY0NdjcJ+PRmZWA3ofhy4UW5sL6ALwppnNv94YSJ9ngKtJGgMsIhkUGdBafS6Q88CXwGIzO9zIAg0HoJekQcBU4BGS4dD+efsoWA/JoNUE4M42rfM4sAv4k+Qmq0yX0zp2k4x6nim5nnJIMkmTJX0h6bKKd0nSCkmTyt5Wdx2SJkraV+DO3yNpfNnb5XKQNFjS1gJ2/hZJZf/q55qRhqCVM8Ee3/ldTsnloJl7gkvy0341KLkxzGtF2XW7gih5Osgrirv93OMA3UjJmz1naXwA6xwwqNk3ibpJFC+FpjvyZI5FTsaw8yGSAKTOBWrb1TwA9Z0PVkWHiSkAOwO1dd1A0ihJtQbu/muSRpVdrwtA0tIGArC07DpdQJKWZJwJapKWlF2fa4P0crBc0u70s9xP+84555xzzjnnnHPOuer6H2XPvfU/3NCQAAAAAElFTkSuQmCC",
    "bot": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAG30lEQVR4nO2dW6xeQxTHf6vUrShxV7dSl7iHqoSGBBUR90QfRIQEiUsiIoLUgxeJW+KFBy9uQRFBE7eESAQRSVsixL0UDQkq7tXq+XuYjzT17XPOnD2z5/u+Wb+nc76z91rrm/mfvffMXrMGHMdxHMdxHMdxHMdxHMdxHMdxHMdxnBHESgcwGSTtBcwBpmUw/yfwvpn9nMG20wZJJ0lapvyslfSYpFmlv7PTQ9IlktZ30PkbskrSnNLfvUsG8hYg6SDgXWDzAu6XAvPMTAV8d06Oe2oKrqVM5wPMBU4u5LtzBlUAJ1buvzMGVQAzC/vfqbD/zhhUAaws7H8gn41yMKgCeLZ0ALUwqAK4F1hROogaGEgBmNmvwOnAF6VjGXU2LR1AE2b2saTDgKuBc2k/FbwV5YaWTmkk3RcxI3hf6Xi7YiBvAU53uAAqxwVQOS6AynEBVI4LoHJcAJXjAqgcF0DldPbaU9I0YB5wJnAwsAuwJzCjoxBipoL/Av7IGMuG/Ex4/b0S+Ah4ysw+7ch3fgFI2hG4CbgA2DW3vxHhLeB+4EEz+zuno2wCkLQlcB1wPbBtLj8jzjvAFWb2di4HWQSgkF//DHBMDvuVMQbcDizKkamcXACS5gJLgN1T266ch4FLzWxdSqNJBSBpP+BtYIeUdp3/WAKcZ2ZjqQwmGwZK2oaQy+edn4+zgVtSGkx2BZD0EHBRKntOIwLOMrPnUhhLIgBJhxOeWH1iqRu+AQ4wsz/bGkrVYXcktOVMzB6EXMnWtL4CSNof+CRBLE4cPwGz2l4FUmQFnz2Fc74lzHQtBX5LEMNUmAHs0/v5S+D3QnFsB5xFmCndJOK87YHTCPMt5ZD0euQa/KclbV006AFE0vGSvo9sy8Wlg54m6a+IgJdK8tz8BiTNk7Qmoj1/KB3wzhHBStIZRQMeAiTdGtmm25QM9oiIQMckbVEs2CFB0qxIARzSxl/boduOEcf+YmZrWvobecxsFSFHYLK0esXeVgAx5yebv66AmLZq1Yc+eVM5LoDKcQFUjgugclwAleMCqBwXQOW4ACpnYItETYSkmcChvV87r/df2v9AIGlBxJz16kQ+ZynU9l+7ge3O6v134V/S6oh2XZDC51QD7VQAkuYo1PRvImu9/678q0MBDM0zgCQDFjP+gpPdgcW9Y0fKfy6GRgCEGv5zJ3Fcrnr/pf1nYZgEEFPDP0e9/9L+szDuKEDSDMYv1To9U1z9iKnhn6Pef2n/TTwtaeP1gmPAZ4SVWveYWWPibaMAJB0IvAjMThFlAmLuqznuwaX9N9GUYLsDcCxwuaTTzeyjfgf1vQUo5Jm9QNrOH5oHowEgZVvNBp6X1LdGQ9MzwFXAvgmDAK/UHUPqttqX0Kf/o0kA5yQOALp9Xhh2crRV3z5tEsDeGQIYphFHaXK0Vd8+bXKUY167io0YE5Gjrfr2aZMAXssQQNZqVyNGjrbq26dNAribUCsvJantjTKp22otoU//R18B9MaMV5A2l399QlujTsq2GiOUmvuw3x8bHzbM7AFgAbA8YTBOt7wDnGpm9zcdMO5UsJm9ChwtaS/6TwUfDdzWNspJElNH4NcR9N/EjcCyjT4bAz43swl3YJ1URpCZfQV8tfHnHb/2fC/TscPiv4nlZvbKVE8eprH5M8BkkkpWk2fr2dL+szA0AjCzX4ArGX+MLODK3rEj5X8gUZmcwIWSfuxj/0dJC1P4KO1fHaaEDV1WsJk9Keklwtz2Ub2PlwPPdvGfV9r/QFHiClADXV4BhuYZwMmDC6ByXACV4wKoHBdA5bgAKscFUDkugMpxAVROWwHEJI9uLa8VPCEKW+zGlNNvlXvQVgBfRxw7HTivpb8aOJy4dQGrcgUyIQr7BcTUt/9cobSK04CkByPac72ksgtuJL0REbAkvayw6tjZCEmXKpTVnywfl44ZSddECkCSPpB0hqSYPXJGFklHKu4//1/uaus7xa5hewArmdrzxDrKbRo1KMwANpviuSeY2ettnKfaOPJRwq5XTnesIGwe2WoNQap5gEX4yp+uWdS28yGRAMzsSxqWHjlZWAY8kcJQys2jpwOvAvNT2XT68jsw38zeTWEs2VSwma0DFhJ2BXXyIODiVJ0Pid8FmNm3wEmEUYGTFgE3mNlTKY1mWdolaTdCkakjc9ivkDXAZWb2SGrDWd4G9q4E84E7CWN9Z+osA47L0fnQQek2SYcBdwGn4K+fY1gB3Aw8bmbZyut0trpXYYn5hcD5wIHAll35HhIEfAo8BywB3kwxzp+IYsUbJe1CqFzVt4BhRfxBeK3+XW8k5TiO4ziO4ziO4ziO4ziO4ziO4ziO4zht+QdI83mFN1tPFAAAAABJRU5ErkJggg==",
    "brain": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAP0klEQVR4nO2debBdVZWHfwsIL8xhBoMSxHYOQWZIFCyBBoMShnJGsRwQiXYVKmC3baWbKifQ2N1Ah2oQNKWlAkHBLkXACZAhSIJAsMFIQkISNEISMxBI3td/7PPk5eXe+87aZ+9zLnn3q8o/L/vstfY+655z9t5rkHr06NGjR48ePXqMOKxpBcoAbC3pCEmTJB0i6UBJ+0raoWiyRtJSSX+UNEfSnZLuM7ONieRvU8if2EH+EknzC/l3SJqdSv6IBTgEmAEsx89fgMuBgxuU/9/Am1LOyYgAOBK4LWLS23ErcIRD/tHA7Qnl/9wjf8QC7ApcC/QnnPwB+oGrgTHDyL+mKfkjGmAisCjDxA9lIXBUC/lHA082JX9EA7wPWF/D5A/wHPDeQfI/0ID8dzc55wM0vgoAPiZphqStahbdL+mjkrZpUr6ZXVOz3E1o1ACA0yVdp/onf4CNCnPQpPwzzOzHDclvzgCAN0q6Ry+upUcqqyUdYWaPNiG8EQMA+iTNljS+CfldyIMKRvB83YKbevRdpN7NH8wESZ9rQnDtTwDgZZIel7R9ZBd/kvQTSQ8obL+apJdJepOkd0g6IIGaZeTPkfTUIPmHSDqlgvw1kl5lZstSKNm1ANMjl05zgZOBtkYLGDAZ+H3lhdrmPFiD/EvzzHqXAOwErIqYmK8RDmTKyhkFfD3yJrTiEmCUU/43IuSsBHaMm92XAMDZEZMS/W4EPh8hbyifrSD/nyPkfShWXtcD3OycjG8lkPmdiJvQpPybqsrsSgiPxdWOiXga2CWB3DGEo1kvTclfheN1V5U6l4EHybfpM93MVlYVamYrJE2PuPQbDcnfSdIbq8rtKoDDgdmOX8FGwnIxlfyx+I54NwL7Nij/XuCwVPIbA3g18GPHwAd4KIMujzvkP5ZB/iPOOegHbgRelVqXwWR5BRDe91+U9HtJ74zoYkVilSRpbaa2ZfGOySRNkfQQ8AUyfRckNwBgf0l3Sfo3SX2R3eycTqO/45nA0mt+BztFXjda0sWS7gBekVAfSYkNAJgo6X5Jh1fsarsE6gzFc9CS41Cm6piOkjQbODqFMgMkMwBgsqRbJe2RoLscv8ClmdqWZdsEfewl6Tbg5AR9SUpkAMBJkmYp3S93H2B0or4GWJyp7bAA20vaO1F320u6ETghRWeVDQA4VNL1SmPhA4xWONlLCZnaluFUxX8PtaJP0iwSxBxUMgBgT0k/Uh6vns/R4eTtpUIxhs9k6HpHST8CKr1yow2gGNi1kvarokAHDpc0OVPfdTJZ0qGZ+n6FpKuqdFDlCXC2pLdXEV6CKZn7r4PcYzgVOCv24igDAHaV9LVYoQ7G1SAjN+NqkHEpkQdXsU+AzytuuXeTpNsc7XPsCNaNZwy3KLibedlL0oUR1/kNANhd0icjZP2XpNMkTXNc8+sIOd2GZwzTFFYMl0fI+RSwW8R1PojzcrlySB83lLhmGZBsSxi4MlbfinJ3LsYyHNcPusaAqxz6DnBRKr3bDWYrYIFTqd8w5CAD2AOY1+GaNcCxiXVvxAAK2ccWY2rHPIYs5wgHanc5dAaYT86lM/Bmp0J/A8a16WsMIYHCukHt+wl5AZLHDNCgARTyxxdjG+wXsK6Yg5Yh48Ar6Ww4rUh6VjBUIa+n7RdK9LkDITT7OGCfjLo3agCD9NinGOvRwLAbaMA055x/JZfuIvjGl2UFEHsEmhy6xAC8ALsQ3MXLcr+n/9KrAMIHmcdXbaaZ/c2jTI/NKfwSv+u4ZAIlniwDeJaBBznbX+do26MzNzjabiNH3KXnhr7G0XatpLsd7Xt05k5Jzznav7psQ48BvNzR9iEze8HRvg48x9Upj7YrY2brJT3iuKT0AZ3HADy7TE842taFJ2o3d4RxDAsdbUvfK48BeMK5n3W0rYvdM7WtC8+ZQul75TEAj5dMU4knOuHxM6wtNMuBZ/5L7wZ6btQaR9tuTIboccnqqm+Agl0dbUvHNXgM4BlH2yOBE4Ec7t2xeOLuuyZGH9gOOFE+V/vlZRt6DGCRo+04hbPtRcC5juuyAGwr3y9oN2qM0G1HMXeLFObSswpL6tU8oMxEx3bkUL6UXCGf7t5DLGg4nSvwpQrzfWQOhXYmRM3G0A9MSq5Ued1j0sV8tUF9JxGfsPoFQhxCFsXmRCoF8IMsSg2v82hCsgcvSwivjiZ0/kGFeZ7tkeVdrt3qbD+Yph6pH1bwmfOyr6QPJtalLFXmyuNz6TaAXzjbD6b2FQHBU/aLFbqYRjNH2lXm6nZPY89x8ChJF7jVeZHkSRdKcImkKk4mY1WP+/tQqszVheRYwQCXVngvAUxNrlRnfT9QUd/BvHd4iUl1n1pR37QfsMDxVCujMoeQILoWCJtQz1WcxMGsA95Wo/59VPvg3gi8NZUy2wNPVFDmHjL6+rXQ90w2dTRNxVrgtBrHsU8xd7HMJ8VOLH6nxAH6CZk6t04wH2X07COklM1R8GnwmL5MTctDYGvgXyqM6V+rKrAnwbXby3oSZrEooec7gD9ETlIM84DcgbGDxzeZuJpGKwmRXNGCL44QuhAYm3D87XTbAziXkEW8KR4AzqFijH7J8e5HXFW1abEC+/DvoD1JwgSLbfR6OfBDYEPEZORiA0Enz4FNzNjHAoudui0j5pVF+Jjy8AI5o1KCTq+hXJxdUywDPM6zMXMwEb/xt/147bQR5K1r900zy+YJDGwl6ftKl2wpB3tL+n6haxbM7C5J/+m87F1t+2v1R8Ku33KVT9j4jKQDzGyVU7HSEDKR/TRX/4k52cx+lqtzQizhEyrvebVS0h5mtmHof7Sz1MPky9Y5I+fNL0gaLZyZrLoWGcg94Wu7KNQ02ox2BuCtdF25sEIJqvoZ9mdq24o6fCK9c97SSaSdAXjCsx82s/lOZWJYUOHahZJudrT/iaQnK8hbUOHaUpjZY5LmOS5peU/bGcArHR3f6WhbhZvkT+DYL+kKhbjGpx3XLSuuuVL+pwEKutbBHY62Le9pOwPwFGtIntu/FUVpVU9OvFskHWpm58V8n5jZSjP7hMLr0ONkcVWNZWAfdrRtuT/TzgA8KceecrStyqck3djh/59XWCoeaWYnmdncqgLN7HdmdoKkYyT9UFKnmMcbCx3rwjP3Lb9L2jkOeBI1ewJGKmFm64EzJJ0h6SOSDi7+a65CytrrzeyvmWTfLeluwrbvmQrZvAbLv1rSDWaWOs9wJzxz33I3sJ0BbLZe9Haci2KCry/+1Y6ZLZc0o/jXNJ6539jqj+1eAZ7MHns62vZIi2dXtOU9bWcAf3Z0nHXvu0dHSieCUJtVUDsD8MSiNxpBM8LxRAAtaPXHdgbgyUZxDAkzevYoByFht+f0teWmUTsD8KQa61OH06Ye2Xi3fB+Bv2v1x3YG8Fu1+Wpsw6dzHoH22JRirj/tuGSD2iTtannTitOmexwCxkt6n6N9j2p8UNLrHO1/2243tNOv1pObTgpFC7L7xo10gL0keQM/2uZs7GQA31Pnbc+h7C1pJjW5gY9EirmdKV+w63qF7fGWtDUAM3ta/qfASZIuZwuo9tVtFHM6Q9KJzkuvK3YvWzLch9sl8h/BniPpWzQUW78lQgir+7akj3ov1TDBrR0NwMweUOfTt3acLWkO4Eku3aMFwASFw6aYymDXmVnH4/oyS7cL5MtTO8DrFUqf/4Iq0SkjFELgyy8Vbv5rI7pYpxKFpIY1gMLd698jFBjgrZKWEvL1v430NYG3GAjpbE4ArlIoYH1che6mmdmC4RqVTSRwiUIFzImRyoyS9PHi3/PAIwp708sVni5zJc0q9h+2eAq37tMV/AlGK5yo7i/pDUpzvP4bSV9P0M+LEGLTljojUjw8C7w/qdKb6t8VFUOA9xdjzcVSoLRLX+ntWzNbrFAGtXQaUidjFPYRshlB0xRjm6l8buNrJJ1qZkvKXuDavzezexUeXTEfhaVESLqMNlW0KtJovYBiTJfJkcjZyTpJU8zsPs9F7gMcM7tF0inyeQ15GHg/pqbpegGnK98vf6WkyWbmShEnRaZ1N7PbJb1ZPl97DwcP38RN0/UCcoxJCquFSWb2y5iLo49wzexBSY0ngnbgefTW6dlblXPMzBMfsAlVz/BzPQEq+/O3YL2jbY5vnBxjknz+m5tR1QBy+OCvkDQrQ7+eDztPdZGyzJKv7EtZStcGaEVVA1istI9LJE1NvSFESJfm8aB9HYnzGhZjmqq089WvipFZlQzAzNYoXXGCFZLOMjNPlcyyvEW+J0CfpOTp7YuxnaV0T4KFZlbpdZUip+xcla9m8ayk8xWCT3eTtJPCuzn3VrA33Y0UHF1diZfLYGbfBf5XL24F90lapTA3SyRNV/nlYq7vivIAFzm3Kt9Ts367A6sjtlRXEVyv69T1LKeOn6lTv3ZKH+JU+hFqrMcDfMWp32AurlHPUcCjTv0m1KVfJ8WNkBzSw/k16fZaqiWNXgv8Q026XuDU7U916FUKQv5c78RmtV7C2fr9Tr1acQ+Z3duAg/Ebam1Pp2EBDsRfUOoJHMeWTn22Ar7n1KcT3yGToyvhmH2BU58NwP459IkGuCFiYv9A4vSqwDbANRG6DMfVJP52AcYBj0Xo0tbNuzGACcSVlVsCHJdIh30JPoi5uBVIkqmU4B4Xk/Z2A93qbEv8L28jcAWRkUWEvPofA/4aKd/DcuAjRAbAEJw9ZxBfA+B/YuTWArBXMUGxrAKmA6XyFBLW+FOJe4xW5f+A8yjp8QwcBPxHMcZY/lxWXlmSf9gA75KUokjkHxVyED6ssN+9UqGo8+4K+/pHKaRwazoUbYOk+xSCaR9TyJu8VmGXc6xC4OwkSQcmkHWmmXmjtTqS68t2hkKEULfzuIIBeRJjNsUVZnZe6k5zGcC2ClVG35Kj/0SsUMj9t7XCk8aTG7FufiXpH83s+dQdZ0nqUCg6RdKDOfpPwFpJ7zSzRwtvmikKTpXdyAOSTstx86VMBiBJZvaspOMVBtBNrJZ0ipn9Pc+umf1KwdF1dVNKteF+SSe+pANmCGXnf1rtgzsZi4CWefMLXQ8DnmpYxwFuBnas815lg7BOv5i4jaJU3EaJTRzybyYNxwZCrcYtL+cSMIl66/tB8AX4J8+EEs4SzgfW1KzrPOCYnPegcYBtCUefOePjIDxtZlLhrAHYn3ColPvJ9QzwWUZSUg1gDKG0bEwxxE6sIxzcvD6hrm8gbHOnrkm8CLgQ6OYlaF4Ip3dvB74N/CVyIp8jHNScS0Y3LmA34JOE74mYUq4UY7wWOJkuSKjVVcmcCO/p8Qo5cMcr7NDtp7AFvJ3CWn29QjDEAoX0p7Ml3Vd4KNep6w6Fnocr5Ow7QCHOv0+h4toqhWXlYknzFSqr3KtQY6lqUaoePXr06NGjR48ePeL5f9gF4kq5Pj3RAAAAAElFTkSuQmCC",
    "building": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAJRklEQVR4nO2dbcwdRRXHf4e+ARZoi0VLKwVata1ETUHlxWIrxGjFD4gm1CICKq39IFbAGJWK8QsxLalGIAimaSIUk4JBQgHfSnlJNFUiaFFJsY2UVlpBxVr7/vfDrNi0u/fZuXd35u7e+SX3w3Pv7Jxzn/3f3dkzZ85AIpFIJBKJRCKRSCQSiUQikWg/FtuBIiSNA84GZgCTgGPieuTNLuBF4PfAL83sH5H9yaWvBCDpGGA+cAVwDnBUVIeq4wDwBLASWGVmuyP7019IGibpc5K2qf1slfQZSX0h7uhXAEmTgbuBc2P7EpjHgMvM7IWYTkQVgKSzgQeA18f0IyLbgYvMbH0sB6IJIDv5jwDHx/KhT3gVuDCWCKIIILvsrwfGx7Dfh7wEnGVmW0IbDj4QkTQMuIt08g/lDcBdMQaGMUaiVwPnRbDb75wPfDq00aC3ALnn/D8Dbwxpt0FsA04PGScIfQWYTzr5nZgAzAtpMPQV4An8L/97cVG05wFV7lQ9GDAVeC8wwvPYR81sTvUuRUbSOEkHPKNmD0iaFNv3bpF0iqSHPL/zfkljYvteOZLmdnHyh8X2u1ckDZe0xvO7fzCUfyHHADM82u4FFprZgbqcCYWZ7QcWAvs8DptekztHEFIAEz3aPm5mL9bmSWDM7C/Akx6H+PyveiKkAHzm8zfV5kU8nvdoO7o2Lw4jpAB8njgO1uZFPHxuZ8GezvpiTjoRjySAAScJYMAZHtuBRC4TJV2Y8/5+YAew0cz2VGEoCaA/+XD2KmK/pGeAn+CSTJ/p1lC6BTST4cBM4MvA05LWSZrdTUdJAOGo89HufGCtpLvl1lOUJgkgHCEmteYBT0l6R9kDkgDCMSGQncnAOrmk2yFpzSBQ0jTgC8Ac6ss33A48Ciw3sz96HjuyencKOQFYI+kcM/tTp4atEICkq4DbqP+fPBZ4K3ClpIVmtsLj2J01+VTEWGC1pHeb2X+KGjX+FiDpAuAOwv7CRgJ3Snq/xzGba/KlE2cASzo1aLwAgG8R53scldkuS6zVwYvl1mHk0mgBZF9sZkQXzpR0SkT7ZRgFfLHow0YLADfijU0dPuwB/p7z8skqOpTLJR2d90HTBdAPRRfq8GGlmY07/IVLFJkL/NqzvzFA7nil6QLYAPwtov0dwLOhjJnZXjN7CLeUfrXn4bmp5o0WQJY0ujSiC0tjJK6a2T7gSlwJmrK8Pe/NRgsgYynw4wh27weWRbALgJntBL7vccipeW9WFgjKBhlTcFG4vH5ryXQ1swOSLgGuAxYDJ9Vh5xC2AzcDy/ogbf03Hm3H5r3ZkwAkzQQuBT4AvK3X/roly72/SdLSzI+6QsE7gA2ZvX6gMMKXQ+656eqEZdkqN9Jny7yzE/N0bD+ahJcAJI0HbgE+Xo87rxG9eNWgUFoA2eX+fsLMazd2QWjTKCUASecCDwPH1evOa4SaOx94hnwMlDQdeJBwJx/CzuwNNB0FIOlY4D5cKDEkoefOB5ahrgDfAKaFcOQwNkewOZAUCkDSacDnA/pyKP0wyTMQdLoCXEu6F7ee3KeA7N5/eZd97iP/Hn4sLjmhFuSKLM6h3qTQHcBaYK2ZtWIJe9Fj4AX4j/p/BdyAq3J1ROKCpNtxRSIrR9JUYBVwVh39H8ZXgfWS5pmZT9GHvqRIALM9+/khrvR58Bi5pJNxpddDxg7eBTwu6Uwz2xbQbuUUjQFKrywBXgCuijhBsow4gaMJRJwOrooiAZzq0cedZrarAl+8kXQ8cEkM2xkfk3RCRPs9UyQAn8DPU1U40iVn4F+Js0pGELCkWx0UCcAnU6iSQgVd0g+lY/vBh65pekrYBrpPla6CfYDvGsG+otECMLNXgXsjurDazP4Z0X7PNFoAGdfi6uyHZltmu9E0XgBmthVXIcN3sUQvrAdmNT0GAC1ZHm5mGyW9hxQK9qYVAgDITsjPs1eiJI2/BSR6IwlgwEkCGHCSAAacJIABJwlgwEkCGHCSAAac1gSCfJB0Iq640/8KK78CbDazV+J5FYeBEICkEcBHgItx+Y65i08lbcGFeu8DHsxLbm0brb4FSBolaTFupdG9wGV0Xnk8Cfgk8CNgk6RrJLV6bURrBSDpPOB3uHIuJ3fRxURgOW5DhlKVt5tIKwUg6RpcVe83V9DdNOAxSYsq6KvvaJ0AJH0T98utcnwzArhF0tcr7LMvaJUAsl/+12o0cWPbrgStEUB2zw9RNHJ5lnzSClohAEmjgBWEeawdAaxoy9NBKwQALKKaAV9ZpgMLAtqrjcYLIAvyXBfB9PWSGh9Ia7wAcBG+bp7zAf6dvbrhTcCHujy2b2iDAC72bH8Qt8HUdDMbbWajceVlv4f/Mq+PerbvO9oggNkebQ8C881s0aHbvpnZs2a2AFcVxUcEuTX4m0SjBZDN6vlUFb3dzO4p+tDMfoBfCfbJknKrcDeFRgsAvzoGAN8p0ebbnn32w75FXdN0AfgUZ9hVcrfPDfiVYQ9dRLNSmi6AYR5tS9UxMDMBuz36jVmgomeaLoBEjyQBDDhJAANOEsCAUyQAn2BI2t4lHj6D4Nx6BkUC8Kn7d7pH20S1TPFom7sHQ5EAtnp0/AmPtolqmefRNneX0SIBbPDoeJakz3q0T1SApAX4bduXu8dxkQCe9PTnNklfyTJzEjUi6WhJS3Db9/mQe05zB3CSJuCKQPsMMgBeBtaRv6P3LMqXVf0D8HiJdpNw26mXYQ+wsmTbT1F+b4M1wJYS7ar4/uOB9/H/JW1lOQBMNLOXDv+gcAQvaQ0tSHhIAG6Z20V5H3SKA9xckzOJ8BSey0IBmNnPgJ/W4k4iJA+b2S+KPuwYxMm2Yvkt8LqqvUoEYSfwzk5b23QMBZvZRmBh1V4lgiDg6qH2NRpyLiBLk1pSlVeJYNxgZquGalQ6ji/peuAm0gRSv3MQ+JKZldrPyGsiR9Jc3BKsk7pwLFE/fwWuMLNHyh7g9Ws2szXADOBWYK+fb4ka2Qt8F5jhc/Khh6lcSZNw6+MuBaZ220+iJ54D7gHuMLMy0cgjqGQuX9IU3K6db8GFK4cKo44ETgNOzP5+GdhEb1cVw4WGJ2T2/4WrDdTrRtRjcOnnx+HCydtwod9eNovq9vvvxu1Z8Byw3sw29eBDIpFIJBKJRCKRSCQSiUQikRgY/gv6kQ1rvjLJUAAAAABJRU5ErkJggg==",
    "building-2": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAJRklEQVR4nO2dbcwdRRXHf4e+ARZoi0VLKwVata1ETUHlxWIrxGjFD4gm1CICKq39IFbAGJWK8QsxLalGIAimaSIUk4JBQgHfSnlJNFUiaFFJsY2UVlpBxVr7/vfDrNi0u/fZuXd35u7e+SX3w3Pv7Jxzn/3f3dkzZ85AIpFIJBKJRCKRSCQSiUQikWg/FtuBIiSNA84GZgCTgGPieuTNLuBF4PfAL83sH5H9yaWvBCDpGGA+cAVwDnBUVIeq4wDwBLASWGVmuyP7019IGibpc5K2qf1slfQZSX0h7uhXAEmTgbuBc2P7EpjHgMvM7IWYTkQVgKSzgQeA18f0IyLbgYvMbH0sB6IJIDv5jwDHx/KhT3gVuDCWCKIIILvsrwfGx7Dfh7wEnGVmW0IbDj4QkTQMuIt08g/lDcBdMQaGMUaiVwPnRbDb75wPfDq00aC3ALnn/D8Dbwxpt0FsA04PGScIfQWYTzr5nZgAzAtpMPQV4An8L/97cVG05wFV7lQ9GDAVeC8wwvPYR81sTvUuRUbSOEkHPKNmD0iaFNv3bpF0iqSHPL/zfkljYvteOZLmdnHyh8X2u1ckDZe0xvO7fzCUfyHHADM82u4FFprZgbqcCYWZ7QcWAvs8DptekztHEFIAEz3aPm5mL9bmSWDM7C/Akx6H+PyveiKkAHzm8zfV5kU8nvdoO7o2Lw4jpAB8njgO1uZFPHxuZ8GezvpiTjoRjySAAScJYMAZHtuBRC4TJV2Y8/5+YAew0cz2VGEoCaA/+XD2KmK/pGeAn+CSTJ/p1lC6BTST4cBM4MvA05LWSZrdTUdJAOGo89HufGCtpLvl1lOUJgkgHCEmteYBT0l6R9kDkgDCMSGQncnAOrmk2yFpzSBQ0jTgC8Ac6ss33A48Ciw3sz96HjuyencKOQFYI+kcM/tTp4atEICkq4DbqP+fPBZ4K3ClpIVmtsLj2J01+VTEWGC1pHeb2X+KGjX+FiDpAuAOwv7CRgJ3Snq/xzGba/KlE2cASzo1aLwAgG8R53scldkuS6zVwYvl1mHk0mgBZF9sZkQXzpR0SkT7ZRgFfLHow0YLADfijU0dPuwB/p7z8skqOpTLJR2d90HTBdAPRRfq8GGlmY07/IVLFJkL/NqzvzFA7nil6QLYAPwtov0dwLOhjJnZXjN7CLeUfrXn4bmp5o0WQJY0ujSiC0tjJK6a2T7gSlwJmrK8Pe/NRgsgYynw4wh27weWRbALgJntBL7vccipeW9WFgjKBhlTcFG4vH5ryXQ1swOSLgGuAxYDJ9Vh5xC2AzcDy/ogbf03Hm3H5r3ZkwAkzQQuBT4AvK3X/roly72/SdLSzI+6QsE7gA2ZvX6gMMKXQ+656eqEZdkqN9Jny7yzE/N0bD+ahJcAJI0HbgE+Xo87rxG9eNWgUFoA2eX+fsLMazd2QWjTKCUASecCDwPH1evOa4SaOx94hnwMlDQdeJBwJx/CzuwNNB0FIOlY4D5cKDEkoefOB5ahrgDfAKaFcOQwNkewOZAUCkDSacDnA/pyKP0wyTMQdLoCXEu6F7ee3KeA7N5/eZd97iP/Hn4sLjmhFuSKLM6h3qTQHcBaYK2ZtWIJe9Fj4AX4j/p/BdyAq3J1ROKCpNtxRSIrR9JUYBVwVh39H8ZXgfWS5pmZT9GHvqRIALM9+/khrvR58Bi5pJNxpddDxg7eBTwu6Uwz2xbQbuUUjQFKrywBXgCuijhBsow4gaMJRJwOrooiAZzq0cedZrarAl+8kXQ8cEkM2xkfk3RCRPs9UyQAn8DPU1U40iVn4F+Js0pGELCkWx0UCcAnU6iSQgVd0g+lY/vBh65pekrYBrpPla6CfYDvGsG+otECMLNXgXsjurDazP4Z0X7PNFoAGdfi6uyHZltmu9E0XgBmthVXIcN3sUQvrAdmNT0GAC1ZHm5mGyW9hxQK9qYVAgDITsjPs1eiJI2/BSR6IwlgwEkCGHCSAAacJIABJwlgwEkCGHCSAAac1gSCfJB0Iq640/8KK78CbDazV+J5FYeBEICkEcBHgItx+Y65i08lbcGFeu8DHsxLbm0brb4FSBolaTFupdG9wGV0Xnk8Cfgk8CNgk6RrJLV6bURrBSDpPOB3uHIuJ3fRxURgOW5DhlKVt5tIKwUg6RpcVe83V9DdNOAxSYsq6KvvaJ0AJH0T98utcnwzArhF0tcr7LMvaJUAsl/+12o0cWPbrgStEUB2zw9RNHJ5lnzSClohAEmjgBWEeawdAaxoy9NBKwQALKKaAV9ZpgMLAtqrjcYLIAvyXBfB9PWSGh9Ia7wAcBG+bp7zAf6dvbrhTcCHujy2b2iDAC72bH8Qt8HUdDMbbWajceVlv4f/Mq+PerbvO9oggNkebQ8C881s0aHbvpnZs2a2AFcVxUcEuTX4m0SjBZDN6vlUFb3dzO4p+tDMfoBfCfbJknKrcDeFRgsAvzoGAN8p0ebbnn32w75FXdN0AfgUZ9hVcrfPDfiVYQ9dRLNSmi6AYR5tS9UxMDMBuz36jVmgomeaLoBEjyQBDDhJAANOEsCAUyQAn2BI2t4lHj6D4Nx6BkUC8Kn7d7pH20S1TPFom7sHQ5EAtnp0/AmPtolqmefRNneX0SIBbPDoeJakz3q0T1SApAX4bduXu8dxkQCe9PTnNklfyTJzEjUi6WhJS3Db9/mQe05zB3CSJuCKQPsMMgBeBtaRv6P3LMqXVf0D8HiJdpNw26mXYQ+wsmTbT1F+b4M1wJYS7ar4/uOB9/H/JW1lOQBMNLOXDv+gcAQvaQ0tSHhIAG6Z20V5H3SKA9xckzOJ8BSey0IBmNnPgJ/W4k4iJA+b2S+KPuwYxMm2Yvkt8LqqvUoEYSfwzk5b23QMBZvZRmBh1V4lgiDg6qH2NRpyLiBLk1pSlVeJYNxgZquGalQ6ji/peuAm0gRSv3MQ+JKZldrPyGsiR9Jc3BKsk7pwLFE/fwWuMLNHyh7g9Ws2szXADOBWYK+fb4ka2Qt8F5jhc/Khh6lcSZNw6+MuBaZ220+iJ54D7gHuMLMy0cgjqGQuX9IU3K6db8GFK4cKo44ETgNOzP5+GdhEb1cVw4WGJ2T2/4WrDdTrRtRjcOnnx+HCydtwod9eNovq9vvvxu1Z8Byw3sw29eBDIpFIJBKJRCKRSCQSiUQikRgY/gv6kQ1rvjLJUAAAAABJRU5ErkJggg==",
    "calendar": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAGP0lEQVR4nO3dX4hUZRjH8e+z5Z/KSCuSstSi0u3K6qIkWAhbKALpoqJiKRIiS4LIC7WLCCLqon90EXVRFkVF/6igUCQkSBLSoCBKWqXSNEzFUkOz7dfFmQWLded9zpwzZ7b3+cDcPe97njPvzzkz7xx3IIQQQgghhF4i6RRJKyRtlLRH0hZJz0ua33Rv/yVpfqu3La1eN7Z6P6Xp3iYkSfMkbdPYjkha0nSPoyQtafU0lm2S5jXd44Qi6VRJW4/zhI4akbSoB3pd1OplPFslndp0rxOGpJVtntBRm3ug182Jva5sutcJQ8X1M9WcBvuc4+hzY1N9jqev6QaO43xH7dy6mqj42J5z6ppeDcAkR+3k2rqo9tiec+qaXg1A6JIIQOYiAJk7se4DSJoGnAvMJP066LleXiZJ7saqcZmjdpKkaxJr/wS2AzvM7Ki/rXRW9YSSpgPXATcAg8CMqo+RkRHgZ2AN8CqwwcwqDXtlAZA0A1gF3AdMrWre8C/bgIeB16oKQiUBkHQP8Cjxr71b1gP3mtl3nU7UUQAknQg8AyzrtJHg9gdws5l91MkkpQOg4mvOj4GBThoIHfkLuMvMXi47QakASDLgLeDGsgcOlRFwk5m9W2Zw2X2Ah4jF7xUGrJbUX3awi6RLgK+BE8ocMNTmW2CBmf3pGVTmFeAxYvF7UT+w1DvI9Qog6Urgc+9BQtfsAS40s99SB3i3goec9cc6RLHFmWI66eE8CNS6XTqOScC0xFoB+xNrpwAnl+jnTOAO4NkSY8cnyST95LgDRpL2Srpf0lnOY+1zHGOw8pNN73PQ0ec+59z9kl53zD/qM89xPK8AC4DzHPW7gAEzG/Y0FApm9i1wm6QvgKccQxdKmmVmP6cUe94Eeu/FvzMWv3Nm9jTwgmNIH3C5pzjVOY7aLWa21lEfxvcgxXudVMk3ytYVgC8dtaENM9sHrHMMmZ1a6AlA6rtdgAOO2rF40t7psTrRzT53O2pPSy3s1VvCvkqs+4tiB6wp35D+ETT1nLqqVwPwUmLdO55Nj6qZ2e9A6pcwq+vspaxeDcD7wNttanYBy7vQSzvLKXoZz9sU59RzejIArdudhoDngL/HKNlEscews6uNjaHVwwBFT//1N8U5DFV9L19Var8ruKzWt1rLJD0DLAYuBPZS3A613szGCkYjzGxY0hXA1a3HGcAw8KGZfd9oc230bABGtZ7AJ5vuo51WID9pPSaMnrwEhO6JAGQuApC5CEDmIgCZiwBkLgKQuQhA5uraCJql9P8LH9LMqmPSugJwfesRelxcAjIXAcicJwC1XINCLZLXyhOAs0s0EpqRvFaeAEwp0UhoRvJfMPUEoMm7b4NP8t3KngD8WKKR0IzktfIEoLG7b4Nb8lrFx8DM1bUTeITiz5iF6pxMDW/E6wrAK2Z2d01zZ0nS80Dlz2lcAjIXAchcBCBzEYDMRQAyFwHIXAQgcxGAzEUAMhcByFwEIHMRgMxFADIXAchcBCBzEYDMRQAyFwHIXAQgcxGAzEUAMhcByFwEIHMRgMxFADIXAchcBCBzEYDMRQAyFwHIXAQgcxGAzEUAMhcByFwEIHMRgMxFADIXAchcBCBzEYDMRQAyFwHIXAQgcxGAzEUAMhcByJwnAKpp3pCmlrXyTPqLo7bfURvSzHTU/ppa6AnADkftQkmXOurDOCSdBAw4huxMLfQE4AdHbR/wYqvx0LkVwHRHfXIAkkk6SdIB+Xwi6fTKm8mIpGWSRhzP+Yik5F8ONWczbwC3OM9hN/AcsInix6RCe5OBC4Ah4Arn2I1mtjC12BuAxcAHzoZCd60ys8dTi70BMGADkJyw0FUHgYvMLPkTm+vzupkJWOntKnTNE57FB+crwChJrwC3lxkbarMLuNjMkn84GsoHYCrwGXB5mfGhckeBQTP71Duw1JatmR0Gbsa3Oxjqs7TM4kMHe/Zmtg24ChguO0fo2AjwgJm9VHaCUpeAY0k6C3iPIgyhe/YDt5rZmk4m6fhbOzPbTbFPfW+rqVAvAW8Cl3a6+FDBK8CxJM0E7qfYwTq3yrkDh4G1wCNmtrmqSSsNwChJfcDVwLXAXGAOMJtiizO0d5TiDfb21mMdsNbMDjXaVQghhBBC+H/4B0QjxbxowPZ9AAAAAElFTkSuQmCC",
    "camera": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAKIElEQVR4nO2da7BWVRnHfw8XReRgKkqmxJgjRAEa3i0CZ+IihpTZRaPJiWGmyWbKsfxQTTMQjtOFbJxSZtSJERVLy7LLCDMmdKNAIAEtKBIUDDXgICIQ5/Dvw3pPkb23td+993rXOes3wxfOWvt59t7/d++1n/WsZ0EikUgkEolEIpFIJBKJRCKRSCQSiUQikUgkeiEW2oG8kDQIeBtwOjCgIDNdwMvA383sUEE2SiVqAUgaD1wHTAPGA/1LMt0NbAAeB5aa2caS7CYAJE2WtFLtwwpJk0Jfl16PpFMkPRD4Ztfjfkknh75OPkTzCpB0HvBTYGRoXxqwDbg6ltdCFAKQdAmwDDgptC9N0glMM7PVoR1pRNsLQNIoYBVwSmhfPNkNXGpmfwvtSD3aWgByn3argXGhfcnIBuCSdv5k7BfagQZ8hXhvPrhP0y+HdqIebfsEkDQC2AIMCu1LixwCRpvZ86EdqUY7PwHmE//NB3cO80M7UYu2fAJIGgv8iWyRvSPAa/l69B+GAAMz9DsKnB/Lp2FwJP08QxBmjaTpko4r0K/jJM2Q9FQG/35RlF+9CkmTMlzcRyRl+WVm9fE4SY9m8DOFi+sh6XhJ6z0v6ouShgTwtUPSTk9f16rAJ1TUSBoq6THPCypJnw7o89wM/j4qqSOUz28kt0GgXNDmHOA0/ObjBwMXA3OAN3ua3QyMNbMuz365IGkAsBF4u2fXfwD3AmuA1z369eQjbDWzw542q9KSACRNAD4GTAXeSXGJGLW41sx+VLLN/0HSB4Efl2y2CxdlXEaIfARJ75P02wyPvzxZJaktPmMl/S7wtVgpaXIZJzpM0sOBT1aSjkp6T+En3CSSJoa+IBUelFTMpJmkCZJeCHyCPSwu5CRbQNKS0Belwja53ImmaOoRKulyXP5bO4xetwIXmNm+0I4ci1wm0Frg7NC+APuAK81sVaOGDQUgaQzwe+BNOTjWKnuAiWb2bGhHqiFpHPBr2uNa7QUuM7PN9RrVnQySNBg3wm2HE3oFmNquNx+gMhqfiksGCc3JwCOSTqjXqNFs4Dz8v3GLYD0uu2ZtaEcaYWZrgEtxn2mhGQt8NVNPSWdLOhxsKON4XdI8Scfne12KR9IgSQskHQx6BaVDkvwTaSV9N6DTuyUtlHRWAfemVCSNkHS7pD0Br+cdtfyrOgiUe/fvItuo33c+/ijwT2A7sA5YATxpZv/KYLttkXuKXVH59y7grcAw/JJysuYjdAJnNJ2bKGlmBpWtkjRVJU7L9jXkpqGvlMt98GWGj6GFngd/SG5iJFECkgbKPyL7TR8DT3gc+Hm5V0aiRCSdKGmHx31aVu04td4/PqPGu83MZ0ozkQNmdgA3pdwsVSOUtR7bPkuw1nu0DY7cDOJZuDoCHYBwg9aXgJ1mpoDu+eITF6kazKslAJ9s3FwSE4pCUn/cyHsKMBkXHKn1yjog6RngSWA5sMLMjpbhZ0YOerRtfowmv2/WKf5+F4+k4ZJu9XxPvpEdcsGc00OfTzUkTfE4lz0+B45WAJKGSPqGXBQxLw5I+roCJJ/WQzkIoJ1XBnkjaRrwZ+CLQN1JEE8GA7cAz6qMzJsS6RUCkNRP0q3AL3EDvKIYASxXwEzkvIleAHKRxyXAlyjnfAYCd0n6nnpB8CtqAciN8B8Erg9g/jPA9wPYzZWoBQDcAVwb0P5sSbcEtN8y0QpA0kdxv8LQ3CZpZmgnshKlAOSKR9wd2o8K/YAlkoaHdiQLUQoA+A75ZCi/BuzP4TgnkTX1KjDRCUDSZcA1Gbu/AnwLeC/QYWYdZjYUGApMAhZW2mRhrqRzM/ZtL9TGkUBJP/PwrYcjcrmFDaet5aZZvyapK4Odh8u4Bsf42rdCwZJGSur28E2Vc5mYwdZ0SZ2etrrlxieloD4YCp6Nn88HcBU7f+NryMweB2YCPrmJ/YCP+9oKSWwC8P3cuqmSp5+JinC+4Nnt6qz2QhCNACQNBS7w6LIauCcH03cCPquRLlIbVQBpRDQCwKVS+8TeF+aR3WNm3cACjy4DgPNbtVsWMQlglEfbg8BjOdr+Ca7iZ7P4+BqUmATgM7pel2eBZjM7iF/uY2lfAq0SkwB83qvbCrD/nEfboQXYL4SYBOCzQLSIRFWfJ0o0NY5jEsABj7ZF7CziUyOhqFrFuROTADo92o4pwP47PNruLcB+IcQkgK0ebcdIOjMvw5Xw7miPLj6+BiUmAfzFo60Bn8zR9g34FdX08TUoMQlgI676VbN8Xjns4SfpVOBzHl32As+0arcsohFAJSK3wqPLacAitVBNtNJ3EXCqR7cVFV+jIBoBVPCdb/8IrW3XMh//pNMftmCvPVD75gMMlrTPw7ce7pFHnX65fQvuzWCnUw3KsuWJ+lo+QKUOwaIMXecAmyR9WFLNc5ZbYTQLt1/RpzLYuasSNo6GGFe23A7cCJzo2e9c3ON5u9z+PetwhbAAzgAmAFfhijdlYT/w7Yx9gxGdAMxsl6QFwG0ZDzGSYtYTzDOzrAmlwYjqFXAMC/GrjlE0q3GrlKIjSgGY2RFc6tXO0L7gXiMfqvgUHVEKAMDMXgRm4VcmJW8OAR8wsx0BfWiJaAUAUCkePQe3h07ZdAE3mNkfA9jOjagFAGBmS4HpuL0EymI3Lt38ByXaLIToBQBgZk/gtp4rIwa/BXi3mf2qBFuF0ysEAGBmW4HLcd/iRWQEdeNSxC9stAtHTPQaAQCY2atmdjNu7v5+XCXyVhEuw/giM7vRzPJYTdw29CoB9GBm283sE8B5uKXkL2U4zC5c1HG8mc0ys6gqojZLdJFAH8xsE3CTpJtxYrgCt8PpaGA4/8003o+74VuATbhKoRvavEpoLvRqAfRQuZHriayucRn0yldAonmSAPo4SQB9nCSAPk4tAfgkNUa3p18vwif9rOp8SS0B+KRf+xRtSOTLhR5tq65WqiUAn5WwcyX5pmclWkSuCskcjy7bqv1nLQE87XHgM4HFSvsFlkYlw3kx8BaPbs3fU0nvz5ASvUZuU8M0JigIuXT1GZKeynB/plc7Zq2tY0/AhUazFDroBl7N0C/RmA6yRW/34raO/b9Z0qoHM7ODku4DPpvBWH/c3vWJ9uG+ajcf6qx4ldtyfDPpMy92DgGjzOyFan+sGQgys+246dBE3CysdfOhwZr3ylhgDW4KNREfG4GL61VMqxsKrqxzu4ZyEy4T+bAHt16hbnGrhnMBZrYFt2bOJzqYCEsncJWZ/bVRw6Ymg8zsD7gNFba36FiieLYBkyr3rCFNzwaa2dO4FbRLs/mVKIEHgAlmtqHZDl7TwWa2x8yux+XWrfR0LlEcK3G/+tlm5lWiLnP9HABJ44DrgKm4pMs+kWPYBnThilgsBx4ys41ZD9SSAI6lMkFxDjCMFDwqisPAy8BzZuazk0kikUgkEolEIpFIJBKJRCKRSCQSiUQikUgk+hb/Bt4BF8hlk9tUAAAAAElFTkSuQmCC",
    "chart-line": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAF9UlEQVR4nO3dTagVZRzH8e9fLQ0KLpRR9EbYwsxAatGmNyVqE2RCRFiCG5HMWoW1aVu0Cy2whTtXib2BpGFCtAmDFHJTFtmbdC28tOjV7q/FzInjueeeM8/MM2dmzvP/gBvvM/M89/x/M8/Mc885A84555xzzjnnnHPOuelnZTaSNANsAtbl/3UCOGhmc7EG5lpK0mZJ57XQeUmbmx6fq1Fe/Pkhxe+Z9xBMKUkzixz5w84EM02P1xWzJKDtJqBIYXvXB64DQgKwbnyTUm1dg0ICsLymtq5BIQFwU8gDkDgPQOI8AInzACTOA5A4D0DiPACJ8wAkzgOQOA9A4jwAifMAJM4DkDgPQOI8AInzACTOA5A4D0DiPACJ8wAkzgOQOA9A4jwAifMAJG5Z0wNwGUlLgPX5vyuB08B7ZvZVowPrkbS3wCeDe/Y2Pd4ukXSLpONDXsd/Jb0u6dK6+vYzQMMkrQWOAlcP+fES4GlgpaTHzUyx+/drgAblxf+I4cXv9xiwsY4xeAAaImk18CGwsuAmW+sYx1RNAZKWAWuAq4DTZvZdw0Maqu/IL1p8qOk7F6biDCBpmaRdwI/ASbI59YykzyRtaHZ0FytZ/Np0PgCSVgBvA6+wcC69EzgiafvEBzZExeKfjDwcoOMByIv/DvDwiGZLgTck7ZjMqIaLcOTvizic/3U2AH3Ff6hAcwN2NxWCCMU/QPa7RtfJAEhaDrxFseL39EKws55RDVfian/QEWBLHWsA0MEA5Ef+u4w+7S/GgNcmdSbIj/yPgWtK7uIw8IiZ/RFvVBfrVABKHvmDJnImiHTkP2pmf8Yb1UKdCUBe/AOUO/IH9c4EtYQgL/4xyh/5R4CNdR75PZ0IQOTi99QSgi4VHzoQgIpz/jhRrwm6MOcPanUAIs3540S5JujKnD+otQGo6bS/mErTQddO+/1aGYAJF7+nVAi6XHxoYQAizPnngFeBMgsnQdcEXZzzS5vEW8IkLZf0fkA/g2Yl3Z7va5tGP91klHmNORNIWi3pbIWxHpZ0WZnXqRGqOQCKWPy+fdYSAqVWfKg3AKqh+H37jhoCpVh8qC8AkpZKOlThBZ1VNheP6mOnqoVgR76ftXl/ZX2g7Bqne1RfAHZVeEHHFr+vn6oheFmpFh/qCYCyo/9cyRe0cPH7+qsSgipaW/ymbwNvI3sDZ6hzwAYz+yJkIzPbDTxHuVvEsg6T3edPdIWvqKYDUOb5gqWK3zPhELS6+NB8AM4Etq9U/J4JhaD1xYeGA2BmZ4DPCzaPUvy+vusMQSeKD82fAQCeB+bHtIla/J6aQtCZ4kMLAmBmR4FtwN+LNPkWuC928fv6jxmCThU/iOpfCr417+NLSb9I+lTSi5Iur+P3GdJ/1VvE1t7qRVF3ANpA5ZeNu7m8SwumgDYxszeB7YRNB43+Pb8qD8CAwBB0uvjgARgqD8GzjL47OURX3swxggdgEWa2B3iQhesUs8ALZMXv/NX+VH1BRGz5Leodkm4CVgG/AqfM7EKzI4vHA1BAvmIZumzdCT4FJM4DkDgPQOI8AInzACTOA5A4D0DiPACJ8wAkzgOQOA9A4jwAifMAJM4DkDgPQOI8AInzACTOA5A4D0DiQgIQ8mEJD1ZHhBTq94C2q0IH4poREoCfAtreLenG0MG4yQsJwKmAtpcAe5U9yNFNA0kzki4Efmr2kKQbmh67W5yFNJZ0DLg/sI9/gE+Arxn/TSAu3F/ACeCgmc2FbhwagK3U9ABDV9kc8IyZ7Q/ZKDQAK4BvgGtDtnMTI+CpkBAE3a/nn4Z9KXRUbmIM2COp8Pcvllmw2Uf2kATXTjPApqKNgwNgZvPAk2Sfk3fttK5ow1JLtmb2PdkjXX4rs71rj9Jr9mZ2HHgA+DnecFwkJ4o2DLoLGEbS9cB+4N6q+3JRzAE3F10TqPxXOzP7AVhP9m2fZ6vuz1UisrWAwgtClc8AF/WerRM8AWwB7gGWxty/G6n+haAQ+b3oXcAa4Drgirr6SlylpWDnnHPOOeecc84551wK/gM9Ck6tFmSuVgAAAABJRU5ErkJggg==",
    "check-circle": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAM4klEQVR4nO2dacxVxRnH/4OA1IoCrriCWtdWrVutFKiiGBOpdvGDGmtqXBtTQ611aWlsWj+1au1iUhGt1thWG8RgUqyKohatWptaqgi4gLjwIoIs+sr264c5t728nHvf+8w598w573t/CYFc5s78Z57nzDn3OTPPSB06dOjQoUOHDh36HS62gHYADJF0sKSDJB0gaV9Je0jaVdIISTtIGtjjaxskrZa0QtJySe9KWizpdUmvSnrZOdddhP4i6RMOABwoaaykEyQdJ2/8ngbOykZJ8yX9Q9JcSXOcc6/m3EbhVNIBkiv8ZEmTJJ0qf4XHYKmkWZJmSnrEOfdxJB3BVMYBgAGSJkg6T9IZ8tN4mVgr6UFJ98g7w6bIelqi9A4A7CbpIkkXKt6VbmWppKmSpjrn3o0tppIAhwB3AN1Ul0+Au4DDYo9nZcAb/k/Apqimy5dNwH3AobHHt7QAuwNTgY1RTdVeNgK3AyNjj3eN6M8AwCBJV0j6kaShkeUUxYeSznPOzYwtJKoDAMdImibp8Jg6IrFe0ljn3HMxRURxgOSqnyLpWuUfsKkSc51zY/KqDBgo6TBJO0la5Jxb0tt3CncAYD9Jf5CP2LWb9ZIWyYdyF0taIh/mXSWpW9LHkraVNEDSdpKGSdpZ0t7yPzk/Ix9OHtRGjXs755ZmqSAx/PckTZYPd9d4UdJVzrnZWerPDeB0YGUbH7LeBe4BLgOOSmaarJoHA0cClwC/B97JWfOJGfUNAR5qUv8m4FtZxyETgAN+CGzOefAA/gNcD3weKGRGwzvEFOClHPR/IYOOIcDDLbTRDRyc5xhYRA7GXzl5shL4FXBklE5t2b/DgZuBFQH9WAN8KrDdVo1f47d5970VkUOBRwMGphELgUuBTxfemV4AtsPfJhYY+nNjYFtW4wPMz7vPvYkcDvzdKLIRC4FzgW0K7UQAwDbAOcCrvfTpeQIcmTDjAxT3TgJv/BcCRPZkBXA5OTzMFQ0wKNG+rEef1gO3AeagF+HGB3i2Hf1MEzkUeDZQZI3NwDRg50JEtxH8M9CJwAXAGcAugfVkMT7ANXn3LU3kYOCv2WzPUmBi28VWiByM/xoBM45VpAPuzmj86cDwtgqtGDkYfxlFvI4GrssgcgNwJQX9jq8KwLY0D/L0RhfwuSKEnk74+/uVwCltF1kxqmT80cAHgSKXFiKyYlTJ+AOBuYEiFwGj2i6yYlTG+InY6zMYf+9CRFaIqhn/KHxQw8rbdK78raia8QcCLwaIXFWYyApRKeMngicHiNxIJ8CzFVU0/q74K9nKlYWJrAiVM34i+tYAoTPoBHm2oKrGPxAftbOwFNipUKElp5LGT4TfGyD2tMKFlpgqG/8g7Lt27ipcaImprPET8VONYj8g8N13X6Tqxt8F+NgoeHIUsSWk0sZPOnCNUfDrwOBogktEXzD+gMSgFs6PJrhEVN74SSdOMop+Db9VqV/TJ4yfdGSaUfjlsTXHpi8ZfzC2xR4fAtvH1h2TxPgzK298SQImGsX/OrbmmPQp40sScIuxA0fH1hyL2MZP2j8W/8yWTzY1et/eVE/ls2aGEtP4+GXjN7D1G9rngXFZOrWPsRM3BDdWYUpg/DlN6t4InB3asXONHfliUEMVhvjT/s9aaOMjYHRI5Zb3/u9TgZ27eUJ8428PrGuxrZsa1TOgSRuWHD5PVCU3bh4A20r6s6TTA6tYLmmCc+7fGWQcK5/XqBVOavQfqQ6Aj+R91iDmaUPZzABnJlff28mfmcCZBbU9RD4pdKjxuySdlNH4kk9o1SqHYnk3g3/3b6GQ+z9+2r2viY77EgO1q/0hwCzj2NSzDLBcWM20HGNsu/V28VdYq2wCWp2KMgH8pgU9s2iDE1Ai49fpsSzP+3paPY2eAfY3aHnDOfeRvQs28KeCXNpC0VMlzcjTCZK6ZiR1h9Alf8+fl5em5Pia1w1f2S/tw0YOsI+h4oWGslmYpOYPrfXk5gRlNH4dlrFPjQ42GtA9DBW/YSibhVHG8pmdoOTGl2xjn2rTRg6wa4PP03jbUDYLqwK+E+wEFTC+ZBv7VJs2cgDLT4zlhrJZmBP4PbMTVMT4km3sU1PvNHIAy4FMHxrKZuExSS8EfrdlJ6iQ8SXbrJi6TqORA1jSl641lA3GOYekcxQ+4/TqBBUzviRZfn2l2rSRA1jW9K03lM2Ec26hpHGS3gus4lRJD5KSnxcf3r1f4cZfLunkAo0v2cY+Ndlmqz+rmlHopk/n3HxJJyrcCSbKzwT/cwLKEdsPwTL2pH3YyAE2GioufP1/nk5QYeNLtoMsUl/WNZrqLUegRsnc7Zybj08vN1tSyBa0iZIeSP5dlXt+Tyxjn/q80MgBVhsq3tFQNlecc/Pwy54el7R7QBWhhpfi3PN7Ysmquibtw0a3AMvPC0vQKHdyuB2EEHPar8cy832Q9mEjB+gyVLynoWxbKNgJymJ8yTb2qTZt5ADvGCoeZSjbNhInOEXtjUwuVz6LOfJilKFsqk0bOcBiQ8UHGcq2leR+nCVO0IzalR/znt8Ty9inniHYyAFeM1Q8ihJtB2vT7aBM074kCX/cjGW1b6pNGzmA5YAhJ+kIQ/m2k7MTlM74CUfIFsh7Je3DRhUskj9Zs1WCz75rFzk9E5Ttnl+PZcy75W26FakOkCzxfsnQwJcMZQsj4zNBGe/59VjGfJ5zLjW622wKsbx6HU9JN4YE3g7KOu1L8kfTSRpv+ErDE8qbOcBcQwMjJB1vKF8oRicotfETjpM/IbxVnmn0H80cwLoCZ5KxfKHUPRM0i3G8o/Le8+v5irH8U0GtAPMN684XUYF8wMCe+Iyn9WcdrE8+ix7VbIVkrFtlQZaGbjY0BBXaIQzsCIxJ/kR7oWUFON5ok1uyNDbB2NhtOfa1Qwr4Y2ctnJylsUH4rd+tsoYKXU1VAz9rrTHY4316OXe5aSTJObdB0nSDxu0lXWgo38HGBWqwurcB0xMbhgOMM3gcwBI6aWJzBz8bLzHaIjxHUF3DDlhgbPjiHPrcoQ7gIqMNFpDXrzL8ub4WltDGffr9DfxWcOvVf1WeAkbQej6aGtfmJqCfA1xtHPt1wIi8RbSSnKGe1YBll3GHFIA9krG0cGs7hIzGfmDU/bkL6WfQPCVOGhuA1GQQeYi50ygG4GttEdMPwJaqp8bv2iloNPCJUVAXELJmv18D7IbPK2ThE8CS3idI2C+MogAeAfLYh9gvwJ/S8peAcQ6P+xvEjcAWHq7x47aL6yMAUwLGdwVFHdAJXBogcBOd54FeASYlY2WllQxquYkcAMwNELkOsKSg7VcAR2N72VPjGYq+xQKHAt0BYruAQwoVWwHw5zJbH/rAP/jlloDSKvqqAMHgD5Q+IIroEgLsjz3UW+PqmMIHAI9lcIJ+PxPg8zK/FTiGjxN7NTYwEngvsAPLgdJtKikK4Dj8LTGELsqyhhH4MvYwcY11NEhk3JcBvgqsDRyzjcCE2H3YAuA7gZ0B2Az8hNjTWQHgb5vXJ30O5bux+5EKtmNm0ngUGBm7H+0C2B14OOMYTY3dj4YA2wAzMnZwOXBW7L7kDfANwu/3NR6i7Ocy41evzM7YUYDpwF6x+5MVYK+kL1mZQ0qCy1KCP9HqqRw6vRb4QWU6XgewHXAdYZG9njwNDI3dJxN4J8hjJgAfM/g2PqljqcGfa3QZ/jCrPHiCEmVgMYG/HTyQ00CAH9TvA5b8eIUADMNHRpfm2N/UvMaVAv9gaF1P2Bvr8KuTxhFxQyp+yfxY4A7si2Z741bK/sBnAbic8GBRM94EbsQHo9o+YMBAYHzS5ptt6M8G4Ip296NGoVcPMF7SHxWW1rUVVkt6Un4//HOS/umcy3SgBbCDpCPlc/KMlc/MYTlQw0KXpLOdc7PbVP9WFD594gM9d0sK37VqY4l8gqQ3Jb0l6X35tKnr9f+UuMPls54Pk099u5f8KVsHqMFpW21gtqRvOueKOoNJUgQHkHw4VNJkST+V1N93EHVLmiLpJufc5qIbj5rRA/86eKqkMTF1RORvki52zr0cS0DUlbrOuVfk07hdJD819xdWSLpY0riYxi8VwHDg54QtM6sK3fhfD6WLY5QGYF/gdrZM4lR11gPTgKIeKKsP3hFuIZ9YeizWAr8ERsUez0ZUIa3bMEnnS7pEUlXWEM6XdJukO51zIUfeFkbpHaAefBq6cyWdpchH1aTQJX/u4L3OOUuW1ahUygFq4JeOjZHPTnqapMMiSZknaZakhyQ9nSTZrhSVdICe4Hcfj5V0gqRj5UO3eR9n95Gkf8mHmOdKetI5V+RBVW2hTzhAT/CRxlGSDpa0v6R9JI2UtJt82Lf2M6z29ypJm5O/V0paJp9YerF8GHmBpDdiROo6dOjQoUOHDh06dMid/wJVkAhXRldmdAAAAABJRU5ErkJggg==",
    "circle-check-big": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAM4klEQVR4nO2dacxVxRnH/4OA1IoCrriCWtdWrVutFKiiGBOpdvGDGmtqXBtTQ611aWlsWj+1au1iUhGt1thWG8RgUqyKohatWptaqgi4gLjwIoIs+sr264c5t728nHvf+8w598w573t/CYFc5s78Z57nzDn3OTPPSB06dOjQoUOHDh36HS62gHYADJF0sKSDJB0gaV9Je0jaVdIISTtIGtjjaxskrZa0QtJySe9KWizpdUmvSnrZOdddhP4i6RMOABwoaaykEyQdJ2/8ngbOykZJ8yX9Q9JcSXOcc6/m3EbhVNIBkiv8ZEmTJJ0qf4XHYKmkWZJmSnrEOfdxJB3BVMYBgAGSJkg6T9IZ8tN4mVgr6UFJ98g7w6bIelqi9A4A7CbpIkkXKt6VbmWppKmSpjrn3o0tppIAhwB3AN1Ul0+Au4DDYo9nZcAb/k/Apqimy5dNwH3AobHHt7QAuwNTgY1RTdVeNgK3AyNjj3eN6M8AwCBJV0j6kaShkeUUxYeSznPOzYwtJKoDAMdImibp8Jg6IrFe0ljn3HMxRURxgOSqnyLpWuUfsKkSc51zY/KqDBgo6TBJO0la5Jxb0tt3CncAYD9Jf5CP2LWb9ZIWyYdyF0taIh/mXSWpW9LHkraVNEDSdpKGSdpZ0t7yPzk/Ix9OHtRGjXs755ZmqSAx/PckTZYPd9d4UdJVzrnZWerPDeB0YGUbH7LeBe4BLgOOSmaarJoHA0cClwC/B97JWfOJGfUNAR5qUv8m4FtZxyETgAN+CGzOefAA/gNcD3weKGRGwzvEFOClHPR/IYOOIcDDLbTRDRyc5xhYRA7GXzl5shL4FXBklE5t2b/DgZuBFQH9WAN8KrDdVo1f47d5970VkUOBRwMGphELgUuBTxfemV4AtsPfJhYY+nNjYFtW4wPMz7vPvYkcDvzdKLIRC4FzgW0K7UQAwDbAOcCrvfTpeQIcmTDjAxT3TgJv/BcCRPZkBXA5OTzMFQ0wKNG+rEef1gO3AeagF+HGB3i2Hf1MEzkUeDZQZI3NwDRg50JEtxH8M9CJwAXAGcAugfVkMT7ANXn3LU3kYOCv2WzPUmBi28VWiByM/xoBM45VpAPuzmj86cDwtgqtGDkYfxlFvI4GrssgcgNwJQX9jq8KwLY0D/L0RhfwuSKEnk74+/uVwCltF1kxqmT80cAHgSKXFiKyYlTJ+AOBuYEiFwGj2i6yYlTG+InY6zMYf+9CRFaIqhn/KHxQw8rbdK78raia8QcCLwaIXFWYyApRKeMngicHiNxIJ8CzFVU0/q74K9nKlYWJrAiVM34i+tYAoTPoBHm2oKrGPxAftbOwFNipUKElp5LGT4TfGyD2tMKFlpgqG/8g7Lt27ipcaImprPET8VONYj8g8N13X6Tqxt8F+NgoeHIUsSWk0sZPOnCNUfDrwOBogktEXzD+gMSgFs6PJrhEVN74SSdOMop+Db9VqV/TJ4yfdGSaUfjlsTXHpi8ZfzC2xR4fAtvH1h2TxPgzK298SQImGsX/OrbmmPQp40sScIuxA0fH1hyL2MZP2j8W/8yWTzY1et/eVE/ls2aGEtP4+GXjN7D1G9rngXFZOrWPsRM3BDdWYUpg/DlN6t4InB3asXONHfliUEMVhvjT/s9aaOMjYHRI5Zb3/u9TgZ27eUJ8428PrGuxrZsa1TOgSRuWHD5PVCU3bh4A20r6s6TTA6tYLmmCc+7fGWQcK5/XqBVOavQfqQ6Aj+R91iDmaUPZzABnJlff28mfmcCZBbU9RD4pdKjxuySdlNH4kk9o1SqHYnk3g3/3b6GQ+z9+2r2viY77EgO1q/0hwCzj2NSzDLBcWM20HGNsu/V28VdYq2wCWp2KMgH8pgU9s2iDE1Ai49fpsSzP+3paPY2eAfY3aHnDOfeRvQs28KeCXNpC0VMlzcjTCZK6ZiR1h9Alf8+fl5em5Pia1w1f2S/tw0YOsI+h4oWGslmYpOYPrfXk5gRlNH4dlrFPjQ42GtA9DBW/YSibhVHG8pmdoOTGl2xjn2rTRg6wa4PP03jbUDYLqwK+E+wEFTC+ZBv7VJs2cgDLT4zlhrJZmBP4PbMTVMT4km3sU1PvNHIAy4FMHxrKZuExSS8EfrdlJ6iQ8SXbrJi6TqORA1jSl641lA3GOYekcxQ+4/TqBBUzviRZfn2l2rSRA1jW9K03lM2Ec26hpHGS3gus4lRJD5KSnxcf3r1f4cZfLunkAo0v2cY+Ndlmqz+rmlHopk/n3HxJJyrcCSbKzwT/cwLKEdsPwTL2pH3YyAE2GioufP1/nk5QYeNLtoMsUl/WNZrqLUegRsnc7Zybj08vN1tSyBa0iZIeSP5dlXt+Tyxjn/q80MgBVhsq3tFQNlecc/Pwy54el7R7QBWhhpfi3PN7Ysmquibtw0a3AMvPC0vQKHdyuB2EEHPar8cy832Q9mEjB+gyVLynoWxbKNgJymJ8yTb2qTZt5ADvGCoeZSjbNhInOEXtjUwuVz6LOfJilKFsqk0bOcBiQ8UHGcq2leR+nCVO0IzalR/znt8Ty9inniHYyAFeM1Q8ihJtB2vT7aBM074kCX/cjGW1b6pNGzmA5YAhJ+kIQ/m2k7MTlM74CUfIFsh7Je3DRhUskj9Zs1WCz75rFzk9E5Ttnl+PZcy75W26FakOkCzxfsnQwJcMZQsj4zNBGe/59VjGfJ5zLjW622wKsbx6HU9JN4YE3g7KOu1L8kfTSRpv+ErDE8qbOcBcQwMjJB1vKF8oRicotfETjpM/IbxVnmn0H80cwLoCZ5KxfKHUPRM0i3G8o/Le8+v5irH8U0GtAPMN684XUYF8wMCe+Iyn9WcdrE8+ix7VbIVkrFtlQZaGbjY0BBXaIQzsCIxJ/kR7oWUFON5ok1uyNDbB2NhtOfa1Qwr4Y2ctnJylsUH4rd+tsoYKXU1VAz9rrTHY4316OXe5aSTJObdB0nSDxu0lXWgo38HGBWqwurcB0xMbhgOMM3gcwBI6aWJzBz8bLzHaIjxHUF3DDlhgbPjiHPrcoQ7gIqMNFpDXrzL8ub4WltDGffr9DfxWcOvVf1WeAkbQej6aGtfmJqCfA1xtHPt1wIi8RbSSnKGe1YBll3GHFIA9krG0cGs7hIzGfmDU/bkL6WfQPCVOGhuA1GQQeYi50ygG4GttEdMPwJaqp8bv2iloNPCJUVAXELJmv18D7IbPK2ThE8CS3idI2C+MogAeAfLYh9gvwJ/S8peAcQ6P+xvEjcAWHq7x47aL6yMAUwLGdwVFHdAJXBogcBOd54FeASYlY2WllQxquYkcAMwNELkOsKSg7VcAR2N72VPjGYq+xQKHAt0BYruAQwoVWwHw5zJbH/rAP/jlloDSKvqqAMHgD5Q+IIroEgLsjz3UW+PqmMIHAI9lcIJ+PxPg8zK/FTiGjxN7NTYwEngvsAPLgdJtKikK4Dj8LTGELsqyhhH4MvYwcY11NEhk3JcBvgqsDRyzjcCE2H3YAuA7gZ0B2Az8hNjTWQHgb5vXJ30O5bux+5EKtmNm0ngUGBm7H+0C2B14OOMYTY3dj4YA2wAzMnZwOXBW7L7kDfANwu/3NR6i7Ocy41evzM7YUYDpwF6x+5MVYK+kL1mZQ0qCy1KCP9HqqRw6vRb4QWU6XgewHXAdYZG9njwNDI3dJxN4J8hjJgAfM/g2PqljqcGfa3QZ/jCrPHiCEmVgMYG/HTyQ00CAH9TvA5b8eIUADMNHRpfm2N/UvMaVAv9gaF1P2Bvr8KuTxhFxQyp+yfxY4A7si2Z741bK/sBnAbic8GBRM94EbsQHo9o+YMBAYHzS5ptt6M8G4Ip296NGoVcPMF7SHxWW1rUVVkt6Un4//HOS/umcy3SgBbCDpCPlc/KMlc/MYTlQw0KXpLOdc7PbVP9WFD594gM9d0sK37VqY4l8gqQ3Jb0l6X35tKnr9f+UuMPls54Pk099u5f8KVsHqMFpW21gtqRvOueKOoNJUgQHkHw4VNJkST+V1N93EHVLmiLpJufc5qIbj5rRA/86eKqkMTF1RORvki52zr0cS0DUlbrOuVfk07hdJD819xdWSLpY0riYxi8VwHDg54QtM6sK3fhfD6WLY5QGYF/gdrZM4lR11gPTgKIeKKsP3hFuIZ9YeizWAr8ERsUez0ZUIa3bMEnnS7pEUlXWEM6XdJukO51zIUfeFkbpHaAefBq6cyWdpchH1aTQJX/u4L3OOUuW1ahUygFq4JeOjZHPTnqapMMiSZknaZakhyQ9nSTZrhSVdICe4Hcfj5V0gqRj5UO3eR9n95Gkf8mHmOdKetI5V+RBVW2hTzhAT/CRxlGSDpa0v6R9JI2UtJt82Lf2M6z29ypJm5O/V0paJp9YerF8GHmBpDdiROo6dOjQoUOHDh06dMid/wJVkAhXRldmdAAAAABJRU5ErkJggg==",
    "clock": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAMv0lEQVR4nO2daaxV1RXH/wsFbR8y2AEqVsVUHCg1pB+kFbDg0Da20jZxVtTaSW1qnY1pHGq1qGhF7WATWwsUU2PjgOIHLJRKjbYxJmqraNICUiZFpgcICL9+2PcJPu60zj3n7HveO7/k5RLuPnevtfe6++699t5rSSUlJSUlJSUlJb0Oiy1AVgCfknSEpBGSDpd0oKQOSf0rr4Mrr5K0SdLaymtn5XWZpEWS3pD0upmtyFP+vOgRBgDsL+k4SRMkjVHo8AEpV7NewRielzRP0gIzW5tyHblTSAMA9pZ0QuVvoqSjJfXJWYydkl6SNF/SXEnzzOz9nGVomUIZADBS0rmSzpc0JK40e/CupEckzTCzhbGFaZa2NwDg45IuVOj4kZHFaZZXJc2Q9ICZrYktTCEBDgDuAjopLhuBqYQJaUkzAAcD04DNMXsuZbYC04HDYrdv2wIMAu4DtkftqmzZBtwDDIzd3l20xRwAOFXSvWq/iV1WrJR0jcKEkZiCRDUAYISkXyos53ojf5N0sZn9K5YAea+dJUmAAdcqzJZ7a+dL0nhJLwFXAVG+jLlXCnxC0nRJX8mpylWSXlPw4i2VtFkfdv1KH3YNd0g6SMGFfKSkT+Yk51OSzjezd3KqT1LOBgCMlzRL0rCMquiU9KyCd+5ZBR/+ulY+EBikYAjjFFzN47RrDyFtlkk6s0iOpKapDHNZzPAXAz8DvkhwEWetR19gLHArsCQDfbYDV2StR24Qfu/vSLmR1hPW1ScAUeYxFd36EIzh/opMaXJPTN1SAegHzEqxUVYBN9JG6+gugP7ApcCyFPX9I9A3tm6JADqAOSk1xGJC434ktl6NIBj9ZODNlHR/Btgvtl4uCN+G51JQvhO4FugXWycvBEO4DtiUQjs8B2Q18UwXwiTp6RSUng0cElufVgGGEeYrrfIMsE9sfepCmBS1+pu/CvhabF3SBpgErG6xbWYSyWHUFITt21b4K3BAbD2yAhgCzG2xje6LrUdVCOv8pOwAbgL2yknWAYSJ2t2Vv8lA2mcIa9W9N8FvsaOF9ro8D1mbBhhPcifPJuDkHGU9DVhTRY41wGk5yjGJ5GcetgPH5iVrXYD9Se4RexcYm6OspwM768izEzg9R3mOAd5J2HZvEY7LxYPg5Xs8oQLLgc/lKOsAqn/zu7OGnH4OKnIdVenMJDxFzEkhcE1Cwd8i5yUecJ5Dvsk5yzac5B7EK/OUdXehRxKOOHlZQzjenbe8nhXKnRHkG0X4SfSyFTgyab2JNhsIw869krx+6i2SJkU6AePxpPXPTIoamNkrkr6qXWcUmqWfpN+Q8Kcg6W7TOQp74x52Sjq1R+51p4SZvSDpLIW28jBe0plJ6nQbAGGCdFuCum42s6cSPNerMLMnJE1J8OhdhMMrLpKMALdK8l50mC/p5gR19VZuUDgw6mGIpJu8FbkMABgu6fvOOlZLOsfMdjif67VULpmernB83MNFOFdX3hHgakneY1cXmtly5zO9HjNbKekHzsf6SnItC5s2AGCowq1cD4+a2ZPOZ0oqmNnjkmY7H7sQx4aaZwS4WtK+jvKbJbXXpkUx+aF8S8N9JV3WbOGmDAD4mKTvOoSQwqx/sfOZkm6Y2VL5VwUXEe5fNKTZEeA78jlHFku6y1G+pD53SHrLUb5D0gXNFGzWAM52VC5Jt5vZNuczJTUws62SpjofO6+ZQg0NAPi8pFGOildKetBRvqQ5fivJE6nsKGB0o0LNjADnOiqVpKlmtsX5TEkDzOw9SXc7H2vYd3UNgHApweNj3iDpfkf5Eh+/lrTRUf4sGlyXazQCHC/f7dhHzKzTUb7EgZltVIhE1ixDFMLo1aSRAZzoqEwK175LsmWms/zx9d5sZACeLd8lCleyS7JlvsIyu1mSjQCE8KtHOyqaaWbefewSJ5WYQg85HhkNDK71Zr0RYEKD97szx1G2pDWedpTdSyGOclUaGUCzdEr6p6N8SWu8IN/+QM2+rGcAxzgqWGhm2x3lS1qg4mX1HK0bU+uNegYwwlHBPEfZknSY7yh7eK03qhpAZT/ZczmiPOiZP54V18DKeY49qDUC1LSYGrzuLF/SOq85y1cd0WsZwBGOD17VEzJnFI1Km7/teKRqn6YxApTf/ngscpSt2qe1DMBz7NsjREm6eNq+ap/WMgDP6Z9ljrIl6eI5JVQ10lgtA/CEJfNsTxaFw+q5T9sIT9uXBuBggqSlwJTKnki70hYG0FP3//srJHVY0saGkJkBeOYAPdUAumhnQ8jMAHoirRpqOxtCYmoZgKexcg+mkJCXU/qcdjKEludqtQzAM7QUxQAeVcjumRZdhvAfQozDGKuGtjCAQkSyNrMNki6WlHaWroGSrpf03wiGUBqABzP7k6QzlO5I0EUMQ2gLAzjQUTY6ZvawpM8oRCxZn0EV3Q0hyzDvn3aUdRmA5wqSd+s4Oma21syulzRcwRA2ZFBNlyG8AmTVRp7PrRqko5YBeDYZPFvHbcVuhnCIsjOE4ZLmkE3WD48BVO3TWgbg2eIdkiQ6VTuRgyEcKumSND+wMsdoKgZABZcBeLd4E0eqbCcyNoRvpPhZkr/N36j2n1UNwMxWyKf8OKcwbU1GhnBwCp+xO+MdZddVgk7tQT1XsOdnwBs1tBCkbAhprzg8bV5zRK9nAC84KhhPATN7NUtKhrAgLXkqbe1JGPF8kkq+6Yxa3R4ZLHIAGAz8lOazhb4HpLZaImRo8TApqZKevDa3pqVgUaA5Q9gBNBWwyVHvFEe/vE/SVRrwoqOiJRQ9121CdjOEVd3a5EWg7vXsBHX1AZY6+uUfrVTmTfr8pfRULR6EjGBHAxOBgzKq43hnnySJPP5BZSc5K/tdirqWVAF40NknJ7RS2d7ASkdl64GinA8oHMB+wAZHfyynQS7Gur/ZlbDlsxwyDpA/wnVJ81wi3xbwrJbD9AOjHRYHYcRo+zTvRQPYt/KN9uAJ8VO38pedFae68VEiAT929sGraVbuzQu8hHZPdV4gCN9+b3LJq9IUYH98kw+A61IToJcD3OBs+07STisLTHUKsZmQY6ikBYCDCQm2PdyehSBD8We7fix1QXoZwGxnm2/BkTLGK8yvnMJAko2IEkkS8K0E7X1vlgIdRMhV62E1MCwzoXoowIHA28623gakffBkD8HucQoFsIAGYctLdkHwwC5M0M7Zp+kBBgD/SyBcmTm0SfBt93axAhiYl4BnJxBwB3BKLgIWGMJBnJ0J2veMvAX9SwIhNwM96gBpmgBj8C/5AObGEPYo/BNCgHeBz+YucJsDjALWJmjPrWR386ih0FcmEBhgGaWT6AOA4SSbVwE0nSU0C8ENeCyh4MtJa7eqwAAj8fv5u3gSsNgKDAYWJ1RgLb14TkD4zX8nYdstJaT0jQ9wLLA9oSKb6YWrA8Jsf0vCNtsGfCG2Dh8CuDyhMhCWiLfQC5xFBCfPFJIt9bq4NLYeVQHubEEpCB7DHus2Jrh3n22xjW6LrUdNCJPCmS0quJoeuIFE2Njx+va7M53Yk75GAH2BOS0qCmEbtPBLReBQwmy9VeZSlPuXQH/guRSU3gz8hAIeLyMc47oe/xmKavydbGMNpQ/QQTojAYSTxtcAH42tVyOAfYDvkXxt3525ZBNeJnsIM94HUmoICPODG8lr18sBYdS7lOQevWrMIGRwLy6EieHtKTYKhGH1YeDrNLj5krFufYCxwP34D802Yho96bItcBnJnUX1WAL8nNARmX9bgH7AOOA20hvmd2cb8KOs9egi1yUFIYjEQ/IFOPSwSSGH4XyFvHqvtZrRjBCN60iFOEgTJY2VlNVcZImkM8zMH9EjIbmvKQn+6z9IOjmnKt9WiHe0SCHHziaFqJnrtCv/boekQQr37joUDPRwhRiInlBsrfCEpAvMLIswtjWJ4lQgODOukHSLpGKsbbNjm6RrJd1dSQ2fK1G9SsBhku6TdFJMOSKyQNLFZvbvWAJEnWWa2Ztm9mVJp8iXAq3orJB0nqQJMTtfapOUMWY2W9IoSdMk9eQ09Nsk/ULSEWY2PcaQ3/YQLp9MIx0XarvwHsFXkNXqp+dBuIs4FdgYtetaYyPBCTYkdnvWor23FvXBsvHbkiZLKspp4lckTZf0ezNbE1uYerS9AewOMFLSuQoTqKGRxenOGkl/ljTDzBbGFqZZCmUAXRCOjk2UdKJC0OTRyn9Cu0PSS5LmSXpG0ryWAzJFoJAG0B2Cu/Y4BaMYI2mEQsqWNFmnEHP/eYVOX2Bm61KuI3d6hAFUAxiq4ModoeDWHaaQ669DweU7SLtyHnYqdPDGyr83SVqm0OGLJL1uZqvylL+kpKSkpKSkpKQkM/4PWQDvZsAfjOoAAAAASUVORK5CYII=",
    "cloud": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAJGklEQVR4nO2dW6ydRRXH/4u2ULC1hZYWSVBa0ZiDVFCCoi2l0tLIpajhYo2RqIlWrUpCRB/E+KBWpUHhwWhiGiNQqKgxFqs9pRx6UeEB7EUaRaQ9JdwptNLS0nN6fj7MVgvZ3z57vm++/e3L+j1+ey5rZtaey5qZNZLjOI7jOI7jOI7jOI7jOI7jOI7jOI7jOI7jOI7jOI7T8VjVAjQCOFnSOZLOlHSGpDdLmi5piqSJksbWgo5I2idpj6RnJQ1K+pekRyRtMbPnWit559BWCgBMk7RQ0nxJsyXNTJT0TkmbJa2TtNYVoo0ATgQ+D2wAjlA+w7W8vgCcVHX5exbgXOA24GALGj2LQ8DtwHlV10fPAFwIDFTY6FkMAPOqrp+uBZgF9Ffbxk3RD8yqur66BmAi8CPC2NspDAO3Am+suv46GmA+MFhtWxZiN7Cg6nosk2PKSBQYB9wkqV9h7d6pnCZpLbAcGFe1MGWQ3A4ATJd0t6Q5qdOumD9JutLMnqlakJQkVQDgLEn3qLP/9Y3YLekyM9tetSCpSKYAwGxJqyVNTpXmUeyV9LCkfyiYeJ+pfXu19vs4BfPwKQrWw3dIOlvSiSXIsk/SIjPbWELanQnwQeBAwsnXYeAPwFKgD4hWVMBqcb8E/BF4NaF8B4CLyqjLjgOYDexPVLHbCY0+pQQ5TwK+WMsjBQeAbpvnxAG8E3gxQWVuBhaS45+eU+6LgU0J5H6JMO/pPYBpwK6CFfhPYFGFZVhUk6EIg4SVT+8AjAU2Fqi0IeA7wPFtUJbxNVmGCpRnE11qJ6gL8L0ClTUIvL/qMrwe4HyK9WjLqy5DSyDM+PPu298PTK26DFkAU4H7cpZtBJhfdRlKhbCxk9e2vxI4tuoyjAZwLOGMQB4G6eYNJOCWnBXzM6CUfYcyAI6pyZyHW6uWvxQI+/l5tnRXdlLj/5eaEtyRo7zDwLuqlj85wNoclbGBDuj2syAMBwM5yt1ftexJAebmqITdhGPdHQ1hYphn3tM9x8uA9ZGFHyJsDnUFhCVirJ1goGq5k0A4vRvLsqrlTg3w3Rz10PanjUe1vQO3SfpERJqPSTrLzA7llqoNAcZL2q5wQ6lZVkvKWhWMSHpJ4SbT02ZGMQlLgLCDFntu//Kq5S4L4PIcvUAzHAT+CqwAPgO8peqySpIIN3Zi2Fy1zGVDml3EZtgCfB04tcrCbogUemFlwrYIYEHypm7MYeBOWm1bAE4mzvDzN1q0n181wLYSGno0RoBVJB4eGlnoLpY0JiKtn7blRKYcflJBnibpakk7gBuAmLaJB/h5hHYepoRjXO0KMKVW5irZDBQ+fd2oB4g573afme0pKkynUCvr+orF+ICkh4ALiiRSVwEIjhpinDOsKSJEh3Jn1QJImiqpH7gmbwJZPUDsjLPqf0MVrJT0UNVCSDpO0h3AJ/NEzlKAvog09krakSfzTsbMhiVdoWAdrJoxklYAH46NODbj+1sj0tjaQ7P/12BmTxLs/Z+VdI2kt6v5ldPxksYnFGeMQk9wgZkV65mA30bMRn+cRv7eA5hE2GxbAtwNvJxgdbCLor6PgL9EZPi1RPXR8wAnAJ8mGNWK8Kuigjwakdm1icrv1CAcR/sU8GwBJbi6iABPRWRU2c2ebodgjl+dUwGeBCaMlkfWKuC4CDkPRoR1IjCz5yUtkvT9HNFPlXR9royB5yM0ret3ANsB4Bs5eoG9wKRG6Wb1AEMRsqVcyjgZmNm3JcWuuCYpLFEzyVKA/RGZlOERxKnPdZIeiIyzhAbb9FkKELOx86Y4eZy8mNmQpGv1f9c4zTBTwfF2XbIUIMYTVnucX+sRzOxRST+MjHZV1g9ZCvBEROIx+wZOGm5S3DB9WdYPWQrwWETisxqNMU56zOxFSbdHRJkBnF7vhywFeCQi8ckKL3o4reUXkeHfW+9jlgJsjUzcXaa1ngclxbx8UveMR10FMLMXFBwyNsslEWGdBJjZiIL72mZ5W72Pjc4EbopIfB5dcBO4A4npqU+r97GRAtwbkfg4SYsjwjtp2BURtu4ftJECrJV0JCKDz/lqoOW8EBG27s5gpgLU5gExw0CfpA9FhHeKE7NnU/eo2mi+e+6KyECSbowM7xTjDRFh65qPR1OAVYrb738f8JGI8E4xTokI++96HxsqgJntVXj9I4bltIEL2B6h7tIug6frfWzGfVus37uZkr4VGcfJx3siwg7mzgVYF3kSZZhe96VfMsAEwsunzXJDkczmRCoAwBOEO4ZOCQCLI9ujmB9jYE0OJdgExBwwdZoEuDeiHYaAiUUzPJN8PvV/SdnODHoM4N2RbZDGdxNwcw4FgOBswpUgEcS77f1qqownADtzKsEqfDgoDPFj/xESeBI5WoB55H8wYiM+McwNMIP4R7rSO64GluVUAAirA18iRgJMJt+Td+kddxIejRoooATDwA+AE5IL14XUGv/POep5G2Xt0BIuLj5eQAmoxf9oKQJ2CcBM8j92eUXZwvUBewoqAcADwKX4eYLXAHyc8DhlHlrjt4ngSz+FZwuAHcCX6fHjZYR1fn+BenwFiNkoKizwXNIpAQSD0zrgOsLztB335lAshCX2YuIsfFksbTbflM/Hny/p9yrvyfYtCs/H71TY2tynuNsx7cZESdMVHEudK+k8xfllyOI3kq5s1nFX0jEX6JN0j6QZKdN1mmarpDlm9nKzEZJ2rWa2Q+EGyoaU6TpNsVPSJTGNLyVWAOl/bk3mS1qm8CyKUz67JF1kZk/FRix12QVcKGmFfEgoky2SLs3T+FIJPcDRmNn9kmZJulnScJl59Si/VhjzczW+VLICSJKZ7Tez6xUuJ/aiV/EyeEXSUklXmVmhlVDLLW+EByW/KWlBq/PuEtZLWmJmMT4cMqnM9AqcI+krCs+g+DHy0dkm6UYz+13KRCu3vQOTFXzYfEzSXMW9U9TtjCj842+RtKYMr+yVK8DREN4dWqjgcGKO4i4+dAvDCq7gVku6y8x2l5lZWynA66kpxNkKF0/PkHS6pGmSpijcdu1kJ5WHFB7beE7S45L+LulhSQ/GGnMcx3Ecx3Ecx3Ecx3Ecx3Ecx3Ecx3Ecx3Ecx3Ecx3GcHuc/r6GJ6QNcniMAAAAASUVORK5CYII=",
    "compass": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAOKklEQVR4nO2de6xdRRWHfwMFtK3lpVBpeYdSqAhVjLwKFnlFECJSQAkPJSEKBuQlhBCFSBGkIAV88AcYC5Y3QQoFLRQrFXlESxR5mpTSQinQB/TeC7SUzz/WOfRyOefcs2bP3vuce8+XNLe5d++91ppZe/bMmpk1UocOHTp06NChQ4dBRyhbgbwAPi9prKQxknaUNFrSMEnDKz83rvyUpG5Jyys/uyo/F0l6QdKLkp4PISwuUv+iGBAOAGwiaT9JEyXtIavwEYnFvC1zhsclzZY0J4SwPLGMwmlLBwCGSDqg8m9/SbtKWqdgNT6UNE/SI5JmSZodQvigYB0y01YOAIyTdLykkyRtXq42n2CZpDsl3RRCmFu2Ms3S8g4AfFbSybKKH1eyOs3yjKSbJN0QQlhatjJtCbAFcBXQRfuyEpiCdUg7NAOwNTAV6Cmz5hLzPjAN2KHs8m1ZgI2A64DVpVZVvqwCrgE2LLu8q7REHwCYJOlatV7HLi9el3SerMNImYqU6gDAGEm/lg3nBiN/k3RqCOG/ZSlQ9NhZkgQE4HxZb3mwVr4k7StpHnAuUMrLWLhQ4HOSpkk6pCCRSyQ9J4vivSKpRx8P/UofDw0Pk7SVLIS8k6TNCtLzfkknhRDeKkiepIIdANhX0nRJo3IS0SXpUVl07lFZDH9FlgcCG8kcYYIs1DxBa+cQUrNI0nfaKZDUNJVmLo8e/svAJcBeWIg4bzvWA/YBLgUW5GDPauDsvO0oDOx7f0XiQnobG1cfAJTSj6nYtg7mDNdXdErJNWXalgRgfWB6wkJZAlxEC42jqwDDgTOARQnt/SOwXtm2RQEMA2YmKoiXscL9dNl29Qfm9CcALyWy/SHgM2Xb5QJ7Gx5LYHwXcD6wftk2ecEc4QKgO0E5PAbk1fFMC9ZJeiCB0TOAbcq2JyvAKKy/kpWHgA3KtqchWKco6zd/CXBY2bakBjgCeCNj2dxMSQGjpsCmb7PwV2CLsu3IC2BzYFbGMrqubDtqgo3zY1kDXAysW7YdeQMMweIWazKU11ll2/ExgH2JD/J0A4eWbUPRYJ+E2DUPq4G9y7ZBkq3IJT4itgzYp2wbygL4KvBWZNktxJbLlWpAAP4UacBrwBdLNaAFAHauVGYM91NmpxA4L1LxhQyAIV4qgG2JjyCeU5bS47AlTl6WYsu7O/QC2AX7JHp5H9ipaGUDMDtC2R4G8Te/P7A+Qcwq6DkU+SkAjo9Qcg2DsLfvBTicuCHid4tScATWgfNyUSEKDgCAyRHl+zq2eCV35a6LUG42gyDIUwUYCXwZ2D7y/iFYs+5lampb+iq2Lf6AzxIGcHi3CrADNvs3r4/984FTcH6jK0602FnWq8hzdAX81qkQDMCJnSoNKr0WN+J3giMiyjufuQLMI991KnN3LsqUCL5K78spEfLudcp4lzxaXPwzfd0MkGAP2Sq9N/PxtwJb4R8aXpG6ADbFdrp6OD+pEgVDukrvy3YRulzolNGF7b9IVhjekO982nMZV16V3pvdI/TaAHjFKecnKQvm307hP0wmPGcoptJ7MzJSz9OdctLsN8TGsx4W0+Krdym+0qvMy6Dzp/AH4Mb399xmNh4c79R1SgjhXec9uQNsgy0tnytL/zZZ0m4Fq3F77I0hhPckXe28rd+6a9gjxTYlLFLzGyTfkTQqhNDV5PW5gkXijpE0ScVXdi3GhBBeir0Z2x/wqqRm9wkskTQ6OnsZcIizybkhSlBisCVqfwE+dOqfJ9HNfx/bbnTKPajR8/r7BBzo1G+a8/qkAJsBt0uaI9O9lZZQ35HoOTc7r/96tCTgXw5Pe5lyN2ruSdwsZVEkSRCFrcWY75D7VKPn1a0wLP3qrg7dbg4hfOi4PhnAIZIeltSq6diezvLt700lp9AtjlvGAxvX+2OjN3ZiP3/vy0zHtckA9pR0t6Qihp7dkm6TdKPzvlTNf5UHHNeuK8ujXJP+HKBZuiQ1bGryANhM0l3Kt/KrlX6UpM1CCMfKnyEktQM8obXpbZrBU5cG8JTjO+PxyGQAtzl09NBVefZRwNA+Mofim5xJ0vuvYfuDDh2eqPecRilVxjj0me24NglYvqGjEz6yW5ao6Q5JM0MIPXWu+4Z8LUDqt7/KI5IObvLaHV1PxvL0etjTrX5GsHF+Vuq+6Q3k3uqUkUt6WCwnkofm5yCAic6H1+1l5gGwPfFBHnel95LbEs1/RZeNnXbvW+s59T4BYx26LCnh5Ixj5A/yPCHpSkn3N2je+6NVmn+FEJYDb0pqdt5/rCwz6ceo5wCeb8bzjmtTMcl5/e8k/SiEsKZgubk5QIUX1LwD1KzTesNAT0DlBce1mcG+qZ6JnSeVoPIrnwvPxpZkwZ8GeMq+Zp3Wc4DhjgcvclybAu9beEWCN19qoea/Fwsd19acQaznAJ60ZCsd16bAM/TrkS9qlkqulGHu34Gn7NvfASrNv2d+YkYIwRMxqyd3qKwFaJZ5IYT/ZZXbBC3hAEUu/vC+hama4UPVes2/lKMDePoARTqA5/vfI+nBEuRKA8ABWo5O858P9RzA81Z7WossdJr/T5K5r1bPATxNS1EO0Gn+P0lLOEDumaw7zX9dBocDyP8W3plI7mHyNf+p5DZLSzjAaMe1sXgcoFsJlqdhx8+c57ytyOZfkrZ0XOtygMWOB/sWGziJiP1nme2ryhwim0D6kuO2ImL/ffGU/Wu1flnPATyTDJ6p4xgK64RV5vuPlk0dn1yU3Ax4HKD5OgW+5lxskFt2KnwbOLuIW+RxFLZIJPak8g+BbfMqgzp657ogxDvFu5Okfzjv6Rf8zf99zTT/rO3dT5J/nF+LB0II8zM+w4s3O+iLtX5Z0wFCCIuBdySNaPLhE5SDAyhh7z+HSv/o0ZIuSfQsDzXf6DqsCCG87no68ISjecllWTgZm3/SNO/9MT0P25somz87dHw8RsA1DgHdJE4JgyVx8HBr5b6hwLex1bt5VXqVhcCmKe1usmzWd9pWN69Ao8mgRxw6DZX0Fcf1zeBt/rsxJ3hD9ik4Rvmd8StZvOHIEMLSHGXUYw/5bPPUpYH1Mj1Jiy91C2ksv+j0LR66AO/W+ZRlc5lD1w+IHaUB/3QIWkCi7eH4m/8iWQikbu08ZbMOvoxhTzZ6Xn8V5tnytZV8PdNGeJv/opguabcQQuEbYXsxUb4QcMM67M8BZjkESdIJzuvr0UoOgGxh6V4hhONK+ub3xpu066FGf+wvSdQQ2bLvzZsUljlJFBb8qRm0KJinZeHdW0oI8tQEf5KoxZK2bLQsvmELUMku5RnnjpD0A8f1tUi549fLs5IuljQ2hDA+hHBpq1R+hdPkmwKennlPBDDe2Ul6nQyJIikheSOWNDKXXbypIC5RpGcRTUPh3lSxp0XKGemUE0tbVHpvgB87bXwmpXDvucALiDjqHNjdKcdD21V6Fezt9x4ueW5KBTYB3nEqcEGEnO2cMvqjbSu9N8DPnHZ3kfpYWWCKU4kenHPk+HPg1WJAVHoVYGtsrsXDL/NQZCT+067viZBzilMGDLBK7w0ww1kW+RwZU1HmNxGVc4RTRqC5fLgDttKrAEdGlPe1eSq0FXZWrYc3gFFOOQFrCeYPtkqvAowG3nSW9Spg67wV86wTqDIHiyrGyNseO7Qi6qSNdgQ7OHJuRDlfVYRyI4BXI5T7ee7KDRDwTfdWWQxsWJSCx0UouAY4vBAF2xjgW8SlwDu2aEUfjlCyB5hQqKJtBLAH/iEfgHfWNomyO+PvEAIsA75QuMItDrALsDyiPN8Hct2d1UjpcyIUBlhEwRspWhnsUO6YfhXAmWUqHoB7IhV/jVSzVW0MMA5/nL/KfTiPos3DgI2x42JiWM4g7hNg3/y3IsvuFUpYkl4TYG9gdaQhPQzC0QHW2/eexl5lFSVkaG8IcFakMWBDxMlEBovaCSzIcxnZjrQ7o2w7agJcmcEosIihK2zcTmDh3UczltHlZdtRF6xTeHNGA9/AOYHUDmATO97Yfl+mUXanrz+A9YCZGQ0FmwZt+6EitsjlvgTlMYvE+y9zAxgOPJbA6B7gQiKWl5UNtozrp/jXUNTi70Ce+xzTAwwjTUsAttL4PJzZP8oA2ACbyo4d2/dlFrYfoP3Aerw3JCoIsP7BRRQ16+UAa/XOID6iV4ubsBPc2xesY/jLhIUC1qzeDnwTWLdE29YB9gGux79otj+mUuJZzMkBziQ+WNSIBcAvsIrI/W3BkjNMAC4nXTPfm1XA6XnbUaXQIQWwt+zgY8/uVg/dkubKEiI8Kum5rCeaYUfi7STLg7S/pH1kCTHyYIGkY0MI/pQukRQ+psTi13+Q7wCmLLwpO9nsBdkZO92yrJkrtPb83WGSNpLtuxsmc9AdZTkQmz2VKyv3SvpeCGFZQfIkleAAkvULJJ0tabKk9hjb5scqSedLurpyNHyhlBpVwlb3XifpoDL1KJE5kk4NITxblgKl9jJDCC+FEA6WdLh8R6C1O4slnShpYpmVL7XIkTEhhBmSdpE0VdLqktXJk1WSfiXLPzCtjCa/5cE2n0wlTQi1VXgPixXkNfoZeGB7EacAK0utumysxIJgzabYKZzWnlrUR8PG78sSULXLauL/SJom6fctkFSqIS3vAL0BxsmyZJ0oqdW2ii2VdJekm0IIc8tWplnaygGqYEvH9pd0oCxv3ngV36FdI2meLA/fQ5JmJzqkulDa0gH6goVr95M5xR6SxkhKPWO4Qpa+7nFZpc8JIaxILKNwBoQD1ALbTTxW5gw7SholO+NwmCzku5HWnnnYJavglZX/d8vyI74oCyE/H0JYUqT+HTp06NChQ4cOHTrkxv8BKIkOb4G91j0AAAAASUVORK5CYII=",
    "cpu": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAKqklEQVR4nO2da6xdRRXHf+u2t01tq4gJYJCKRGhLUGnhi+GVUEDlFYiUEKhRHkYqBktMAB+JIX4xDUoTBTFpFA0QobykaJBCohQNTYD6LC3Y24K0FkN4ltjb9p6/H2Zfae7d+9zZZ88+s88580vuh3v37FlrZq279541s2YgkUgkEolEIpGIhaQFkm6XtEXSa5KelnSDpNk1yLpA0lpJO7KftZIuqEHO7KwNT2dt2pK1cUFoWT2NpCskjSqfEUnzA8mZKeneAjnKrs0MJGt+pnseo5KuCCGn55G0RNJYG6NI0lZJcwPIunUKOZJ0awA5czOd2zEmaUlVWT2PpGc9jCJJN1aUc4ymdrRxwxxTUdaNnm16toqcEAzFFC7po8Biz+JV39Hn4dfeoaxsFXx1XZz1QTSiOgBwZImyH6soq8z93ZR1ZEVZlYjtADNKlB2uKKvM/d2UVaYPghPbARKRSQ4w4CQHGHCmxxAqaQZwKHBsiduGJZ1RQezhZcpWlFXmG2C+pL+a2asV5HWMdUOIpGnAKcD5wLnA0d2Q22PsATYDa4A7zezlbgit1QEkDQGXAt8j8nCnx2gBTwDfMLO/1SmoNgeQdAKwGji+LhkDwD5gFXCTmb1bh4BaHEDSUuAO4H111D+A/B34nJm9Erri4KOALGZ/D8n4ITkOWF91jiKPoE8ASZcAd4euN/F/dgAnmtmuUBUGM5SkE4EngVmh6kzk8iRwhpntC1FZkFdA9rW/mmT8bnAqcFOoyoI8ASR9AfhliLoSXuwFFprZSNWKKjuApOnAi6Rxfre5y8yWVa0khAOcBvy+w9v3AuPj22nATFwY1YAx3Dh4FFA1LScxnMmalv0+lskJ8l49AOO9Nk3DtWO8TWNZmTl0Nv3cAuaZ2Y4AenaOpB96Ln8apyVptaRPRlW8IUiaJulcSZtK9qMkXRtbfyS9WNL4X4qtcxORNEfSupIOsD620jMyo/qyOqrCDUfSQSpeSp7HfyXFi7lImldCWUn6RDRlewRJl5Xs00OryPNaDyBpHvBxJscNFpaQtRcX00605wFgP/5rNZZK2jzhby1gq5m9VEkTSafLf93+VLxeSZkBQtLrgfr8OXWafCLpcvklUviSHMAThXMAydmwMA0t9wNCLnnxz7gxbCjeNrMPBKyvb5H0FvD+gFXuBY43s+cnXiiaC7iOsManhvr6mdB9NQNn00kUOcBpgRWASAtQe5Q6+irXpkWC6nhUhwg7zwI+g8vdWwgcAQTfP6AkbwEvAZuAXwHrzaxq6LqOsb2/TeU2NAhNx3F2SdMlXS1pZw16heZ5SWd12tasvftq0GtDnqyiV8BDVRpQQEcOIOlg4FHgJ8CHg2pUDwuA38ntBNLpozz0pBQU2LRoFDAXNwo4KqACb5rZB8vcIOlDwFO4Tu1FHgYuKrt6R9IbwEEB9RgBFpnZ2xMv5D4BzOwd4GxgW0AlSr0Xs/+e++hd44NLhFnZwX0hp7+3AefkGR+m+NiQ26Dpa8CF5IeCh3Hz2T68YWYHe5ZF0leA233LN5wzzexx38JyQTPfp+VuJr8yWsA/cY/9H5vZ7qKbK31tSjoTeMyzuLcDyH3tb6U33vk+bMSt5m35FC7pAGeZ2bpOFWtqdvBn6R/jAywCToqtRB5NdYBzYytQAxfFViCPpjrAcbEVqIFPxVYgj6Y6QD89/seJuhtYEU2Nz/uOLABuAJ6rS5EpOAH4vmfZRs6ENtUByrCxzBArJIq5Hi8QTX0FJLpEcoABJznAgJMcYMBJDjDgJAcYcJIDDDjJAQacKQNBko7GLWzIWw/wkTqUSpRihaSJE03j6wEeNrMX291c6ABy+/neAlxNelI0mbPbXFsp6XbgOjPbm1cg1wGyEOedwNLq+iUiMgR8FThE0sV5y9WL/rMvILzxq57CMUiE7quLcMv6JlHkAHWcaZdSw/ypo68uz/tjkQPUscHztKmLJDLq6KtcmxY5QOhduRLxybVpkQP8pQYFxqYuksioo69ybVrkAD+rQYHRGursV+roq5/n/bFdbuCawArUke/Wr4Tuq/uAB/MuFKWGCVgG3IaLKiV6kxbOhpcVpawXRgKzyNE1klbRPhTcLhKVqJ/fAhNPEqkeCh4nq+AHedey1LDkAHFZ1Y+pYYkukRxgwEkOMOAkBxhwkgMMOMkBBpzkAANOcoABpx+ygxdLijV9vSiS3GA01QHewn+TJN/8/NjkbtMWm6a+Av4VW4EaiHu8WwFNdYCNsRWogX/EViCPpjrAfbEVqIFfx1Ygj6Y6wB/pr6fATuCJ2Erk0fYj0HOr2OCYWUvS9UDH05wN47tmtqemuh/Q5K34vbeKLUTSfJU7xHAqSh8aJemWgPJj8aikUsu8FfbQqBG5M6Byafp28cO474HzA+rRTTYBJ5nZm2VuUuzt4oFrCGt86CDbJdtn//P05q7hjwEnlzV+RujMoKNwNp1Eu9zA0HT0vWBm+81sOXAWMPGEzCbyCnAVcLaZvdFhHXV8W+XatOgV8G/gsMAKtMysUsqTXNbyycAlwLG47VdDPio74R1c4Op53Akh66p+8EkaI/wI7VUzm2TTIgfYDMwPrMCYmTU19NwoJO0nfH7gC2Y2yaZFXvaHwMLBHYic8KOOvsq1aZED3EL49KSUGuZP6L7ai7PpJIoygzYDywmbFZSSQ/0J2VctYHneucFTonR8fBTUxePjvbY7lzSP/FDwAuBHnu3aD8ytMSTaF0gaAvbgPxS8FjcCOZAWsNXMXgqp2yQkHVbSIy+tVaE+QNLxJfv0iJjKDknaU0LZEUmNPDmjKUi6o0R/jsmFy6Mq/FQJhSVpndwsY2ICkr5csi+3xNYZSV8vqbQkbZJ0njo/XLmvkLRI0i866Mebq8qufOaNpMOBl+ksdLkfF0oFF/maiVujMIQbCu3DjYlDr/odzmSNR9vGMjmhd+awTM5wJquFa/Mo7w31ZgMzOqz/VDNbX1XByki6E7gsRF0Jb0aAY8ysUswg1ITDd0iRvm7zrarGh0AOYGbbKQg1JmrhGeDeEBUFO/dObjjyBHBKqDoTubyLW2UUZC/HYHPO2eqdi3ErYBP1IOCLoYwPgRcdmNku4HRge8h6E4Az/vVmdn/ISms5+lTSYbjty3o+ebIh7AGuMrO7QldcS2JI9iQ4BVhJ2iG0Ks8An67D+FDTE+BAJB0H3AycSXMzkZrICPBt4J6iXT5D0LXTr+VmrZbhTiJZAMzqluweYTybZy0uj/BPIcb5UxHt+HNJh+BGDb7rCXZTcOyJJyuAczzL/gZYVUHWg8Acz7IrgPuBXWbW9XWT0SZjzOw/JWez9pnZ453Kk1TmDKQdFWWV+e7ZZGYT9/rtGumdPOAkBxhwkgMMOLEdIPc0ywKqxhPK3N9NWWX6IDixHWB7ibLbKsoqc383ZW2vKKsSUR0gW7bsuxXMQxXFrcUv0aWVla2Cr64ba1+63XQkLZFb3dqOrXKbVlSVdZvHOrvbAsiZm+ncjjFNkbQxMEi6UtJoQUdtU5stTkrKmSlpTRujrJEUZHMGSQsy3fMYlXRlCDl9g6SFkn4q6QVJr0naIOmbknwjamVkXSjpEUk7s59HJFWJMhbJmZO1YUPWpheyNi4MLSuRSCQSiUQikfDmf1XF/7nUVd3FAAAAAElFTkSuQmCC",
    "credit-card": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAEhUlEQVR4nO3dT2gcZRzG8e8vtYgiRaRQFDVBEKSikFARPKut+B8FiwhaPShEPHhQhIKCglo8qQj+IUKVSMGDFHuQ6lVQtAd7UDyobdVqbQ1YobTVPB42pzjv7rw7s7NJ3ucDucy88+5veZ+deWd2MwNmZmZmZmZmZmZmZmZmZmZmZmZmZrbqRdMOJE0CdwDXA5cCm5b+rF0CjgKHlv6+Aj6MiL+adDpUACStAx4CZoHpJgVYI6eAj4DXI+LzYTrIDoCk24CXgc3DvKCNhIB3gaci4kTOhrUDsPSp3wU8mVebdeg4sD0iPqu7Qa0ASNoA7AG2DVmYdec08GBE7KnTeGAAlj75+4CtDQuz7iwCOyJi96CGEzU6ewUP/mozAbwl6bpBDfvuASTdDuxtqyrr3GFgpt/EMLkHWNr1vziKqqwzlwPP9mvQ7xCwA7i61XJsHB6TdGVqZb8AzI6gGOveemBnamXlHEDSFPBj5gudBT4A9gO/ZW5rgwVwBfAIMHByt8xJYFNEnKrVWtITynNI0jWZRdkQJE1I2ilpMXOM7s55kfmMjs9IunaE79kqSHohMwBvVPWTmgNcklHLfER8k/8WrKHngO8y2k9VLWwjAPsz2lpLIuIf4M2MTSarFqYCsDGj42MZba1d32a0rfyNRioAOV8TK6OttWsxo23lWNf5LsDWMAegcA5A4RyAwjkAhXMACucAFM4BKNw5LfQxI/la0JjMNO0g9XuABeDCpp3birIQERctX+hDQOFSAVjXaRXWhcoxTQXg3BEWYuNROaapALQxObSVpXJMPQcoR+WEPxUAn9etPZW/HUgF4OwIC7HxqBzTVABOj7AQG4/KMfUhoByVY9rGbP9vfMgYl/XABa33Kmkh4x8Obmy9AKtF0k0Z4/RnVR8+DSycA1A4B6BwDkDhHIDCOQCFcwAK5wAUzgEonANQOAegcA5A4RyAwjkAhXMACucAFM4BKJwDUDgHoHAOQOEcgMI5AIVzAArnABTOASicA1A4B6BwDkDhHIDCOQCFSwWg8bNorBONn+2UGrzfMzq+KqOttavySWAJf1QtTAXgSEbHj6r3qHnr3taMtr9WLUwF4KeMjjcz4Bn11j5J08B9GZv8ktP5PRm3HpF6DzJ+RpLnAx2QNC3pSOYYzVb1lbpd/Pn05gG5NyD6EpgDfsB3GhuFi4Gbge3k3eBLwGRE/O/QnpxFSnoPeCC3QluRvo6ILVUr+u2yn8e3f1srXk2tSAYgIr4H3h5JOdalg8D7qZV9LyRI2ggcAC5ruSjrzraI+CS1su+sPSKOA/fiewevVi/1G3yoeSlR0sPAO3Xb24rwMXBnRPS9rF/rvD0i5uidEZxpoTAbvX3A/YMGHzK+yImIeeBW4ESDwmz0dtH75J+s0zjryl1EfErvy585fKFnpTkI3BIRT0fEv3U3GvqYLukG4HHgLuC8Yfuxxg4ArwG76+zyl2s8qZO0gd6ZwhZgEpiid8nS2ncMOAr8DHwB7I2Iw+MtyczMzMzMzMzMzMzMzMzMzMzMzMzMxu4/1upWe9CuEfwAAAAASUVORK5CYII=",
    "database": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAMTElEQVR4nO2daaxV1RXHfwtUHAABRcCiDIqPQnFEa4ljqRitVdtYG4dSjEPTJk36wVYT2nRQ2jR+avqhqU3T1rHWqaDiSCyoFQuioKiIlcEyFUWGhwwK/35Y+/oez/fgnnvPOffcd/Yv2bmXxxnW3ed/z917nbXWhkgkEolEIpFI6bBGG1APknoAQ4CjgIGhDQ6v/YDe4bVPeH9Q2NXC39uzEVB4vw1oBbaEv7eG1/XA2vC6HngPWGNmu9P/dPlQeAFI6gmMAEYDnwdawr+H4Rf+gMZZB8BOXAgrgGXAEuDN0Jab2a4G2rZPCiUASfsBY4Hx7do4oFcj7aqDHcAiYH67trhIomioAMK3+1TgbOAs4Ez8dt2d2QI8B8wObX4jBZG7ACQdDEwELgIuAQblbUPB2ADMAh4FppvZpjxPnosAJB0AXAB8G7/wzXpLz5rtuBDuAJ4ws4+zPmGmApB0NPADYApweJbn6oasB/4C/M7M3svqJJkIQNKJwE3AZcB+WZyjRHwMPAD8xswWpn3wVAUgaSjwU+BaoGeax44gXAg3mdmytA6aigDC9O1mYCpwYBrHjHTJduAW/I5Q9+yhbgFIGgH8FZ/CRfLjJWCymb1dz0HqEoCkLwMPA33rOU5KbGJPN+268LfNuCt3a3hfob3rt6NruC/uOj4kvD8Un662dzcfmtHnSMJm4FIze7bWA9QsAEmXAXeR35ROwErgLdzN+g5t7teVZrYlJzvcGKkP7o4eHl6PBcbgruqjyc/HsgO40sweqmXnmoyUdB7wONkN9IRf5Ir79GVgoZltzeh8qSKpN3A8e7q0R5OdKHYB55vZrKQ7JjYojPQX4LfCNFkKPAn8E5hjZutTPn5DkTQQd3mfDZwPjEr5FP8DTjazVSkfd08kzVZ6zJX0Y0mjMzW6gEgaHT773BT7s+axQLVGn5OCkasl3SKpJVNjmwhJLZJulbQmhf49K0tDZ9Zh2OuSviN/LhDpBEm9JE2RtLiOfn40K+OG1mjQKknXyB/9RqpAUk9J18rvlknZLenILIy6vgZj7pXUP3VjSoKkAZLuq6Hfr8vCmOkJjfhF6kaUFPmYKQnTszBiYQIDsp2KlBAl+zmo+qlhjwQ27Eiw7SBJExJsH9kLks4Ajkiwy/ZqN0wigCShSj2Bp+Tu4kgdSPom8ATJvK5VX6skAtiYYFvwByn3S7pT0pCE+5YeSUdKuhv4O96XSaj6WiURwIcJjahwNbBUPpBJ233c7ZA0UNKtuGv8yhoPU/W1SiIA7XuTLjkE+AmwQtIf5CFjkXZIOknS7fgTz6nAwXmcN4kA0uAg4AbgFUmLJN0oaWTONhQGSSNDHyzCH7BdT84RVY0M2BwH3AbcJmkx8AjwFDDXzLY10K7MkHQQ8CVgEh4eP7axFhUnYndsaDcDOyS9BMwhxAKY2X8baVytyB+dn4LHA5wNnEbBciKKIoD29MLTxD59qiVpHbAQWExb4uVSM1vXEAs7IGkQ/nx/DJ7AOgY4gSbIeiqiADpjEH7bnNT+j5K2A8vx0LBVwBraYgLX42lXraFtNbOqpkeS+uED196hDcAdMZWYwCHA5/BQsGE0cSR0swigKw7EQ62qCiiR9pjIbAIqef09KEaQZ+40uwDqoZQXvCN5TwMjBSMrASygds9h5LN8iPdp6mQlgPl4bPx38Tj+SG28C/wQL4UzP4sTZPYTYGatZnY7Pr8/D7gTz86J7J1WvD7AV4BRZvbbLPMhMh8EhgpazwDPSPo+XhXkYjw2Pg7EnI14TsQMYIaZteZ14lxnAeGD3Q3cLWl/3NkzCTgHODlvexrIJ3i202zc/T0nj2ogndGwDg8feFZolVy7M4DTcdfpKTSBJ61K1tGW4vYi8Hye3/K9UZhvXEjufDw0ACQdhbtUx+DOnsprUX86NuGD3sXh9Q08p7GwzzIKI4DOCLVx3sMLJ32KPNR8GD7TGA4Mpc1VOzi8r1QKTYMt+MVdF1r7SqErKs3Mmm7qW2gBdEXo6A+BV/e1bfDrV3z6ldCqfZWK3Up4hlDt84NmpSkFkIRwAbv1RayH6AouOVEAJScKoOREAZScKICSEwVQcqIASk4UQMmJAig5UQAlJwqg5CQRQJLybrEUXPpk0v9JBDAio20j1ZFJ/ycRwGEZbRupjkz6P4kAkmx7tKQirCHQLQgxDcMS7FJ1EfAkF/WjBNv2xYtEJalsFekESYOBp0m2oGbV1yqJADYk2Bbgi8Crkr6WcL9IQNIlwCt4kGwSPqh2wyQCWJ3QCPA06hmSHpV0fA37lxJJJ0qaCfwDj3FMyppqN0wigJ01GFLhq/jdYLqkcyUVatHqIiDJJE2U9AieB3hBHYerOscgz5hAwzOCLgbeknQHcI+ZrcjRhsIhaThwFTAZOC7v8zcqKHQ08CtgmqR5+K3ucTyGvp5ydIUn3P1OBC4ELsUTYBp2R2x0VLDhhZNOwwWxVr7syXPAC8DiRi6tngbyRTXH4FlPZwLnUqCMp0YLoCODgStCA2iVtAAfCS8CXgOWmNnmLvZvKMH3MRovgTcOz3c8ifQSVFKnaALoSG86VAwDkLQGL6W6DC8StRKfpazCV896P+07h3zFk8PxrKOhwJF43v6I0EZR24i9oRRdAF0xJLQuF0iSVMke2oSndu2gLfunq0SRyuom/fBydX3wPMT+7f6vW9GsAqiGbnvR0iTGA5ScrATQdFmyTUAmfZqVAB4ALseXOI/Ux4vAZXifpk5WApCZ3W9mp+Nz/D/hA7FIdWwG/giMN7MJZvYg9a3X0CWZjwHMbJ6ZXYdPka4EpgPdshx8nWzDPaJXAIPN7AYzeznrk+Y2CzCzj4B7gXslHYIXh7ogvCYJduhOLMeLRM0Eng59lCsNmQaGuncPh4akY/F6+pUiUS000D+eEcLrBs0Fnscrg73TWJMK4gcIHfEOPlaohECNx92o44Av4HX4m6Us+zb8Yr+Ou68X4AtfFK5SSSEE0JHQUc+EBoCkHvhPxShgJF4cahjujh2KjzHyWo1jO7AWdz1XCkUtw0u7LgVWhgKZhaeQAuiM0KHLQuuUcOc4Ao+KHYDHJvbBPYL70/ZQph+f/Ylp7yLeggdVbKStQtiG0NaZWZJFNAtN0wigGmJBqOREV3DJiQIoOVEAJScKoOREAZScKICSEwVQcqIASk4UQMmJAig5UQAlJwqg5EQBlJwogJKTRABJolKHJ7Qjsm+OSbBt1cEoSQSQ5Dn7JEk/T7B9ZC9ImgZMTLBL+kkkkq5Rcu4JUTqRGpDUX9J9NfT75CyMOaYGQyRptaTJ8pi+SBVI6iFpiqQ1NfT3bknZhNlLeqlGEUjS65Kuki8aHekESfuHL8sbdfTzC1kaeHkdhlVYLemXkpIMaro1kkZJmqbavvEd+XqWhpqkF1IwssKLkm6UJ4aUCvlF/5GkuSn252wlLMGXOPtG0nH4UuhJSpdWw9vAk8CzwHNm9n7Kx28okgbimU8T8XS4USmfYjNwspn9J8lONaVfSboYeAjoWcv+1ZwCWAL8G5iHZ9a8FpaYLzyS+tBWJOpUvGxuS4an/AS41MweS7pjzfl3kq4G/kx+uQXCM3CWhLYUT65chmfi5CoOeUWwyvL1lSJRLaHlmez6MTDZzP5Wy851JWBKuhDP+C1CafgttKsShhdM3oBn9bTiRaIqDpLOHCXtl5TvT1uRqL54ltFh7FklrAil3zYB3zKzJ2s9QN0ZuJJacBGcVO+xIomYD1xRb4Zx3c4ZM1uC/8b9DE+ajGTLNmAqMCGN9PJUc/AljQB+jdcH6m75/Y1mN3AfMNXMukyQTUomF0nSCbhKv0F2M4WysAt4EJhmZovSPnim31JJI4HvAVPwAVSketbjs6zfm9nyrE6Sy21aUi/gIuBqvC5QXoUcmo3teL2gu4DHzKyeRTqqIvff6TB/vhC4BPeIDcjbhoLxAe4BnQHMzNuf0dCBmrwC93jcPXoWMIH0XcxFYzPwL2AOMAuvHdSwNREKNVIPghiLTytPxVfWGAsc3Ei76uAjYDG+3sE8vHLqG0VaBKNQAugMeSDJSHzVjRbc5XoM7oI9Cq/900h24oWiltNWJOot4E3g3aIXiyq8APZGEMdg3DU7qF07DHfnHor/pPTDhVL5eenLZ6enu/DbM7hbeSdtaw1sxN3HG4B1oVWqhK0t+kWORCKRSCQSiUT24P9hGpDy7eZ/HwAAAABJRU5ErkJggg==",
    "dollar-sign": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAHsUlEQVR4nO2dXYycVRnHf89SW1P5sBhFrdWo0AZDWxrF6AXW+IGVhqS0YEOIxli9sirhRgl4ZbBpIlbUJtI0aUiMjaUxSrW2ASJESLhAKBCoIWgblY+LZkWtQGm3fy/OrN3ddHfmzM55zjv7Pr9k7p4z5z9z/jPv+XwOBEEQBEEQBEEbkLRO0j5JL3Re+yStq60rKIykBZL2aHr2SFpQW2dQCEnbZ2j8cbbX1hkUQNJSSWM9GGBM0tLaer0YqS3AkWvo7fOOdGJbQZsM8P5CsUNNmwzwpkKxQ02bDBCchTBAywkDtJwwQMsJA7ScMEDLCQO0nDBAywkDtByrLSAXSRcDVwIfBi4BlgAXAvO7FF0I9LrUewJ4tV+NMzAKPAhsM7NnCrx/NkNhgM7q3JeBLwAfrKtmIJwENpvZjtpCGm0ASR8HbgM+T8O19sFp4GozO1hTRCO/VElLgDuA62trKcxTZraypoDGGUDSJuCHwPm1tThxsZn9pVbljRkFdPbr7QJ20p7GB3hfzcrn1ax8HEnnA/cCq2trqcArNSuv/giQdC5wH/Cx2loqMAq808xO1hJQ9REgaR5wD+1sfEjzAdUaH+r3AbYAayprqMXvga21RVR7BEhaA+yvqaESo8A2YGvtXz9U+vI7z/1nSdO4/XKM1Hd4DDgC/KdL/E3A2h7f+3fAj/qXNi3/JI39qzd8VSTd3sMBjel4TNL6Tv8hp867Muq4q9Rnbxruw0BJF5F+jbkcB24GdpqZBquqvdSYB/gWaWUuh7+R5s0bsYI2l3A1gKT5wKbMYi8Cq83s6OAVBd7DwDXAOzLiTwIbovHL4W2A9ZnxPzCzR4soCQB/A3wmI/YY8P1SQoKEmwEkvRdYnFFkh5kdL6UnSHj+A+RufPhFERXBJDwN8IGM2JdjyOeDpwHelRH7ZDEVwSQ8DXBeRuxLxVQEk/A0QM6kU7sXSxzxNMCJjNhzi6kIJuFpgNGM2NYkaaqNpwH+kRG7QpGx0wVPAzyXEbsQ+GQhHcEEPA1wiHQcqle+UkpIcAY3A5jZv0km6JUNki4tpSdIeC8GHciIPQf4iaS2bRp1xdsA92TGfxr4dgkhQcLVAGZ2CPhTZrHbJX2phJ6gzsGQOzPjR4Bdkm6Jx8HgqWGA3cDzmWVGSJtD7pf0ocFLai/uBjCzU8AtfRb/FPC0pL2S1sZk0eypeTRsPyn1y2x4HXiadDKo2zHrK4Feh5WHgT92iXmDNL39d9Ik16HOUHeoqGmAxaR1/7fV0jBgTpPmOQ4Ae83sicp6eqJqp6pzQPS3pDH/XONxUod3d+vPAs6EpG/M4pzgMPC8pOtqf8/T0YhhlaTbgO/V1lGYA8DXzCxnVbQ4jTAApH8C0rn5ufg4GGcU+KKZ7a8tZJzGGAD+3yf4OXOnY3g2xoCbzezHtYVA/RQxkzCzA6TzA435hRTgHOBOSbfWFtJolJJAPFex8+bB5trfc6MeAVNRygKykZRT4IrKckowBqytmS+40QaYiKTlwHXA1cAq5k5n8Riw0sxerFH50BhgIkpJplYCyzhzX8CbuxQb9FTweaTdyyvIz3gylf1m1msCq6AfVChJlKT5kj4n6ZeSTs2iP7Ch5OefjkaNAoYRM3vDzA6a2UZgOfBAn2+1RZmZzwZBGGCAmNlh4LPAd8jbAQ3p+puNAxfVhTDAgDEzmdlW0rb2XBN8s4CkGQkDFMLM7ga+m1nso5IuK6FnOsIAZdkC/CGzjOs1OWGAgnQymm4mTfj0ymx3SWURBiiMmT0L/DqjyCpJbymlZyphAB92ZcTOIz+hVt+EAXx4gLwEGctKCZlKGMABMxvfvdwr7ymlZSphAD/+mhHrtiEmDOBHtxtNJuJ24CUM4EdOo54qpmIKYQA/cvIk/6uYiimEARxQOtW8IqOIW6LMMIAPl5PXscs9Pd03YQAfbsiMd8uVHAYojKQLgK9mFDliZi+X0jOVMEB5bgUWZcTfX0rI2QgDFETSJ0h3HebwqxJapiMMUAhJlwB7ydu+/hLxDzD8SFoFPAi8PbPojk4KHTfCAANE0ojSKedHgHdnFj8O/HTwqmamxtWxcw6lZFXrScmvlvf5NtvM7NjgVPVGNQNIWgRcht/CR85U7GJJ3e44XEQ6GXQFaSv4Bf0KA44CW2dRvm/cj4ZJWgLcQfrFzJXzfbPhNHCVmfV7oGRWeF8evQx4CLjIs96Gs6VW44PjP4CkEVKe4Mu96hwCfgOsN7PcAyQDw3MUcBXR+BN5CLihZuODrwFWO9bVdA6SEkO8VluIpwHe6lhXk/kZcI2Z/be2EPDtBB51rKuJvAJ83cwadSm25z/AvYAc62sKIt2EfmnTGh98L406DOz0qq8BnCat7H3EzG70XOPPwXUiqDNluhu41rNeZ54h3Y10t5kdraylKzVmAg3YAGwiDQu9poIXZtR1Ani1S8xrnLkv4M+k7OAPD0Ojt5JSSaKGnVgObjlhgJYTBmg5YYCWEwZoOWGAlhMGaDlhgJYTBmg5bTJAzuWNrbnosU0GOFIoNhgGJC2VNNbDOsCYpKW19QYFkLS9BwNsr60zKISkBZL2zND4e5T2LARzGUnrJO2T9ELntU/Sutq6giAIgiAIgsCF/wHgw4IxwYWhmgAAAABJRU5ErkJggg==",
    "eye": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAMTklEQVR4nO2dedBXVRnHv0cRwdRkCwQXwqVlSkVQQHGpsRHNcWhME5cS09IKyVzKykxtRCe3EQTcLUtLm2kxZ7BBc6NyIlwxFgEFURBwQxR4l09/PL/XXl9/v/d9z7n33Pv76fnMOM68/O59vs+55557z3Oe81wpkUgkEolEIpFIJBKJRCKRSCQSiUQikUgkEolEIpFIJBKJRCLR8LiyBcQA+LikT0vaU9IwSYMl7SipX+W/3pK263DYeknvSHpN0jpJr0haKWmppEWSFjjn3ipCf5E0fAcAekkaJekASSMl7StpaCRzyyTNkzRX0hxJ/3bObYxkqxAasgMAe0o6WtI4SQdK6lWSlI2SHpM0S9K9zrlFJekIpmE6QOWinyDpWEmfLVlOLeZLukfSnc65xWWL6Q513QGA3pKOl3SabIhvJOZIuknS3c65d8sWU4u67ADATpImyS5835LlZOU1WUeY6pxbWbaYugYYBtwMbOLDxybgRmBo2e1cdwA7AjOAzaVeomLYBEwHBpXd7lLJjwDsGX+upB9K+liZWkrgbUlTJF1d5lSytA4AHClpqixQ81HmeUnfdc79rQzjhXcAoK+k6ySdWLTtOucOSZOdc68XabTQDgAcIekWWVi2CDbIwrhLJK2Q9LKktZW/d2zoPrLH0ICKvp0l7SYLKfcuSO9KSacWORoU0gGAnpIul/T9iDZbZGHaxyQ9Luk/kpY651qznBTYQtYRRshCzgdJ2kfSlpnUdmJS0pWSfuKca4pkoziAIcA/Ir1RvwH8Bvga0KdAn/oCxwO/rWiIwaNAUSNlHIDRwMs5N8xm4I/AeGDrOvBxa+ArwJ+Appx9fQkYWbaPQQDHAe/m2BirgJ8Dg8v2rRbAYOBiYHWOfm8AxpftmxfAJKAlpwZ4EfgOtuzbEAC9ge8By3Nqg2bgzLL96hbAz3Jyeg0wmToY5kPBHg9nA2tzapMLyvapU4BLcnCyCbgG2KFsf/IC6ANch93JWbmobH+qQj53/lxg77J9iQUwHJiXQzvV10iAPe+ysBm4EOhRti+xAbbCXmazzhi+nYeezEEZ4KuSfi9pi8BTvCDpeOfc41m1dAWwq6TPyZJFB0ratvJPb0taLYsaPuOcW16AlgMk3SVpl8BTtEg6xjn35/xUeQLsD7yToRfPwtYGYulzwEHATGCph64l2PL0WCBatBToD8zO0H4bgBGx9HUlfjDZgjxXA1HCqUBP4FRgYQZ9bSwAJgJbRdLaA5iWQd8KYGAMbZ2J7kl4eLcFOCuitkOB5zI0aC0WA4dH1H0O0Bqo7SGKfH8CrgoU2gycHElTbyzlKjbXEqmxsZEmdKo4JYamaiLHEdZTm4EJkTQNAZ4KbLgQ7ifS4hNwEmGdoAU4LIam9uL6EvbcbwFOiqRpD+CFAE1ZWQjsFsmniYTdZCuIGUDDll5DmBxJzxBgWaCmPFhIvJHg/EBNt8XQI+DIQEHXRtLTm2KH/VrMIt5sZnqAnlbyfhRgje0zj27j/oiNU8QLX3e5OpKPPYAHA/QsIs9FNOCiABHLiBTkAb4UoCc2B0bydQD2bPelW+sFXUa5sASMRfLL22+SdLBz7l8ex3QLLL/wGVk4N5RVkp6o/F+yJNB9JGXZrPGoc+7gDMfXBBgr6SH55SGul7SHc251VuM3BPS+CzMZ7VzPqQF6wBZfbgf27+Tcoyq/CZ2LHx3R75Cl9ulZje6O/6rVXOIFShwWmvVlPrCPh519gecD7DxLpLUDbBXxSU89m8iyFxG41dPgZmCv/Nz+gJ6DPPUAPAxsH2CrHzAnwF60bezACPxHpxtCje2E/y7dq3L2uaOmmZ565hNw8dvZ64+tDPpwfZ4+V9E01VPPRkISaYFfehp6FSvOFA38pqJNwPAcbI7CLyoXtTIIFo1d56EHfNcJgG2A1zyNTIrkc5umXT315BYRA+72tL1zXrZr6DnbU886bCd2tw2c5GlgGTY9iwZwlKem/XK0PdrT9pF52a6hpxf+sYETqp2rVhrXEZ6apjjnNnse44vPvH+VrJRbXjwu6VWP32eJUXRJpZ7AFZ6Hjav2x1odwCczd5WkX3uKCaG/x2+fcM6Rl+HKuZ7wOGRAXrY74Rb5dcqq17RWB/BJMZpZUIWLjpU9OyNb9Ks6q7r+yXv4aA2iUnlspschVaOcoZm879OSwzkSYWRu+1odwOcOOoNi9u297fHbGEmSPusE6yPYfx/ANpLO8Dik6ghWqwM86XHigZImevw+FJ/n3XByDMlWzrWvxyE+WkM5TX7vGk9V+2OtDjDLU8yPiL+J06cO7yBJuU0DJY2WX2NHrRlcGXHP9zys+9cUSwDxDQRFS/WuaNrFU8/tOdr2TYXbKS/bNfSc46lnLT6BoIqRKzyNrCHyjl784vLNgM+wXcvmfvjVOogdCu6H/815WYihIfgvBkXJ/2unaYanniVAvwz2+uOfCjctT5+raPLNE9xIaK0hrG6vD7kswHSiZ6ynHrAlXe9OgF38kOXgMTF8r2jaH//KKzOyGByGf/3eecTbR+eA/3rqARsJRnnY2Q//ZWCwLWmxEkJ64p8FvREI3X38nmHfYRfg4pz8rqZnYoAesCXdu7GFnQ9cJKxzjQHuILy+UZRtbxV9lwXomdrVebuTFLqjbFqzbVe/bUezpEOdc3M8jukW2OjyjKRPZTjNq7LY/iuyNmhLCv1EhnM+J2lv51xzhnNUBThU0gPyi9y+JUsKzR6TwKp3+LIc8FnA8dFzCOE7aWPxxUi+DiJsO55vnKBTEb0IS5KcTbwE0WsD9MQiSioclgj6cICeBeSdnwEcHtg4UXLkgC2B+wI15UnM3U83BehpBb4QQ4+AXwU20rmR9GyPJX6WxQIiBb+AHwdqujmGnjZRfbD6tb60AqdE0rQb+ZSC8WUBEOVjF8DphL3jvEjs+orYvryQaVIzNfLSctC0HfCXAE2hzCLenf91wtq3BZstxAf/dYI2mok3EmyJFZ6KPTu4knjP/NMJj0FcGkNTLaE9gEcChbYS6Z2gom0MVms/b+YT6+XKdF9AeOedTaRO2ZngQYS9D7QxjYiVrYCjsb16WZkPnBxLKzbVy1Lr4EWgiCTUquJHYMUKQ3kgpngsvHsAcD1+cYxFWAeN+rlaYCBh8/w21uOx6bUaeZSKHS/pDwr/hs5LkiY45x7LqqUrsB07baViB+n/4e31sjzIhZKedc69VICWQyTdKSn0AxjNksY75+7LT1UgwJkZejHYy+ElRFpFrCewIf8yspWObwW+WbYv74PwwEV7nqSs2rcFgK3nP51DO51Xti9VIayWUEeasedvo381/D2wNK7p5PMZnZ+W7U+nYNOZPFgH/IAG+lZQR7Dk2vOA13Noj6jT51zB3gny+DwK2FRzErYRoiHAttdPBlbm1AZNwGll++UF9i29LFPEjqwBLqW+Pxs3pKJxTY5+rwe+XLZvQQAjyRYsqkYTFvc/hjp4PGDD/LHAveQ36rWxnIzz/NLBIoZZgh2d8SZwFzCBSNlHNXwaAJwI/A54K5JvDwJZUtS6RVEfj+4h6ReSzlM+O5Kr0SrpaUmPygo6zJO0yDnXkuWkFe17yD4ePVr28ejPK17btUqaIumirNq7Q9Gfjz9M0m2Som6dascmSYslLZV9nGq1pDWS3pD0rqS2ugZbS9pG9gn5AbINr0MlfVIWNYxa/qYdyyWd4pz7e0H2it/bj62jXyPpG2XYr1OQdKukc51zbxRpuLQLUBkNpsuG148yCyWdWeRd355Yz+Mucc7Nlj1LL1ABBRXqkDdlW7z3KuviS3UyBFfedi+U9C0V97wti02SbpR0qXNuTdli6gqsBsAMbE/bh42NWF5C1CKSvtTFCNARYJCks2QjQvD27jphreyOn+qc86k0Vgh12QHawKpaHCurhzNWda63HcjiEbdIuqdS0q0uaZQGFfaJtgmSjpO9PNYjT0m6R9JdzrmlZYvpDg3TAdoD7C7pKEmHyyJzPp+zyZMNkh6RdL+kvzrnlpSkI5iG7ADtwTZBjpQ0RtIoScMlDVP+U9xWWURxnizU/E9Jc51zTTnbKZSG7wDVwHIHPiMLMg2VhZ4Hy+oN7yAbMbaT1Jbq3SQrRLlB0uuS1kl6WdIKWQh5saQFzrl3ivIhkUgkEolEIpFIJBKJRCKRSCQSiUQikUgkEolEIpFIJBKJRCKR6Db/Aw8b/wQD2easAAAAAElFTkSuQmCC",
    "factory": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAIj0lEQVR4nO2dbYxdRRnHf0+3VYqFkioiVlsCDdrWYgytxVhRqBFTfIlfREJibGwp+KHxpRISEyvGDxjjG4nRDQkJjVCNL4GggKZkqy0KdosirUW0aw3dvtAUoa0NtNv+/TD3Yrt07z1z75x77r3z/JL9tHOe/5wz/ztnzpyZ54DjOI7jOI7jOI7jOI7jOP2PlRVY0gxgMTAfmAmcHRniKDAKbAMeM7MX0tbQSY6kqZJWSNos6YTSMSZpo6Tlks6q+jydcUgakHSzpL0JG30i9khaKWlS1eftAJJmS3q0Aw0/nt9JekvV5581kq6QdKCCxq+zX9Kiqq9DltQa/1CFjV/nRbkJWqalpwBJs4EtwPlpq9My+4GFZra76or0GtEDKUkDwD10T+MDXADcIx8YRtPKBbsReG/qiiTgSuCzVVei14i6BUiaCowAbyqnOm2zF7jYzF6quiK9QmwPcAPd2/gAFwLXV12JXiK2B9hMfPd/DNhE6DkUUa85wBJgSqTeRjO7KvIYpxmSZih+evcBSTPb0Jwl6aFIzTFJ56U8dweQtKyFxh9IoDtZ0oOR2tekOOcciBkDzIsoewy4ycxORNbnVZjZGLAKOB5xWExdsybGAG+OKLvJzEZjKzMRZvYs8GjEIS3fdnIjxgCviyj7r9iKFGBnRNlpJej3JWXNnJ0sIWbM7aS0hS79hk+dZo4bIHPcAJnjBsgcN0DmuAEyxw2QOW6AzHEDZI4boACSpktaLelhSbslPS9ph6R1kq6VVOrMY9X69UoMRryOHewHfUmmsNXtYBO9P0mam0Kz0/reA0yApEuBIeBOYEaT4ouAxyQlWyzbKX03wDgkTZH0FeBJ4P0Rh54L3Ccp5rV55fpugFOQtBjYCnwDaGUX8huAb/WSvhsAkDRN0veBPwAL2gx3XeyvUNI5VelnbwBJ1wLbgdWkuR4DwIcj9bdVpZ+tASRdIGk98CtgVuLwTUfkNf2fVKVfJzsD1B6tlgN/Az5Vksy5BfWv67T+eCaXVIGuRNIcYBC4Okf9M5FFD6Cwt+BW4K9UcPFrj3aV6Tei73sAheQRdwLvzFG/GX3bA0g6W9LtwB+p4OJXrV+Ufu0B3kV4tLsoUTwBzwOvr1D/IGGiJyn92gMsIt3FHwE+BPyiYv1fJop3Gv1qgBSMAXcAl5nZhn7V79dbQLtsAVaa2ZP9ru89wOn8F/gC8J6KGv9Ip/W9B/g/DwGfM7NdFerfbGb/7qRoL/UAYyXFPQDcYGbLmjT+sZL0nztFv6OND71lgF2J4wm4G5hrZvcWKJ+6cer68wrql0Iv3QI2Jow1AqyKHF2n1N9JyKBSxdPFafRSDzBMmEtvhzHCipkFsRffzIYJy7RS6Ff1aPkqesYAZibgSxRPNTeeJ4B3m9ktZna0xRjt6G9NoJ+cnjEAQO1Xc1vkYUeBNcBiM/tzm/qPAF9rUf+KdvXLoJfGAACY2W2SDgG3A69pUvy3hEerkYT6X6/pf7Og/k1mVkbOpCT0VA9Qx8y+S3jDdjdweNy/TwK/Bz5uZtekbPxT9L8HXDaB/gnCgLGu37WNDz3YA9Qxs6eBz0haCVxKyBN8BNhhZi92QP/vNf0VwNsIOZQPA0+b2aGy9VPRswaoY2bHCa9et1ekP1alfrv05C3ASYcbIHPcAJnjBsgcN0DmuAEyxw2QOW6AzHEDZI4bIHPcAJnjBsgcN0Dm9MTbQEmzCF8SfQF4qvYGsN2Y04F3EH4E28zsPwliTgbmEzZxjnT7WgDo8h5A0tWSthKWZD9CWFe3T9JaSa9tMeZMSfcS9gNsJiweOSDp55IuajFmPQHFKPAXYAMwIumpWhKo3kcdTtUqabkaf6p2g6RmS7LGx5wjabRBzOckRaVpkzQg6f4GMU9KWh139mfUKeX6d2UPIOntwA9pXL+lwNqImAasp/EHMM8HfqrQlRdlDfCxBv834DuSLo+I2TG60gCEDZJFuvjVkqYWjLkUWFig3FwaN+grKHwbeU2BokXLdZxuNUDRHLnTgKK/rJi8u0XL1gd8qfU7RrcaYHpE2aKfio/5pHwZMWPOqWN0qwFiNmIWLbsrImbRsmXUs6N0qwHuK1juH4Q8u0V4gOLfNC6kX9vO/UTKmJ2mWw3wA8IO3kYI+HJtz2BTzOwZ4EcFiv44cgvXLTQ31h7g2xExO0ZXGsDMDgPLmPgz9GPA583s/sjQXwR+1uD/vwZWxQSs7RdcycQJJEaBZWZ2MCZup+hKA8ArO28WALcCjxPy5P0TuAu43MzuaCHmy2b2SeAThCzde4B9wMOExNEfbWXnrpndRdiqNgg8Q8gpOAx8FZhfYbKpdJQ1E+UUI6uZQKdzuAEyxw2QOW6AzHEDZI4bIHPcAJnjBsgcN0DmxBggJkGiGys9McvUCrdVTEPFzJFfElHWKcbFEWWPFC0YY4A9EWWXKKzldxIgaTawJOKQ3UULxhggJg3aFGAwcnWtcwZq13CQuFvAjjIqcp6ksYg3UpL0oKS3Jq9MJkiaLek3kdf8uKTCaxUtskJDwAciz+M4YQfOToovycqdSYStcEuI3743ZGaFP08bG3wd8QaYAlxV+3PKZ11M4dge4CzCWr0LY45zOsYocImZvVz0gKjndTN7ibDMyelO1sY0PkT2AACSJgFDwJWxxzqlshFYamZR46xoAwDURvbDwBtbOd5Jzj5goZmNxh7Y0pStmT0LfATombz4fcwhwmrm6MaHNubszWwL8EFgf6sxnLbZR+j2h1sN0NZLm5oJFhKybDidZYiwP6LlxocEb+3MbDfhGf9GYG+78ZymjAIrCL/8mPczZ6SlQeBE1OYJrgc+DbyPkBjBaZ8xYBNhkmd97KNeI5Ia4FRq89GLgXnATOCcsrT6lMOEX/t24PFOfAjLcRzHcRzHcRzHcRzHcZx+5X8o3A/Pj9aJLAAAAABJRU5ErkJggg==",
    "file": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAFh0lEQVR4nO3dT4hVZRjH8e/j5DiVUqSbotCoRW1Ks1ArIhtBhhb90UVtXARt22gUhFKbaNWihW2ihSQk1WiQFhWaJlmkDRpSEvbPjZCpkaU25a/FzATpHT3nvO+559x5ns/6nPd9597v3HvuOWfuQAghhBBCcMe6NZGky4FlwP3A7cA84JpuzV+RATOAfmAacA74HTgI7AZ2ATvM7HRjK0xUewCSbgRWA6uAWXXP14BTwNvAK2Y20vRiyqotAElXAC8ATzH2G+TBFmCNmR1ueiFF1RKApFuBYeCWOsZvuT+Bp81sfdMLKSJ7AJLuBrYCV+ceu8dsAJ40s7+aXsjFZA1A0h3AJ0zN9/oqtgIrzOxs0wuZTLYAJM0BRoDrc405RbwPPNLWCKZlHGs98eR3MgRsljSj6YV0kuUVQNIy4KMcY01hrXwlyBXAbuCeHGNNca2LIDmA8QO/fRnW4kWrIrgswxiPVdjnCPAaY+G04oE4z2ZgZk1jTxwTtCaCJJJGVM6b42cJW0vS8ZI/UxXb1NIDw8IkDUgaLfFD75GU41WnVioXwMkS256v8QhSPwbeRLm3kXVm9nfinG3zDHCs4r6Nf0RMDWBOiW1HgZ2J87XR98AgPRpBagBlrvKdavt58arM7AA9GkHOM4GujUdwH3C04hBDwLuSBvKt6tIigIzM7BvgAapHsBzY0s0IIoDMei2CCKAGvRRBBFCTXokgAqhRL0QQAdRsPILlVP+IuBwYrusjYgTQBW0+TxABdEmmCN6QlPU5iwC6KEMEK4G1+VYUAXRdhgjWSlqcaz0RQAMSI+gDXpXUl2MtEUBDEiOYT7U7sS4QATQoMYI1OdYQATQsIYL5khamzt/627MaohLbDksazTBn3/i8Ze7UfpTEO7IjgM5OUPzLK+q6e7iIQeC5lAHiLaCzQ00voKAFqTfZRgCd7Wh6AQX1AzekDBABdLYJ6JW7l5MOBCOADszsCLCx6XUUNC9l5whgcs9S/XRtN/2RsnMEMAkzOwo8DrT9VvYfU3aOAC7CzD4GVpD4W1azpGOVCOASzOw94E7gs6bXUocIoAAz+xa4F3iQsT8dP9XsivKJM4EFmZmAbcC28UuxNwPXUs9jWOf3E/xPBFCBmf3D2NnCWs4YZrq2UEi8BTgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXATgXAbSTldhWKROlBlDmX6vPlNSfON+UJ2kGMKvELmdS5ksN4JcS204HlibO58Eg0Fdi+zLPQV6SBiSNqrgvJE1vbMEtJ6lf0t4Sj+do46+qkkZKLFiS3pJ0ZaOLbiFJMyW9U/Kx3Js6b47/Hv4hML/E9iuBJZJeB74ETmdYQy8bABYBTwDXldz3g9TJyxxtdiRpAfBV6jihktvM7OuUAZI/BprZCPBp6jihtJ2pTz7kOw/wfKZxQnHrcgySJQAz2w5syjFWKGSjme3KMVDyMcAESbOBfcDcXGOGjn4AFprZiRyDZTsVbGa/Ag8Bv+UaM1zgJPBwricfMl8LMLP9wBCQbYHhP8eBITM7kHPQ7BeDzGwPsAQ4mHtsxw4Ai83s89wD13I10MwOAXcBLwFn65jDiTPAi8AiM/uujgmyHQRORtJcYDWwCriq7vmmiBPABuBlM/u5zolqD2CCpAHGrgYuZezU8VxgNnFPwjngGPATsB/YDuwws3jlDCGEEEII9fgX+SxGKvAT4BsAAAAASUVORK5CYII=",
    "file-text": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAHdklEQVR4nO2dX6wdVRWHv9XK7VXbaIQHNGopQsQHtbQaQI0KpWkuJgoCiSYEA4kxFUNiwPhgIPpiNMY/MQIvBhJSEkm1VCPFiKEUCWCE3rSGKAGpygsJ2GIof+q19+fDzDW13srZZ+85e82Z9SXnbWbttWe+OTOzZ/+BIAiCIAiCYHDYpAqS9HrgQuDjwPuB04C3TKr8MTFgFTADrAAWgReBx4EHgQeA3Wb2SrUMM+lcAEnrgOuAK4E1XZdXgcPAT4Efmtl87WRS6UwASW8AvgFcS3MFDYGdwPVm9ufaiYxKJwJIeg+wAziri/jOeRn4ipndXDuRUSgugKQPAXcDby4du2fcDnzezP5ZO5H/R1EBJG0A7mc67/XjcDdwqZkdqZ3IiSgmgKRTgHng7aViTgn3AJd4lWBFwVg3Eyd/OeaAuyStqp3IchT5B5B0IXBviVhTjMt/glICPAh8uESsKcedBNkCtA9+jxXIZSi4kuB1BWJ8Zox9ngF+TCOOiwNxHHcBqzuKvfRM4EaCLCTNK42ftK2EbpF0MLFO47BLTh8MR0bSrKSFhEo/LKnEv06nKE2AFxK2PZ7qEuS+Br6LtNvIjWb2r8wyvfFV4Pkx963+ipgrwCkJ2y4AezLL88jTwCZ6KkGuAClf+Q57bxcfFzPbT08lKNkSOGhaCT4KPDtmiDng55Jmy2X12oQABTGzPwIXML4EW4Cdk5QgBChM3yQIATqgTxKEAB3RFwlCgA7pgwQhQMe0Emxh/FfELcCOrl4RQ4AJ4LmdIASYEIUk2Cap6DkLASZIAQkuA24ol1EIMHEKSHCDpHNL5RMCVCBTgpXALZJWlsglBKhEpgTrGa8n1v8QAlQkU4LrS+QQAlQmQ4L1kjbmlu++e1YllLDtDkkLBcpc2Zab0lP702T2yA4BlucQo09e0VXv4VHYBHwtJ0DcApbnidoJjMjZuZ1sQ4Dl2V07gRGZAd6REyAEWJ47gb70Xs56EHT/DCDpTOCTwBmMLuwi8BTwCzN7MrVMM3tG0h3A51L3rcBp1UqWtDlhEMTBxNgzkm6SdDRj4MXRNkbyHEWSTpX0XEbZk2Jrat2OxeUtQJIB24AvkpfjijbGtjbmyJjZs8BnAe9d2f+Ss7NLAYCLgcsLxru8jZmEmf0GuBR4qWAupcl6VvEqwNUdxLxqnJ3M7JfAB4CHyqbjA68CrPcU08z+BHwE+ATN0PHDpZKqjde3gJSm2InENDMBu4Bdaj7FngG8lW6OYZfzE/wXXgXYR2YDxwliFsHMjtK0FnbSYqgy3xZGwust4NYOYt7WQcze41WAncD2gvG2tzGD43ApQHu/vYJm7sHFjFCLbYwr2pjBcXh9BqCdS+AaST9gwk3BQ8KtAEu0J/C7tfOYVlzeAoLJEQIMnBBg4IQAAycEGDghwMAJAQZOCDBw3DcEAUh6J2ktgbVZBJ4ys7/VTuS1cC2ApAuA7wAbaucyDpL20qwheF/tXE6E2ytK0lU06xD18uS3bADubeviEpcCSDoLuAWn+SWygmZCB5erqHo9wF+mWbV7WlhFUyd3eBXgY7UT6ACXdfIqwJtqJ9ABLuvkVYC/1k6gA1zWyasA09h/z2WdvApwE81aPNPC0zR1codLAczsReAi4EDtXApwALiorZM73LYEmtkTkt4LfAm4hB42BdOM8PmRmbkdXOpWAID2wH27/QUd0JcrKuiIEGDghAADJwQYOCHAwAkBBk4IMHBCgIETAgwc1y2Bkt5INAV3ilsBJL0buAdYVzuXMTkZOAf4gqQ5M3M5Bb3LK0rSGpop2fp68o9lHc3UcmtqJ7IcLgUArgFOr51EQU6nqZM7vAqQPK9vD3BZJ68CrK2dQAe4rJNXAf5RO4EOcFknrwLsqZ1AB7isk1cBvg8cqZ1EQY7Q1MkdLgVop2ffSt4soV5YBLa2dXKHSwEAzOw2YDOwt3YuGewFNrd1cYnblkCAdlz9xpggojtcC7BEeyDdH8w+0pcrKuiIEGDghAADJwQYOCHAwAkBBk4IMHBCgIETAvgkZaXzrNXQcgVIWVp9taSZzPKmHkmrgJT+g6/mlJcrwHMJ254EnJ9Z3hDYBKxM2D7lHJRF0qykBY3O7ySdVC1h50iakfRowvFcqP6vKmk+IWFJ2q5mwEdwDJJWS/pZ4rF8NLfcEl8Dfw2sT9j+MuA8SbcCvwdeKZBDn5mlGUByNfC2xH1/lVt4ytPmskg6m3532ugz7zOzP+QEyH4NNLN54Le5cYJk9uSefCjXDvD1QnGC0bmxRJAiArRdt+4sESsYiTvM7IESgbKfAZaQdDLwGE5HwEwRB4CNZnaoRLBiTcFm9nfgUzgdATMlvABcXOrkQ+FvAWa2D5gDiiUY/IeDwJyZ7S8ZtPjHIDN7GDgPeLx07AGzHzjXzB4pHbiTr4HtbBgfBL7FdA3xmjSvAt8EzjGzJ7sooNhD4ImQtBa4DrgSp+vmOOQQcDvwva4Hl3QuwBKSZmm+Bp5P03S8lmYenaH3SVgEnqdZU2gfcB+w28zinzMIgiAIgiDohn8DzuH96H67e7gAAAAASUVORK5CYII=",
    "filter": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAHW0lEQVR4nO2daahexRnH/5O4Ra0aXKN1q8YoCC64ROtWF7BaqhYjaBH1ixIVQVwQUdyKQhWFtiJKqWibKHRJKI2KGkyDW1qtK+5LWo1rFOWamOXm/vxwbmrQ9733nZlzzrzzzvP7mDvnzP955p/znjPznDmSYRiGYRiGYRiGYRiGYRiDj/M9ANhc0lGS9pW0naQJNWsy/BiR9IGkFyQ95pwb8jm4ZwMAUyVdLek0SRv6dGK0xgpJ90v6lXPu7V4OGNcAgJN0qaQbZAOfCyskXemcu228hmMaAJgg6feSzqlJmNEud0qa6ZyjW4Pxfr9vkg1+zpwn6bqxGnS9AgA/kTR/rDZGFoxIOsI590SnP45lgGcl7d+UKqNVFjnnpnf6Q8efAOAw2eAPEgcDB3f6Q7d7gF80KMZIwymd/rGbATpeLoys6f0nQNLODQox0tBxTNfr0niSx4k/kTTXW45RB6dI2rrHtr1P4gHv4McFQfKNYIALPcfoFZ+TP+B58jXAqQ3Ga6wDcNpozn34W6dzdbsHWOCpaYKkP1FNHhkNAhwt6V75r8Iu8OlkJ2DY02EAXwD7egozegTYD/gyYFxWAdv5dnZPQEcAHwC7NpSDYgF+BHwYOCZ3hnQ4BVga2OEbwDYN5KFIgG2ANwPH4iOg1yeF73V8LNXlI4R/A5vWnIviAH4APBM4BiuBo2IFnI7/Heda5gNWRBIIsD7wUGDuR4Az6xJyQaAIgPuoCksMDwBH+H0YwMV1C7o5QszttYopAODWiHzf2IQgB/whQtTltYsaUIArIvL8R6o6zkaErY//LOFaRgArLxsH4JejuQrhH0C39Z3aBG4MPBEocBg4uVGBGQOcCKwOzO0iYJO2hG4JvBoodDlwaCtCMwI4CPgqMKftz7sAPwT+Fyh4KbBXq4L7GGAq8HFgLpcAaeo3gL2BzwOFvwfslER4HwFsDywOzOEXwD6pAzgC+DowgJeByUkDSAiwOfB8YO6WA4enjkGSBPyc8JuXp4CNU8fQNsAkYGFgzoaB/ircBc4NDAbg7zT9+NJHABOBv0Tk68LUMXQEuDYiqLtS628L4LcReboqtf4xAX4TEdyY77ENAsD1Efm5I7X+cQEmAH+OCPKi1DE0BXBeRF7mAhNTx9ATwAbAI4GBrgFmpI6hbqhulEPK7AAWABuljsELYDPgucCAVwLHpY6hLoAjCX9UfpFcH5WpJjneDQz8S2C/1DHEQtxk2TvAlNQxRAHsTlWXFkLWBabAjoRPl38KTEsdQy0ABwJDgYl4C9g2dQy+AFsRvmC2DDgkdQy1AhxD9dseQlYFplRL5k8GxroKOD51DI0AnMGAF5gSXzRzVuoYGgW4LDA5ALPp4wJTqrK5uyPiuyR1DK0A3BKRpN+l1t+NyLh+nVp/axD/P+XS1DF8F/xf116XWfTxla0RqH4rHwxM2AhwduoY1kLcyzOPAhukjiEJxN8t/7QPYjgGWBEYw7/I6OmmEaiel18LTOAyEhaYAgcQPr/xJhnObzQC8TNmeybQHDPDuQTYpW3NfQ3xBaY7tqh1CnFrHLaJRieIWzV7iRZWzahWOf8TqHElcGzTGrOGuHXzf9LgujlVncPDgdpsI61eAc4PTDLAX2mgcoaqkDOm0mlm3ZoGGvqsdg64PULPwNc6NgJwV0TSr6lRx9UROvw3aDIqqC67cwITPwKcUYOGMyMGfw65FHL2K8S9QTME7BbR91TC39hdCPjsuWx0A5hM9ZgXwkMR/Ybe8b8EbFFnDooH2AH4b+CAeG9hS7VFXgiLgR2ayEHxAHsStnnlvIC+QlYql5JgWroogOlUi0A+DOOxVy5VObvv8u4yILsvrWRXhOCce1rSDEmrPQ6bKMln6fhE+eVmWNKMUW1ZkZ0BJMk594Ak31ekf+zR1neJeeaoJqNNgMc9LtFPeZx3kcd5FzYZY9NkeQVYh9kebbf3aOtzFz/Lo23fkbsBevpE+ig+e+j5TOC869G278jdACMebX1i9dlqteuXuXMgdwMYkZgBCscMUDhmgMIxAxSOGaBwzACFYwYoHDNA4ZgBCscMUDhmgMIxAxSOGaBwzACFYwYoHDNA4ZgBCscMUDhmgMIxAxSOGaBwzACFYwYoHDNA4ZgBCscMUDhmgMIxAxSOGaBwcjeAz6vZGwEHNaYkU3I3wAqPtpMkLQLmmRG+JXcDLA445gSZEf5P1gZwzr0v6f3Aw80IytwAo8yJPN6MkDPANMI/MdOJecByj/bHpc5BDNlfAZxzr0uq8wshJ8hvl7BVNfbdOtkbYJQrJL2YqO8PE/VbCz7bofU1VFu0z5c0rcVuhyRNds6tabHPWhmUK4Ccc0tU7fE7t8VuH8558AcW4GeEf9zRh+NTx2p0AXDASQ0aIetNoouhISMMAXukjs3woEYjrAR8Pj5h9BORRvgUODp1DEYNrGOEZ3sY+NXA3Xh8dygXBmYeIAbgMEknS5ouaVdVM4GfSXpd0mOS7h99zDQMwzAMwzCM/PkGWKXpc9jMtL4AAAAASUVORK5CYII=",
    "fingerprint": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAQq0lEQVR4nO2de7CeVXXGn01uInipQoSCmHCJkVrBlKCCGARrbbAoMlJqvTSV2o6DICO0eGmxg7YUHaagFNEgLVEQdKQyCgkyQMotigoBC8SQEKWBQCQXMCDJSX79Y3+BY/jec9bzXr73O+F7ZvJPznrXXnvv9e3Lum1pgAEGGGCAAQYYYIDnHVLbAjQBYIKkaZKmS9pb0isl7SlpV0m/J2lnSS/e5rN1kp6UtFbSo5JWSnpQ0v2SlkhamlLa1Av5e4ntQgGAqZJmSXqTpJmSXitpQs3NbJR0t6SfSLpN0sKU0oqa2+g5xqQCABMlHSHpaElvl7RPS6Isk7RA0lWSbkgpbWxJjtIYMwoA7CDpcEnvl/Ru5aW8n7BO0pWSviHpxpTSlpblCaHvFQCYLOkESX8jaUq70oTxgKSvSpqbUvp128KMSQCvBr4G/JaxiyeBC4H92h7PMQPyxF8KbG516urFEDAP2Lft8e1bAJOBC4BNrU5Vs9gInA/s2vZ49w2AccDJwNpWp6a3WAt8DBjX9vi3eggEDpR0kaQZbcrRIn4i6YSU0uK2BNihjUaB8cAZkn6s5+/kS9JBkm4HPt3WatDzFQB4laRLJR3Sg+Y2K1/JlkhaIekhSY9Ielz53k6HbpyyafjFknaXtJukqcrm5KmdvzeNmyT9ZUrpwR609Qx6qgDAbEnzJL2soSZWS7pB0i2SFkm6K6X02yoMgR0lvU7SGyQdpmxybuoQ95iyEixoiH97AE6nmavdEuBzwEFka2HT/dgBmAl8Hrivgf4MAX/fdD96BmACcHHNg7QO+DLwR33Qv4PIV7v1NffxIrJXc+wC2Am4psZBWQacCOzcdt+2BfAi4CTggRr7ezWwU9t9KwXgpcAtNQ3EMuCDwPi2+zUayCveHGBFTX2/Gdg2dqG/QZ7822vo/DrgFLL7d0wBmASc2ulDVSxirCgBsDNwWw2dnkf2BI5pAK8AvlnDeNwEvLDt/owI8vK3oGJHV5Kvi9sVgHcCD1ccm6vp14MhkIBLKnbwKuDlbfelKQC7Aj+oOEZz65SpNkMQ8GlJnyv5+WZJn5L0hZQSoxG7AF4iaX/l0LHdJU3Wc4NC1ytbCR9WDvW6N6X0eAOyJEmflHSmypviT0spfbE+qSoCOIryRp71wDtqlmcf4O+Ay4DlFX5ty8j799+SA0/rlPEo4PGScg0Bb69TntIApgBrSnbkIeCAmuSYCvwjcHdJWSK4C/gM2Z9Rh8wHUv5csBp4ZR1yVOnAeODWkh1YDuxdsf0EzCYfPHsZQbSZbOD6E/KSXqUP+1DeZnBFlbYrA/hsScGXA3tVbPtdwJ0l268TPwPeWbEvUyinBJuAdqKjgRnkECcXK6mwn5JjBheWaLdp3EiF7QzYl3LbwWFl2ywN8tL/sxLCrgNeV6HdA+nv0LFNwNnAC0r27/X4B8NWFOCUEoNT6eRKVrp7S7TbBv6XkoqOd6Pq/RZANm2WcX2eWrHdo0u02SaeBD5Usq+fCrbR+0MgOdHBxfeoflo+u0S7/YCzMQNVyLebK0fhu4peXwPJBzA3bn8lsEsNbf+HO/J9hG9i2vDJ3sSvAlu68LsDeHXVMZUk179+RolvTqgpP255DTx+I2mpct7/amUTtJT7tIukvZQDQev2ur1P0iTg+JTSUOSDlNLTkj4CnCPpXcrBqesk3Sjp2p4nn5J//UOm5l9SY/v7lGh/PdkcPAeYRmAbIsf8TQc+DFxOeXNtN8yLyNCXICdqOlhDzSlQwDnBtq8H3kvJ69g2be4IHE99tod/q2MsegqyG/Mps6MnNyDHeODrI7Q5H5hZd7vD2j+Y6vEOAB9oSsZGQPxashX302DgAnAkOYP4ns6/S4G3NdVel/ZnU83LuAF4ba/krQTynuh29oNty900yKFv7rY4HHcBk9rux6gAjjA7towxEL1bF8jRymWLWPxL2/KPCvykjhPblrnXAN5CuejfjfTzVgBMxHO+rKMPkzZ6AXK6WBkluK5t2QsBvMPszJfblrlNALMotx3UGhJXG8g5bw7sXD1gD+CtwBvpUdw78HJyPMPbgMOo0aNGNjq5+FEN7U7s9OmtwB519EXkA10U95m8Z/Bc48oG4FzgRbV04Nm2JgDvIVviHuwi+xBwBTU5VxjZVlGEI0q2NQn4Z54bl3kL8KYqnZhqdiAcEk6OpRvJsHQH8NLSwj/bziRy7MLKYB9WUYOThZwo+ktz/L5Xsn/Xj8BzI3Bs2U58yOxAyAIH7AI8FuD3X6UEf7adQ4FfmH2ArHyV6wyQYxYdDGEu3cBZAb5PAHuW6cBXDOEfjQ4a8MmmBmRYGydTrdxcLYcyRv51dkO4MATwQvKWGcFZRXxGmjTHpu7Uxv3zIN04SW82ZJAkAWdK+nf5buvhmFXh2+H4J5P+OIN2puJu68K4wa4KQLbjOwaKmyNE5IQKJ3LWOgeQHVCfcb6po90ipJRuVi4tH8UM4queI+PvF/2haAWYJsnJyY9eY9yg0F9GCYFDJdWVL7eiJj6SdIFBmyT9cZD2UYPvbkV/KFKA6QbzzZLuCtI6S+sTkhZGCMlOlYtVbdl/hp1y/f+68F3ll0iiiI6RowAvoOBWVaQAzgMMy1NKTwVp32DwXWDw/aikuipyz00p3VsTL6WUNkj6ofFJdIxWm6Ls3u0/ixTAMYgsiRCRjTuOYl0b5DtBUqWQ82G4UtLHauI1HKG+dDCNgEW0k7r+tMG36zZQpADO9Su6T0+XV4/gliDdn2mEQ84o2CBplaT5kt4r6dhOMGbdiPZFyrefaUHaxwy+L+n2n0V7plOb5/+CdE4m8FMKriySjjH4StIWSRdKOi+lZJmvK+Ae5UenogfrvSXdGaBzqqB2bbtIARznSHQvclaVpSmlzaOTScrvCEWxRbkU67eMbyojpbQJWK744To6Vs5q1TX6qGgLcHz664N0TnJIaFUh1xNyzJwX9nryh+FXBm00mroxBXDCqTcE6RwPXzSRZIrBU5LOM+nrhJMcE60J6DxTZymAc5+OvqbpGJaie1vXg00Bnuzhnt8N0SutVDBZXeCsANYZwGEcpXVuANFKYU7oeamy8cD+yu8Uvkp5u1soab5xRtkKhz7qjay8BRQpwGIVGA62wZCknwcFcFys0cFyHnII5eRtRce6eL6kv9bvKu9pkhaT8/zaXFGkBreArweZXplSWhukdbaVqAI4K0D44Wdy/t63JH1Y3VeuAyQtpIyfvV5U3gKKFOA7ylaxkfCIpE9UFaAA0clyDqvOYB2rvOyPhMmSzjF4NgFnVe3qru/KoFOt8y8kfU3d9+M7Jc0y37dxsmCik7WjwdM5A8wJ0h1T5GTpEZwtsOuqWrgsN5Cf7vxao5PlPKYQva5K8ZiF8ZJeI8/nXyeaU4Ct6Bx06jjsOMalqPvUsS38xqB1ZG2zhHtlBejlu4HOoEYnqykFqLy39ghjSgGcyXoiSOcYgtYZtI7NYqAAQTiHpWiZdodn1GchbZ8rQFc7SL8qQHSynMclovYKaftUgPZWgE40kGO0iU6WowBO8EQTCrDG4BmlHRsKIP+p1egAOArgTEATChAKcDVpHeNaq1uAE2GE4kEmToxB2yvAtYpF+dypeAyhYwfperXulQK8wqBdn1Ia1clBLgHnvKXn+ONrV4CO4ex4ZRN6ER6RdLxhZKt8tY7m8+1FrhU0g3LVvwoTE7pgpAEaDvdNwVAcPX5iaPgQmFJaopzS9W397p68ufN/Mzs0UVS2hI5oCSTnrH9B0oxh/70GOE/SWUYErRO1+3CQzlEqKZ5I0ZgCSFLHf3IcuSjF1vS7nxte1eFoTgGAOZLm6rkD8jJJn5V0GDA7slzLyzNYGaRztpXHU0pR87Jbus3xyT+DzoTfVOZbSSI/Jeu42Lsa14qSQ6cr57SN9Gs4Url4dASOAjwUpIsErGzFKoPWLS9bKtKoBrgV2LserIsm+BTFfgknARGXrPNAVNTF7ChAdFuRfAVwYv3qhKsAXa/BRQoQTVDcWdKIhaE6hyrnnb1oplET5wrJizGQ2lsBLNN60VZdpACOk2U0QfaQt6+uCNI5V6DotiKNnS3AUdTCK3CRAoTz8gO0+xq8JOmBIN0Kg+f9Bq2jAMgLNasTzhwtLfpDkQL8t8F4tKjgaKKjJD2aUoq6gqMyDkn6viGD47Z+qonHroNYrPgrKoXxnUUKcH6AOcqvWI82AE6xiUJN3RYppdslXR4gPdeMXXT2VifGoFZ0xv20AOndkv7TboD8RExRmfhNwElBPs4DCxeZMu5EfiSiCPMwK5fjVfq8x+HdBICPU1wR7R6qPHTdGeB/ABYBvwaWAhdhPIpIvEgjQESjt+W/A/B+4IfkQo8PAVcBR7u8Ovw+bshrB4OS6yR+BPgS8EXgOCo+bQMcQK5Oej/5VfGbO/2o/GROJZBr8jo4qlWBs8xnGPJeY/I+ke61/R4EjmyqTyOhaW/ggSZ9NM2sSThX4HCYGfAJSV9S9yjiPSVdDbzFaLsWNK0Arzdo18vLoW8KjoUtpADAFEmjvQ4yUdJcwInyqYymFeAgg3Zxi1eq4XAsjFEP4wcUi97ZTyNU9WwCTSvAwQbtHY1J4aEJE7NTdteppFoZjSkAsJtyOlkUP21KFhOOkylqYnaMYT1FkyuAu5Td3ogUBsieTccQNKoCkO0QToU0x8RbGU0qgFMWdo3iZeGahFtvMLIF7C0vJP4XpgyV0KQCOE+g3NYnB0AncAXF4hedF0i2qJ5X0sNoRAHI5sfXGJ+UDo2qGc5krQ6Gwzn7/4qUUk/dy02tAK5F7wa3AfKTMJcBK8gvltxINrFWqRjuKEB0qXYUoKfLf2MArjHMqWvdSQP+FdhSwO9HlHy2Hvi+IXfIcdVRzCjOLSN3FTS1Ajh32etSSuEKXsBHJZ2u4hTugyV9m1zoyYWzAkQPrQ7PsDu8LvQyO7gIP4gSksuonxkgnSXJ8gYCE+VVHh11uQYmy8tf6PkW0JQCLA7SudE6hyvnJUQwWpWvbfEH8uLsI5PlWvVqe6giiqYUwKkz6OTsHWrQTjFoJekQg3azpGUBOscbusaMXKoFTSlAE3UGJemNBq0bruXwvi+YFhcOnFH83aVa0YgCNFFnkBzd4kyS61xyVoBbg3TOFhDdNmtFHa9sdUUDdQZnySvJtihK2DmsOfb6UUPBOodK5wawfSnAVtRYZ/BPDdotkn5s0DtnCylWGHJ/eRU8WlGAfrgGjorOnd55G+inKSXnDDDboH1MMRuAkxAzpBZuANIYUQDl/EMnwXR+lLCjXM7qEnVcWR5A443EWjFWFMB9GWyBQXuIvAetok/AOTkD4fPK8xLAvYY9fRVGYCVwgcEb4A+DfBOwOMjz8NKDs70D2M+coAsN3uPICS9RrDBlPwTYOArPS+xBqRFjYQt4s0n/XYN2f3m1Bq1HpVNKtyr7JLpFD6NsJznB4Vk3Gr8G1gA3WdOJLXAfe7jMpFdKaT4wTdL7lA1ZOyqbka9IKfVLJHT/AjjGWKLnmbz3MngvpZyLua8xFraA6xWvHm4pQErpV4qHo1/cJ3GLzz8ApwZ+oVai5jDeRwCbR+G9jFyWbYA2QL5SnTPCBP0PufBiWf5zgKcLeC8nl80boG0As4DLyYGgjwDXAX9FDQmVwHTgK8B95KvhIuB0wClGNcAAAwwwwAADDDDAAP2P/wdp5W0ZH5ZlJAAAAABJRU5ErkJggg==",
    "flag": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAH2ElEQVR4nO2dW6wdVRmAvx8KlgJCEW8UKEFuBYoaiRAsSMBKuOgTJMYLBBU1oqKJ0WA0JkpifNAQxJgQ4yX4AhGlBlSg0hSpCCjVQoUiWAoFqkILNqVAL58PazdUUtuz98zes89Z/5fM09lrrX/P+vbMmn/mzA9JkiRJkiRJdUTXAWxD3Rd4BzAXmAMcChwE7AfszyuxbgDWA88ATwKPAQ8D9wMPRMSLIw18ktOZAOruwDzgXOA9wAnA7g273QQsBe4AfgcsjoiNDftM2kSdq16pPu3w2aAuUC9S9+v6u1eLGuo56qIRTPr/40X15+p56rSu90k1qGeq93Q48TviSfWb6qFd75+uGdoaQJ0NXAW8f1hjtMAW4FfA1RFxe1dBWNZDJwEnA7OB6cBzwApgUUSs7Cq2vrEc7i9V13f4Cx+E+9VPqnuPcF/trV6urt5FbHeq540qroFRD1BvGuo0DZ916nfUI4a8r+apK/uM7Ub1gGHGNTDqweqKVqeiW7aqt6kfUKe3uJ8OVa9RtwwY10PqrLbiaWUNoO4D3AUc30Z/Y8jzwALgRmBhRKzvp7H6BuC9wAXAOUDTq5C/AqdExAsN+2lNgO8Dn26jr0nAZuBPwH3AQ8BqSmYSYAawL3AgMAs4gpLgessQ4rgqIi5r2kljAdSjgeU0z+JtYyMlm7cM+DuwBlhHyfJBSQ2/jpImPpKSNj4e2KOl8ScLm4E5EfFIk07aSIh8huaTvxq4nnJJdldEvNxPY3Uv4J3AmcDZlHsKY3OfY0hMAy4FvtCkkzaOAKsph7tBWAh8F7glIrY2jWW7mGYB5wMfpIgxVVkVEYc16aCRAOohwOMDNF0CfDkiljQZfyKoxwIfBS6inJunGrMi4qlBG+/WcPDZfX7+OcpknDqKyQeIiL9FxBeBg4ELgbtHMe4IOaxJ46YC9HN9/AJwQkT8OCJsOG7fRMRLEXFtRJxMSbv+DOhrrTEitva2iTKjyWBNBejnFPJSRDzRcLxWiIh7IuIjlCPY14GBD6EtswI4g5J3mCiNTuNNBZjURMSaiPgG5TB6PnAr/f362uIp4HPA3IhYPMqBqxZgGxGxKSJuiIizgMOBrwEPDnlYKYvhi4DDI+J7EbFpF21aJx+MeBURsQq4ArhCPY5yO/tsyrphz4bdPw/cSTnSLOiN1SkpwE6IiOWULOe31BnAiZQk0zGUI8VBlAdW9+o1eZny0OpaykOrTwMrKef2ZcDDbeY72iAFmCC9Gy939LYpQ64BKicFqJwUoHJSgMpJASonBaicFKByUoDKSQEqJwWonBSgclKAykkBKicFqJwUoHJSgMpJASonBaicFKByUoDKSQEqJwWonBSgclKAykkBKicFqJwUoHJSgMpJASonBaicFKByUoDKSQEqJwWonBSgclKAykkBKicFqJwUoHJSgMpJASpnbN4Uqu5LeQ3rXOBoyhu830x5FevM7T66EfgP8G/gSWAV8DDwAHB/lovvj84EsNTLnUepozefUl5tIsWnZlLe0XvMDv62Wf0L8Hvgdkrd3Q3tRDw1GbkAvTdwf4xS0OmNLXc/jfJC5xMp1bReUhcDvwB+GRH/anm8ulHn91HydKO6aMByqW2wWb1VvdARFogeBHVtH99rfpOxRrkInA6cPsLxXs3ulFPNT4E16o/UUzqMZyyo9SpgH+BiYIm6XL1MnbmrRqNAPYz+inE1KsBVqwDbcyxwJfCUeq16mjryqqPqbuqnKIWh99rV57ejUQHpFOAVpgMfBhYDK9TLLYUxh456OvBH4AfAa/tsvrLJ2E0rh86n1L+ZqmylXFJeR7mKWNNWx+o0Sj2izwOnDtjN4xHRb/HO/2FsEkFjym7Au3vb1eq9wC3AIuDefnMMliLX84D3ARcAb2oY3w0N2+cRoAGbgYcoGchHKVnJZyjlcaWcx/ejlKw9HHgb8FaaVx7bfvyxKB/fNhuBpZQqW49SKm89S/nCUFbwB1Iqlh8JzKGkj/cYcZzTgON7Wxdc3XTyYXwEWE05z94E/CEi+qrpq06nZP/OAM6i1PibSFp5srIU+ErXQfSbCdwRt6nnWu4LtBnX69VL1IWWDOBUYrnadO3QDg4uwBJHlIVTD1K/pD7Y0gR0yQLHJGEFDCTAWvViO0i09OJ9l/oT9YV252XoPK9e0sU+2yn2J8AG9eCuYwZQ91c/qy4bwmS1zXXqrK732Q6xPwHWdh3vjlBPUq+x/MrGiSXqvK73z05xCgiwDXWG+iH11+qmIUzoRNiq/kY9o+v9MSGcQgJsj3qg+nH1ZstzDMPmCfXb6lGj/q6jzASui4gDmozXBZay8adRcgynAm8HXtOw283An4GFwM3A3V2VlR+XRNDY0isb/9vehronr2QAjwJmA4dQspN7U9K/UCZ5PeXh1X8C/6A8vLoMWNrrt3NSgD7pZSnv622TnnweoHJSgMpJASonBaicFKByUoDKSQEqJwWonBSgclKAykkBKicFqJwUoHJSgMpJASonBaicFKByUoDKSQEqJwWonBSgclKAykkBKicFqJwUoHKaCrClj8928r9vyc5pKsCqPj77WMOxkiHQSICIeJTyz44T4cYmYyVjinqWumUX///+mNrvO3CTyYL6CfXlnUz+sV3HmAwZ9Tj1h+oj6rPqvepX85efJEmSJEmSJGPGfwGse6ut/pWuDQAAAABJRU5ErkJggg==",
    "folder": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAFmElEQVR4nO3db8jdYxzH8fdXcw9zJ9byt4mMIvcDppUUEWOTJ1KEZP41xQMSyVISpfBACLXWIp55wAP/Q4v5UzIRwtjWmsk2zZ+Nm48H51Zav9855/rd55zrnHN9Xg9/v+u6zrff9Tm/P9fvvjtgZmZmZmZmZmZm4y96MYikRcBi4CRgATCnF+NW+B3YAnwOrI+IXX36HOtE0kJJD0jaqDymJb0l6WpJc3Mfj2JIWiDpKUl/Zpr4KpskXSOpJ2c0qyHpEknbs051e29IOiL3cRpLku6W9E/mCe7GFklTuY/XWJF0b+ZJTfWzpFNzH7dR0PGaKekaYE3/S+m5zcDiiNieu5Bh1jYAkk4EPgEOGkw5PfcKsCwilLuQYbVfh/1PMrqTD3AhcEPuIoZZ7RlA0gXAqwOspV92A1MR8X3uQoZRuzPA7QOror8mgdWSOp3tilR5BpB0NLCJzpeIfW0H3gF2zrKuKgEsAs6i2VLzrRHxWG9LGlOSrk987PpL0p2SJgZQ23GS3mzwaPibWu8srBNJqxMP7rUDrm9C0tsNQrBOvhR0JumDhIP6dqYaF0ra1SAEq3LUO1IkbU44oDdlrHNFgwD8I+lhSZO56h4mdTeBPwOHdTnG0oh4rXclpZH0MrC8QdfdtG5Yt/a2or6aBn4CvgQ+iohvZztg3d10ymvV3KtsNwKf0X1g/zMJXNz7cgZH0jfA88DTEbGlyRgjf0MUEVuBW3PXkckJwCrgW0mPS0r9Eox+AAAi4jngxdx1ZDQB3Ax8IWlZSsexCMCMlbSujyU7HHhJ0h3ddhibAETEj8BVwN+5a8lsP+Chbh93xyYAADNPIyvJf2M6DO6TdGWnRmMVAICIeAZYAfyVu5Yh8JSkE9o1GLsAAETEGuAcYGPeSrKbBzzRrsFYBgAgIt4DpoAHgd8yl5PT+ZLOq9s5tgEAiIhfI+Ju4DjgLmBD5pJyua1uR91S8A7g0C4HvyAiXm9SVQ5q/a3DEuBE4EjggLwVJZlP69I2P7Hf38AxEbGtq9aSdiS8XDk/sRibBUlzJd2j1r/GpVhRNd5YXwLGUUTsjYj7aT3upjizaqMDMKJmHnfXJXQ5uWqjAzDank9oe1TVRgdgtH2X0Pbgqo0OwGhLee9ROdcOQOEcgMI5AIVzAArnABTOASicA1A4B6BwDkDhHIDCOQCFcwAK5wAUzgEonANQOAegcA5A4RyAwjkAhXMACucAFM4BKJwDUDgHoHAOQOEcgMI5AIVzAArnABTOASicA1A4B6BwDkDhHIDCOQCFcwAK5wAUzgEonANQOAegcA5A4RyAwjkAhXMACucAFM4BKJwDUDgHoHAOQOEcgMI5AIVzAApXF4CUHyOa6EUh1siBCW2nqzbWBeCXhIHPSGhrvbU4oe3Oqo11AfghYeDrJc1LaG89IGkSuC6hy/dVG+sC8EnCwEcDayTtn9DHZkHSBLCGml8DrfFpygcsT/xlakn6UNJFM8VZH8z8cvgySR83mJ+LqsaMmg86ANgGHNKgzmlgd4N+1tkkMKdBv53AkRGxd98dlYNFxB5Ja4FbGnzYHODQBv2sf9ZWTT7UnAEAJB0LfAXM7VdVNhB7gJMiYlPVztqFoIj4AXi0X1XZwDxSN/nQ5gwAIOlA4CPglF5XZQOxAVgSEXvqGrRdCo6IP4BLqVlEsKG2A7is3eRDF+8CIuIrYDlpq4OW1y5geUR83alhVy+DIuJ94GzSVggtj43A2RGxvpvGXb8NjIhPgdOAFxoWZv33LHB6RGzotkPS6+CI2BERVwDnAu8mFmf98w6tb/3VEZF0v9b2KaATSVPA5cBSYIpmq1SWbprW2v6rwAsR8VnTgWYVgP+beQdwPLAALx71y17gJ+C7iPgzdzFmZmZmZmZmZmZmZmZmZja0/gUcFq7BhjcYPwAAAABJRU5ErkJggg==",
    "gauge": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAJ7klEQVR4nO2da7BVZRnH/+/hKqJB4QUvKZJJ5QUvQzKJTRhlpWE55mTjaDY2KaVdxtFvTTU5k1njTI3pWM0YmtKkZVcdICwQyuygCYMVh4AE5MBB5CbX8+vDu1OMtfbZz7rsd6193t/Hs9+znv9z2Wuv9V6lSCQSiUQikUgkEolEIpFIJBKJRCKRSCQSiXQgLrSAdgEMkzRO0psljZQ0Rq/7j6StkvZI6pO02Tm3L4TOdtNRBQA4SRMlTZZ0uqRJkk6RdKKkY42Xe0nSWkk9kv4paZmkpZJWOecoSnNoal0AQJekcyS9X9I0SefLf8PLZIukxZIWSpovaalzrr9km6VRuwIARkiaIelySR+RdFRYReqV9BtJj0ia55zbG1iPidoUADBV0jWSPiFpbGA5aWyRNEfS/c65v4QWU3uAw4DrgeeoH0uB6/B3rMpSyTsA8CZJn5d0k6SjA8vJy0uS7pJ0t3Nue2gx/0+lCgAYJelmSbeourf5rPRJ+pak7zvnXg0tplIADrgaeDHgLbtdrAWuwr+yBie4COB0SfdIek9oLW1moaQbnHPLQ4roCmUYGA58XVK3Bl/yJd9v0Q18Fd9LGYQgdwDgnZIekHR2CPsNtsp3AUs+DmMCavmbpE855/7RbsNtLwDgOknfkzSqZFOrJT0nabl8V+6/JW2QtNE5ty1F25HyXcbjJU2QdKp8l/JkSW8tWe9OSbOcc/eXbCcMwAjgvhIfrnqAu4HLgWNK0D8euAL4AbCqRD/upeJ9B2aAo4GnSgjWSuBrwBkBfDoL+Aa+8IpmEVD3/g8PcBrFfmP2AA8C06jAqxT+FfZC4GFgb4F+9gBvD+1fLoApwKaCArINuAM4LrRfaQAnAN8BdhTkcy9wbmi/MoH/hm4rIAi7gTuBt4T2qVXwP3l34e9WedmKHwirD8AFwPYCnP8FcEpof7ICvA34VQFx2AacH9qflgDOxldtHtYCl4b2pSiAy4B1OWPyMnBWaF+aAkwENuZ0dDZ+NLCjAMYCP80Zm/XAyaF9SaTh4As5nNsBXB3aj7IBrgV25ojTciBkr+WhAEOBuTmcWkmAd/lQAGeSr//g98CQ0H68BvDtHM4sBMaF9qHdAOPwHT5Z+WZoHyRJwEygP6MTjwIjQ/sQCvyUt8cyxq4f+HBoB44H+jI6MJsq3cYCgf/5fChjDHsB63qHwoQ74ImMwh+MyX8dYAgwJ2Msfx1K9PVZBQNDg4iuMMAw/MNdFq5pt9hj8R0TVp7GT/yMJACMBrozxLUPaN8CGeAnGUSuo8IDOVUBP6C0IUN8f9wugVOxP/XvoS592RUAP5ayzxjjfmBKO8QtNgoD+FLpwjoM4LYMcf5j2aJmZhA1lwpM3GgH+Fe6M4HpQK45hEAXsCBDvMvpG8C/9j1rFPMKcGIpgipEI/G3cuhA2DPA9BzXnYB9cslfi/TtYDGXGIUAfKEUMRUCGIl/tU1jP3Btjut/JUPcLy7QxdeE/Mko4lk6vLOnkfzHW4jFbmBSRhtDgeeNsf9D0Y6eZxQA8L5CRVQMWk/+/7gnh60ZGeI/uUhnf2Q0/rvCjFcQ7MkHeCGnTetw+71FOXsE9geRes5kbQGyJR9gQ067U432tgOji3D400bDT+Q2WlHInnyAPxdgf77RZv4xAmCe0eiM3EYrCPmSD3BbARo+ZLT5eF6DR+FfY1plBR3Y6UP+5PcARxSgowv4l8HuXgZYTzHQ/gCXSLK8yv2wkzZRlHzyJf1S0gczXmKTpI8WsT9QYz/C+wz/Mkx+K71s4Kdstco+Qs1OKQn8iuZmnTwD0UvBk12B47DdledkNTQU2wKPuUU6GhoqmPyDtC0w6NhCkw65Zj8B50qyLNB41NC20uDX5/9c/icwC5skXeSce744VW/gEUPbsZJSVxQ1K4BpBiNICjM3rWBqkHzJb01rITWXzQrg3QYDK5xzLxraV5KaJF/OudXy2960SupknGYFcI7BwAJD20pSl+QfxJOGtuelfZBYAPjNkiYYDCwxtK0cNUy+JD1laDuRlH6ItDvAu2TbQewZQ9tKUdPkS5Jl4oeT9I6kD9IK4FTDxXdJWmloXxlqnHzJPwNY9hxOzGlaAVhu/yuccwcM7StBzZOvRswtG0sm7rSSVgCWCY21+/bXPfkHYYn9CUl/TCuA8YYLrzW0DU4HJV+yxT4xp2lr9CzLjNYb2galkfzHlG9gZ7pzbllxqnJhmWSSuAdD2h3A0gW82dA2NN9V5yRf8mcUtUpiTtMK4HDDhRM3Xq4a+F03P5fx36uYfEl6xdA2MadpBWDZrHi3oW1ILlW28xGqmnzJFvvEnKYFxDIJpC6vgCdn+J8qJ1+SLMfbJh5KUcSJIXWZArbV2L7qyZcK+KKmFYDl9Mu67G1vWTlbh+RLttgn5jStAHYZLpx/7nl7mK/WxizqknxJOtLQNjGnaQVgebKvxS7ejcmqn1Tzfov1qk/yJdvZiolvDGkFYHm3r83JFs65lZKmSHpIb3yA2tf425QaJV+SLEfjJPYZpPUEbjRcuOzDlArFObdO0lXADfIHQknSMuec5Z26KiT276eQmNO0ArBM76rlfv6NhFsmVVQRS+wTc5r2E7DacOFM694jhXCaoe3qpD+mFYBlmPEYOuWEqxoBHK+UAZ4UEnOaVgArjHo6djl4hbHGPHF/gsQCcM79R7aRJssU8kgxWPZd7HXOJQ4dN+sK7jYYeK+hbaQYLjS0Tc1lswKwzDqdCliGkCM5aEzbt+wK+nTaB80KwPKKNELSRYb2kXx8QCmjeyksSvugWQEskrTfYOTjhraRfHzM0HafpMWZrGDbF/hlBvHxL+0Cf8yM5UTWJ5tdb6D5AJbt3sYo+0zbSOvMlGTZbua3mS0BZxgqDfJuShQZEOybdll6CxMNrjAYOwBYlpVFDACTsJ3V8PeBrtnKlLCHDRq7JH3R0D5i42bZpuD9LLdF/HnAlqrbCVjGqSMtAIwHXjXk4QBw0kDXHfAO4JzrkW0DiFGSbjG0j7TGrZIsb1nznHNrCrEMXGGoPIBdDIKDItoFcBJ+y3kLlxUpYCiwxihgdmECBjnAA8bYr6LosxqALxtF7MePWUdygD+e17IxJMBNZQgZDWw2CrmycCGDDOBKY8x7MRzO2fLKIOfcDkl3GvWPMbaPHIo1hnc45yzrOloHOBzbqZbFH140yAAuNsR7neXbn1XQZ1sU0wccVqqYQQB+8KevxZh/ph2ChgBLWhBzY+liBgnAjS3EewntOqUNf3p42inX/cDtbREyiABuJ71Htpt2b9UPDAdmAYvwT55rgDnABW0VMojAHyo9pxHr3kbsZwHDQ2uLRCKRSCQSiUQikUgkEolEIpFIJBKJRCKRSCQSiUQikUgkEokE5L9ahhK8S27oMQAAAABJRU5ErkJggg==",
    "git-branch": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAKB0lEQVR4nO2da4weVRnH/4+VmnJpMBQoabhFWWIFWghLL9xDoKaEEoxGhDaGGgNC1BiN0S8mRCMJaOXSUuADGDGpXKIGKKF8kGAt0NhQDN1tt6XEKIZuo6np4gfo5eeHM9uuZXc7z5mZd2be9/ySzSa75/KcOf933pnnPOc5UiKRSCQSiUQikUgkEolEIpHofqyOToETJX1R0tzsT29J+p2Z/acOezoFcKak8yT1STpV0vHZvz6QNCxpu6S3zezv9VjYAYBbgT18nD3ArXXbVyaAAZcDjwDvjjPmidgJrAYuA2r5kFZCNvkHJxn4wW4QATAVWA4MOSZ9IrYBtwHH1D2uQgAnMv4n/0j2EL4iWglwFTBYwsQfyQ5gUd3ji4bwicjL8rrt9QJMAx6rYOKP5H7gk2XZ/YmyGsrB3KMXiSpbO8AsSW9I+kYHuvuOpLXAp8torJMC+FRFZWsFOEfSBkkXdLDb6yS9AXymaEOdFEDXkX3yX5Z0Zg3d90l6seidIAkgEmCapBclnVWjGX2S1gBTYhtIAojnAXX2tj8RiyTdF1s5CSAC4Fp15oEvL98FLo2pWNrrRK8ATJW0smAzuyRtzn5L0mkKbz4zC7R5j6QrCtpVHcCjjnfdR+u2dyLw+TPGsg/4FXDJJG3Py8rsj+xjSSevhQu6QAAE3/62iIkZAHL7NoCLgHci+tmCc+0gPQP4uEzSuc46f5K0wMzeylvBzN6UNE/Sa86+Pi9pgadCEoAP70LVoKQbzGyvtyMz+7ekGyW966zqsjE9BB4F4GRJZ0uaJelLjqr7JS2NmfxRzOxfwC2SXlf+2I3rPH0kAYwBOFvSpZL6JV2oELwR62n7jZltLmqTmW0EnpX05ZxVPgucbmb/yFO4pwUAHKfwiVmc/T6jxOYfLrGtFcovAEk6X1ISwHgQAiuul7RUYeKnVdDNbkmbSmxvY9bmKTnL9ym4qY9KzwiAEI93h6Tlyn8hYxk2M8pqzMwANiu4ffNwct62u14AwBxJP1R4gOvUeKvoZ9fRixzihLwFu1YAwGxJP5F0k2qKfm4DXSeA7LXtp5K+Lil6mbQg+yto07NOMJK3YNcIIHOB3i7pZ4p/dSuLUwEr6zkgG9tFjiq78xbsCk8gITTqFUmrVf/kS+Ehs7/E9ubL8WCnsMEkF60XAHCbws6iK+u25QjuLLGtu5zl385bsLUCAI4Ffi3pcR3eYlUlI5I+cpRfCnhu2+MC9Ev6qqPKO2b2Xt7CrRQAcIZCJO6yirrYJekZSd+XdI2kmWY2XUFseZki6RngpFgjgBmSnpJvntbF9lcplBQPAPQD70eslR+NvwA/AuYwwZo6Ya+elw1EiACYkdX14loO7hiUIABgEfBBxEWZiGHgHqAv5xgM2BrRz05gnuNa9Wd1vAzS1M2kFBQAcBPwYcRFGY8B4GuAewMKYaNmDAeBp4H5jDNJBHEtAJ4EDkT2UdVXYnEoIADC5H8UeVHGsh24GYh+9gGOIS4sbCzDwEvAE4QYwHXZ34owQMSewcY7goAvSPqtpCLbo/dKulvSQ2a2r4g9ZrYPuF3B7xB7uz1F+Rd28vItM3N7IBv9FkB4BXpW0tQCzbwgabaZrSg6+aOY2auSHiyjrZJYYWZ/jKnYWAEQXvWek3RcZBP/lbTczG4ws3+WZ9khvqeca+4V87KkH8RWbqoAjpX0B8VvlBiUdLGZPVGeSf+PmR1QcNAMVtVHDoYkfSWzJYqmCuBmhZi8GNYqhGFvK9GecckCPpfI4XsvkSFJi4sm1mqqAGIfTh+VdGORSFwvZrZT0sWSnu9Unwrevvlm5g0Z/xhNFUAM95nZHUVuh7GY2YhC4MkvJZUWCjYBv5B0fetS6uHzA3i5t+7xjUJw5qyvYIwDwNV1jy8aqhNAU/cRLiHs1SvKALCMEhND1QLVCOAFCmTHqBqCe3chsArfZs/twEpgYdU2tllVWyXdUsd3fl6ykLDXsh8Bp+twqtiZOhzHMKKQKnZI0hbPen5roNw7wAjwubrH1A209S3g22a2tW4juoE2CuD5Kj18vUbbBLBX0jfrNqKbaJsA7q5oYadnaZMAdkh6qG4juo02CeDHZa3nJw7TFgEMSnq6biO6kbYI4F4zO1i3Ed1IGwSwW9Kauo3oVtoggMfNzLMlK+GgDQJITp8KaboANplZHeFWPUPTBfD7ug3odpougBfqNqDbabIAhuVIdJCIo8kCWF9mrr3E+DRZABvrNqAXaLIA3qzbgF6gyQIYqNuAXqCpAthrZsN1G5EoSBYWfTnwCLDXEfT5IbCakJOnmSlPEhMDTCWcrjVUQgTwNkJaliIJIhKdAriKkKyobHYAZWfVSJQFMA14rIKJP5L7afs2qW4DmAX8tQOTP8o6Cp6anSgJ4Bzgbx2c/FGGCEmiEwUo9IQNzJL0Z9V3hPp2hUQJe2rqv/UUyZc3TSFJ0lmlWeOnT9IaGrxDuOkUcQQ9IOmCsgwpwCJJ99VtRFuJ+goArlVIT9YkLjOzDXUb0TbcAgCmKqzT50qwPAG7Fdb7R1/nDihkzyxynNt6M7uiQP1EHggevhj2EfLiXjJJ2/OyMvsj+1jSyWvRcxB8+zGJkgeAuY5+LsKXUmWULaS1g+ogLOx4eRWYHtHXScQdmFB5Xp2ehbCq52EgZvLH9DcD/8EJq8occ2IMwLuOidgHxKZ7HdvnPMJhC3nZUcZYE0cAnOmYBIDSdvQQTtrwcHpZfXc7HkfQ+c62H3aWn4wVzvJeW3sWjwA87/27JG1y2jIZG+U4DlXFfBQ9hUcAMxxlN5cZ05+1tdlRxXPMak/jEUDuM+kVvHxls8tR1mNrT9PUqOBEh/AI4ANH2VO9huTAc3zMSAX9dyUeAXgewi6kRJds1pbnIGaPrT2NRwCeRA0zJfU7bZmM+fI92KWkEjnxCMC7VftOZ/nJuMtZPm0rrwJ8fvn9gOe2PVGf/fjO0k2uYAfetwBPFNAUSc8QcXT6KMAMSU/JZ+e62P4SR4GwV8/LhhgREFYCY5aDF1Qx9oQOBYRsjZiUncA8Rz/9+JeBIWxJSwEhVULYqBnDQcKq3vzxJokgrgXAk/i+88eyrI5r0mZigkKPUXjKPrdAv7sVfPvvZzacJmmuigWFDkqaE3OEei8TGxZ+paRXYutXxDWxR6j3MlFrAWb2qqQHS7alCCvS5HcYYAqwNvK7ukzWkbaGRVN0c+h0Sa9Lml2OOW6GFDaHtusg5QZRaDk4O6Z9ierxvQ9JWpwmvwEAJwDPdfC2/xJwYt3jToyB8EywAl8Idww/J33nNxeCM2d9BRM/AFxd9/gSOQGWEPbqlTHxy0iJodoHwb27EFiFb7PndmAlaZ9f5XTUk0fYsXOeQtz+TEnHZ/8aUYgkHpK0xcze66RdiUQikUgkEolEIpFIJBKJRCLR3fwPNeitkxdNG+sAAAAASUVORK5CYII=",
    "globe": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAQa0lEQVR4nO2dfbBV1XnGn4WKRAh+pUKQ4EdHQrUatelICsWQmLTT1hgz02kSq3HiTJomnTo1KoyTWtPWSUJJWohtYhrHDFo67ZiPRtFpMRAMWhOt1qpVtFFB5EtRCQgIXn79Y5073lz2Pns/a69zzr1wfjPMMPesvZ53nbXO2uvzfaU+ffr06dOnT58+Bx2h1wZ0CuDtkmZImi7pnZKmShovaULr/5MljWsl3y1pk6T1knZIeq31/zWSnpL0ZAhhYzft7xYHRAMAjpF0rqS5kmYqVvjEzDLbFBvD/ZJWSFoVQngls0bXGZUNADhU0nmtf++T9C5JY7psxj5JD0taKWm5pBUhhDe6bENjRlUDAE6TdLGkSyVN6q01+/GypNsk3RJCWN1rY+oy4hsA8DZJlylW/Gk9Nqcuj0m6RdJNIYStvTZmVAJMAb4K7GD0sh1YSByQ9qkDcAKwCNjZy5rLzOvAEuCUXn+/IxbgKOAGYG9Pq6qz7AEWA0f2+vseZESMAYDfl/Q1jbyBXafYJGme4oCRXhrS0wYATJf094rTuYOReyR9JoTweK8M6PbcWZIEBGC+4mj5YK18SZoj6WHgKqAnP8auiwK/JGmJpN/utvYIZ5mkS0MIL3VTtKsNAJgjaamk4zsk8YakQzuUdzc01kv6WDcXkrr2CgCukvRD5a/8tZKul3SdOl/5amlcK+mLktZlznuqpJXA5zLn2zta7/u/yTyd2kacV58HjAHeBmxJyOcu4D8SntsEHN3Sng3c2LIpJ4uBnozRsgGMBZZm/FI2A9cxbB4N3JyQ17PAMcTGsy7h+W8Os2ECcDmwPiGvMv4JOKy7tZYJYDxwZ6Yv4jnil/uWAp3ZwD4zvz3AOUPymIW/ADUwNI8heY0FLgGeNvMr427grZ2ur6wQfw33ZSj8DmA+MLZEZwzwUEK+1xTkdW1CPj+lpJsmNoRrgNcS8h3OfcD43PXUEYDDiO/WptwOnFihdVlCvquBQwryOhS4PyG/iytsPJ44XmnK3cDhDaunsxB/kU3f+ZuB36uhdQTwgpn3TtpsyAAzgF1mnuuAcWV5Dsn7AtIGqkO5lR4tGNWCuH3bhB8BU2pqzU/I/+oa+V6TkO8VNW2eBCxPyH8oN9TR6jrE5cxUBoAvUNA1l2hNBLaaGo8Qj5JV5T0WeMzMewswoabthwJ/3SpzKrUaXNcA5pC+jfsa8Lum3udNjX3AbCP/uQnlmGeW4QLSzzzsBWY5eh2DOJdem1iQlzEqpqU3HnjR1FmaUK7bTI1NFExTKzTOAV4ydQZ5nnhcrncQV/n+LbEAG4AzEjT/1NTZBUxL0Pll4kkehz9O0DmVWJkpLKOXg0JgXqLhz1MxxSvROwR4xtT6SoPyfc3UeoqE5VvgJNJXEK9MLV8jgNOIK2ouW4nHu1M0P2xqbSduP6eWcTL+Qk7lFLZE63TiK9HldeBXUsuYBLHrX5Fg7E7Md/4wXXfj5ssZyupObe9soHUOaaegV9HNVwFwcYKRA5ij/WGaJ+NNnXYBjc8YElfydpvlPLGB3ofMcg7y8aZlrWvgROIAzuW6hrrXm3o3ZiqygJtGeFkhzkKOylTktsbdkGDcCmou8pRojsHbst1HxvcicIZZ3mdo0CUTF4tWmZoAi3KVucywk/AXfDZTc3m3je57Tc27c5V5iA0rTRuSxzotvcnARlNzD+brx52yXC3/2NVlIYQN5jPD+QMz/Terk9i4r5SPNhELIWyS9GnzscMkdWZa2GqR7k7ZdzPoHkLsReryIiXnBxracTje/sMGMhzpAn5gaEKso9o9rmPg1XrTo0YddkrKsWkxS9JxRvqlIYQ9GXR/gRDC65L+xXjk7ZL2OzGUwJ8oeiypyzhJf5ZB902AY4mLKg7zM2kvMHVzfOlltswybfliJl1382sHDRbAigxwl3yfJVM3DDxu6nZsQYS4AOas2/93Jt3D8Q+uVp59kOq/Ai4ybV6QoxsGpko61Xjku528bNnK+3vGI2eQYTGq9fpZaD72iaa6kiTg18yWtxFzW7SN9iWm9twcuhU2nWfa9LFMuuPwF+DOqsq3Tg/Q9tBjAQtDCLvMZ8o410i7TVI3rlTdo+hKri7vzSEaQtgt6e/Mx9y6+0WIJ3ydKdg2ah6Nqqm/xtBuPOU07LrdsCvb1W/grcDPDe1NVByDq+oB3i9vCnZbCMH5dZRCPO3iuFRZmUO3A1ozyLRGH0LYruiJrC6TFN3olVLVAD5giEnx2ncu3i3v9vI9GbWrWGWkHSPp1zNq32qmf3+7D6sagDOoWivpx0b6KioHMEPYJunRjNpVPCJvcebsjNorJT1npE/rAYjuV99lCN0aQthnpK/iTCPtg5m129LyCPpfxiP22cc22kj6Z+ORs4Cjyz5s1wPMrfh8OMknYUr4VSPtQ5m16/CgkdYpSx3uMtIeojazqaoGUJcdkh4w0relNXJ1BoCP5NI2+B8j7XTy3vX/ibxXUGldtjPKWVNfHULYa6Sv4kTFrc26PJZRuy7O9G6cpHfkEm6tsjprHjPLPmjXAKYbAiuMtHU42UiLpKcz69dhTUu7Lk6Z6uBMRd9Z9kHZ3fYp8vzt516BO8FI+0IIYWdm/Upac/ItxiMnZjbBmXEdCUwu+qCsByhtMSU8aaavwnEk9VxmbYdnjbS5nWM9YaYv7NHLGsAMI+PNHYicUdhaS1ifWdvheSNtVo/hre/8ReORwjrN0QPk/vVL3vJzL2P5ONr5Dmi8yRojbWGdlm0UOK31ODKewW9ROmotYFYH9OvirPB1wk7ndnBhnRautQPLJP1OikV9Rix3hBDOH/7HslfA6HJL1qcOhXVa1gCmdtCQPr2hcCGqrAEcLIEbDiYK67SsAYxsf3R9Uii80zG6nRH3aUxZA3i9q1b06Qa7i/5Y1gA2d9CQPr2hsE7LGkAvl1f7dIbCZeuyBrC9g4b06Q2FdVq2FOw0gCeU9zCoJJ2v+svRP5WU5Q5eAmcrnl6uwwZJd2TWn6P6G3dWA3A2ObaEEP7ISF8J8WbrhTWT3xdCyHsduibAYtVvAPd24HtyfniFTjrKXgHOLpOzdVwX56CFs3WcG2fTzNm6rYuza1tYp2UNwNninZTr5ssQnB6ol8vWjnbWbevWUW9ni9lqAE4PIEm5PVW+YKQ9KbO2g3POL/fMyv3Onyr6Y2EDCCFslPRzI/PfNI2pwonHNwU4IrN+JcBEeQdX1mY2YY6R9tWW06n9aLcU7LwGct/L/5mRNsg7wZwL99zkM5n1ne+8tEdv1wB+YgjMIa9nrrWSHA8juW/e1MFxeL1L3vnBtrS+aydgxP1lH7RrAM658yOU8QZs6+6dc9bfuUeYC+e+31OZ7y7OlOSEkSuty3YN4EeSHKOTHUGX4Nz2cW4S56Lu/F/Kf3PJOa43oDbX2UsbQOvYsbPCdlHm+2/Ofb9308AXsUvr7qJzINS5R1ilPUaS4xn8oRDCq2UfVlWYc+VrmryRaRUPG2knSjo9o3YVZ8rrgp2r5FXMlXfPsG0dVjWA5YaQJF1ipm/HA/Lu3jkOpZriNPR98q6SV+E6fkp3nE10W77JcEqU20nUk4b293Pp1rBrmWFXNs8l+E6iNlDxamzbA7RG407ItYnyPVy3w9nsmEsXQq0T4/c6vc29GeU/K+/I/tIQwkAjReAso8VBQvy8Ntp/aGq39YeTyaYPmja5ru7LdFMcRVa6+KkctYcQHpbngGmSpE8a6duxQt44YL+bLx3gQ0ZalM993afl7T4+HkLI4zkFPy7wWjKFOgceNXU76Sx6DF7E8iy+i4i/fje45FV18q47b79J3imhaZI+Z6Rvh+N8apq8i6UusyQ54W8cZ07tmCdv6/k1STdn0o4AC80WuBNovFULzDZ1OxZeHfiGaUvj2AXACfjBKxfkKO9wQybjR7tuPDUjhoxxpqIv0ZmQMePwInuuJ0/IGMcvMXQqZExrP/nbpv0XABeYzwzXHZD0HeORY1X/PKHDhZJKHS4W8J2mG0DARyS5oWi/lSFIVzHANPyI2luARv5xgDmmZm6vZQLuMW34jYZ6U4kBsBz2AI6DrSTDFptGQQyC6IabG6o5BnjO0NtHYoDqEv0zzfL+jOaBI1ebmgBfzVXmdsZNxJsKDfJXDXX/0tT7x4xl/rap/RcN9b5k6kGM1HJkrjJXGXhRgoEDgLOIMlzzJLygyruBxp65gHfgvfYGaNANAxcSezCXRoEqUwz9YYKRO4HkA6TAXaaeG2ipSNN95S1roDUTf8oH4O7aNgc4FX9ACHEqlXSGjxha3WEHDaJ2AVPwp75JzrWA04FXTC2IdeAeUM0DcGWCwRDnyPYiEXEw+H+mlhtoaajeP5haT5Aw+CO+3lLGVQA9uRY3aHgAvp9o+AZq7FYVaH7W1NlNWmM7hTitcvhUgs5p+Ov8g9xBB/c+6hbgaLwp2lBewRwTAEfgRTIDcGL+Dup8z9TYADixlQff+S+ZOoOsA451y9URiDF19yYWZCfm7ACYn6BT+xAHfnBIAGvzizjad6OxD7IHeI+j13GAKxILA3HqdD01F4uIR6PcVbJHqXFiCBhLfJc7bAJqHRIlLvJ8ibSp3iCX19HqOsBXGhQK4ophrWVj/DMKUCOiOX6kbqhZIcTl3R8n5D+UL9fR6gnEQeGtDQu4hRobSMBb8AdPu2gzZSIOyHabeT5LjcMvwEfwe63hLKHXg74qiOFm72xYUIjboG1H78ClCfneR8FJ2ZbdDyTk1zY4NHAycbTelOV0YJu7IwATWl90U3YSu+TCXxhxXeDBhHz/vCAvd68B4D8p+UUSzw9ci7+QVMS91BxjjBiA8eTpCSAOsuZR4AcAeA/+gGovQ7ZridvNb5h5DAD73Q8EDgc+RfrcfjjLgdHpvZ044r0p0xcBcXxwHcN2vYBvJeS1FjgWOI64Muny9WE2TAAuJ31Fr4hb6MJdh45CHBguyPilQOxW/xU4n3hk7Bi8Y2OD/Dtpm1obgaOIr6DZwI14N3bqsIi8l20L6dqIkrhmvUDlrulSWad4e2mXpC9kzruMzyvegvq48jup2ivpyhDC4sz5FtLVKQUwSzHwcbYomsN4Q/kb2HAGFOPxdoK1kj4aQij16JGbrrqLDyHcq+jMIXnPvIJOV77Uucr/gaSzu1n5Ug/iBYQQtipe4bpKnh+gA5U9kq6Q9OEQwsvdFu/pqhJwiqQbJH2wl3b0kFWSPhNC+N9eGdDTiCEhhKdDCL+leOEymxetUcBGSZ+QNLeXlS+NkJAxIYTbFV28LFIcBR+o7JH0t5JmhBCWhBCcm88HB8TLJ4vIs4Q6UthNXCvo1OznwIN4F3EhsL2nVdeM7cRFsBEbhm9kby1KIh57+qSiA6peeARN4VFJSyTd3Jr1jFhGfAMYCvG618WKA6hexgkoYqviJdZbQgire21MXUZVAxiEeHTsfZI+oOg37yx1f0A7oOjLcIWiK7YVjR0y9YBR2QCGQwyecK5io5ip6D089z25VxV97t+vWOmr2nngHC0cEA2gCGCyYjib6Yqu3Y+XNEHRw+c0xVfI4DHu3ZI2KW4s7VB0sbJescLXSHoyhNCPpdinT58+ffr06dPnwOD/AV/IfjPKTCW5AAAAAElFTkSuQmCC",
    "graduation-cap": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAKDklEQVR4nO2de6xdRRWHf6u1tKBQKWAhUktbWhEM5dVWsMhDULQQ8EEIJlSpSoImEJXKKyQlUVMTDVBJwAgoiE2BpIAGRSoQoRrkVaBQaQu0tqXQB719kb77+cc6x9tc76Vnn7NnZp9750vuH03P2bNmzTp7fvNYM1Imk8lkMplMJpPJZDKZTCaTyWQymUwmk8lkMplMJpPJZDKZTKbtsdQGxAIYJOlsSWdIGitpuKQhtf9eK2mFpHmSnpD0uJltTWFnpmSAI4FbgfU0zvrad0altj/TJMAJwD3AjgIN35VdwJ+AU1LXJ9MAQD/gPGBOC43eE88Dk4H+qeuZ6QKwT61xXgvQ8F15A7gS2Dd1vfs8wOBaY7wdoeG7sgqYBhyU2g99DmAkLtLeT9DwXdkMzABGpPZLr4dyhF0o6oLx5NR+6lUABpxVc267MBe4kCwYm4e4wi4Ui6m4YKzcTCBwgKRLJU2V9PHE5pTFakm3SZphZutSG7MnlQkAXET9QNIUSR8OXNwySXdKel7ug3HyoPtE4HI3S7pD0s1m9p/AZbUHwPHEE3Yv4d3KgG7sqE8i/SuCHXXBOCGFz5NDfGE3F2/cht54wMSabburZltbQ6ewezWCY7cD9wPjWrD3WODXwJYI9i7CBeOgMn1eCYADapVbEcGRG4FbgGEl2n8oPuu3LoL979bKGrJ3yyoOcAQwnWJLsc3yTs1xBwasz/54IC+LUJ9NeCCHFqblAxxHPGH3MnAZMDBi/QbgEz3PRqhfXTCOj1W/pqFTPMWgEuKJvi4Y6RR28yM4oGVhFwpgLP7W2x7BD/W3XjrBSKewWx6hwnVhV/n+EDgM1yIdEfwSXPd0V8HhuLDrnRUsCdL8QEob+XRXoaOAmcQTdpOBfYJVKBJ4F3kpceY+tgG/A0aXWQEDridO3zYH+CJVEjklUfPjl4EnIvhxG3Bty34E+gO/D2zsduBe4PiSfF15gJOAWYR/m95LK3sSgF8GNG5j7fmVF3ahAEbgfffmgH6+tVnjziTM2PZt4GrgoyX7s20BhgDX4aI3BOcWNciAeSUbMR/4Fr1A2IUCGAh8B/h3yb5fSJGuADitxMIfB75ELxR2oaBzb8LfS2yHxt8CwE0tFrYDHzKeENBPfQJgPPAAsLPFNnmwSKH/bLKQTXjwDA/okz4JMIrWciE6ihRWdN1+JXANWdgFBzgYuAHPTCrC7iKFrCnw4PXAV8l9fFTwVcii0/GNCUE8AbIorwFTiLgu39fAR2eTaG42sdAboJU1/ZX4NGTbLeBUFXx4OIXWkmR2FCnwRy0UVCcLwhah3AmiLUUKPpzyFn/ykLAgdE4RbyqpDQA2FDXijhILr/M34ByyYOwWYBxwH62P+bujWEoa8DGKDzUa5RV6yXp/qxBm1q87iuckAqcDWwMatRyYiieE9imAQcB3KX/evyeaS0rFN2hsCGzcDuBhQm5pqgj4RM4sfNNGTJrPSsbP2nsygpG7gD8Ax5bo80qAT+XeRpj+vRFaT0sHzgYei2DsbuBR4KwSfJ8UYALlLOa0SnnnEuBJk7H2wPeY0l1VCHteYbOUfzAFnUmTMbaIr6yVVdkFJ3zGbjKwIII/tlNs5BDuZBLiJk1uwCdJDg9WoYLgwu5q4pxX+L8kGbxLbpTwR9MQN01sG94NHRO8Yj3XdyTeGDHOK/y/JBmqFgBdnBM9UTRi3U7Egy+GsOsx+5kqB8AeRtZTxXcFcE5XXsTfQB8KUI/Ywm6vmcC0QwDsYezMIG7qniW4Jmn5tDHiCrs6sxu0reUA6NeqgwqwKWJZR0i6WdLbeB9d+LxB4BDgaklLJN0t6VPlmviBrIlVUMwASMFgSVdIegvvho7e2xfwGbtbJC2VNF3SYWFNTEtvD4A6+0i6RNJ84CFgYtcP4KL1QUmL5EGzX2Qbk1C6WCoJFOYU036Szpd0PvC6pOfUeVLoJwOUV3mqGgD3SeqQ9E2F+yUeVfsLzXZJMyV9RNLXI5RXiKp2ARvN7Hvyq92mKaIoKpH1cg0xwswulVSpQ6LrVDUAJElmttbMbpQHwuWS3khsUiMslR96PczMrjWzlYnt+UAqHQB1zGyLmd0u76e/JumZxCZ1xwuSLpY02sxuNrPNqQ1qhLYIgDpmttvMZpvZyZJOlfSwpMYTHsoHSY9IOtPMTjKzWWa2M6E9hWmrANgTM5trZhdIOlrSbyTFvOp1m6S7JH3azM41sycjll0qbRsAdcxsoZldJp/9+4mk9wIWt07SzyQdYWbfNrMFAcuKQtsHQB0zW2VmN8gF4xXyKdyyWCLpSknDzex6M3u3xGcnpdcEQB0ze9/MfiVptKSL5NfCNMtztWeMNrMZ7SLsitDrAqCOme0ys/vNbJyk0+QLOusb+GpH7bOnmtn42jN2hbQ1JTFnAinw2VLtMrOnJD2F7xM4QdJYeVdxkHwqeI38IqmXJM0LpOSLZEE16qvtBZ7ZbXZwzAAoshx8GtC/7F9erWGfrf1FoxZ4pxf4ysYGP7e0wDO71UQxu4BlBT47StIPQxmSgB/LRymNsrSRD9WunnuxwWc+VKD88gFOKbgrZgdwYVKjSwC4mOJ7Bxu+Tg74PHvfbvcmsH/IejZi6ACKX7q0E08ebbt0cvw4l2so3vhrKbinET89pKdcw7eAaix141ewNcNjwJGp7W8UYAzNnwre1Nm++HH+t+Ongq4FnsFzFULfwto4wDE0vzt4K37kzKGp69ET+G0hM2g+83cnMCZ1PYIC3N2kc+psAe6iQteuAp8BfkvrZyncGdv26H0rMFTSq5IOLuFxCyU9IOmPkl4wsygrg0A/SSdKukC+y6eMX+1q+eJS1M0vScQVfnDxwyp3GPqepKcl/UO+Nv9yWVe147d4Hidv9M/Kl6LLvNlzt6RzzewvJT6zIZKpa3zP/fTAxayStFg+CfJO7d8d8mvcN6hzL0F/SQfI9+0dKGmofDv4CPmawtDAdk41s18ELqNbUl+6+FNJ16W0oQLcaGbTUhWefHwNXCXp5+rFC1M9sEv+y78ppRHJA0CSgEmS7lG5/WqVWSvpEjN7NLUhlfjVmdkj8hW6v6a2JQJ/ljS2Co0vVSQAJMnMVpjZOZK+IWlFansCsFzSRWY2qepbxZMD7IevAaxucWKlCqwCrgL2Te3XtgMPhO8DrydtwuZYAFxOxRu+EiJwb+CrgZ+T5wp+RVJVTwrrkDRbLmifNrMiu6CS0BYBsCf4AdNnSjpP0hckpV4lXCxpjnw6+kkzK7JNKzltFwBdwc8Xnihpgnyq9lj5rF4INkh6Rb4L5xlJc82srQVr2wdAd+D3EY+RNFLSMEmHSzpEvgl0sKSB8infPemQZxdtlK8rrJGPRpZLelPSIjNbHsP+TCaTyWQymUwmk8lkMplMJpPJZDKZTCaTyWQymUwmk8lkMplMphD/BUqcGggUizckAAAAAElFTkSuQmCC",
    "handshake": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAALPUlEQVR4nO2de6xdRRWHfwM0FKqAFaEYHqLhIVgFURHEQAkaGwVLQ0CD+ChPCdREE00IhpaIYOKD8AjSQk3BSIhgNGjViEFNFaTFIIjR0kSQVJB3oYW2tPfzjzkHLtdz7tlrHnvvc+58yUmam+mseaw9e9aaNWtLhUKhUCgUCoVCoVAoFAqFQmH0cSH/CdhN0nxJh3X+dL+knzjnnk/VsKbp9PFgSftLequk3SXtKmn7TpHNkp6X9JSkRyWtkbTWObe1hrbtImmepPd2/vQXST91zr2QW7aA04Hn+H+eA07P3oBMAG8DzgNuAf7Vo39V2ATcDXwHmAvslKGdpwLP9JD9DHBqankThZ8OjE0yAGPDpATA3sBFwP2BEz6IjcDtwMnAtATtPa3C+J+WYmx6Cd+N3k/+RJ7DL5+tBTgO+BmwNWZ2jTwBLAb2CGzzLvR+8ifyDP4VkRZggaGzC5I3IAHAHGCldeYSsxH/itjd2PbPGWR8tmq9OxjacNjgIq9yIvBvQ/nc7CXpi5KOarohknaW9GVJC4DFkq6puHF8t0HG4ZJuqlLQogB7GcrO6/wK/dlN0vcknQlc6Jz7Xb+CwExJbzfUvXfVghYFmGUoW6jOuyTdBdwh6RZJayW9Sf4pfp+kI+RNUQt7Vi1Y2Q8ArJF0gLEhhWZY45w7qEpBiwJslH9/FdrPS865GVUKVlIAvEPjpagmFepmZ+fcy4MKVd0DVH6njCAb5N/Lj0l6WtIrnb/PkB+X/Tu/7RppXX/2kHdRT0pVBZhKG8AXJf2q81sp6WHnHJP9B2CG/IZtjqSPd/7dNLOUUAGmwgpwj6RrJd1eZekcj3Nuo6Tfd36LgAMkndX5zUzd0IpUmrOqy9YoK8DdkuY4545yzv3QOvm9cM497Jz7mqT9JF0kf2pYN0UBBrBB0hmSPjSZEyYG59wG59zlkg6StEzSpK+SxFSasxyvgC2SNhrKh7Cj4kzSP0o6xTn3RKL2TIpz7kl5j99SSVernj1CuocWuM1wEHF1MsG923IBkx+JDuJbQGM7dmA74BTgTwFt32wo++Mq7cmxAvzXUNYEcIGkqxQWyYSkLznnsiroIJxzY5Juk3QbcKCkEyUdLekQ+cijXSRtkjc518hHW62WdJ+kT0laXFFU0hVgjUHzzkom+PVtOIfwJ38MuDBHu+oEONvQ539WqbPqUmjxAyRfAYBzJH1fQ/zkJ8Iytml8N8BOxqftA0kEvyZ/yj/5XYAjjf2Pj0nEB0ta2C9BX7uyy+SPo5G5AD5oFDo9QV/L5PcAmG4chyMH1VllD2DZTa53zm0ylO8J5Z3fk87YWmL/B85dagWIdqyUyR+IZYxrV4AoCwA4T3GTv3DEJ1+yjfHwKADwUUnXKG7yrwmVP0TUrgDZfQD4yxI367V7d6b/rqkz+VJiX0BbVoBvy0ewWJlqky+1/BVg3gQCsyWF3CecipMvJd4EVjkMyr0CfEX2eLqkkw/sKOnDko6R9B758Pc9Jb2xI+tFSU/KH878VT5UbKVzbnMK+UaSrgCTQmY3MPAG/F05C2P4U8FogGOAm4AXjG0AWA8sB45O0RZDm+tzB5PZ9QjMM9afZPLxl0RDzuP7sRI4LrZdFdtenzuYzG5g4LvG+i8O7oyX9xbgVqNMCz/CeOs3oA9J3cGD3r253cCWG6//kHS5sf5XAeZIelBSziwan5b0AHBsLgGp3cEpFSBkA7iPoezNzrltATKEz1fwa9UT3LqXpN8An88oI9lGMKUChJwDWDJZPBRQfzeM7AZJ0SlaDEyTtIxEm9UeJDMFm14BskJcDGEsTtJVmZSgthUgtxt4vaHsbEvFDU9+l1xKkMwd3PQKYEkjcwZQKYq5JZPfJYcSjMwe4EFD2QPlr1lNSssmv0tqJUi2Bxj0ROVeAf4gnzCpKosAJ+myiYmV8Hn4vi7pYsVN/t8k3SGffbPbp1nyqVo+IenQwHq7SqAELuz87mBqiAYGZgAbjHIAHgYuAeZ3fouBtQH1jOfPVLDf8V7EeyPkRHszqcMdTE0uR2CZUU5qxoBLgcqxCMD2wDciZQYrAXXMDXYtC4oGBg4FthllpSJ2Ii4kLnI5SDZ2d7D9rgbwSYOAqPvvwA+MHUpBqoOlmMuqMUqw3iDnpH71TGYF1OkE+mqCOiwkiyfo1LFQYXf/Y6yDJBvBVAoQFQ7unHtK0mckZc+1rwyRRA0pQRJTsC0rgJxzd0o6X3mzaGQLI2tACbKvALXfCHbOLZV0rvKsBFslnZ0zhrBmJUjiDm7NCtClowQfU4JbRuN4XNJHnHM3JqyzJzUqQav2AEk3cM6538onUV4maSyiqm2SlkqanSsRVC9qUoK83kCjmXGiWUD1drwTWIotcHM9cD1QKWFyxrZnMxGBkwx19TXTe/rMsecGPtI5d6+hvJlOm46XdKx8KNk+8jn3JX+s/IikB+STNd6V4pZyCojPa9Rz04qP9bvHUFel3MHdyhtLCjGKEJ/rYGGPOvPNETW5gacSJFYCErmD+20CLSZgkqQQo45zbomk8xS+MbxyvBIERAf3nNN+CjDUsYBtpaMEMdbBlcDZ4/4WbQmkUIBa0q02Af5bfRcDq/Df41sL3ACEBoWkMBGX4C2iZ2X7kFT1OQWuNrxbKqUkHTaAQ4BH+vR5C3BuZP2xKW+t9MyckmIPMHIrAP7Lmyvk0733Ypqk64g4To5cCULI9gqwfHpuWFio/pPfJTrQs2YlqP5QY8sN3NNOHWaA1cb+x8b41fE6qJQ7uNugp42VJ7uz3waa6D/5leAxS2NCImxHRgma6j95laD63UrgxkAhI6EETfYfH+6eg+stjTgUb+qEMPRK0GT/gR2ARwNl92MTcLC1IefSQLRrW2iy/6RdBbYBXwhtSPJTrGGC5kK+5wfKnMh9wPGxg1CUIG5jdm2ATIsCPA2cMOF3PLBvykEoShCnBLca5S021B2UOcUMRQlilWBpRTnTsJmhK3L3fXzjihLEKcGiCjIuNdZ5RQ1df10DixLEKcEl9Mh0gn/yLw2oe24TgzAUSkCG8/xOvbFK0M1vcDJx+Q3W01Q4Hi1XAkbvPL8X1b17OYgchGzOIvyT32/yk8nv9L8pttLwvYfxg9AqJcAv+1Ux2+kTZF0b2PdYlqQar2homRJgO88Ho53eQ97tgX0PZR3w5lTjlQRapATYz/Ohop0+icyHAvtuZQs1paY3Q0uUgPCMYYsiZJ4QKNPCGHBmijHKBi1QAsLP8wmVj0+rl9MqeIW2T34XGlYCGjrPB14OlDmI/9DWZb8fNK8EtZ7nA7MCZU3GVmAJMDNmLBqD5pUgxk43ycc7xlLRivwGSaB5JYix0yvJB3bFFsa1Bfg73pxbh7cgVgBXAHMZtdvWNK8EMXb6GNA3oTV+8u801vmLmP4MJTSvBLF2+i/xpt70Tn2z8Mt+SADn+WlGdcigQSUgrZ3+UsT/3Yz/WPbUhOYCLXPb6VVZnnpMhw6aU4JcdnpVtgAHpB7PoYSalYA8drqVb+Yaz6GEGpWAtHZ6CKvwXykvjIcalAC7nZ6adYDlK6lTC+KVILWdnpJ1wCF1judQQnyMXUo7PRWrSXlDZ9QhXaBljJ2e4vtFW4DLKO98OzQfbbscmIc9rAz8xC+nmHpx0JwSvM5OB47AH9Cson98wbPAz4HzmcoevtTQjBL0tdPxN3feAbwfOA44HLCk0ytYqVkJip3eRmpSgmKntxn8+zXX10WLnT4M4C9SWj4bU4Vipw8T+E3YXQkmvtjpwwzFTi9IxU5PQciXrFoJME3SvpJmSpoh/yWxx51zI5fOvlAoFAqFQqFQKBQKhUKhUDDzPxRtK89FcCChAAAAAElFTkSuQmCC",
    "heart": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAKt0lEQVR4nO2de6xdRRWHf0MppTwLCDZgeZVSEbFNRZSCQKqgCGgBjUEwPGJ8BQRj1ATF+AiCJgQDCgGpRGsVVIRWKdpQ2grIqzy0PGpNG5ACBfoE2nLb2/v5x5yr5XTve8+amX32uffO91e77z4zv5lZZ5/ZM2vWkjKZTCaTyWQymUwmk8lkMplMJpPJZDKZTCYzCHF1C9gaYDtJoyXtLWl3STs0/tQlaZ2kVyWtcM5Rj8K+AUZJGiOvf6+mP6+T9LKkF5xzK9utrYzaDADYXtJRko6RdKSkd0s6RP8f9DK6JC2VtEjSI5Lul7TQOdddndptAXaVdLykY+X1T5D0thY/vlpe/2OS7pU0zzm3tgqdHQUwAjgTuBVYSzrWAr8FzgBGVKh/FPA5YA7QlVB/NzAfuBBo1YgGDsBY4GpgVcJOK2MVcBVwUEL9RwC/ADa0QX8XcAtwdCr9tQGMB2bgLbzdbAamA+Mi9E8EZgE9NegHuAc4JuWYtAVgD+CnjUGom03AT4DdDfr3AW4GttQpfCt+D+xf5ZglA/gk8FLNHVbEi8DUFvSfA6yuWWsRr+PnCB315vY/gJ2BX9XbRy0xDdipQP9u+N/eTucvwD6pxi2JNQGHSLpD0uEpymsDiyRNdc4tkyTgMHn9h9aqqnWWSzrdObcwtqBoAwAmS5qlbRc+Op2Vkk6TX3eYKWlUvXLMbJB0lnNuVkwhUQYAnCj/zdnmkTpAWC9pmKQd6xYSSLek85xzM0ILCDYA4MOS/qTqOm+T/ABJ0s7qf4Wwk+iSX7be1Pj/cPlVwpEV1LVF0tnOuVtDPhxkAMCRkubLD0wKnpI0V9LDjX8vc8691lTnKEkHys8zPiBpiqR3Jao/hi2SHpB0t6QH5fW/ULRfAbxdb9V/nLxxxLJJ0inOubsTlNU3wBj8K1UszwPfa0wgQ7WMA74LLE+gx8pi4GIiZuTAnsAXgMcT6FkDjA/V0qrg4cADkUKXAucCKSy/V9cOwPnAs5HaWuEp/J7Ddqn0N9rwEeCRBNpSPZULRV4ZIW4j8C2q3azZEfgO8GZkRxbxBnAJfhezKv0OuIC4hagbqxI3mfA1/X8Cbfu9xm/ePBXRic08SsSeQoD+/YAFEXo/mlrQcODJQDG3UeVjqVzzLsDtEZ3Yy2+Atr8mAtsDPwvUvAxI98YBfCVQyA0k/q006h4G3BSoHfyGVq1r7/hJbgiXpRKwG/BqgICf1915Df3b4Xf2rEzrBP2SBFweoP81UjiXAJcGVH4XFU6WrOCfBDcY9F8HDKtbdy/4yeGMgHH4YWzFI4AVxkqfBfZI1PakAKcDz/Sh/UngE3XrLALYCfvEdg2wS1/l9vmIA86RNN2gs0fSFOfcAsNn2gr+sf5++VW4AxqXn5O0QNLDnepxLEnAJEkPSbI8Xb/onLshtMJ5Rou7KaiiTMvg/R0tPBha0RhsvnBvAKMTtzfTBN4z2epcO7asvL5e0abKtll0vXNuheH+TACN8wPXGD92urki/Ey+VboZKE6LgwC806rlXIJtlxC/CvW6oYI7K2prpgTgD4bx2QAU+lOU/QQcIanP14cmgpwRMlH8znDvSEkTi/5QZgCTDIX3SJptuD+ThjnyziitMqHoYpkBWHbunu6k065DhcZkcJHhI4cVXSwzgNLXhgIeM9ybScsThnsPLrpYZgD7GQpebLg3k5Ylhnv3LbpYZgAWH//nDfdm0rLccG/hmJYZgMWBY43h3kxaLH1fOKZlBmDZbNjU/y2ZirD0faETbpkB9BgK7ph98yGIpe8LXxnLDGCDoeCBdqZuMGHp+8IxLTOA1YaCC2eXmbbwDsO9q4oulhnAS4aC32m4N5MWy0mgwjEtM4BlhoLfa7g3kxZL3y8tulhmAE8bCp6AP7iZaSPAXpLeY/jIM0UXywzgcUPBwySdbLg/k4aT1bdDTzOtL9njvYE3GvabZ5rlZ6IA7kzhD9BXBXMNFWwGLPsHmQiA/bGd0/xrWVl9PUIsXj7bS7rQcH8mjotkWwSy+2sAB2HzCl4H7GmuKGMC2Bubu14Pof6awH2GigCuStzeTBPANcYxmR9T2XnGyrqAvDBUEcDh+NC3Fs6JqXAk8IqxwgXUeCR8sII/5Xy/cSxWEBuRBR9yxcrFidqdaQB8LWAcLk1R8Sj8KVMLGwHLKlWmD/Bh661xj1YBu6USEBIjYHEyAUMYYHdgSUD/fz2liJGEhWC7nTwfCAb/uz8roN+XkjoaGz42XgiXJxUyhACuCOzzj1cl6I5AQRdUImgQg09OFcJtVYral7DET5uBUyoTNsgATiUs3c6r+HjElYr7VIAwgPXAsZWKGwQAHyQ8O5k9DkCgyBsDBa4FsgdRCcCRhOdTvL6dQkcSHuF6JVB4UnUoA0wgPK/iQtodzRQ4kLAAkjQ+d0RbBXcw+NjGoX35CnBA/7VUI/x4wlOovkJ+EvR+8637Lb10AcfV3YDzCM+suQqffWRIAryP8Md+D3Bu3W2QJAGXBTYC/KRnct1taDf48PsxCbS/XXcb3gLhoc3Be7icUHcb2gVwAjavnmaurbsN24Bftw4JaNzLeuCkuttRNcBJjbaG8ms6dX8FH14uJknDRuDUuttRFcBpxKWz+SMdFIW9EHwCJ4vPejNdwBl1tyM1wJmEvzEB/Bmrb39d4JM3WSKNNrMZOKvudqQC+Axha/u93EWFybYqAf8kiPk56GYQ7CICZ0cO/mxqyFmUhIYRWEKaNtMDXFR3O0IBPg9siWj/LAbaN78ZfLqWX0YawVfrbocV4MuEL5AB3ELC5Jq10jCCaRGdAamyYLUB4JuRbZ1Bp8/2reATH1lPtTRzZd3t6I8Eg19rqr1KaRjB1ZEd9OO621EG8IPItl1Hh6SqqxTg+4Opo0hj2D+qux1thfhH5Y10wKOyMfjXRral43/aKgH4RmTH1TpZwk9ub45sw4CZ3FYC8CUG4OtSY/CnR+juAS5pt+6OhAG2YIJf4LotcvBzRJWtIX69vC1LpvgAWjMjdHYD51etc0ACfBp7AIStmUc/uXIj9e0EzIkc/M9WpW9QQPye+d+AXSvQtTO2yGnNDMpt7koAPoYtTmEzD5MwYBX+iPbfI/S8SYdmJe9YiPebexQfTjVWxx7AQxE61gMnpuiTIQf+nNxrEZ3/OLB3RP37AP+IqP8NYErKPhlyEOc7D/AMYM5pAIwGFkXUuxY4uoo+GXIAk/BnCkP5F9BycgV8KNZ/R9S3Gjiqyj4ZcuADJYUeoQIf4qbfhJj4M49LI+p5mXzkrRrwgRJfihic54CD+ih/LPB8RPkvAIUpWjOJAMZHDtJyCg6g4F89X4wo9z/AuDr6JIaO2VO3ABwsaa6kAyOKWSzpEfk+OErSoRFlLZP0IefcsxFl1MKANABJAsbIG0Hd37ol8oNvSeOaSQE+aNXTEY/tWJ4ERtfdD0Ma4hdrQolaZMokhPjlWisLSbDMnEkIPqj1A20Y/PvIMZA7E2BX4nbt+uNeKvQ3yCSgYQQxThtlzM6DP0DAB6q4AltqtTI2488xWLJ0ZToBfOTNBRGDfw8wse52ZCIBjsWfUG4l88lqvK//kIhcNmBXAkPAHyKZJGmipP0l7SnfByslPSefM/kJ59yW2kRmMplMJpPJZDKZTCaTyWQymUxy/gsxrL5ZpHMSuQAAAABJRU5ErkJggg==",
    "home": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAGsUlEQVR4nO3dUahlVRnA8f/ScianCEYnKcYxpzGNHnqZMBiFUUNQGKIXBR+0QSOxoqcoiNKQiUQfjCw0hKQwRXyzR8FxNMfKHkJmVLApw4lJccCR0jszzt+HfS5Nl7l39tprn7PPvuv7wX243HX2+vbe3znn+86+Z20IIYQQQgjVSUMHMCvqR4Crge3AF4BNwPrJnw8DrwF/BZ4CnkwpvTtAmKFv6mb1l+oR23tbvU+9cOj4Q0fq2erd6kLGiV9qQb1r8uoRxkK9RN1fcOKX2qdePPR+hRbUbeqbPZ78RYfVK4fev7AC9WuWveSfzjH11qH3MyyhJvWOKZ74pX6mnjH0fgdAXaM+PMOTv+hxozgclrpBfW6Ak7/oOXXD0MehSuoW9ZUBT/6iA+rnhj4eVXF6lX5Xo+0QRvdRsLoTuB84q2Az/wQeBP4y+X0rcAtwfsE2jwK3ppR+XbCNsBybSn9XD8/W36lnn2L76yZ/K7VLHd0Ta67ZX6W/Yvvm/9rJE4XzRIfQF/UcdU/hCVlQb8yY83r13cI596rnTfPYrHr2U+kfVq/oMPc29Y3CuaND6Mp+Kv2/qZcUxLBFfbkwhtF2CIOxn8/096qf6CGW9erThbHENYQ27O8z/cfssQizKUJ/20NccQ1hOc6o0i+ILzqEaXGASr8g1ugQ+uSAlX5BzNEh9ME5qPQLYo8OoYRzVOkX7EN0CLmc00q/YH+iQ2jLOa/0uzI6hNNzRJV+V0aHcGqOsNLvyugQ/p8jrvS7MjqEhqug0u/KmjsEV1ml35U1dgiu0kq/K2vqEKyg0u/K1d4hWFGl35WrtUOwwkq/K1dbh6DutLzSf1Y9d+h9mRX13Mk+l1iw+Y7EoDtye+FOaPO/+GsH3ZEBqGvt53sIPxxqB75fGPgJ9U4r/hKFTYdwp+UdwndnHfiOwqAX1JtmGvQcU2+y7G30ffXaWQW7wbJK9i11+0yCHRF1++TYdHVIXX/6mcoDvb8gyFeNxZaWpV48OUZd/WLaAW5Sj3YMrqpKvyvLOoQFdeM0g/tJx8CqrPS7sqxD+PE0AzuQGUz1lX5Xdu8QXplWQBdlBnLMqPSL2VxSP5557Ptf4la9ITOI7/UeRKXUH2Qe++vabjvncutFGWPfAO7NGB9Wdg/NiuZtfbbtwJwEyKng96SUFjLGhxVMjuWejIe0XrouJwFyFmXKydbQzpsZY1t3XKP/j5uKOI2NRgJULhKgcpEAlfvQ0AH0QT0TuHTy82kyiqBM7wH/APYCf0opnZjSPDMz6gRQ1wHfAb4JfGrG07+u3gf8PKX03xnP3ZvRvgWolwH7gF3M/uQDbAR+CryofmmA+XsxygRQvwo8CVwwdCzAZmC3+pWhA+lidAmgbgMeAdYMHctJ1gCPqpcOHUiuUSWA+lHgYebr5C9aC5xyJfJ5NqoEoCn45uFlfzmbgW8PHUSO0STApNW7beg4WviWI/pS62gCpenxh6j2c20Evjh0EG2NLQHGYjRt4ZgSYJ7f+5caTaxj+iQwZ3GEl4Bnep7/cqDt17LX9Tz31IwpAXI8k1L6Rp8bVB+gfQKMxpjeAsIURAJULhKgcpEAlYsEqFwkQOUiASoXCVC5SIDKRQJULhKgcpEAlYsEqFwkQOUiASoXCVC5SIDKRQJULhKgcpEAlYsEqFwkQOUiASoXCVC5SIDKRQJULhKgcpEAlYsEqFwkQOWmlQCRWP07M2Ns66Xlc07UfzLGbs4YG9rZkjH2nbYDcxLgXxljL1fPzxgfVqBeAGzLeEjrc5WTAPsyxn4YeEBdrSuQzMzkGP6KvNVc9rcdmJMAzwPvZ4y/BnhC3ZTxmHCSybH7PXB1xsOOA39sOzjrjp7qHprFknIcA/4AvAqUrK+fs0jT0ItElc5/Bs17/mXkr+P0VErpyraDcxPgZuDBzIDCbO1MKT3UdnBuAqwFDgCfzAwqzMZB4DM592zM6tdTSu8Bt+dGFWbmR7k37My+q/dkIeTd5NcCYbp2A1fl3seo023d1Y3AC8B5XR4fencI2JpSOpj7wE4f2aaUXgd2AEe6PD706giwo8vJh4LP7FNKfwa+DPy76zZCsUM0L/svdN1A0UWbSRJsJe/O1qEfu2le9juffOjhqt3k7eAK4OvkXS8I3RwEbqF55nd62T9ZpyJwOeoa4AbgRpouIecSZljecZpPFn8DPJLb6q2k1wQ4mfpxmjtnfJ7mVi8f63HzZwEXAudMfn8L+DtwtMc5hp7/HZpn+37g+ZTS2z1uO4QQQgghhBBCCCGEUJUPAGrCIezrqgP2AAAAAElFTkSuQmCC",
    "image": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAJAUlEQVR4nO2dW6xeRRXHf6tXWo2JwQgGKiAWFJtCL/HShkjUhBexaRtrEZTSlOKN4gUJMeFBTYzxglxMFKKlEFGgWK34IKmJgm15oFYFayJiD6kgYAQjoZWelvP3YfaWpvL17Nl79uXb3/ol56HnzKyZfuv/zcyeWXsNOI7jOI7jOI7jOI7jOI7Tf6wuw5JmAUuBk4E5wAnA9Lra6xnjwL7s53Fgt5kdqqOhpAKQNB34ELACOB+YndL+CPMP4E7gdjP7bUrDyQQg6QLgm8DcVDadV+TnwBVm9ngKY5UFIOl1wF3Ae6p3xynIAeBaM7uuqqFKApB0BvAz4MyqHXFKsQm4zMwOlzVQWgCS5gE7gNeUteEkYSuwyszGy1SeUqaSpOOBn+LO7wLLgG+VrRwtAEkGbAZOL9uok5xPSLqkTMXoKUDSB4G7yzTm1MoB4AwzezKmUtQIIGka8KWYOk5jzAa+GFspdgpYDbwlthGnMdZIOiumwrTIBlZElgc4THhaGCNscTqTMxNYDLwtst5U4FLg88l7JGmWpBcUxy8knZK8MyOCpBWSnov8zMeyhXryzrwvsiP3KawZnApImi/p+cjPfkFR+zFrgDdGlD0MXF5lh8oJmNnDwJWR1QrvzMYI4A0RZXemOqxwALiNcCxclMLTbl0C2BtR1pkEM5sA7o+oMqdowRgBxARz+Go/PS9GlJ1RtGCpswCnP7gARhwXwIgzlM/pkmYQgkxPzn71BPBM2TPxUWYoBCBpKnAu8AHg/QyIO5T0F0LM3FZgu5m91Fgnh5ROTwGSTOH4+c/Ar4DPcOyg07lZmV8DY5LWZ+JxBtBZAUhaBOwmxB6UCT6ZA9wMPBSzNTpqdFIA2bf+AeCcBOYWADslXZzAVu/onAAkXUMIM0/5UslxwO2Srk5osxd0SgCSVgNfoZ5X1gz4aja6OBmdEYCkxcBGanxfMbO9ydcEL9MJAUiaAtwCzGqgudnArVmbI09XPoSLCIu1pjibEN848rQuALUXafxl3yPogAAIOQROLVl3nPJHz2/K2h5puiCAZZHlBXwfONvMZprZTMKQfmv2tzrb7h1dEMAFEWUFrDWzdVmsHBDi5sxsLbCOOBHEtN1LWhVAdqoXs8270cw2DfqjmW0kvDJdlNMVspqMLG2PACcS99x/Q4EyN0bYm5L1YWRpWwAnRZQdB/5YoNzDxC0MY/rQO9oWQMx8bRQfLWJGldiFY69oWwB/jyg7HZhXoNx84iKYY/rQO9oWwNPEfQM/nahMzkTWh5GlVQFkMXx/jaiyRtLaQX+UdBkQkynjsboSMA4LbY8AAPdGlDXge5I2SjonfwtW0gJJtxEOlOpqu5d0QQBbI8sb4R343wEHJR0khI59tIG2e0cXooK3E5JHnFaibpVNnL3Azgr1B5IdcJ0IvDb71b+Ap7oYpdy6AMzsJUnXAj9ouOkvpHSIpKWEDCrnEZ5Wjn4/b1zSI4To5i1m9mCqtqvQhSkA4IdA0iTIk7CLBJnOJE2RdJGkPYSR7LPAQl755cwZwCLgKkKQ6iOSVquObB4RdEIAZiZgPSHVWd3sB9ZlbZZGIRnTDsLIFZWYKWMe8CNgu6TWEm91QgAAZrabsLirc2dOwCVm9odKRqQLgYeAdybo0xJgl6RVCWxF0xkBAJjZ3cA11CMCAVeb2Y8rGZE+BdxB2rD1VwF3SvpkQpuF6JQAAMzsa8Aq0k4HLwIfMbNvVDEiaT3htLGusPWbJG2owfZAOicAADO7h/Ay6O4E5nYB7zKzO6oYyZz/XeoPW7++SRF0UgDwvzXBYsJo8FgJE/uAy4F3mNnvq/SlIefnNCqC1vcBjkW2Ut8saQthsbSMEMY1l/93xgRBKPcSdvh2pnjOz+b8uob9QeQimDCzb9fZUKcFkJM58jfZz1VZGNcJvJwN62/A06nzEtY850+GATdKmmJmMVFOUQyFAI4mO8F7IvuphYaH/UHkI0FtDXR2DdAm2bBfxflbgHcDryY84p1HuGGlDAZcT7FgmGhcAEdRcc4XcKWZrTSzB8xsv5kdMLP7zWw5IVilzNfZqOklFhfAEVSc83PnD5yvzewG4GN0KA7RBZBRcc7PnX/TZAXN7BY6JAIXAM05P6dLIhh5ATTt/JyuiGCkBdCW83O6IIKRFUDbzs9pWwQjKYCuOD+nTRGMnAC65vyctkQwUgKouMMnYEMdzs/JRLCBBkUwMgJIsMO3oe6TOYCsjaoi8BtDjmRYnJ+TQASF37HovQCGzfk5FUVwfNGCvRbAsDo/J9F0cEx6K4Bhd37OESKI4dmiBXspgL44Pyfry46IKmNFC/ZOAH1z/hEUyY+UUzhHUq8E0GPn10ZvBODOL0cvBJAokmfknA89EEBX9/aHhaEWgDu/OkMrAHd+GoZSAO78dAydANz5aRkqAbjz0zM0AnDn18NQCMCdXx+dF4A7v146LQB3fv10VgDu/GbopADc+c3ROQG485ulUwJw5zdPjABibuI6LrYj7vxJKRzrDxS+BSVGADEJmc6NuZ7dnV+IN0eUfSp565I+rDjWFLS7XtJEpO2cCUlXJP/PdgxJp0k6FPG5XFpHJxZGOud5SfMnsenOnwRJ0yRti/xs3ltHR0zSWGRHnpO0fIA9d/4kKHzzfxn52bwgqfAaLGrOlfR1wo0XsewhJG0+mP37rYRk0GXZSVyY9LAxnXCp9hLik3luMbOVRQvHCuAswt28UyM75TTHxTGZ0aP2AczsT8Rdz+40y6NE3oUU/dgl6aSsoZQ3ZjhpWGlmW2IqRO8EmtmTQONXmziTsg34SWyl0pmwJX2HkNPGaZ8x4O1m9s/YilUEMAO4h3CBg9Me/waWZOuzaEofBmU3fy8Hbi5rw6nMXmBpWedDxdPA7CaPjwOfI1zI6DTHNsKwv6eKkcrHwWYmM7sOOBPYXNWeMymPEi7SOt/MCmcCGUTy61AkLSRc5X4h8PrU9keU/cB9hJtI7kp5N1Jt9+EoXKG+CDgVOIVwwVN+pj2DkMosz2b1LGElGxNz0GcOAc8QLsPaBzxoZv9pt0uO4ziO4ziO4ziO4ziO4wwv/wXgFOifIhh2EgAAAABJRU5ErkJggg==",
    "info": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAMeUlEQVR4nO2deazdRRXHv8MuD9pS1FaKBYyUpVaCRqhCS4qARGU1CkpYNQYxkbAoDSGKEZSlAoUCwQSVFqs2YpD1j9ZCbSUQEqpBloIJhRa6QBfoe4W+8vj4x7m1r+Vu57fN7973+yTNa95v7ptzzpw7v5kzM2ekioqKioqKioqKIUeILUBeAJ+QdLCkcZIOkrSvpB5Je9R+7lX7KUl9ktbVfvbWfi6XtETSS5JeDCGsKFL+ougKBwBGSjpG0hRJE2UNPizjat6WOcOTkuZLWhBCWJdxHYXTkQ4A7CTpuNq/YyUdJmmHgsX4QNJiSY9Jmitpfgjh/YJlSE1HOQAwXtLZks6TNCquNB9iraS/SJoVQlgUW5h2Kb0DAB+V9F1Zw4+PLE67/EfSLEl3hxDWxBamIwH2AW4CeulcNgDTsAFpRTsA+wHTgY0xWy5jNgEzgQNj27e0ACOAGcDmqE2VL/3ArcDw2PbeQinGAMA3Jd2m8g3s8mKlpCtkA0ZiChLVAYBxkm6XTeeGIv+QdFEI4blYAhQ9d5YkAQGYKhstD9XGl6TJkhYDPwaifBkLrxT4mKSZkk4sqMpVkl6QRfFek7RR24Z+pW1Dwz2SxspCyIdI+nhBcj4s6bwQwlsF1SepYAcAJkuaLWlMTlX0Slooi84tlMXw16f5g8AImSNMkoWaJ2nrGkLWLJf07U4KJLVNrZvLY4S/FLgG+BIWIs5bj52Bo4FfAq/moM9m4LK89SgM7H1/Y8ZGehubVx8HRBnH1HTbAXOGu2oyZcmtMXXLBGAXYHaGRlkFXE2J5tFbAPYALgaWZ6jvH4CdY+uWCKAHeCQjQyzFjPuR2Hq1AnP6c4CXM9J9HrBnbL1cYN+GJzJQvheYCuwSWycvmCNcCfRlYIcngLwGntmCDZIezUDpB4H9Y+uTFmAMNl5Jyzxg19j6NAUbFKV9568Cvh5bl6wBTgFWp7TNvUQKGLUFtnybhseBfWLrkRfAKGBuShvNiK1HXbB5flIGgJ8DO8bWI2+AnbC4xUAKe10aW49tACaTPMjTB3wttg5Fg70Sku552AwcFVsHSbYjl+QRsbXA0bF1iAVwJPBWQtstw7bLRVUgAH9LqMAbwGejKlACgENrjZmEh4k5KASuSCj4MrpgipcVwAEkjyBeHkvo8dgWJy9rsO3dFYMAJmCvRC+bgEOKFjYA8xMIu5Eh/M5vBTYmSLILegFFvgqAsxMIOcAQHO17AU4m2RTxO0UJOAwbwHm5uhABuwDg2gT2XYltXslduBkJhJvPEAjyZAUWLFqQwM7T8xbsAPwBn1V0cXg3L4DRwAqnrfvJc3YF3OkUCEq4sAOciq02vl4z8qPAGZRsoQWLFnrJZ60A88h3ncL8NRdhEgLsCsxpIu9DwO6x5RwM8IDT5u+SR4+Lf6Wvj5IFe4Db25B7Vmw5BwOMxT81vDFrIfbGTrp6mJqpECkBxtH+9Orw2PIOBrjKafte7PxFS9rddfo9WW6ddlkq6SZH+SI4Se3re2qegiTgRknLHOV7JJ3fTsF2DXKWo3JJuiGE0O/8TN7sn1PZ3AkhbJI0zfmxc9sp1NIBgM9LmuCoeKWk3zvKF4XnhFCq00Q58RtJnkxlh7bzKmunBzjbUakkTQshvOv8TBEsyKlsIYQQ3pN0i/Nj3rbbFmyH7yrH4ONtwDNWKAxsAevpNnR4ngKOmCUB2BN4x9EeK1PpApzoqAzg7gz1zRzg01jwpxGrAc/rrnCA3zrb5IRmf6/VK+B4p3wzneULJYTwX0lHSPqjpM2DHg1Iuk/SESGEZ2PI5uBeZ/kvJ64JeMbhaUvpoMOMwHDgKGASsFdsedoFe5W94miXp5NWNBLfuvQ1Geta0QDsaHq7vE8TB2/2jZ3S4vn2POIoW5GORx1ld5TlUa5LKwdol15JybqaiiQ8pa3pbdqhYVs2c4AjHRUsCiFsbl2sIgtqUVZPGpmJjR40c4BxjgrmO8pWZMNjjrIHNXpQ1wGw9WRPvv3uS2pUfhY6yg4HRtd70KgHaOgxDXjRWb4iPS84y9ft0Rs5wMGOP7yqG27O6DRqNn/T8ZG6bZpFD1B9++OxxFG2bps2Wijw5Lf3CFEKgGGyTR+fq/3qGUn3hxDeiSdVIpZIavekVd02beQAnhW95Y6y0QG+JelOSSO3e7QW+EEIYU4EsZLi2SVUN9NYIwfwpCXb4CgbFeAM2UJQve3fIyX9CQghhD8XK1liPLav26aNxgBd5wC1bv8ONc+PHCTdUSvbCZTCAXodZWNymj7c7ddjpMq3KbQRuTmAZwzQKQ5wWE5lY5KbA3QjnkybpdzWlgeNHMDzrR4yxiohqcdqjRzA07VUDhCPUjhAZ2Wy7i4qBxjilMIB9nWUrciWTzrKuhzAcwTJu3RckR0e279R75eNHMCzwONZOq7IFo8D1G3TRg7gWeIdRRHZqSq2obbVu60cADVcDuBd4i02U2WF5Lf5S/V+WdcBQggrJHnWxic5halIz2RH2fUhhJX1HjQLBXteA54zBBXZ4LF5wx69mQM85ahgMh14s1enUrO158KIJxs9aOYAnn3nu0v6gqN8RTomyre41bAtmznA45I+cFRSJYIujq86yg6oScaThg5Q23b8L0dFZ9FBx8M7lZqNPZnBn2l2g3qrBvMc+Ror38i0IhlT5AsBN23DVg4w11GRJJ3jLF/hx5v4aV6zh+30AKsclX2DkiaJ6gawS6RPd3xkhVoM5ps6QAjhfUmzHRUOk3Sho3yFjx/KtwQ8O4Qw0KxAO4M2b/Lky+mAa947DWA3ST9yfqxl27V0gBDCYkmezFmjJF3gKF8UngQWZUx2caF8R/aeCyH8u1Whdqdt3l7gJ5TvqvNXciqbO7Vv/2XOj92TpQAj8WWoBLgyMwEygPbTxQ8AnuwouQP8zGn7XrK+VhaY5hRiI3BApkKkhPYujLg9tpyDAfbDLt/wcEMegozGf9v1/ZkLkgJaXxkzh5K9urC7jTzkc2VMTZg7nMIAnJKLMClg20ujXq/9v3TnAYHTE9j7tjwFGovdVethNTAmN6G6FGBf4E2nrfuB/fIW7FanUGCXIJYyBXsZwS6OXJTAzvlf04NdHdss5XojfpG7cF0CcF0C+64Ahhcl4FkJBBwATi5EwA4GOA34IIF9zyxa0L8nEHIjUG0gbQAwEf+UD8C7apuJsIfiHxACrAU+U7jAJQeYAKxLYM9NQJzTWcDlCQQGWE7JgkQxwS7lTjKuArgkpuABuD+h4G8AnZKKJTeA8cCyhDZ8iNgXXgN7YdfFJGEdQ3hMgL3z30pou9eAvWPrIEnC7t7ZnFCRjQzB2QE22vfexr6FfuCLsXXYBuDShMqATRGvZQgEi7Agz3Ukm+pt4eLYetQF+HUKpcAihl0bNsbCuwtT2uj62Ho0BBsU3ptSwdWUcAEpLdjCjje2vz0ziT3oawV23ewjKRUFW6Hr+Kki8ClstJ6WuXTK+UtgD+CJDJTeCFxFydbo2wHYDfgp/j0U9fgn4DkLGB+gh2x6ArBLkK8Ado+tVyuwTSffJ/ncfnvmYucBOg9sxHt3RoYAGx9cTVGrXg6wXu9ikkf06jEL2Dm2bqnABoY3ZGgUsG51DnASsGNE3XYAjgbuwr9pthXT6abDtsAlJA8WNeNV4FdYQ+T+bQF2wS6cvp7suvnB9APeAyCJKXRKARwlu7HDc7rVQ5/sDsPHZPfqvZD2RjMsG9chsjxIx8ru6MlrLPKqpDNDCA0zemRN4XNKLH59j4pLKPGmLN/REtkdO32yrJnrtfX+3R5JI2Tn7npkDnqQLAeiJxVbGh6QdH4IYW1B9UmK4ACSjQtkJ12uldQZc9v86Jc0VdItIQSKrjxqVAk4UNIMSSfElCMiCyRdFEJ4PpYAUUeZIYSXQwhfkXSyfFegdTorJJ0raUrMxpdKcmVMCOFBSRMkTVc5T+ZmRb+kmyUdHEKYGaPLLz3Y4ZPpZBNCLQvvYbGCvGY/3Qd2FnEasCFq06VjAxYEGxXbno0o99Ki/j9tvECWgKpTdhM/K2mmpN+FENbEFqYZpXeAwQDjZVmyzpU0OrI427NG0n2SZoUQFsUWpl06ygG2gG0dO1bS8bK8eYer+AHtgKTFskxq8yTNb5WQqYx0pANsDxauPUbmFBMljZOU9YrhelnO/Sdljb6gWQbOTqErHKAewGhZKHecLKw7RnbHYY8s5DtCW+887JU18Iba//skLZc1+BJJL4YQPPkSKyoqKioqKioqKsrL/wBPDooeMlQTGAAAAABJRU5ErkJggg==",
    "key": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAJQ0lEQVR4nO2dfcieVR3Hv9ec29Rn01KzTG00fFm6Ek0oEsPlCtLCBhUhhGIopb0gVAS9iJgVvYCUlsQIBmlYijZJc7QN+6PBsiKV9gQ65zY3n2eYe3Fzz+bz6Y9z33X3tOt+7vN2nXPd9/nAw2Bc1znf8/v9znVf5+13SYVCoVAoFAqFQiF/gKuBNcCOzt8a4OrUugqRARYA91PP/cD81DoLEeg4/7E+zu9yV2qthcBYOB/gdeCcQcueE1N4wR9ggaSHJH1owFvmSPrIoOWXAMgYB+d3WTzohSUAMsXD+ZL0SmA5hSax/M0/GlekbkPBkQDO3wRUqdtRcCCA8yeAs1O3o+AAMB94xNP5y1K3o+AAMAZsLs4fQTrOf6k4fwTBPPZLzx9FKL/5o0vH+WuK80eQAM4HuCN1OwoOBHI+wDRwc+r2FCwI6PwSBG0jgvNLELSFiM4vQZA7DTi/Nwi+kLq9hR4adH4JgtxI4PwSBLkQwPkTwB0dZ7oGwU2p7TCSBHL+sk5ZN3sGwcrU9hgpQjq/p0yfIHgZeHMqe4wUMZzfU7ZPENzdtC1GjpjO76nDNQgOACc2ZYuRo+P89TGd31OXaxB8MrYdRhJgEbC1Cef31OkSBD+KZYORBbN7d2eTzu+p2zYIHgzd/pEG89j/cwrnd+pfADxjUd/v+pVXjoZZgDl7/xtJFzsWMSnpA1VVPeVY/wJJD0t6h8Vt/3KpqzAD/N/2Xw7Q83/vUO83+5VbngAD0NPzr/ItyrH+bs//oMPtG1zqLHQI0PN7mQAusKzftecDbAGOiWWboYc4q3oDBwF+zge4MbaNhhbiLunOGgT4O38jpfe7QTPr+bVBgL/zdwJva9puQwHNbub4vyDA3/kvzSyzMCCk2cnznyAggfOzyyQBnCZpTNJJnX8lab9M3pt9VVVNRKrXd6g3JWme472Tkq6UdLvchnqSNCEzyfS0zU3JAgCYI2mZpMslvVvSuZLOkbRollv3ShqX9E9Jm2TGuX+vqsppjN3R4uv8SUkflnGgS1InycwRuPrDyfmNA8wFrgJ+Bez2enD+L5PAvcCVwFxLTSG3cc3zLMur/mwBzgJ+iF8ihEHZBfwAOHMAXTG2cTUZBHk7HzgbWAUcasggvRwCfg4sqdEWcxtXE0GQr/OB44FbSeP4mRwG7gQW9uhrYhuXb8avfuQ71MPkst8WqeE+vAB8lAaTMxDnSZBnz8cY9s7AjY2BT3C6bOMKGQTZOn8J8NdAjcwVn21cIYIgW+dfgPlsyTDjbXzMO8HDjvVvBZaG8lkwgMuBPYGMnCvWa/h97DUH+Bpw0KL+B4FTQ9QfFOB9mIMHMdkP7ItcRz+COX+G7c7EzIvU7Sw+CDwAXBa67pk4TT12jPKEpDcE0jEuab2kP0naLGlLVVWTM+o8RdLbJS2V9F6ZKeSBP43iwKSk5TGnVzFZvc+TacepMmseW2Smtg/GqtcL4C3A9gC9aztwO3Cuh5almKPUod9BovT81gMcA6zzNO5zwGcA15Wzo+maD9wIPO+pDYrz68HM7rkyBXwXs8M1lr7jOhpf89D541j6Wg1wCXDE0aibgXc2qPVCYNxR6xHgoqa0tgLMo3+To0F/C4zNXktwzQuBPzhq3ojZr1CQJOBzjob8BZbr84F1z+0404UbUunOCszK3oSDAe8lg14EVMAfHfTvAo5LrT85wJccjLcWODa19i6YufgNDu34fGrtSQGOxT4Jwk4yTE4EvAn7+YJtBByutg7gU5YGmwaWp9ZdB7DCsj0wyilWgEctjbU6tebZAO6zbNMjqTUnATPle9jCUK+S4aN/JsDp2K3GHcacVRhK+r2lr5RkM4RbVVXVLk890amq6kVJqyxumStji9ECsw49KK/TosOIwOKO5kH5dWrNjYLZuDBpYaDHU2u2BbtFrd1kMKcRg7pGLZN0ikU5vwygpWlsNJ8s6fxYQlLSLwBsWOsrJAGPWV4/UgFgs9NmvPNi1Sqqqtohc8B0UM6LpSUldQFg09i/hBCSiL9ZXOu8cyln6gLgrRZl2PSi3Nhsce2sh03bSF0ALKz5/6PxfAAdqdhqca2NTVpDXQDMlqShl30hhCRir8W1IxUAJ1iU8WoIIYnYb3Ft47uamqAuAA5ZlDE/hJBE2GxQtbFJa6gLgFHpGTaPdZufi9ZQFwA2jc1+BbAPNtrb/K5TS10ATNb8/9Fo8/jYRruNTVpDXQDYjO1tPl6QGzbTu22e76ilLgDGLcq4BLAZNWQB5qyCzZc/bCaNWkNdAPzDoox5ki4NoKVpLpNks3N5pAJgk+yGPR8PoKVpPmFx7ZSkJ2MJyRLgCYsNE3uA41NrHhTgBGCvRfs2pNYci367XNZblLNI0rV+UhrletnNAayLJSRbgIstegiYs/nZnAaqA3NK6AXLtl2YWncSgKcsDXVLas2zAXzFsk15Z+COCfBVS2PtBWz2EjQKJjmTbdKpL6fWnQzM4RDbbBtryXAHLW7pbQ7SgsMuUQF+amk0gG+l1j0T4DaHdvwkte7kYHL9T1kabhq4LrX2LsA12H92fQpYnFp7FgD3WBqva8CPZaB9JfYBDHBXau3ZALwRtywhR0j49UrgOuwOuHbZDZycSneWYHL7ubKaBtOt4J+6/tqmtLYGzHnBxz2M+iQNpF/DpLPzSV3/KCaFa2EmwGnAix7GPQLcDZweQdsZwM+wO/U7k+3kmJk7JzAp4l0TRnZ5reMs7ycCZsr6Hvy/T3SYBrJzDwXADZ7G7uVpTMLnFfR82KlP3WOda78DPBNIwzRwfRO2yw3n3zrgG5JuC6ily3ZJz8rsweueORiTOa6+RNIZEer8elVV345Q7nADfD9QD0zJ91LbsdUAX8TvxSsV08Ctqe03FACfxi7zVmoOANektttQAZxPuJeymIwD70ptr6EEk6J9dWIH1zGNyV7e5qNs7QB4P2Z4lwvjwIrUdhkpMHvvbsFv5tCXHZiX1Oz3Kg4tmC9m3oT5WFRTPAt8Fmjz0fXhArOYdClm2jbG10b3YN4/riDDLWm5kmTVC/PlsPdIWt75u0iS7XLxAZnTOutkzjBsrKpqKJM4xCSLZU/M8utZMvkJl0g6UdJJ+m/yif2SXun8PSdzeHVbVVU0r7ZQKBQKhUKhUCgUCoVCoVAoFNrGvwErSbta+6aJvAAAAABJRU5ErkJggg==",
    "layers": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAALgklEQVR4nO2de7BXVRXH17k+wDcCY5QpDzEJU1DwkckEiqJZ0UzR9Bp8VDTTH1KGlUXTY6YpR8AHltnojCE9hmmcmmaiJCVHjRBIwTQUCBTkjYE3BO713k9/rN/p3rne372/fc7ea5/f73c+M/zF3LPXWvt3zvmevddaW6SkpKSkpKSkpKTpSGIbYAUwUESuEpEpIjJORIaLyODKf+8VkW0i8qyIPC4ijyVJcjiGnSWeAUYD9wL7qZ39lb8ZHdv+kowAFwKLgHaHie9JB/AH4LLY/pTUANACfARYlmPSq7EamAkcFdvPkh4Ax1Ym54UAE9+TjcBs4PjYfjc9wCmVyXjNYOJ7shv4HjAkdhyaDmAUKtIORpj4nvwXWAiMih2Xhgc/wi4UqWB8f+w4NRRAAkytBLdeeAqYQSkYs4OtsAvFBlSjHBc7ntUo3EogcLKI3Cgit4rI6ZHN8cVuEblPRBYmSbIvtjHdKcwPABgpIl8VkZtE5ITAw70qIg+KyGrRGFwk+qM7M/C4Byvj3pUkyebAY9UHwAXYCbvn0NfKMb3YkS4irTSwIxWMl8aIeXSwF3ZPoZNb0xMPuLxiW2fRbKtr6BJ2/zQIbBuwBLgoh73nAfcDhwzsLbxgzAxwcsW5bQaBfAO4GzjDo/3D0FW/1w3s31UZa3D/lhUcYATwY9y2YrOyoxK4UwP6cxL6Q37FwJ9W9Olzdih/ggGMx07YrQVmAQMM/TsGXeh5xsC/VDBeYuVfZugSTxYUQjzR7IKRLmH3vEEAcgu7UADno0+9NoM4vIy+igbGdDgVdlsNHE6FXeiFmtzQJRj/YxCXnVgLRmA4KuwsHAwu7EJBl2B81SBOregNMjykQ2OAX2En7GYCxwZzyAj0FXk9sM4gbm3AYmCMTwcS4DbgiIEDy4BpFEnkeKISx2uAvxjEsQ2YC7TkNboF+IWBsYuBCzzFuvCgiS0WT9NfA0fnMfSOgMa9AcynDoRdKFA9dWclFqH4aVbjJhPm2/Y14BvAIM/xrFuAQcA3ge0B4g0wPYtRazwb8TxwAw0g7EKBCsYb8b9Z9jIu6Wno6pYvHgOupQGFXShQwfgh4HGP8/BRFwPm5RysHRU5FwaMU1MATAB+Q37BuMhl0CczDtKKippwCxJNCrqzenclxll4yWUw1+3O7aiIKYVdYIBTgW+hq6QuHHEZZKfDhQ8AnyTvokOJE8Akx3mCWoUg8JLjhan8zSxi7lA1OHQlrj6RYX5aXQb6XYYBUnYB36EsmvQGMAD4AvCvHPOy2mXA2TkGSjlIWTSZC2Aw8G3cH/W9cYfLwO/EX2LDW2gSx8UBY9VQoNXPC9HKY1+MdzXifo+DpzwBfJhyUahXgIvRm+Utz3H/YxZjhuLn0dMbLwKfxzCRs6iQT9jVwkHgPVmNm0TYoogd6Ddt3WX85AU/wq4/OoBP5TV0KuHz+luBu4ARfsJbXPAr7PriCHCDL6NHEabzVk/a0SSGCV4MLxCEEXbVWEeIzGngCuBP2OTAN8QuIuGEXW88B3yO0J1JsC2ajJ8D7wgFr3726ahl0aR9Drwj2La1KU6RDPZFk4UqEqGrX6Fl9XNh/P8/xCmajLbCCIysTIaFsKuvIhkauGgS27Y25tXPXsG2aDINlnfBSLMIu1BgWzTp7XFJswq7UGBbNJm5XQylsAsLBe0rQCnsbEHfrdfiNwe+Gp3AUuDKXuyIsWL3tn6F1vQrMNBS4+kiMlJE9ovIX0Xk0SRJOn0bg+4BzBGRT4hI9qLG2tgoIn8XkU4RmSAi5wYeT0TkURGZlyTJMt8XRpNyrxaRySIySEQ2i8jvkyRZn/WCA4Cf0/tn3LPAOd6sf/vYI9Ddwaw58EWiDf0KGhcwXuegc9KTTnQO3T4h0cfyI/04thOPvfiq2JE1B74IHECrq98dOEZn0P/W8iO4fE6iq3m1sCSgb93tGYBmEL2YcTIs2QrMQbueW8RmSY12zXC56NIaL9qOoYJFn0wh06fyYC7s0CdkrauRS10u7FKnPjWgj33ZaKna++LPwFWRYjDFwc7tLhfe43DhQ2huQDBR2I+tMQ6QCi7s+vE5XatwEck7XQbI0qCgA60oujyg733ZPAStSNqVwfZaMRF2ffg4keyl4mtdBlqcM1ArgI8ToWAUGIhuCGWpb6yGqbDr4U8CXAcsz+nDwy6DfjbnYCkbgS8T4ZRNNN/+Y+iuWlairdjR9eXjaxPq0y6DHw/s9TQwqKb4AXBawJj15c8lwAM1+nQA+CVwRSRbQ6x97KHKTVh1cQD4mojM8+zfYRFZJCILkiSpvWOFJ9Bs2fGVfyNEZKiItIjIHtGDpNaJyJokSdoj2DZSRL4iemjWiZ4vPydJkvmuBh0NPO3xV9idqIKxSJBP2NXCCrI2i0STODYEMqy7gVEEYyzwJ+z6YxMwLK+xwwj3JOhONMFoBf6FXV+sBPwcvIm+Dm7BrzCsxh7g+0QSjCHAdlNrP9qwy/+XC3AC+o293sCRw+hq21jvjhhBgx2a1d2xtKb9bwaOpfn/UfYcsoBt2vg6YjbnQvP/rTZl1qB5hKGzhZyh2dPGgdHoBsWbBs7/G83YDX3IdC1+W56GGr0aql+A09B3kYVg3I/+6N4Vwc8Yp6HWT9o4XYLR56ZMNcwEI40q7EJBgwhGmknYhYI6E4w0u7ALBQUXjJTCzgYKJhgphV0cUMF4M3q3huYQ8DP0nd5S+Xc++Q5jcGEb8HXglNhx7xf0jpgLrAL2oRs2DwBByqiAo9CzB1YZTATo18Nho7GCnoYKnIvOzUZ0rlahc5ctjQ0YC2yp4kwbMMuzDz3H/yB2HUZCEvw0VPSLoVoDji24fh6jd361yU/pAKYF8qm7Le9Ff9lWd6oP0tNQ3Tp0Z4vPNHQu+mILLk8C9NFRC7WnGucEzUv4Ifp4KyoH0BPXgtZM9ojL2hptm+ty0dUOTp8V0L/ebDsRO8FYK1GEHXCWg41OJ4a4fJZ9lwgdrLAXjL0R7Zh7tCXfbQ627nG5uOsZtlHbvqPnHFsKxmjH3KM9lm7BvSnnVpdBVmQMTNrFc0S4EPRp91jCCUYzYVfFt9OB28m+CfW0y2B35gxWO5rqPDFgTPqy36dgNBd2PXw5D3gI7f+fhwUug07KOVh3lqMp0DEel6lg3JzB7q1EXLEDrkT7NPh6rX3AZfAE+IengVNeAG4inmC8Gu2X01edwyb0FXIdEdLOUGH3mQCxr/0LoJsxUwgjqnag6jVakgP6ZHgfuto4GX3Mmlf+drMnq7CrhU5gSlbDfhTAoJSogrEIkF/Y1cLteQxsAR4MaBxEFowxwJ+w649F5C25Q/XArdisxS8nkmC0AP/Crhpt6GvWXxyBs4GHsWn7Hk0w+oZwwq43OtCUuzEhHbI8J2hXZay6O4kcFZpWx+ikWdHhJr4XBy3bvreinciyHYFqCPAO9EdrsWu5uzLW0JgOxzgn6NJoDlcB26P0NqI3X7FK6bE/J2gGkRtLGPu8Gt19DHswZF6AcdidE7QBvRuOM/QvLXZZaeBf+tS7zMo/b2B7TlBwwUijC7tQUOeCkWYTdqGgzgQjpbALBwU+WNLYtvoQdqGgIIKRJhV2hVlzR2vkZovIF0XkpMDD7RCRh0TkGRFBRCaKyEwRCV2n92Zl3AVJkmwKPFZNFOYHkIJm4HxJRG4WET997uKzW0TuFZH7kiTZG9uYugAt0b4ebZZQr6ynERs9WIJuR09DU7HrhSeB6TRR+1sT0JLuxdgIRlc6gN9SwH2KhgM4E5iPNlmIzUHgJxiXyJWICkY0dduio0dPdqFnFDXeil29ga1gLIVdUUEF4zWEEYylsKsn8CMYS2FX76CCcR5uJe+vA/fQBMKucCuBoUAPUJgsIlOk69CoIaIx2Ccir4jIWhFZLiLLkyQ5EsXQkpKSkpKSkpKSktD8D+NCLFQZzJjdAAAAAElFTkSuQmCC",
    "lightbulb": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAIj0lEQVR4nO2dW6xeRRXH/1MKpdxaaqtWW6BCobQ1ogYVIiQITySUkhhKIN4aRe0LeEkxanww8d0H7xBjojFaQoKXWqHQIFBpArZEblXk0IulPRQOttjWtqfn58PsNoeTc+i35vtmZs935vfWnrX3/Gfttfc3e2bN2lKlUqlUKpVKZdLhcguICTBV0lJJH5R0qaQFkuZJmi1phqSpjemwpH2S9kraJWlA0guSnpb0rHNuOK3ydPRdAAALJS2TdJ2kqySd2eUp/yvpMUnrJf3eOfdSl+er9BpgFvAVYDPxeQq4Azg3d78nPcBFwE+Bgwku/FgOAD8BLszth0kHMB/4OXA0w4Ufy1HgHmBebr/0PcA04Nv4u69tHAC+BZyW2099CXA58Fzea9wRzwAfzu2vvgFwwF3AkbzX1cQR4GtA69+yWi0QOFvSLyXdmFtLIPdJ+oxz7kBuIRPR2gAA3iPpT5I+kFtLl2yRdL1zbk9uIePRygAA3ifpIfmZu35gQNK1zrltuYWMpXUBAFwg6RFJ5/fytJKel5/a3Sppp6QhSYeav0+XNEvSeZIWSbpMfuq4l/7ZJulq59zOHp6zvwDeCbzYo4HYAeA3wApgToCWOcAtwBp6N9G0FZgdw3fFA0wHNvXAyQP4qdoZPdQ2Ez/VvK0H+h4HpvVKW98A/KJLxw4CXwJOjajxVODLwKtdar07lsYiAVZ24cwR/HrAzIR6ZwJ3N22H8qlUelsNcCHwZqATXwOuz6h9GfB6oPZ9+AHv5AU/y/dwoAOfaYMD8QH8fGAfHsitPyvArYGO20jCR/7JwOcjhA5gV+TWnwXgDGBngMM24aeIWwUwA3gyoD/bgdNz608OsDrAWf8E3pFb+0Tg5zFeCujXV3NrTwr+7re+Sr0JLM6t/WQA78eer7CbyfQUwL9LW/lsbt2dAtwe0L/bc+tOBvbEjrW5NVvAv908YOzj33PrTgJwpdEx/6PAxEtgIXDY2NePpNY5JXWDkm4z2v+4xFx859yLkn5mPOzWGFpaAzAFP+DplEPA3Ny6QwHea3wK7CRxGlnqJ8CHJL3bYL/GObc7lpjYOOd2SbrXcMg8+VyEZKQOgOuM9v2wanaP0f7aKComIHUAXGWw3SFpYywhCXlU0isG+4/HEjIeqQPgcoPtWucc0ZQkwjk3IumPhkOSvgkkC4BmMGdJzXoolpYMbDDYziUghS2UlE+Ai432m6KoyIO1L4uiqBiHlAFgSfF+zTln+d1sNc657ZL+YziklxnRb0vKALC8/hU38dMB/zLYJpv7SBkAloIKfXP3j8IynzErmooxpAyAswy2+6KpyMcbBttuy9p0TMoAmHpykxMcjaYiH4cNtslqDKQMgGMG2xyLVLGx7Few+KorUjracgecEU1FPixJrMm2k6cMgP0G23dFU5GP+QbboWgqxpAyAAYNthdFU5GBZon3EsMhu2JpGUvKANhmsJ1Pf+2iXSzbW9BALCFjSRkA/zDYOkmfiCUkA5Yl3uO1DJKQMgBelu39/pOxhGTgJoPtgHPOMmfQFckCoFkW/avhkGVA8YNBYIGkqw2HPB5Ly3ikft9+2GA7TVI/7Ji5QzY/99My+FvBp0pbOEgLdgCHAszDVl7mKC3e+tYTsG+g/DMFFFwcD+BeY1/X5dYcHeCLRqcA3JlbtxXCtr73/1Zx4Ex8ZQ8LR4BrcmvvFOAy7FVPdjFZCk3jK35beQNofdVQ4BJsm1+O8/Xc2pMBnBXopEFgSW79EwEsBfYE9Gs3kCwHoBUAnw5w1HFnLcytfyzAoiZAQ/h8bv3JwW+hXh/osAFatGcQ/wWTHYF9eQLox/yHk9M4bijQcU8C2fMGgHPwFctCOAhcmrsPWQGWE15w8deZtU8B1gZqB/hCTv2tAfhuF05clVH3d7rQba0d0L/gxwO/CnTkQSDZTppRmq8AhgM1ryNiTeMiAU4j/HG6kYQDKfzXy14I1LqBFoxdWglwOvBgoGNXJtT5zUCN64DpqXQWCb6G4GMBzn2FBHcWMBvYH6BvQ734HQKcS9gjNnr+APC9AF1PA+fE1tZXABfjS6pb2IH/ZHwsTdOxl4cfBM6LpamvwX+zx8oNEfXcZtRyDLDWRaqMBrjP6PTfRtTyB6OWH8XSMmkAzsfXC+yU/URYV8c//i3pXUOAZUt8Flq/CNFU17CUWjtb0sciSLlC/vuCnfKDlOndobQ+ABp+aLS/MoKGjxpsR2QvE5uFIgLAObdVkqWadowaO5aR/FPOuX9H0NBzigiABku+vKUgU4xzPhqh/SiUFADrDbZ/idC+5ZyWDTCVTsCvvW/pYPS9hQgLQ7nbr+hExu3bJV3uASz78Itqv6ITaWRreOta/HDzf5YqHEW232uK3HIl+QUjSUubfz6b+p07d/uVSqVSqVQqlUqlUqlUKpWKiSJnAvEp1svlv0SaO9f+kKTNku53zlkKYldCAG7GnpqdgteBm3P7p68BVhC+lTwFIxRW6auYnwD8Y/9lJfygUiBDkhaU8nNQUuLCTWr/xZe8xuW5RXRKSQHQ+hJxoyhGa0kBUNKHpIvRWlIAWNLCc1OM1joI7D11EBiDxqGr1O7HK5JWlXLxi4Q6EdRTivkJGA11KrhSqVQqlUqlUqlUKpVKpVKxUeRM4ETgy8QukTQnUhN7JT3nnBuOdP5KCMBU4BuEf7XLwiBwFxFrElcMAKcAv0tw4cdyP3BK7v5Pepq7MRerc/e/W4oeAzR34B5JszNJ2CtprnPuWKb2u6aYhJAJWKJ8F1/yg83FGdvvmtIDYGZuAWqHhmBKD4DtuQWoHRqCKToAmlLymzNK+JtzbkfG9rum6ABoWC1fnj01I03bldwAK4HDCV//DgOfy93vXlD0a+Bo8J+PvVPSNYo3FfyqpEckfb/5hkGlUqlUKpVKpVIe/wdo82tEYrtdhgAAAABJRU5ErkJggg==",
    "line-chart": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAF9UlEQVR4nO3dTagVZRzH8e9fLQ0KLpRR9EbYwsxAatGmNyVqE2RCRFiCG5HMWoW1aVu0Cy2whTtXib2BpGFCtAmDFHJTFtmbdC28tOjV7q/FzInjueeeM8/MM2dmzvP/gBvvM/M89/x/M8/Mc885A84555xzzjnnnHPOuelnZTaSNANsAtbl/3UCOGhmc7EG5lpK0mZJ57XQeUmbmx6fq1Fe/Pkhxe+Z9xBMKUkzixz5w84EM02P1xWzJKDtJqBIYXvXB64DQgKwbnyTUm1dg0ICsLymtq5BIQFwU8gDkDgPQOI8AInzACTOA5A4D0DiPACJ8wAkzgOQOA9A4jwAifMAJM4DkDgPQOI8AInzACTOA5A4D0DiPACJ8wAkzgOQOA9A4jwAifMAJG5Z0wNwGUlLgPX5vyuB08B7ZvZVowPrkbS3wCeDe/Y2Pd4ukXSLpONDXsd/Jb0u6dK6+vYzQMMkrQWOAlcP+fES4GlgpaTHzUyx+/drgAblxf+I4cXv9xiwsY4xeAAaImk18CGwsuAmW+sYx1RNAZKWAWuAq4DTZvZdw0Maqu/IL1p8qOk7F6biDCBpmaRdwI/ASbI59YykzyRtaHZ0FytZ/Np0PgCSVgBvA6+wcC69EzgiafvEBzZExeKfjDwcoOMByIv/DvDwiGZLgTck7ZjMqIaLcOTvizic/3U2AH3Ff6hAcwN2NxWCCMU/QPa7RtfJAEhaDrxFseL39EKws55RDVfian/QEWBLHWsA0MEA5Ef+u4w+7S/GgNcmdSbIj/yPgWtK7uIw8IiZ/RFvVBfrVABKHvmDJnImiHTkP2pmf8Yb1UKdCUBe/AOUO/IH9c4EtYQgL/4xyh/5R4CNdR75PZ0IQOTi99QSgi4VHzoQgIpz/jhRrwm6MOcPanUAIs3540S5JujKnD+otQGo6bS/mErTQddO+/1aGYAJF7+nVAi6XHxoYQAizPnngFeBMgsnQdcEXZzzS5vEW8IkLZf0fkA/g2Yl3Z7va5tGP91klHmNORNIWi3pbIWxHpZ0WZnXqRGqOQCKWPy+fdYSAqVWfKg3AKqh+H37jhoCpVh8qC8AkpZKOlThBZ1VNheP6mOnqoVgR76ftXl/ZX2g7Bqne1RfAHZVeEHHFr+vn6oheFmpFh/qCYCyo/9cyRe0cPH7+qsSgipaW/ymbwNvI3sDZ6hzwAYz+yJkIzPbDTxHuVvEsg6T3edPdIWvqKYDUOb5gqWK3zPhELS6+NB8AM4Etq9U/J4JhaD1xYeGA2BmZ4DPCzaPUvy+vusMQSeKD82fAQCeB+bHtIla/J6aQtCZ4kMLAmBmR4FtwN+LNPkWuC928fv6jxmCThU/iOpfCr417+NLSb9I+lTSi5Iur+P3GdJ/1VvE1t7qRVF3ANpA5ZeNu7m8SwumgDYxszeB7YRNB43+Pb8qD8CAwBB0uvjgARgqD8GzjL47OURX3swxggdgEWa2B3iQhesUs8ALZMXv/NX+VH1BRGz5Leodkm4CVgG/AqfM7EKzI4vHA1BAvmIZumzdCT4FJM4DkDgPQOI8AInzACTOA5A4D0DiPACJ8wAkzgOQOA9A4jwAifMAJM4DkDgPQOI8AInzACTOA5A4D0DiQgIQ8mEJD1ZHhBTq94C2q0IH4poREoCfAtreLenG0MG4yQsJwKmAtpcAe5U9yNFNA0kzki4Efmr2kKQbmh67W5yFNJZ0DLg/sI9/gE+Arxn/TSAu3F/ACeCgmc2FbhwagK3U9ABDV9kc8IyZ7Q/ZKDQAK4BvgGtDtnMTI+CpkBAE3a/nn4Z9KXRUbmIM2COp8Pcvllmw2Uf2kATXTjPApqKNgwNgZvPAk2Sfk3fttK5ow1JLtmb2PdkjXX4rs71rj9Jr9mZ2HHgA+DnecFwkJ4o2DLoLGEbS9cB+4N6q+3JRzAE3F10TqPxXOzP7AVhP9m2fZ6vuz1UisrWAwgtClc8AF/WerRM8AWwB7gGWxty/G6n+haAQ+b3oXcAa4Drgirr6SlylpWDnnHPOOeecc84551wK/gM9Ck6tFmSuVgAAAABJRU5ErkJggg==",
    "lock": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAHaElEQVR4nO3df8idZR3H8fd3bXPm/DVqs7F+qNNhv3AtW5tlFpjLYiCRWFhaWknrjxL6RSFBKCERQSwsSLAsCzSpyBrhj2itzFIyiAa2JZubW7Umns398NmnP+57IE9zz/W9f5zrfp7r+4Lnr+e6r/M91/XhnPvc5z7XBSGEEEIIoTiWu4C+SToNeDNwPrAMOBNYBJwOzK2bHQL2ALuArcBm4FHgITN7etw1j9OMDICk84H3AmuANwCzGnY1QRWEe4G7zeyxbioMnZN0sqRPSvqr+vMXSeskzc/9fENN0imSvixpT48TP9l/JN0o6eTcz79YkmZJ+pik3WOc+Ml2SbpOUtO3mNCEpGWSNmWc+Mk2Sjon97gUQdJHJO3LPOHHMpJ0de7xmbEkzZH07cyTnGK9pNm5xyvVtPgYKOkk4G7g0ty1JPoFcIWZ7c9dyFQGHwBVH7nuBd6auxanB4F3Dz0Egz57lXQC8DOm3+QDXAzcI2nuVA1zGuwrgCQDfghc2VGX+4GHgEeAfwA7gVH9v/nAYmApsBxYCZzY0ePebmbXdNRXOSR9toMTsmclfU/SGknzHI89T9Jlkr5f99HWp/scqxlH0mpJh1sM+EjSTZJe0kEtCyXdXPfZ1CFJb+pibGY8SSdJerzFYN8laUkPdS2RdE+LujZL6uptZeaSdEvDAR5J+uAY6rtGzS9E3dR3fdOapKWSDjYY2CdVfQU8rjpXSNrZoM4Dks4aV53TjqQfNxjUbZKWZqh1maQdDeq9Y9y1TguSXitpwjmYeySdl7Hm10l62lnzc5LOzVXzYEm6zTmQE5LeNYC610o64qz91tx1D4qk0yXtdw7i13PXfZSkbzprH0k6JXfdgyHpeucAPiHpxbnrPkrSfEnbnc/h2tx1D4ak+52Dd1XumieT9GHnc9iQu+ZBkHSqfFf9Nkt6Ue66J5M0W9IWx/M4qAHcUziEbwMvBjw3UKw3s4meamnMzJ4D1jsOmQu8radykg0hABc62h6m+oZwqO6g+i1BqtV9FZJqCAF4o6PtRjP7d2+VtGRmu4BNjkNW9FVLqiEE4NWOtvf3VkV3PDW+prcqEmUNgKrbvRY5DvljX7V06GFH28XK/A1h7leAxc72f++lim79zdHWgJf1VUiK3AFY4Gh7BNjRVyEdehKQo33rm1bayB0Az9W8Uf1Ra9DM7BDV/Yepsl7RzB0AzwWdwX32P45DjrZzeqsiQe4AhMwiAIWLABQuAlC4CEDhIgCFiwAULgJQuAhA4SIAhYsAFC4CULgIQOEiAIWLABSul0Wi6nv9LgUuAV4OLKG692/yillzqBZoSuoW2NtVjT07jfSxHVHd7v58B4FtwBP13y+BB8zsSGcV1joNgKQLgC8B7wSSF2UKSbYDPwC+YWZPddVpJwGQdDZwM/C+rvoML2gv8EXg1i5eEVpPlqrf6N8JnNq2r+CyCbjczHa36aTVSaCkG4CfE5Ofw2pgo1quOdT4FUDSR4HvtHnw0ImngNVmtrXJwY0CIOktwH38/1l9yOMxYFWThandbwGqlja5i5j8IXk98K0mBzY5B/gMvt/zhfG4WtI7vAe53gIkLQQeB7KvbBGO6WFgpZkl/zTN+wpwHTH5Q3YB8B7PAd69bdY62wP8nuoK1lam/snUbKqtXedTnd3ubPB4Q3RG/bcP2MLUP3M7gWrhjHXAS52P9QGqj+bdkrRYvgURJyR9vPNCCiJpgfwrqD3Ty5oDki53FnJj50UUSNKJkh5xjn3yyaDnHMCzkMF/gVsc7cMLMLNngWup1kdIdWZqQ08APKt5/NbMDjjah+Mws0epzqVSvSK1oScAnpUsWn1BEY7Js/TMGakNPQHwXDPo/MaF4BrT5LmKW8IKFwEoXASgcBGAwkUAChcBKFwEoHARgMJFAAoXAShcBKBwEYDCRQAKFwEoXASgcBGAwkUAChcBKFwEoHARgMJFAAoXAShcBKBwEYDCRQAKFwEoXASgcBGAwkUAChcBKFwEoHARgMJFAAoXAShcBKBwngAkrz/r7Dek6WWuPJ16Nio6z9E2pPGs0P6v1IaeAGx3tF0labmjfTiOeunXixyH7Eht6AnAPx1tZwHf7WXN2jJ9jmovwlTJAUhWr1n7jHPN2vskLei8mIJIWqdq4e1UE5KSl/X1bhhxJ3Cl8znsptrO5E9UO2KGqc0FzgKuAlY6j/2Dma1KbewNwFrgp86Cwnh9wcy+mtrYGwADfgckJyyM1Qg4x7O1rOvzer0Xzee9VYWx+Zp3X+Gm+wbeDnyoybGhNzuBc81s5DmoaQDmARuBFU2OD507DFxiZr/xHtjokm29GcQV+K4Ohv5c32TyocU1ezPbAlxItY9gyGMCuMHMbmvaQRfbxy8EfkIVhjA+e4H3m9mv2nTS+lu7ev/6i4BP1EWFfgn4EbC87eRDB68AzydpEfApqitYS7rsO3AA2AB8xcz+3FWnnQbgKEmzgLcDa4BXAa+k2skqdhxPc5jqBHtb/fdrYIOZ7ctaVQghhBBCmBn+B+79XwxlLTbtAAAAAElFTkSuQmCC",
    "mail": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAISElEQVR4nO3dW6wdVR3H8e8PAYuGchFtrfUSLMaqCb0EjEVBRUOiQNM+iKKUGEwUEy2JTVFiAj4YatQoBoMvNZFgEJtQtBAfWhHboqGI1bSFGFBDpaWmYgstl9rLz4fZJz097Musvdfs2/w/yX7ZZ9ZlZv5nzfxn9qyBEEIIIYQQQgghhBBCCCGEEEIIY0g5K7M9A5gBvClnvQEAA3uApyUdzFVpTwFg+zxgcePzfuCUHJ0KHT0HPArcBayV9FK3FXUVALbnAauAy7ptOGRzAFgN3CzphdTCSQFgezpwO/BZ4KTUxkKldgM3SFqTUqh0ANg+F/g18N7EjoX+uh1YLulYmYVLBYDtC4DfAG/ooWOhf34JLJN0qNOCHQPA9mxgC/DmDB0L/XO3pKs7LdT2OG77NOA+YuePos/YvqHTQp1O5FYAC/P0JwzAdxuH75ZaHgJsnwP8HZieu1ehrzZJurjVH9uNADcRO38cfMj2la3+2HQEsH0SRV45o8tG93VZLrQm4Mwuy7YdBV7F9iKne9T2J22/tstOhg5sn2V7pe2DifvmqO1ZKQ3dmtjAvbZPrXDdwyS259vem7iPvtKsrlbnAO9O6M9u4BpJ/0tek9AVSVuBqynuEJbVNBtoFQDlhwtYLenFhOVDBpLWA+sSiryt2ZetAiDlws+fE5YNed2TsOzbm33ZKgBOT6g4248TQrK9Ccue0ezLuKVbc60C4DUJdSy33TS6QnUa2/yrCUXK71PbhxJTjF22lyavReiK7SWNbZ7ilZQGjiZWPmGt7bdUuO61ZnuWi2su3Tia0tCxLhux7f22r3dxOTlkYFu2v9TYtt0q9QuhiQaP9NDQhM2231PhdqkF23Ntb8qwP46kNPpyhgbt4lziFsf9gWS2T7V9s+1XMu2L8j8dd29DTTOP2/5ghdtrrNi+yPaOzPug6R3aVsfplGvMZcwFNtq+w5EytmR7uu0fAxuB3IfP8vvU9r7M0TfZLttLMq/cyLO92PYzFW73/zZrdxBn6rOAe22vs/3WAbQ/VGzPtH0nxY9v+55CDzJVuxzYbnu5a5gyukjtlgHbgWsG1Y9Bb/jpwA+BTbZr88SR7TnABuBnDPhhm0EHwIRFwFbbqzzGKaPtU2zfCGwDPjro/kCeANiQoQ4oHi2/keKwMBQbJyfbHwC2UjxVPS1TtY9kqudETssCPmb7Cts785ys2i4uRd9p++xKVrCPbL+uMbLluLo64Vnby2x/PKFM0yygVaeTAqBR5gzbt7n7G0nNPOviRGkkufiV9NMZt8cJ/xgepgCYVHaR7e0J5csYqZTR9gwXOyqnJz3l0OgMAZD9JFDSH4D5wNeBjo8nl3Q5sM1DnjL6eGq3g3yp3WHgO8D7JD2Yqc723MMIMKWeObYfTKirjIc9hCmj7Xfa3tDPdfUwjgCTSXoKuBS4Fih/DGpvqFJG2ye7SO22U6xrDi9SjKAXS9qRqc7ynGkEmFLnTPfhuNhPthfYfizzOq2z3fQ3/E3aH76TwBJ1j3zK6ApTu8R+jF4ANOof2ZTR9idcTWqXfEnYoxoAk9oZmZTR1aR2T9nu+rzBw34S2MkopIyuPrX7baY683GfRoApbQ5dyugitVufuU+P2V6QaZuN9ggw2TCljC5Su+XAX4EsAQ68RDHSXShpuB+o9QBGgCntV5UyfqRE2/Nt/ylz2/e7ZGqXuJ1G+ySwRD/6ljK6mtRujyvMTJwhAE6uqnM5SFpneyNwK/BFej9kieJE7jLbqylmQDXF7BmfJ21ijHZMMYP3SkmjN2GWh2QEmNKnKn4rX4W/2b6kT9tkfE4CO5H0MLAAuIV8KWNOh4FvA+dL+v2gO1PWyAQAgKRDkr5Fce1g86D7M8kjwEJJ35RU/jHsITBSATBB0hPAJcD1wPMD7MoBYDmwSNK2AfajayMZAACSjkn6CcUjVGsH0IX7Ka7k/ajsyxmG0cgGwARJuyUtBZZSzFlYtX8Dn5Z0haSdfWivUiMfABMkraUYDe4AqviPNPBTYK6klOnZhtrYBACApOclfZni/ODxjFU/CVwq6bqRzOvbGKsAmCBpM8dTxpd7qOoQxUWo8yX9LkPXhs5YBgCckDKeB3yf4mWLZe0DbgPeJekmSb0E0VAb6kvBOUjaBayw/Q3gw43PAuAdwNkUl4efA3ZSTHv7EPBQmTdujYOxD4AJkg4D6xuf0DC2h4BQTgRAzUUA1FwEQM1FANRcBEDNRQDUXARAzUUA1FwEQM1FANRcBEDNRQDUXARAzUUA1FwEQM1FANRcBEDNRQDUXARAzbUKgJQHLt+YoyOhK6cnLHug2ZetAuBfCRV/KmHZkNfChGWbPjfZKgBSHnpc7D7NEhKOsz0duC6hyK5mX7YKgC0JFQu42/a8hDKhB7anAXcBMxKKPZHSwGynz+N7wPZKj8F7foaV7Wku3jD6l8R9Y9sXNqtTbRrbDFzUZV/3k//9wwHO6rLcbmC2pFftk3aPhn2P7gPgzC7LhWr8vNnOhzYjAIDtTUC89n207QfmSGr6dHSnC0FfA45k71Lop1Wtdj50CABJWyje5hlG0x8p3s3cUttDwATbvwCuytGj0DfPABdI2tNuobL3Aq4F1vTcpdAv/wGu7LTzoWQANGbLuIpiqpUw3LZRvJNga5mFS98NlGRJKyiu/fdjPr6Q5ijFDOWLJP2zbKHk28GS1gBzgR8AL6SWD5V4AJgn6QuSDqYULHUS2Irt04AlwOco5tw/p5f6QmmHKCao/hVwn6R/dFtRTwEwle3XU8y+NTN33QEopqndI2nvoDsSQgghhBBCCCGEEEIIIYQQQghhqP0fUJHlcKdF8uoAAAAASUVORK5CYII=",
    "map-pin": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAANYklEQVR4nO2deZBdRRWHf00SloQkbEnYZZElLEIAWWRHESgKEAFZBEkJWIgRSCGbhYDsKOCGkhJQSmQLRRAphdKygFCAAUxEArIkERAkKwnJJCHLzOcf/UbjOCTv9On77nsz9/tz5t57zunu1/f26dPnSBUVFRUVFRUVFb2OULYCRQH0k7STpOGStpW0paSNJW0gabCkdbrcMlfSfEmzJP1L0jRJb0p6VdKrIYRljdG8sfSYAQAMkHSgpM9K2lfSCEmrZ3r8EkmTJD0j6U+SngohLMr07FJp6QEArCPpi5KOl3SIpDUaJPojxYHwkKRxIYQPGyS3AgjAwcB9wGLKZzFwD3Bg2W3TowH6AacDL5fZ26vgJeA0oG/Z7dVjAPoAI4FppXatjanAV4DVym6/loY41b9Ubl+6mAQcXHY7thzAMODekjsvJ/cAQ8tu1+5oulUAcJKkWyWtX7YumZkt6ZwQwoNlK7IiTTMAiOv42ySdVrYuBXOXpFEhhIVlKyI1yQAAtpU0TtKOZevSICZLOjaEMKVsRUofAMChksbq/12zuVgu6Q1Fl+5USe8quns/rP1PkvoquoeHSNpM0laSdpC0naQ+Bek1V9JxIYQnCnp+80NcKi3N/MHVDjwHXAEcAPR36DcAOBC4Evhz7dk5WQKckrNNWwbgPKAjY2P+BTgX2KhAnTeu6T0po97twDlF6dyUkK/zlxOXi3uWYMNewP3kmRU6gFGNtqEUgLPwd34HseO3aQJ7tiMOhBw2jSzbnkIBjib+aj08D+xdti1dAfYGXnTatgw4smxbCgHYDWhzNM4iYDRQ1Fe5G+K+xQX4dinnA58q25asABsAbzsa5WVgeNl21AuwIzDZYe80YL2y7cgCsBrwmKMx7sWxlCsL4hLyAYfdjwKl+2ncEKfEVK5s5UYgBq9c67C/tVcGwA6kvQ876EFrY9KXvYuA7crWPwni1P9sYud/vWz9cwOMSmgLgPG04iwIfC3R4IvL1r0ogMsS2+SrZetuAhgMzEww9LaydS8a4M6EdpkBDCpb97oBbkgw8mniYY4eDbA6aa/Ga4rQJ/u7BdhQ8VTNWobb5kjaJYTwXm59ukJ8n24qaaikgZKQ1CZphqT3Qgg0QIfNJf1V0rqG2xZK2iqEMLMYrTIB3Jwwuo8vUJ8+wOeAG4EJwMKV6NFWu+YG4BAKjOoFTk5opxuK0icLwLrAAqNRhcTIEQNLrwXeTWjoTt4FrqGggE7gYaM+82jmbwHsTp82YNPMOqwNfI+4hs7FQuIMsnZmXTdn5TNSd5yXU4dsEL1ebxqNuSKzDocB/zTqYOEd4KDMOl9l1OHvNKNfADjIaMhMMv2iiE6na8kfstUdS4Gzc+hd030QMNuow7655GcDuN1oxCWZ5PYjHrxoND8l0xlA4HKj7ObylxA7YY7BgAXEo91euX2AB42Nl5O7M7Xf+ti+WWaQKS4i1zJnP0mW/et7QgjzMsj9sWJugLI4FbjI+5AQwhxJ9xtuGSrpM165Ur4BcLjx+ju9AoETJTXDjuH1wFEZnnOH8fojMsjMAzEku17cp2GAzYihU83CPGCY06YAvGWQOcHbjlKGGQAYLGlXwy0PeWVK+qGiG9dLm6QFGZ4zWNLlngfUXNDjDLfsRoZVVI5XwB7G5zzmEQbso5gXKIVZkm6SdICkgSGEgSGEQZIGKSaYurl2TQpn4Q9T/53h2r6SdnfK8wNcZJi2FgGuRE7EWDkry4DvUkdsITGW72rSwtddbm1gLeAjg7zRHnlZAH5lUHi8U9YnsDt7PgD2T5B1OPHdbqEd2Mxpo2Wr+BceWVKeV4AlZu1Fp6xTZdN5oaTDQghPWwWFEB6XdJSkpYbbVpP0ZausLljaaHunrCwDYEvDtZOdsqzLrdEhhBdShdUGzreMtx2dKq+GpY22csryAayJLdr1AIesQcR3eb1MIMOmCdHb+IpB7jIgeYVCTIxVLx2AKxuqdwYYJltU0TsOWSMUv3zr5eYc0T0hhHZJlnCsvrIti7vytuHaoOgVTMY7AKzHl2Y4ZG1ruHaxpN86ZHXlN4rpYevFomtXrCFfriNk3gEwwHDtkhDCYocsy9f1xBCCpcNWSk3vSYZbklcCIYQ2/Td1TT24nEHeAWB5/3g7xPJefcspqzv+YbjWG7Zl+aG4/CreAWDZkuxwyrIYusQpqzssA3hNp6x2w7WubWHvALBMVd6Yf0tevcFOWd1hiV9oc8qyzKyuQhbeAWApmjAAXxCDJX6giFwCOxiunZsqhHg4xnKmwpVw0jsALJ0S5Ev/OtVw7XBgE4es/4Ho3rV4PC26dmWIbEtrV2CNdwBYlyybO2S9Zrg2SDrdIasrI2XrFIuuXbGuIFwnhVwDIIQwV7bXwNYOcS8rZvesl/MBy9GrbgHWl2SJxZ8r6RWHSMuW8vwQwnyHrCx7AdMM1+6UKqTmkXvScMsQSWM87uDavWNke3U9WdM1FUsbWZam3ZJjALxuuNYbwGDdb/+SpKsc8q6SPeh0rEOeJO1muNbzqskDMdCiXubgOHAJ9Ac+NMjr5A7LpgmwBmnn+OcBli/4rnJXwxaD8J1UWZ3kmAEmGq5dT9LOqYJqtfrGJNx6hqTJwAmsZADWOuAYxaPbKVk5bnO6u0fI5sOwuKe7Jcd26UaKlTbr5cIQwk0OeRtKmiLbPsSKvK0YezdR0vTa3zZSnHqPVPpKZYGkrUMIqTGFAi6VdF29l0sa5pGXDWyHQp/KIO8Sg7xGcUEGuyzhYOW//zsBfm5QfDn+GPp++PPy5mQCzvQ2wCbY4h2znA/MdTLoD4Zr+yh+nSdTK+R8tKTCU8rUwXTFyh/e4tInytYfjzvl5YOYFWyJYfR6g0M75e5O3kQQVhYDe2WyxVIfcTGZk1W4AR43Np4nbGpFuSdjixXMxTLi+cQcNuxhlP1oDrlSvleAZHfSZEmyEEK4T/Fw6gc5nlcncxTDzR/I9DzrIdemqj0o6T8JoiynWhYS/ey55G+NL0V7vbxOxvy9xFT6lnzKC8mYKCrbDFDbGHrEcEt/SdlyAocQpiqemb9FxUQEtUv6maQ9QggW9/eqGCVbBNHD3g2gwgA+b/w1zaSAWgDEI2R3k6+g0yPAiAL0HIgtswo0czFqoit1itEgtwNlJfrsBPwAmG7UCeB94BYgeQezDv0uNer0GpkzhBWRKna04jRcL7MUXag5zul3C9H/v4ukgxXL026neKilM9J4geJ6/g3Fo1lPSPpbCMEbyLoynQYrbqVb4vq/GUK4tSCV8kD0CVh37LLmC2wFsFcSmYvjyFlDIWbqtDCfAit+NhvEFDfWDKHXlq133QAbYS8V4z7r3ioAvza2TRswpGy9TQA/MhrZTia3ajMD7I+9ftD3y9bbDHEWsE5zE2niwpBegL7YfP4QX4+t9evvBLjeaCw0Q96bgsCWT6kTT0xjuQDrYE+EvADYomzdcwNsg33ncjqt8uX/cZBWLu2PNGNK9ESIDrLxCe1wVtm6uyG+915OML4Z0sBmATg/wf6e8z1ErCVg/fJtAzyZNpoCYDj2qb+DZqwJ4AH72hfgBZxJkMqEeL5gYoLdvyxb9+wAQ7HvfEErroFrEDeirMwkY5xEUwGMTGiQDuDIsnW3AhxNWsHok8vWvVCAxxIaZTax2GJLAGxB2mxnCahpTYBNiTtbVp7HmWi6ERCTZ1rqJ3Qym96yIQacmtBAANaqGg0HuCvRtiwRxi0D8EBiQzWtf4C09T5kKj7VUhDdxJYSKZ0spQnj4oBDSTufMJVmLgdbJMA+iY02G391jmwA25P2XbMU2LNs/UsFuDCh4SDG57vy5GbSfwPsgbCdFBYQ2zIQq2U9ktiA4wFvVk6P7mthO9K9IuPoQRteLojfA9bC052MxZF2xqFzH+zl3zt5jd763v84gJ2JG0Ap/KQEfcck6jofsGQd7T0Qc/ikuE8BLmugntaS7510AF9olJ4tiaNxAb7RAP1S1/oAruKSvQLiR+HYxAZuB3KmiO2q25mkz1D3Un301Qfx63pCYkMvpwC3KtF9nXrY9BlKXK20JMAwYFpigy8DUsvLdqfLCaRnIZlCq4Z1lw3Rw5ayrQrRy3ZsBh08nT+LHhDWVirAfqQngloKHOeQfZKj89voBSedGgJwlKMjlpEQZQOcRloBaYgD7/Ai2qLXApxO+hf4cuAMg6yzSf/gawdOKbItei3AuYmdAnHwrPLYGXFzKnWgdQDZ8h5VdAPwbccggI85a0f0P1znfPaFjW6PXglwjbOjxrDCqRvixs7tzmdWXr5GAtzg7LCHiQ6n/qRvR3fSOtk7ehLYU9F05VnguarzWxjgRmcHeri6bPsr5N5BTMVdr6ciI8DFpC/fLNS1nKwoAXwOnHpYDpxZtp0VKwE4heiKzc0SetvpnVYFOIL0+MLuWAAcWrZdFQaAvYjbsV5mAJ8u256KBIiZuVIPbUAMVf9k2XZUOCBmJnk6ofPHU0Xy9AyIvv7zielXVsUM4q5jz8jQtQp6VZQqsbDzMYpFpnZVLBkrSe8r1uH9vaRHQwgflaNhRUVFRUVFRUVFRQP4N664EzhExoxPAAAAAElFTkSuQmCC",
    "message-circle": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAMRUlEQVR4nO2daaxV1RXH/1tUhCqKIs4TFWjrmDqidagmRqlGar/USoxTSRNNHdpqtdomdkqTJn6xCtShcWqbaKxFUWNbFBQVooloY5ker1pA5fkeIDI8Hvz6YV0nOOe+u849w7vv7l9yv9y7z95r773uHtZZe20pEolEIpFIJNJ2hKoFKAJgJ0ljJX1F0hhJB0s6SNJoSSMlDZO0+zaPrZP0kaRuSV2S/ifpHUlLJS2UtCiEsLkM+ctkUCgAcJik0yVNkHSCpCMl7ZxzMb2S3pI0X9JcSbNDCJ05l1E6LakAwFBJZ0m6QNI5kr5ckShLJT0j6e+SZrXiCNEyCgDsIOmbkiZLmiRpj2ol2o4eSY9LekA2OlCxPA0x4BUAGC3pSklTJB1arTQN0yHpj5LuCSF0VS1MSwKMB+4BNtK6rAemAmOrbs+WAev4PwNbKu26fOkDHiIqQjrAPsDdwOZKu6pYeoG7gL2rbu8BA7AjcD2wptKuKZce4IfAkKrbv9JFIHCspHslfb1KOSrkNUlXhBAWVCXADlUUCgwBbpM0T+3b+ZJ0nKT5wK1VjQaljwDAwZIekXRqCcVtkdQpM+V2SFoh6X2ZyXe1pM/v1UdKGiFp39rnMEnjZKbkMjrnRUnfCyG8W0JZn1KqAgATJT0oac+CiuiSNEvSHNnosiCEsKGZDIHhko6WdKLM3HyGpFFNypnGh5IuCSE8W1D+1QAE4BaK2dq9DdwOHI9ZC4uuyw7AScCvgYUF1KcPuLHoepQGsDPwQM6NtAb4A1D5+gE4Edu+5r2LuQ97q9m6ALsCz+bYKEuBq4Fdq67btgC7YVu7ZTnWdyY2BbUewB7A3JwaYglwKbBj1fXqD2An4DLyU4Q5wIiq6+UC6/z5OVS+G7gOyPvdfuFgU98NwOoc2uEVWkUJsGH/5Rwq/SD2JrClAfYFHs6hPeYw0KcDbPhrds5fgW0XBxXA+bW6NcPTDNSFIbbVu7/JCs4AitpjVw4wCniyyTa6t+p6JALc3ESl+oCbgAHvoNIs5GMT+UnV9fgCwMQmKrQGOLfqOpQN8C1gbcY26wPOqboOkiTgEODDjBVZDhxddR2qAjgWWJmx7VYBB1VdgR2BlzJWYCnmzt3WAIcDnRnbcDZV+hQAtzXR+dVq7wACG0WzKsEtVQl9LObi5GU58Z+/HdhIkGU62ETZ0yg29L+WQdg1wDGlCttCYH+qLAvDeZQ5FQDXZhCyjzZc7XvBdgdZdlTXlCXgaLLZt28qRcBBAGYn8NJNGUY0zK3ZyxO0gZEnLzBjURaL4Z1FC3Y4/oXfcmCvQgUbhAB7418U9gLFHZTFTrd4GXQvdsoCuCBDez9YlDDjsIVc9cK0EcAjzjbvo4gjaMB0pyA9wD65C9JmYP4EXn/DqXkLMRrY4BTihlyFaGOAHzvbfj15rrvwb0s6sSgekRwAhuI3Feez7cb84Jc5C78il8IjnwJc4eyDJeSx9QbOchbcQQt477YamLvdMmdfnNlfvo2cpJnslPWOEEKf85lIP9QCUN3hfMzbd18Ec23ucWjcagbgoY3BAjAC38uibvpxIu1vBDhLvmhcD4cQ1jnSRxyEENbKTlY3ykhZZLVU+lOA8x2FSdL9zvQRP/c503v78DMwz51GeTtzQREX+E4lL6qXV+oIABwqC47QKI860kaaw9PWY7GgHInUmwLOcBQiWbjUSDnMcKY/Pe2HegpwiqOALlnAo0g5zJdFE2mUCWk/1FOAExwFPB9C2OpIH2mCEMIWSS84Hknty0QFwI5jH+EoYI4jbSQfPG1+FCnW2bQRYKx88fZfdaSN5MM8R9pdJB2e9EOaAox3ZL5FUmWBDtuYBZI8025in6YpgMevrKPZUGwRPzWLa6fjkcQ+TVOAQxwZL3akjeRLXSPPNiTaAtIUYH9HxsscaSP50ulIe0DSl2kK4AlnvtyRNpIvnrZP7NM0BfCEcv3AkTaSL6scaUcmfZmmAJ53+qsdaSP5ssaRNrFP0xRgF0fGcQdQHR870ib2aZoCeEKR9TrSRvLF0/aJfZpHdO2WuB8vkkyaAmxx5DEwAxe2B562T3TUTVMAz7z+JUfaSL542n5j0pdpCuBx7BxoV7i2E562/yjpyzQF6HZk3PIBnVsYT9snbtfTFMBj3Ek0MUZKwWOyfz/pyzQFWOHI2PPiKJIvnrZfmfRlmgK848jY4zsQyRdP2/836cs0BVjqyHgMMMyRPpID2BG8Qx2PdCR9maYA/3FkPETSUY70kXw4Rj5DXmKfpmWwWD4z48mOtJF8OMmRdqOkJUk/JCpACKFX0r8dBZzmSBvJh2840r6ZdmS/3hDi8fQ9kxJu7YwYWFxgz8mtVA/iep32sqOAUbKbsCPlcKJ8Tjtz036opwCzHQVI0oXO9JHseNs628EdLNBQo3h2DpGMYHGEPf2ysF5+/c3bzzhkGw8c70gfycZJ8p3beLrej/0pgPcYcgwPVzyXO9N7+/AzsCBR3Y7hZg0xSFRhALsDHzn6o4t+QvbVHQFq9oC/OWQcoTgKFMmV8nlsP9Z0yD78gSI7Gaj327Yw2Gj8jrMvmjfQYaFiO5wFT8qhzpHPAUxx9sFiGggV26/1rhb5Y7pT3u8400fqAOwi6WfOx6aHEPLx2MauL1nv0L54XjBHgJ86//3rAI+lsCEhpjmFiJdF5ACwP/67BPO/PAr4mlOI1MhUkcYB/ups980UdTMrvsDR5xUiRBsBXOjsfABXuF7vK1zPiaH4ergJsCl0mvOxXkm3ex5ouJOwm6g899B4ji5HPgfmW/GAJO86amoIoZiILcBNzqEoXg+fEeDnGYb+Loq8oBN41SFMN/Gq2EwAk8h2efQPihTqYGCrQ5i/FCbMIAY4DtvDe3mZjC55jT50kSTPP/qxDLK0NcA4STPlP229SdL3C43VDMx2aON64ithF8AY4N0M/3yAm4sWbl98c9LjhQo0yMDuZM7a+c9jHsKZaWQKmNRguk+Iw3+DAMfJHDYPzPD4B5IuqYWOLw7gOYdGbgJiwIgGwKx8WRZ8YObes8sQcq9aYY0ys3ChWhzMv+I2sm31PuG6soS93CnYVaUI1qJgt7A/00THA3jNw00JPMMhWB/giTHcVmAGnvea7PwnKOteZuya0o0O4f5VimAtBvY+3/tKN4lZlBmHAbjYKeDVpQnXAgBDgRvxO3Mk8RKwWxFy1htOPH59WyVl2v8DB0iaKLun+AhZ4KM9JK2VRStbLDvd+qykVwb67WTYhVuXSrpV+cRPekHS+aXeyQwMBz72aKgj7yHABOBXwOv43jEsBa4t6t/QDNiUeT1+1+16PEkV4XeAi5yC/qif/PYEvgs8BKzKoWHWAlOBSiOTYAc1TwGm4zux0wjTKGvBl1CxPzmF3c4HDTga82adg+0QimIR8FvgZJo0izbYNkOAU4Hf4T8v0Qh9wPVF1+MTEt/wAW9KOrLBPF4LIRwPDJd0tmw+n6iUS4oKplsW12COLMLJG83Ondh0c4zsVO5ptU++LtefsUrSxSGEfxaU/3akKcAqWdSPRpgrW7CdKd9FE2WA7GKlRbLLrVbKbOg9soDY62vphspew+4lC7+6vywE2zjZQq4M55bnJU0OIZR6piJNAT6Q7+KoSHY2SfqFpN8X/mIngbRFxnuKClAGr0q6KoTwVlUCpL3mfaNUKdqPHklXSzq1ys6X0hWgbliRSGY2S7pT0tgQwl1VDPnbkrYGGCYLGN3oQjBSnz5JD0v6ZQjBE4e5cNIihW6Q9JuSZVkni0YyRbbVmiQ7lt7K19Ktl3S3pPEhhMsGWudLdbY3mBVqlnwhSb0skfRU7TM7hLApQY5RsrAzU+SLjlUli2XKe38I4cOqhalH3f0tdj5ttmw/nAe9tfxmSnoqhNDw7dfYQZPTJU2WuakXZYzJSpekRyU9IunF3IIzFEy/Bo6aEjyq7CPBSlmHz5T0XAgh8fIiD1gMojMkXSDpXOWnoF4WyhbMM2QjWHMBmSqgIQtXbTq4RhampL+F4VZJ82XD+kxJrxf9bwAOlCnEBFkc3aOUv1Vyg6QFslfTr0h6oWyrXRG4TJxYrJpvSzpPZh/fr5bHCpnt4B+Sng4heG61zp2awo6R9FVJh8nMuvvJTtvuLgtnN1xmApbs2Pta2UK0R3bB0gqZGblDdtlCx0DYtkUikUgkEolEIk3zf+m/uUnA/LM3AAAAAElFTkSuQmCC",
    "message-square": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAEiUlEQVR4nO3bTWgcdRjH8e9TmlgsIlrWghURweZSREHwrNjoxbebiiCGehBRQRARjCKCB8GbKNgSvJQeetEWtBGUUEGQHkQQ6ktVKh6USAVbtdrWn4fNRTqz2f9mZieb5/c5zv7nZft8s9mdbMHMzMzMzMzMzMw2vmjiIJJuAG4BZoAesLmJ49pFzgHLwNfAsYj4bq0HHDkASdcCjwMPANet9UJsJCeAA8DbEfHTKAcoDkBSD3gFmAOmRjmpNe4fYB8wHxGnSnYsCkDSPSsn6pXsZ2PzCzAXEe8Pu8OmYRdKeh54Fw9/PdsOHJb07LA7DBWApJeAV2noTaO1ahPwmqT5YRavOlBJjwDvrPGirBsPR8T+QQsGBiBpJ/A5cGmTV2Vj8wdwU0ScqFuw2q+At/DwJ9lW4M1BC2oDkDQL3N70FdnY7ZZ0R92Dg14BnmnhYqwbtbOsfA8gaQfwIwUfE1csA0vAb4X72XCuBG4DthXudwG4JiJ+Hmq1pD0qc07Sc5KmCy/MCkm6RNILks4Xzmiu5CQLhQd/tMXnbBUkPVY4o30lB/+s4MBL7T1NG0TSJwVz+rTqGHW/468uuI4D5ZduDSn5t99RtbEugJLP/j8UrLVmfV+wdmvVxroASu75q2CtNetCwdrKWZd+zLMNxgEk5wCScwDJOYDkHEByDiA5B5CcA0jOASTnAJJzAMk5gOQcQHIOIDkHkJwDSM4BJOcAknMAyTmA5BxAcg4gOQeQnANIzgEk5wCScwDJOYDkHEByDiA5B5CcA0jOASTnAJJzAMk5gOQcQHIOIDkHkJwDSM4BJOcAknMAyTmA5BxAcg4gOQeQnANIzgEk5wCScwDJOYDkHEByDiA5B5CcA0jOASTnAJJzAMk5gOQcQHIOIDkHMNl6BWv/qtroACbbbMHa5aqNDmBCSboReLBgl2+rNjqACbQy/MPAdMFux6o2bm7kitZI0jRwPbAdmOr4ctazHnAn/Z/8kuEDLFZt7CwASTPAQ8BdwM148G06HhFfVD0w9gAk3Qq8TP8NTIz7/Em9UffA2AKQdBnwOrAHD36cTgILdQ+OJQBJO4FDwMw4zmf/81REnK17sPUAJO0CPgKuavtcdpG9EXFo0IJWA1gZ/seU3bGyZhwFnlxtUWv3ATz8Th0F7o6Iv1db2EoAHn6n9gKzEfH7MIsb/xXg4XfmJPB0RLxXslOjAXj4nThO/3P+wqB3+3UaC6Ch4Z8HTjdzRRvSWfp/1fuG/r39xbo7fMNqJIAGhv8lMA8cGaViG10TAewC9jP68BeB+yOi8gsL1q7KW7KSTgFXDHmMfxn908QicJ9/6rvTxMfAUY/xIf2ffA+/Q119IWQRuNcv+93rIgC/7K8jdQGopfN5+OtMXQBD3UYs5OGvQ3UBfNXweTz8daougMovEI7Iw580knqSzmjtjkja0vXzsRFIetHDT0zSlKQlDz8xSZePEMFBD38DWXklmJd0epXB/yrpCUn+yvcEGXpYkrbR/y9Ju+l/vXsL8Cf9P+V+AByMiDNtXKSZmZmZmTXoP8oLREgjsokPAAAAAElFTkSuQmCC",
    "network": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAHgUlEQVR4nO2dXagd1RmGny/JicZE/CFFJceoKAaLUBVFaxqp1EiLokaIiqCIKL1QxAoWUdALCwYFK/iH2N6oiBdBa1VQY/xLSykEjYjmohFjTIx/SUwUmxPNeb1YW4xxZu+Zs2fN2uP6HjgXZ2b2rHe+9e41e89e37fAcRzHcRzHyQ5LLaAJJJ0CnA8sBA4FDonRDLAZ+AD4D/CYma2P0I5TFUmnS1qlNExKWi7p8NRxyA5JJulmSbsTdf6efCnpwtQxyQpJ96fu9b34VtKVqeOSBZKuS93bJeySdEbq+NSlUx8CJc0D/gfMSq2lhI3AAjP7OrWQqkxLLaAmtzG6nQ8wDlyfWkQdOjMCSNoH+AzYP7WWAWwGxs1sMrWQKnRpBDiT0e98gMOAX6cWUZUuGeD41AJq8KvUAqrSJQMcllpADY5ILaAqM1ILqMGcGsc+B9zTcPs3AH+oeOwBDbcdjS4ZoA6bzOylJk8o6aImzzcqdOkW4ETADZA5boDMcQNkjhsgc9wAmeMGyBw3QOa4ATLHDZA5boDMcQNkjhsgc9wAmeMGyBw3QOa4ATLHDZA5boDM6dKcwF01jr1C0tI++8eAfYDpvf93AxPAN31es1+N9utoTUqXDLCxxrEze39VmUYwRVNsbvBcUenSLWBdagE1WJ9aQFW6lBs4B/iU0U4OhXA7OdTMPk8tpAqdGQHM7CvgqdQ6KvBaVzofOjQCAEg6BniHevf3tjnDzFalFlGVzowAAGa2Drg9tY4+PN6lzoeOjQAAkqYBy4ElqbXsxZvAb7pUHaSzSBqT9GiqYkAF/EtSjNqEThkKpeKulrQlYcd/LWmZQvUSJwWS5kq6QdKaFjt+raQ79DMoEtm5zwD9kDQLOIpQTGLva7seOKfiqYrqCwj4hJB6vm0YnU4CJD1U4x3+UGq9bdGpr4FO87gBMscNkDlugMxxA2SOGyBz3ACZ4wbInNafBEo6APgtcAJhgae2TLgIOK7isWuBtn7W3U2YQ7gGeNXMvmypXaBFA0haANwCXESYkev8lJ3AE8BfzOy9NhqMbgBJBtxImMgxyjN5RomdhDfLX81MMRuKagCFyRsPA76g0tT4G/DHmItPxL7/LsM7fxiuIvIUuGgjgKTfAStitpEJk8CZZvZ6jJNH6Zzeff8Nwid9Z3hWm9kpMU4c6xawCO/8JjlZ0sIYJ45lgFGbsftzIEpMYxng1EjnzZnTYpw0lgE6s2hSh4gS01jp4XUSOJ8DNkXSUcZM4Ehgbu//zwkZvW3n9c+j+kTV2TGFNIqkrTUmYC5OrTcVkhbXiNPWGBr818DMcQNkjhsgc9wAmeMGyBw3QOa4ATLHDZA5U3oSKGkGIQV7nOInVHWKLp4kae9pT7uADYRU7H7VO0ceSfsB84FD+GlcTqpxqjFJZxVs3w5sMLNPpiixOpLOlvS4pG01nmANwy5J/5S0VNL0wQpHA0mzJF0r6RVJ37QUq88k3SvpxBgXtEDSypYupIzVkqJMimgSSRdL2pAwTpOSHpZ0cFMXtFjtveMHMSHp0kYurGEkTVOoFzQqrJN09LAXtUgh6KPEbkmXNdRvjSHprtSBKWCzpPn9dJfOCVQoe/Y28Iumg9UA/wdON7M1qYVAGPYJCR2jyGpC/cKJop39vgbeymh2PoT5Bg8qTD5NiqR9gTtT6+jDycC1ZTsLAyhpHvA+zdbQj8ESM/tHSgGSrgHuS6mhAluBI3oFt39E2QhwAaPf+QCXpxYA9FuZZFQ4GDivaEeZAc6Op6VRfq/wUCoJkmYDUaZrR+Dcoo1lBjgyno5GmUWYV5eKcbqz7M4xRRvLxNcpfPwV5YstGSEVfCbBbOodO0HIiy/jQKpnLc0HPqh4bNPUiZOAL/rsn06I1Rjh2icJj8Qneq8tYgyYU7H9wrK2ZQaok8Z9oZmtqHH8QCRtI5igCilrDdSJ03Yza+bpXA+FCbUvVjy8ME7+a2DmuAEyxw2QOW6AzHEDZI4bIHPcAJnjBsgcN0DmuAEyxw2QOW6AzHEDZI4bIHPcAJnjBsgcN0DmuAEyxw2QOU3MaC3K7x+Wrsy0rUNZfv8wDJ0KXhbo7cBBFc+xbFgRQ7K9I23PJiygkYodRRvLbgEfRhTSNB8lbHtDwrbrUliPucwAb0YU0iRbgI8Ttv8p7Re6nirvFG0sM8DyiEKa5Fkz65dgEpXekm5Ppmq/Jk8XbSwzwL/pxijwSGoBwIPAt6lFDGATsLJoR6EBeuvU/TmmogZ4wcxeTi3CzNYCf0+tYwC3mdnOoh198+8k3Q38KYqk4dgCnNrW8qqDkLQ/8F+qr03cJs8D55bdKgcZYAbh88D5EYRNlZ3AOaPw7t8TSccShtnx1Fr24F1goZn1S0rtj6Tpkh5ot7ZRKR9JirJ4UhNIOlzSW4lj9D3PS6r6LKfSxZ0l6d1EFzMh6R5JcwcrTYukMUk3SdqRKFYfSrpSFQtr1iqypFCUaSFwCfBLwkpWzbnsByYID3jWE9KfnzazlN/3a6NQInYJoTLHUYRYxUhl3wFsJAz3zwAryj7wOY7jOI7jOA4A3wFOAzEiIjqjwgAAAABJRU5ErkJggg==",
    "package": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAKwUlEQVR4nO2da6weRRnHf38EhFYxMSSILVdBbhGRywEqlatAIvgJiH6gKqA1MbFqDGAiikZRQAn1g4ZKjNQEIsSQCGhMOGlruLYUjKQtF8ECvaHFAtVeoO3fD7PLKaen7bv7zuzuOe/8kn7oOXtmn5l53vf5zzO780Amk8lkMplMZuBQ2wY0he19gHOAc4GPAwcB+xe/Xgu8DPwNGAbmSdrchp2ZyNg+xPZs2+vcO+ts32r74Lbtz9TE9n62f2p7U4WJH82moo392u5Ppkds72H7Kttr+pj40awp2tyj7f5ldoHts2w/FXHiR7PE9gVt9zMzCtsH2Z6bcOJHc5/tj7Td74HH9mTb19ve2ODkl7zlIC4/0PY4DBxOE+frkvVBkzh9nK9L1gcpcfNxvi5ZH8TE7cb5uowbfdDZVLBDTL0C+BFwQOTmVwG3A4uK/58KXAkcGPk+rwLfBX4jaVvkticutr9qe22CT+YGh8ze+8e45yTb19hen+C+WR/0gkPefmmCCbBDbD60BxumOGiNbYlsyPpgNB6J828nGPQnbJ9Rw6ZTbD+cwJ5xow+S47Cen2F7dYKBXmn7K+5jfW5bti+1vTyBfWttz7L9nphjOm5wuvX8TuN8H7ZmfRALp13PD7uHON+H7Vkf1MXNrOc32r7BET/9O+lL1ge94rRxfmessv0lJ8zPO+uD3eP28/ZP2J6euI9ZH4zG3cvb/94JtUHR56wP3O28/UbbP7b9vsRjMHj6wN3an98dqxzyA8niq4M+mGF7RQL719i+0nY39nEcvu4XJOjoBts/sX2d08TXxU6vDybb/mHRl9jMsz01pf29dHDI9quRO7bN9p3e7pl8p4+vhyUep4Nt35XA/tW2T0pp+646NWT7jcgdetz2tF3cM1V83ewQX5M+/297mu2FkW1/3fbJKe0eqyMHFzeOxQrbl7uHuOaR9fdLEe9fUu4fjDd9sM5NhQMHwRdr8Dc4xMjJNeyY5LDiSBFfm8gfxNYH/3QTD6ba/nkEY3eI833YM9VZH5TcnNJWbO/r/t61s3cT5/uwbcj2I/2O4BhE313cif0x9MFG2+9NaeSVfRi31fZvbe+Z0L4m9EHK/YU9izHa2oedX0xlH7YfijCQf7d9bjIjaUQfVH7CqAebzy3Gpl/mx7atNPCDtrdEMLDkXttHJDF2xObU+uDQCDYeUYxFLN52inSx7c9ENLJks+2bnH793Tl94HB+wU3FGMTm/BSD+O0EhpYkf7/OI+vvVQns71kfuJl9k2+mGMBbEhpc8qTtM6Mb/+5+pNyxXORd6APbZxZ9TE385aDt2xowvOQep9+/L59ZiK0Pttm+e3v7bR9a9KkpftXrOHT1leZLgGedMD8v6RVJM4DTgcdiNg1cCixxSJrdACwj9Kkpep7XrjoAwN7A14FnnDA/L+lxYBrwBWBNxKYnAd8CvgPsE7HdqKRygLcitnUgcBuw0Iny85IsaS7wUeBGoItnBG5J0WgqB7gbmAPEfCP2ROCvjrT+HgtJ6yVdCxwL/CHFPWpg4B7CmEYnlQNskDQTGAIeitz2RcBSJ8zPS3pR0iXA2YTTQ9tiEXCGpMuA/6a4QVINIGmxpOnAZ4HlEZveF7gGWOaE+XlJ84GTgC8T3vVvipUETXKqpEdS3qgRESjpPuA44FrievIUgj54zAl2GAEkbZN0O0Ef3ERafbCRoEGOkTRXkhPeC2hwFSBpg6QbgaOB3xFiWyxOAR4q1t+HRGz3HSS9CcwG7iWu7SX3A8dKulbS+gTtj0njy0BJK4v196lAzK+3cv291CHTt2+shm3vbXsW8AzwOeIerbMYmC7pYknLI7bbE63lASQtAs4ALiMc1R6LScD3gecccv99TZbtiwkTfysQU3SuBmYCQ5JiC+WeaTURVKy/7wGOAX5AiIGxmArcQdAHp1X9Y9sn2l4A/BGI+VjYVuAXwNGS5rR9eFQnMoGFPrieILRi64Mh4BGHvP+Hdnex7f1tzwYWAp+KaAeEldAxkmYVmqJ1OuEAJZJWJMzPXw78o9AHO6Rmbe9VxPkXCCnomKnn14CvSTpM0vMR2+2bTjlASZGf/yTx8/OT2U4flD8s4vwyQpyPufn0GvAN4ABJv4zYbjQ66QDwzvp7LnAUcDNx198HAXc4HCkzTIjzMV+/3kyw+XBJsyVtjdh2VDrrACWS3pR0NSP6ICbnFP9icj9wnKSruxLnd0XnHaBE0suFPjgPeLpte8bgaeC8Yj3/QtvG9Mq4cYASScPACQR98K+WzQH4DyHOn1jYNq4Ydw4AO+iDtvbv3yZseR9VxPkk+/WpGZcOUCLp9WL//mOEPfOmeBD4hKSZktY2eN/ojGsHKJH0fLFnfiFhOZeKZcCFkj4taUnC+zTGhHCAEkl/AY4n5Nj/HbHpdYSt7BOKe0wYJpQDAEjaImkOI/qgn+cTtzAS52+UFPNZx04w4RygRNK6Qh8cDzxQo4lhgrKfKSnmt0mnmLAOUCLpWUkXUc0JHpR0nqQu5huiMuEdYDtWVrj2xWRWdIxBcoDMGGQHGHCyAww42QEGnOwAA052gAEnO8CAkx1gwMkOMOBkBxhwsgMMONkBBpzsAANOdoABJzvAgJMdYMBJ5QDZseJTpdBGz2cOVJmo/1W49vAK12Z648gK1/Z8xlAVB1hV4drpjlAQKhMoDsY8vcKf9DxXVRxgaYVr9wLmOGF9oEHB9l7Ar6kWAnqeqyoO8CjhfJteuQB4IH8T1KcYuz8R3ojulS3A471e3LMDSFoHPFzBEIDzCad91yoOOai4KCpJOJ2syuQDLJD0Rq8XV1Xrd1S8HsKxrtcRzv/v+9i2iYxDWZvLgWcJY1bnrMM6c9QbtvdxqI/TDwtTHeu6G9urVDy5rQX7pjkU1eyHV1yxcGSlbwBJm4DvVerZjpTHut7lrA/K0rF3Ek5VH+qzueskpT0rwaHq1fw+PbWkdvHoGnZ36hvA8YtHD7up8OpQkHF1JMPtUEo9qT5wRxzARZx33PLxK21/uI49tVK2klYQCjf0rDZ3wxRGjnVtXB80RdG3x4C5hD7H4HXgIklVEnXvUDtnL2kx4Yi11XXbGIMhJqA+cNw4vz2rgLMlPVW3gb42bSQ9SRB18/ppZ3SzhCPZx33+wO9ez3+euMfMDwMnS2qzpE2giGtXOE051Cj6wA1qAKeJ8yVrHMa6e/kUpy3L+oT7KBvnhhzA9im2H07Q/7ccCmnGrwweG6cry2qHsnGVz+93YgewPSVxn2OeZdwMtk+z/WiCAalctt2JHMD2JNvX2F6foJ9LbF9Qb/Q7gkPiaIbj5g1KqpRtj+oADnH+UtvLE/Rrre1ZTlQutxXcsj5wRAdwjvP1cUv6wBEcwDnOx8MN6wP34QDOcT4NblAfuIYDOMf5ZnAD+sAVHcA5zjeP0+qDlypc+1wiGwYrztfF6fRBWwxunK+Lgz64ymn2F5pijUMf8ptSdXFafZCKHOdj47T6ICY5zqfE3dUHOc43hbulD3Kcbwvb+zlk/Ta1MPGbinvHrDWcqYOb1wc5zncR22fafjLhxOc433WcZn8h5+3HGx7RB/3kDzY6x/nxjYM+uMX2axU/8T+zPbVt+1PTvUeLE2F7b+Bs4CzC4RVHApOKX28Angf+DMwnvGM/4YpEZjKZTCaTyWQyAPwffKRm0g6K4rMAAAAASUVORK5CYII=",
    "phone": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAALwElEQVR4nO2deazdRRXHv4eWstSyt6HKKotQBWpoK7IIKKWAYFFAUFyI1CUgYAhLJRATojEBARH9R4yAjUZEISxCCoJaFoHKUpFNylK7UOhKgS60fR//mPfq8/l+t/fMb373d9+780maJu/e35kzM+fOb+bMmTNSJpPJZDKZTKbjsDIPAybpQElHd/+/j6SRkoaWV03vSlog6VlJD0u628xeTyA3UxZgS+Bc4CVaxzrgbuCouuvf0QCnAHNb2PH9cS+wZ91t0VEAWwA31dzxvXkH+FLd7dIRAFsDD9fc4UVcVnf7DGoIv/x27fweLqy7nQYttNewX8R64Li622og0nAZCJwi6Xct0qUsb0oaY2ZL6lZkILFJ0QfAlpKubqEuZRkl6fK6lRhoFBqApCmSdmqVIomYAgw0nWulXwMgePi+3WJdUjBM0jfrVmIg0e8cABgnaaZT1hpJ0yTdJ2lpSb0kaS9JX5B0mPO52Wa2V4LyOxfgEucs/BVg34p0+TrBDexhjyp06RiAWx2NvRoYU7E+XoM8tUp9BhNFk8C9HTKmmdlzKZRpwFWSPMu7/ApokiID2MEh474UijTCzNZImuF4ZGRVugw2igxgmEPGshSKNMFix3c3r0yLQUaRAeCQUSqoxIFHp0yTFBnAWoeMTVMokqmHIgN4zyHD87rItBlFBrDKIWN4CkUy9VBkAO84ZIxIoUimHooMYIVDxtYpFMnUQ5EBeJZ226dQJFMPRQbg2cwZlUKRTD0UGcAbDhmjUyiSqYciA1jokPGBFIpk6qHIAOY7ZOycQpFMPRQZwFyHjBGAZ/Mo00YUGcBrTjl5+3WA0q8BmNkiSW875HwojTqZVtMoKvhfDjmVRgRlqqORAbzgkLNfWUWawJNzIG8dN0kjA/CEeY0tq0gTeI6De15fHU0jA3jGIWdHoDJ/ALCbpIMdjyyoSJVBRyMDeNopa3wZRYoANpX0c/leAc9WoctgpNAAzGyuwoHLZvl4eXX+F2B7SXdImuh4bJ2kx1LrMlhpNAJI0t8dsg4to0hfgIMlPSnpGOejM8zsrZS6DGY2ZgCeX9I4IEl0EHCepL9I2iXi8RtT6JCRBBzlPJEzKUGZ33OW2Zt5wGYp6p6RBAwH3nN0wI9Klvf5Ep0PcEaiqmd6AB5ydED07BvYDlhUovPvJxxrzzjY2BxAkh5wyBsDfDBSl3PlO5LWm4WSvmxm2QPopBkD8J79m+xVovuXe6b3uW6WS/q0mWXnTxUAQ4FljqH4wYgyPhI57C8ADqyi3p3CRkcAM1snabpD5sH43cIfdX5fkh6RNMHMnoh4NtNNM68AKXjjPDJPc+rhiSxG0pWSjjCzec5yMn1o1gDulu+8YJU5fFdKusLMPAdYMwU0ZQBmtly+yeBYwLNF7IlCHi7pYsf3Mw1odgSQ/BlDv+b47lNO2WcDORq5lQAjgJWOGfpSQrbRZuW/6lwBTKuyvp1C0yOAmb0t6XaH7G0V8vw1y/WO70rS6UDyLehMA4BJzl9p00ElhDsJljrlPwkMqbLOmV4AmwBznJ30SYd8bz5AgPOrrHOmD8Blzg66xyF7OMG75+Fd8v1BrQMYjW+LuAvHkpCQGtbLg+RXQesAfuvsoD84ZA8BZkUYwSVV1jnTC+AgZ+d4R4HDu5/xsJYQR5hpBfgvkrrTKf9XTvkA/yafVG4NwOSIDmo69z8wClgcUcZ95PlA9QAG/NPZOY/iCNsCTo8wAIBrq6x7phvg1IjOOd1Zxm2RRnBWVfXOdENwDHlHgXnA+xxljAIWRhjAOsAdnpZxAnwuonOudJZxLP5VAYTNq09UVfeMNswFHnN2zFrgAGc5V0QYAMBbwEFV1T+jDet2L4/jmK0TglNnRBrBcvLOYbUQN1lzRfYQ3NDzI41gBXB4VfXveIA9CTeIeVgFfNhZzsfwBab0Le/Eqtqg4wF+ENEpT+E80ElYfsZMCiGsDs6uqg0GGknP0hFCwJ6VtJvz0WvMzLWvD0yV9ENnOb25TtL53ecekgNsp5A0Y1+FdLpNh8c1YK2kRZJelDTTzF5OIDMtwHERv8ou4PiIsq6NHAV6eABIdsUcsDlwJmGy6r3tNIaXgMtptwuzgd9EVGYJIRmUp5xNiNs06s1coFR2k249voU/mCUVa4CfEUad+gFGAm9EVOQJYAtnWUOBm0s24FrgUiI2kYBd8R2hr5KFwHHeOlQCcFJkJdzh3qQxAoBHgKavzSXERbyZoNyUrAcu9LZhJRA/PE+NKGsIcFOCBlwJXAQ0TEtHWI6uSFBeVVwW33OJIIR6vxKh/Hrg5IjyDPhxogacBRxSUM7OxL3iWo1r57USCMOkJ4i0h5VETs6Ai4n3E/SmC/g1sEsv2ZsCf0sguxW8QztETAMXRFZgKRCViJqQcCrWY9iXVcBVhMltqhGmVdzbqJ1aklSJEAV0q6QYN+zrkg6LcXoA47vLTbVOXqk0Dp1WM9HM/tTfB57TwdF0J286Q747CHoYLel+YNeIcmdKGifprxHl9sdA7HxJKvSytjStGrCvpEclbRXx+KuSjjSzORHlDpX0fUkXqcV1Vkis8ZCkl1X+HoPtJB0p/2Wd6yXtZGaePAzVQIjuiXWTvgbsUaLsScDrCd6rzXInid20wGYEp5W3DT35GqoFOKdEoy4gcmLYXfZI4Pclym+Wu6gwPB3/EbpfVKVLFMDVJRp3KeX996dRnRfvVVrglyeciWyWR6rWxwVhA8V7xrA3K4GTSuqwA3ADaXwGPawmrD4qBzjLoZd77lQ5wDBgeonGXo8zrKxAj0MJgSkpiM14GqP3MQ69FrdKLxeEnACeoaw/puHcRexHjyHAFOK3dNcT7jloGcBEh36eG+FbC2HPwBta3pcncMYTFOiyJTCVEJ/QLPMA780mpWGwGIAkAdtQ3giWACck0mcr4LvA7AblvUhwcye5KSVCx9IG0Fb59YGtJd2lcvcPIeknkqaa2epEeu0t6QBJO0rqUrhM6ykzm51Cfgm9Jkpq6OvvxTIza4+IoUYQ5gRlJoY9zAL2r7s+VZJiBGjJXoAHM3tX0gmSbi4pan9JjxOGcc+dg5l2gOAnuCbBSABhghiTkr6tSTECtD3AuaQJsV5L8D6OqLtOqegIA5A2bCAtT2AEENb5XwHa7vXnJYUBDIhGMLN7JB2kcCKmLKMl3STpMeCIBPIGNAPCACTJzF6QNEHSbYlEjpP0Z+CPDML5waCFEPl7PnGBpkV0EY63D6gLqFK8AgYswATg5YRG0MN0goOl7eloA5A2XGJxQwVGAPAPQtBF28YB0ukG0ANwInGZxJphGXAdvjuQWgLZAP4LsD1hW7hKniZs/sRca58csgH8P4QgiSrmBr3pIhwkvYAaT94AX3ToPL8uPVsOsAUhccKq9H3fLy8SklUcTws9jcCNDh2bvr5n0ADsDtxSQYc3Yi0hruFq4GQqut4O2J+QGKJZbulPTlvFA1QFIYL4SgVvYh0skjRL0jOSnlM4ITU79sZzwjb3nZI8c5GLzeyKvn/sCAOQNpxPnCzpcknR5woSs1rSHEnzJS2QtFjSEknLJL0taVX3v57AlpGSJilcxzfMWdZYM5vV948dYwA9EDaBTpJ0qULMQCfwvJmN6e+DAbMXkAoz6zKzWySNVQg8ebBmlVrBT4s+6LgRoD8IBzm+I+lk+YfWdmeOpH2K4iOzAfQC2FHSmZKmyJ/ssl2ZbGZ3FH2YDaAfuucJn5L0VUmf1cDNC3C9mX2j0ReyAWwEwu0mkyWdJmmiJFde4xqZIeloM1vT6EvZABwAW0k6VtJnFJZj3kQNrWKGpBPMbMXGvpgNIBLC2f/xko6WdJRCtFI7jA7XSzpnY7/8HrIBJIJwOHWCpEMUPI7jFU4StYo5ks4zs9s9D2UDqBDg/QpHyvaTNEYhdfxekrZNWMzzCuv8X8YchcsGUAPAtpJ2V0hft4ukUQpu3h0kbaNgIMMVfBI9/0vBJbxIYS9hpqTp/bl3M5lMJpPJZDKZhvwHzyn8vdmzstsAAAAASUVORK5CYII=",
    "pie-chart": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAMsElEQVR4nO2de7DVVRXHvwvkIUpKkA8SMUMKTW0qUQsFGs3R8dWYllo6vjUzUvPZhDo6Zk5ZouUj1NRKCx1fmfkoEAyFnDBLE6+agIo80yvy5n76Y587cwfO79zf+p19zu+cc3+fvxju/u299t7r7Mfaa+8lFRQUFBQUFBQU9DgsbwEaAWBnSYdI2kPS1pI6JC2StEDS65JekdRmZutzE7IgPsBewBNAB92zGnge+CVwArBT3vIXZAQw4HJgfYqOr0QbMAk4EOiTd70KUlDq/Fuq7PhyLAduA8YBTTO1NrygwDYK8/OOkt6TNN3MXqgivwskXRtJvCTekHSrpNvMbGmNy2pNgN2AXwNryvzaHgeGZshzJGEurxcrgZuAT9aijVoSQsdPATZ007ivA0Oced9V2/5OZB1wJ4UiJAMMBW6n+47vyu2O/Lemvr/+cqwFbgAG17ItmwqgD3AR8EGGBl0NDExZzlfj9WPVLAfOBnrXun27o1eehQN7S5oj6RpJW2bIop+k3VOm3TND/rVikKQbJc0E0spfE3JRAKAv8CNJf5O0W5XZpd1/b1tlObVgtKTngUvyGg3qrgDACEkzJV0sqdpKd0j6T8q0mznyfUXSYwpm4A6vUE76Srpa0lRgWI3L2oS6KgBwhKTnJX0+UpaPm9niSHl1ZZqZHWJmIyQNlLSvpAmS7pO0rAblSdJ+kv4BfKVG+ecHwfp2Bels7ml5H/iUQ4ZfOfK+qUI+vYB9gKuBuRHr08l64PtxWr4BAPoB90ZupDeB0U45oihAmXxHE4w97ZHrOBnwTFuNB7AVMC1io8wFzgA2zyBLTRSgS/4DgQkE5YzFo8AArywNATCYcHwag7nAcUDmNQs1VoAu5fQBTiGeIkwHPpJVnlwgdP4LESq/FDiLCEMhdVKALuX1By4krFWqZSbNogSEYb/aX34HocM+GlGuuipAl3K3A+6psj0AZlCD6SDqNhDoJ+lBVbfNmy/pQDM7zcyWx5EsP8zsXTM7VtLhkt6tIqsxkqYQ2fEkmgIQnCDukDSuimymSNrTzP4SRagGwsweUfA5fLSKbA6RFG10igowsYrhbR0wocby5TIFlJHDgEvxnXpuTGPZCYDDqqjQcmB8HWRsCAXoIs+hZLcdrAcOjCFH1VMAwaX6rox5zZf0JTObWq0czYaZ/VHSWGVbF/SW9Dtgh2rlqEoBCAuSexV86b20SRpjZmkPc1oOM5ujcAYwL8PnQyT9lipPEasdAS6TtFeG716VNM7MFlRZftNjZq8pLJznZ/h8f0kXRhUoLQQbeBa/+nnkcOxJg60Bysi3C/BuhvZcQxVOJZlGAMLQP1n+8/z/STq4+OVvipm1STpY0gfOT/tKmkzGqSDrFPA9pXfF6mS9pGPM7OWMZbY8pTXBcfI7oYyWdGZ8icpA8N7N4sB5Xl0ETJa7oaeAjWT9QYb2XY7TTV7KNgJcKb8D5/2SfpahrJ7K1ZL+5PxmkMKi3IVLAYBdJZ3oLGOBpFPNDOd3PZZSW52kcEXdwxkEu0xqvCPA5fIt/JB0spm95yynx1PydTzd+VkfSRM9H6RWAGCUpKOcAt1hZk85vykoYWYPS/q987PjcVxB84wAFzjTL5N0kSN9QXnOldTuSL+ZpPPTJk7VocB2CtsTD5cVV6Orx8wWSrrK+dmJpHSmSfuLPkPhGlZa2hTuxxfE4Qb5TMUDJJ2SJmG3CkCwMKXKrAtXmNk65zcFCZjZaoXtt4fTSPFSSZoR4ABJHtt9m8IJYUFc7pRvFNhFwY2sImkU4HhHoZJ0nZltcH5T0A2lEfXnzs++2V2CigpAcPI8wlHgcgXnkILacJt8h0VH0Y0TaXcjwAGSPP7od5vZSkf6Agdm1i7pHscngxUcThLpTgEOcxQmBa/ggtribePDK/2xOwU4yFHQS2b2T0f6ggyY2XMK7xak5eBKf0xUgJI5cSdHQfc70hZUh6etR1LBA6vSCLC/oxBJesiZviA73rZO7MtKCrCvo4BFCo89FdSHWQo7rrQk9mUlBfiCo4BpxXl//SjZWZ52fJLouV1WAUp7x10dBTzjSFsQB0+b70HCFfukEWCkfIc/sx1pC+Iwy5G2v6QR5f6QpACjHJlvkPQvR/qCOLwgn/dw2Qe1khTA86jx62a2ypG+IAJm9qGkNx2flPUVTFKA4Y6M5zrSFsSlzZG2bJ8mvb2zvSPj/zrSNgtjgVvyFiIFn3Ck3a7cfyYpwDaOjN9xpM0TzyHVKPnWQc1AWW/upCnAc917iV+WXMhyBbuVKPukbpICbOXIuFl8/p/LW4CcKXtQl6QA/R0Zf+iXJRdmqXmmq9h0SPpzuT8kKYDnYcamcP4smU9/kbccOfGYmZX1J8w1YkgOXC/f3rkVWCfp0qQ/JimAJ0Zu00TMLBlPjpe0Jm9Z6shEM3sx6Y9JCuCx7GWJ9ZMbZjZT0rHqGUowSdKPKyVIUoD3HYUMcqRtCMzsAUnj1ZpGLCn8gM8xswndHdMnKYBna+d+laIRMLNnJX1G0iWS3s5ZnFisVHi7aZSZ3Zjmg6TVvicOT9WPFeZFyYX9GuBahXd29lHwg3QHpMiRdZIWKuzzp5bWOalJUgDPfnknT4GNiJl1KBiKepyxKGkK8JhNR8YQpCAfkhTgNUceO9OscW0KEhXAc8bfW+Ed/IImJEkB2uTbJ+8dQZaCHCirAKWryP925FPxAmJB41LpLODvjnzG0QCh0Av8VFKAZx35DFa2Z+MLcqaSAnhunki+hyQKmgGgzfFY8at5y9uTAD4H3ArMAV4GpgBHpnkYylPIJOeL1a6AzgXZIERiTwrS9QiwRayCDnAqwM1RCi5IBDgnRT9MiVVYH2CZQwHagYFRCi/YBEJY3rSxiFNtzSu6hJXsAQ84ZBwo6WRH+gIfpyn9o11HRykRGOcYASCETW8aN7FmAegHvOXoh1QhatM4hT4t6Q2HrMMlneBIX5COkyV93JF+RbSSgYuco8A8wHO3oKACwADgbWcffDemAEOAlU4B8glo2IIQgk17WAHE9dUEbnIK0Q4MjSpEDwQYhj9K2w21EGQE/kihxavhVQI86GzztYDnfQeXMHc6hQE4sibC9ACAYzK09+RaCrQzIVath0XAtjUTqkUBtgeWONt6NbX69XcRzHs+APAY0NPuIWYG6AU8maGdf1oP4QbjMw938sOaC9ciAFdmaN8lxF75VxDwrAwCbgAKn4FuAL4GdGRo31PrKWQvYFYGIVcAnidoexTAPsCHGdp1BvWeYoHdCYsOL4uA4jLJRgC7AksztOcq4NN5Ce01EXcyH2eQ41YGGInf1NvJeXkK3huYllHwBUDZ50t7EsBuwDsZ2/BJ8t5dATsAizNWYBHQY72JgX3JNuwDLAQ8D3rWDoLrmNdM3MkKeqC1kLDa9x6wdbIW8EZ0qS3AuRkrA2GLeBl5D2d1gDBtXkW2rV4n3867HmUBbqmiUgBPECKVtyTAUOCvVbZR/JO+WACbEdySq2ExcFTedYkN8A2yz/ed3E+jX8EDNif7zmDjynpcoBoSYEfgoQjt8RQhlG/jAwwEnolQ6Q8I3jDN9GaPJAnYAphIWORWy1Sa7REOYMuS4DF4i3D+0DfvenUHwXv3O4RtWgyeoNk6vxPCdPBwpIaAYDy6APA8Z18XCKekF5PdoleO+2iWYT8JwsLw5oiNAuHA5HZgDDEvQ/rr1gsYS/CUWhW5jtfT6As+D8AEYF3kRoJwEeUnwP4kxMaLXI8+wHjgOoL7e2zWAmfXuh6d1PXXA4yXdK98IWk8tEuarhBUcbakOWZWVUALwnTzWYVHJMcoxOGt1f3HhZK+bmYzapT/JtR9+CS4it8t6ct1KvIthUev5pX+vUQh7u7Gj2D1VwiV8zFJwyTtqPAGYr1eQn1S0rfMbFGdypOUgwJIYe6UdL6kK+WLUNqKrFJ4r3hSHvGXc1tAScEJQuFxY0+k8lZihqRTzSy311VyPXwxs5cV5tUzJS3LU5Y6s1jSKZLG5tn5DQUwqLSyzuJm1iysAq4FPFHZehbAcGAyYTvUKqwh2EKa9mn9ukNQhOvxX4xsJNoJo9qwvNsziVwXgWkg7MNPknS6pHy8X/28JOlmSXeZWXvewlSi4RWgK8AXFaJ+Ha2wX28kFkn6g6TfmNnsvIVJS1MpQCcEG/l+kg6VdJBC7J88eFHS45IekTSzFJyyqWhKBdgYgnfsGAV7wl4KptvY4exWKHT4bEkzJU2vt9WuFrSEAmwMwdI4XMGUO0LBrLu9whnEEIXg2J02kAEK0bbWKXTyUgVz8TuSFihET5kr6c08LHUFBQUFBQUFBQUF0fk/HZ6A9wqjMqQAAAAASUVORK5CYII=",
    "puzzle": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAPg0lEQVR4nO2dabBcRRXH/w0hBEQRBREiEEQQArgBpUQRWQSRRShRsBDZ3MWtLD4o5Qe1tEBLtFBxAUtAlNJSUUzIRgKiBQmCG4YHJUsgAfIeUYxKSF5e3s8P504cnjPz5nT3vXMnmV9Vvrzc2326z5nb2zmnpQEDBgwYMGDAgAFbHKHXAlQFMF3SmyUdLmmmpOmSnl/891OSHpM0JOkOSQtCCI9XLN+ehXyvlXSgpD0k7VD8978lrZR0r6SlkhaGEB6pUr6+BJgCnAn8Bhine8aBW4EzgK1LlG8acAGwxCFbgyXFu9uWJV9fA5wK3B/RsRO5Dzgls2wBOA9YkUG+FcC5wBbzNe8IsBPw0wwdO5HrgR0zyLcX9kXKzS3YMLLlAuwH/K2Ezm1wP7BvgnxHAU+WKN+TwJE5+7RvAA4Ehkvs3AargIMi5DsZWFeBfOuAk8ro49pCdcpvMIzDCIA3UY3yGzwDvLEb2fp+4gAcKGmxpBdVXPWIpGNCCH/t9BCwl6Q/SHpBJVL9j9WSDgkhPNrpoa0qEqYUil9hL5Svos5Fnb4EwFaSrlX1ypeknSVdzSSrg741gKLjF6k3ym8wmRGcK6mrT3FJHCXpPZ0e6MshoCbKb2ZE0tEhhGWNP2AbNA9IeknPpDJWSNo3hLC+1X/2nQFkUP64pLmS5skUFCTtI+mtko5X/FfxWUYAXCDpqohyNkiaLelmSQ8W8u0t2yY+SdI2EWWeF0K4OuK9egEcRNps/85i0tiu/IOBuxLKH26UD9wR8f5sbNLYTr4ZwJyIcm8vQx+VQrry5wHTuqhnu+LZWIaBE/CdPQB8jS62c7Ft5MudZY8De+TRRA/AlD/ibHQzXSm/qb5pwNyE+jY6n5+NYy8fM4KbnHWcH9f7PQbYBXjU2dhmXMpvqjfVCLplFJgRId/exbvd8l1vHbUAO4SJZS4Rym+quwoj+EWCfL901HNbqzJqvQ8AHCLpjMjX50k6LYSwLrb+4t3TJM2PLaML5ia8u9Dx7PRWf6y1AUj6gOKWqsnKb1CUcarKMYINkn6V8P7Djmef1+qPtTUAbFJ0csSr2ZTfoEQjuD6EMJK5zM0DYI+I8TRpzO9CpmmkLRGbWUPi0gy40FHfQ63KqO0XQP4t1Oy//Ilk/BIg6YMhhBWJ5RzreHZlqz/W2QCmOJ+fKelS4NVlCNOgyQhuiiwCSZ8IIVyfIgewj6QTHa8MpdRXOcDMhM/rYuANJcs3BbgEGHPINQLEzGsm1r0V/qHoghztrgxgKubZksJ1QKln8cArMCfU9R3kGAG+BOyUob6tgG86+2GcNg6jtT4NBObJTuhSWCHp9BDCnRlEagvwPElvknSwzAFkTNLjMm+gO0IIYxnq2EfSFZKOc756ewjh9an1Vw5wltPS27EW8IyXtQHYDXg/8CtgQ2T7z+l1O6LAxtmhPDbAesD7y+kZ2F7/jxKU3mA5MLXX7YkGeAO+iVYn1gAH9LpNkwGcA/wnU5s7uoT1BcBHMnUGwF8pcbMoFeCijG29mc0lZAz4AL7jz058vtftaQXwdvyOJO0YxiKiNx+AQ4ClGTpnbd06B9gZWJ2hbWDL5/6b9XcD5g1zLPAD4ImETvp6r9vSDPDldL0Dpvy39ro9zwLYE5gFvBEL48oyK8U2Rd4C3B3RUWuA5+SQIxVgG+DvGZQ/Qsk7oF0DvAb4Hq1/peuA+dg637vn36quKcC3IjrszBxtTQWLHUxlIdDrGIRNY9mP6H4ycx9wRKa6r3V22g9z1JsK8Cmn3M08ApxNHWb72Of9kYhGjAEXZqj/ufjmBS3PyKsGuMzZX+PA74D3UJdNHvKEaH8wgxyfcXbkc3O0P1Fmj5//BiD1fGQTWfwByBeifTlwaGIZNzieDZLaRuJUyGrHs1MkXUOHCCcPyQaQUfmSxb1dQdp4dr+kUcfzvQjdnkjHHAMt2FUWlZxsBEkGkFn5DQ6TBUJGEUIYl+RxC0tehWTgFpmHsIddJS1ONYJoA6Dc5Azvjn0R2F6SZ1x/JrauXIQQnpJ0Y8SrL1IGI3BDeqDmZKzCsmvEyHaUs66X5u6fGIBXEX/quSkquQpBy1Z+g6jtTOBqRx3ryLARlQvg0oT+Kt8IqE75AH8EXMkQsF+Rx4GiVDcxL9iOpjfqtxlX9jKvcDni85c73+l6RYBFEXtTw15aSmclgAWfzHe2o5n8RgC8nAzx+VgSBC8/xBwuO8n3KuLyAr8ua0dlgvQIpFUkZDSdKMxOpKVf3RSfD7wyoUEXAwdQTA6B52ATvmuI85u7lzrsn7eBdCO4nwy5jUVa4uX/S86AuSmlMAo8lVgGwPuSO6dkSDeCH6cKcGpC5S0zc2B+AbncnmJ5kLocokwC6UYQ5w6PzUhj8+13TMuC/9g2N/XymJkE0ozgXmIuvcBu2ohh0hBtbF4Rc2ycgyujNdFDSDOCd8RUGHOpQdfx+cChwNORDYrlTmA7d2fUBOJzFi3yVrQ7/nHanZwBy6dXVRr1IaAuqWWjIe5LsBF4saeS85wVLCUy2AJ4M+agWSZL2QyU3wBLZunNaHq2p4LvOgoeI3EPGlvf/8XZoG65kj7+7LcD21PxJKS8wlP4bx0Fz87UoGnA57CgjRw8QJ/N9r3gGwpu8RT8kKPgZGfOCXXvjm0Zxw4LQ8D76JN1fgrAxx398kCrMtodhXocKh70i96e4sbOTwIXy9LEnShpliyleyvWS/qTpN9IuiGEsCSnPJ0ojGyGzDtnW1kq+tWSlocQ/lWBCC2V2oaW28LtDMCzR17KfnoIYa2knxT/BOwg6+wXSJoq6WlJqyStyJF9oxuwc4gjJJ0iu43jILXJ3w8sl/RbWTKpX4cQni5BpHLOMrCt0m7JOgTUEWx+8lFsXhHDGsz1O2vKduATDhn+5inYMwmck7NRdQM4Db8fQzueAT5Ppvt+8fkNuCaB33EUvBE4OEeD6gR23FzWmcU9wP6J8r0a3zLwm57Cz3U26C42o7U28GLgD84+8PJPwJPps1m+7TGXOQ/de1pjSzHvrRdRFzPUDUz59znbHss64ASnfDFbwWPArt6OuDWiQaUmay4b7LNf9i9/Ik9j9yJ0I1/sYdDNMZ3xzsgG9e2XgN75KSwHnj+JbCm3l5we0xlbU5JDSB3BZvu95AcdZEvxBVhGjENIUfEpCQ3qGyMoOjjXUi+WceDwNrKluIS55hitOifl0qa+MAJ8Fy+0YiPwD9KTWy+cIFeq8tMzoAA7Ej8UQM2NAEs0FbPDN4odmx9BsbGDZTDbH/g08Hhkf72iKCtV+UPkSn4BvAzzzY+ltkaAZS7zsoxJfCCwdDXXRZT9VdKV/wSWWTxrR6WGhs2nhkaAPz/PMmDnLssOwLed5T9M+tW1tY0PvIkaReNKEr51/yhO7ycsB+CfE/rMQ3nKb2pQajKoS0oV0AGwLRVcvwqcmNBf3VK+8psalGIEY8ArKxF0ErCgVw9RWTixiWaZYfXRyo/KwhFCWCbpaEkxlx5uLekzMfWWgGd/fKOk38dUUuQtWhDzbheMSDomhOBNNCUpIUdQohGcyiRbnxXhmZSuCSGsT6jruoR32zEi6ehY5UuJWcKajGDY+epUSUem1J0JjyvZ9qSFky+QdFfC+xNpKH9ZSiHJeQILAY6R3wiqzWzVmn84np0mab/YikIISPqw/OngWpFF+VKmTKGFIOfLxslu2SVH3Yl4bt+W7Cr5aEIIv5f08ZQylFH5yWDLqPcCSyJmrl/utfySee86ZH6MDNur2NlDTFaTh0h0JcsCtqw5H1gZ0YgGqb+ELGA5iDxcS4bUMsDrsV3FbhjHbkhJvnk0GWAv4DZnp7Uiyx0BqQDvipD9Cpwp7NrUPQXLwzCP1iFxK7CLMOrhdAscTb6rTWqxJQzsAPwrog1/xMLbc2Vc3wbbmJqFXY7VfTh3FQBvI18s/xd73Z5mgG8ktGUVlqnsOGqcdSwJ7Jff6XZsb4f1fixrAtiDdGcOsAwkqfcd1AtgBvnus9sInNzrNrUC+EKmNo6S4eaTWoDN9nNM+MBmsh/rdZvagTlh3JOprQAf6XWbksGWejlYQ02uaOsEMJN86Wo2ALN63aZosF/EY4mdMIqtX7NGxZYJNpnLNd9ZRk1WO26wLBsxrAd+ju0Q1mG71w3mwJErhd0ZvW5PJ9ouW7Bc+oc5y5st6cIQwiNJUmlT9o1Zkg6VNF2WgWNY0j2Sbg0heA5yYuo/TNLPJO2ZWNS8EEKab37VYDN/L5eRZ4t0Z+ArdN5wGgV+RsmeRVhG02tIy228lgy7hpWC//N/I3mUfzq+ncYx4BJKHmex3blFzj5pJvoYuScA33c0bj2Q+pkUcBHxv7Q5VOByjl1M8TX8YWT9tRoAbnc07ucZ6vNmJm1FJcEnxEXpeudSvQV41NG4DyXW9VLyzbhLNQLiQ7R3K0umUsB3I8dxiXX9OKJDO1GKERCv/JW5ZclJjqPM6MkfsLskfy77zhwv6Zc5jaAo6wZJb4l4PeZG0MpoZwCeLJd7J9R/ssq5uzebESQqH0nfTpWhTNoZwGOOMqIvepZ0UsK7k3G8pBtSjCBR+ZJ0bQjhntj6ewZwlWOMGwVmRNTxMuLvyvUQlbiK9BDth6mZ30PX4N8ImoNjIwg7ZvZeH5diLC4jIF35I8ABcb1fA4jbCr6cLowAU743Zn4cO6VLvbp20mSWWBLGesbnVwl2zYqXm4C2k0Lss78gotzbivdTQ9PvpsP5AZZ+1ZuBs5m+U36n08D3Soq5Ym1M0hxJN0t6qPjbPrLJ4gmKm/WfHUK4rpDrQEmLJcXeATQuaaGkuZIaGbT3K2Q7VvFL46Qo3V7RyQC2lV0GMb06cVqyXNJ+IYRNMXUZjCA3fal8qYO1F6HQn61QlnZc3Kx86VkBqTGh6bnpW+VPCpbkKCZncC7m02FiSXrOolT6bsyfSDez9j0l3S2pq8xYGRmW9JriDqG2FApYpOqHg/pE6ZYNcCTV3fAJ5kXT9Rk61X8Jhkm8K7HvAE6iGiNYC7i3XjEjSElm2S1PsKUpvwGWVfPJEjt3FS0SJjvk25e0tLaTMUTuDJz9BhZHd0sJnbsAOx5OlW9H8vsYgOURyJN7t9/BVgfn4PMcasfDwFlkjqzFfPuHMsi3jIghaYsAmIr5892Oz6FzHPgdcDYlukxjl168A1iM7w6kjZgX8Olkiv+vK9l+dVj417GSXifpAEm7SXph8d9/l/S4pCFJSyUtDCFU6iqFJVw4TtJrJc2U9JIm+VZLWlnIt0TSghCCN+vZgAEDBgwYMGDAgAH9wX8B33kWzQSee/wAAAAASUVORK5CYII=",
    "refresh-cw": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAL5ElEQVR4nO2deaxfRRWAvwNI2cQCggmggLggVFCCiYbgghENoBYlVEBRNBJqMC5xgWhEoola0BgVAXEJIURERVkqigtRVFRCpWoRQUVAlLZ2oaViC7zPP+b35FG73Ln7+/V+yQtJOffOmbnnN8uZM2dgYGBgYGBgYGBgYGBgYGBgYGBgYPyJrhVoGnV3YD9gL2BPYBdgp9F//ycGrAQeBFYA/wD+CdwVEYtbVbgA6s7AbODQ0T8tAL4bEaty3zVWBqA+HTgcOAx4PjCLx3/oMqwCbgcWArcANwGLIuLRiu8thXoCcAGw63r/azkwNyKuaF+rjlB3VI9TL1bvtj1WqteoZ4yMrq36zlEnNqHXhDqnLX06QX2COlv9lvrvhj90URaqH1b3b7DeO6vLCuiyzDREjBfqPuon1MVNfcUamFB/qr5RnVFz/d+coccpRd+7VZ1KNoF6iHo58GfgTGCPjlXaFAG8GLgUuEf9qGkSWgeHNCHbWwNQD1S/DfwWmANs07FKuewBnA38Tf20WtVwd8yQ3amoYO8MQN1dvQj4HfA6pv9KZQfgvcBfRj1CzodsnN4YgBrq6cCfgNOArTtWqW52IvUId6gnqb0w7F4YgGn2fANpfVt13d539gQuA35oi0vIjdG5AainArcCL+lal5Z5OfA79R1d9gadTazUHYALgTe1VORq4B5gCfAvkodv0pu3LamL3hV4CvBUoI219I7A+cAx6lsiYmkLZT6OTgxAfRpwFfC8hopYDNwI/JrkJ1+U69NXdwOeQ1pSHQa8EDigZj0nORr4rXpCRPyyoTL6gfoC9Z9lPS2b4Gb1LPXgBnV/isnJ8w11VQN1WKuetpGyL8p4z0VNtUEl1FepD9bYYItN3sFndVCX7dTjTXsCj9RYJ9Xz1K3WK296G4Bp02ZtTQ20yOQa3bbrekEa0tRPqstrqp/q16fWz+lsAKaPv66GRrlTfYPr/Tr6gvpE9Uz1XzXUVXW+ut3o3dPTANRXWv2X/4D6XvUJXdenCOqTTD3CQxXrrfp903Az/QzANOGrOuZfre7VdV3KoO6vXlex/pqM4MsZ8t0bgGlcvL9CpVebnETTHvVN6ooKbaF6b4Zstwag7qAuqFDZ29Sm1tydYPpB/KxCm+RQ2ACamkxdSIrJK8N84IURcXuN+nRORNwDHAmcSwpCHU/UU6tYrjpuu4D/h3qi9UwQN9qOXVVsf9PYXYZ5nSjdEeoR1us3mEr7Q4Bpbf4VMqJRpnBuRHygLl2mAxFxI/BSoPUNoEZQTy9prRd2rXuXqLPUJbX99hPtDgGmMK4y3dm1bgFj/uawfiMobAB1bQd/nPxIntuAk7o6YVM3pmXru4GXAWUigbcmrQ5aDQ6pbADqgcDbMh9bA7y+zFm2PqK+lRTO1ovNqRzq6AE+Rn4A5zvHZZ2vvhy4mB6E15WhktLqIcBxmY9dExFfq1Juz5hH/z7+zKKCVRU/k7wxaxUwt2KZvUHdh8eOaPeJfYsKljaAUeWPz3zsnIi4r2yZPWSfrhXYCK2cDDqdvDnEncDnK5TXR1Z2rcBGKKxXKQMwBWa8NfOxj0TEw2XK6zGLSCHmfePaRt9uOp+fwyJ7GsZVFfWD9flvamGNGc61sh/ljZny8yJiomRZfec84OqulRgxAcxu1LlmSsuyJsMiF1tzsoS+oW5jCgbtMnnF3eoLcnXPdjuqs4HvZDzyyYg4K7ec6Yi6DXAQ5VzBuexEyn62FpgfEXe3UCaYEjLl0PqhjYEGMS8b181d6zuwabImgep+wNMyHrkyT52BtsldBRyeKT8/U36gZXINIGeWeT/w+8z3D7RMrgHkbHzcGBHjG/48JuQawEEZsr/JfPdABxQ2AFOeu5ywrwX56gy0TU4PsG/muxdlyg90QM527t4Zsqv7mGd/XFBnkpJoTuZYuhW4MiKa2542pTMryh8aU2QLRz3ZDZ80XqGenPu+nB5gtwzZ+3MVGdg8ow98KRvew5kJXKoSEZcVfWfOHCAnb96KDNmBAoy6/S+w6Q28AL4wki1EjgFsnyH7YIbsQDFeR7Fo38n5QSFyDCBnuBi30K8+kJNUs7BsjgGsy5D9T4bsQDFyeuDCybRyDODWhmQHipFz7KzwjzXHAK6kWLjxSoZt4CZ4YobsQ0UFCxvAyMlwBpvObyNwRqMOiS2X9e8J3BQPNKZF3Y6IgWKot2c44jaYcHpDlDqLbheuyC0Y04USayg+EXxtRPQlVH2gKureGb9+1cIp+sbytM4YMitT/q9FBQcDmB7kOIGWREThSeBgANODF2XIZu3EDgbQc0yHanOisRfmvH8wgP5zKHlb8b/OeflgAP3nmEz5XzSixUA3qAszln935r5/6AF6jDoLyLkG74e5ZZTOE2hKEnUMMAO4i3aCQJaSLoF8pIWy+sBbMuWva0KJx2G6ByjnhHDdLDYlY+js2ts2ULdXl2a0yypHN4w1qdQr1Efr/6aluMoxTjStnpbZHoUDQcsqtLV5qWHa4IONVrojTG19R2ZbHF2mrJxJ4AeAHcoU0iDvczx7gZOBZ2bI/wO4vkxBOQZwbJkCGubJ5B1Y7T2mcfyczMe+WnZinGMAhWPNW6avepXlPeSdw1xHuqWtFDkG0NdY/3ayY7WAui/w4czHLq+SfznHAP5WtpAGWdBaerSGMUX9XEjePGsC+FSVcnMMoG/hXhOkiem4MBd4ZeYzV0TEbVUKna6u4HXA2yPix10rUgfqwaSUszk8DJxdtewuvGkCq4Ey+WyXAjcAnx2jK2dmAt8i7+QPwAURcUfV8rswgCClNz0yIrboPAKmtPvfIG/NDylF/Ufr0KGrIWB34Cem3a4tktGk70vAUSUef19EtHsE33Sxc90sGY1/WxzqZ0q22fUm42ld4SYMQNONo0e0XqGOUMPyH3+lmpOqd7P0YRWwC3C9elLXijSNacz/KsnbV4bTI+KeGlUqjs31AJNMqPMc031+dab6gwrt0+6F0BuoQNMGMMnPrLmb6xr1YPO3d6fyGxsK9mhqCPh7hWePABaqp9SlTFeYxvt3AL8if6k3yX3AcRHRSNaVpgzgB6O/sswELlG/rz6jJp1axbSxcx1wPvlOnklWA8c2edlmUwbwKDCbakYAyTf+B/VTZqQ+6xJ1O/UsUqrcXN/+VNaSbgDrR7od8+YAF42e2U6dX2Hsm8oy9UNqTr7C1jCFcZ2i3lVDXdepr+m6To/DEgYwem5b9es1NMokK0w9Qi/u7TVF755mtUneVP5j3z4+lDeA0bNbqefW1ECTPKJ+T52jth6rqB6knmde6PbmWKUe2XZdCmEFA5jyjrera2tssEkeVL+pnqru1VD9t1IPUz9i3nGtovxdzckDUAutOl0i4mJTJvFvAnV+qB1JV9kfD2A6I3cTcAvp3qI/RkThBNYmX/vewIHA80nn8w8n75RuDjeTlnqNzfY3Rutet4i4yZTD5mvkn3wtyjNHf//zJagPkeIHlwLLSUmXJuMcZ5CSYe8G7EEKymz2lM1jfAl4V1Pr/M3Rids1IpaqryaFQc0j/YKbZnvggNFfH3iA5Nu/vEslOtsMiggj4ouk06/Zp1qnOT8Cntv1x4ce7AZGxF8j4ijgRJLbc5xZSjrxe1RE3NuxLkAPDGCS0a/h2aRQp76eQSjLw8DngAMi4pI+3afYGwMAiIg1EXEO8HRSlOyajlWqygRwOTArIt4VEcu7Vmh9mjoZtDpXkalExNKIeD9pNn42MN1uIFsHXEL68CfWEb3bOeqbM5watW7lqjNMSapvMAWO9JX71I+pe9ZZ/16g7mzakNkcy2xww0bdTz1LvaWxz5jHKvUy9WjH86j6Y5j87pv6BU6oc1rUZx91rnqlKbi0Le5Uv6gea9NpWRomO7xYPQG4gP+/wGA5MDcirqhDsVxMGTUPJLltDyX5F55D3n3HG2IJyZ28kHQh9s+7cNk2Rdn7AnYmBXxMXie/APhuRKyqS7G6UHcF9iP59vcgGe4upFO4M0Zi60grjgeAZaSLL+8F7hruQBgYGBgYGBgYGBgYGBgYGBgYGBiY/vwXLbwihNGd63sAAAAASUVORK5CYII=",
    "rocket": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAANh0lEQVR4nO2debCXVRnHv0+ipimuuaKokLtMmaOklIG5i04u2ahkmhiRMmOpYWo0pZELNQOOqA2iTrinI5OQpjSmJe5jhORGsqggm6yy3k9/nHvpcvlxf7/nfc+7cHk/M/ef333POc8553nPe5bneY5UUVFRUVFRUVGxyWFFC1BRH2B7SftL2lfS3pK6SNpF0q6StpO0vaQtJW3dkkTSVElPSBpuZos2lHelACUD6CKpp6QjJB0u6TBJu6XIcpqkU8zsrVr/rBSgYIDdJB0vqY+k3pK6ZlDMNEk92hsJKnIEOBi4HngVaCIfrqslSzUC5ASwp6QLJJ2vMKznzWtmdkTbHzsVIMgmA/A5SadKulTSyZI2K1CcfWr9WClABgDbSLpE0iCFmXsZoNaPlQJEBNhBodMvl7RTweK0ZVWtHysFiEDzG/9TSVcorMvLyIpaP1YKkAKgk8JQP0Tp1up5sLTWj5UCJAQ4VtIIFTOjT8LiWj9WCuAE2FnSMEn9lP8yeqnCFu8HkqYr7BYe1WDaBbV+rBTAAXCupNsk7ZxDcXMkvSjpFUlvSJpkZtPbyDNKjSvAnFo/VgrQAM2z+5GSzs2wmGWSnpX0V0kTJL1lZjWXbq3Y05H/rFo/VgpQB6CXpDEKp3CxWaJwYveopKfM7DNnes+5wcfOvDdtAAOuBlZlsC//HNAP2Lq+JBuUbzNguaPMM2O2T4cG2AZ4JHKnLwNGAgdHknE/Z/lfjlFuhwfYG3gjYscvBG4AvhhZzr4OGZoIm1XrUc0BWgF8RdKTknaPkN1nkoZLutnM5kfIry2eN/pDM1uSgQwdB6APsCjCG98E3EOw7MlS3iccMo3PUpaNHuAM4LMInf860DMnmWc55Bqah0wbJcA5wMqUHb8MuIpwNpCHzPs75TsnD7k2OoCzSL/Mexk4MGe5L3XKmMUexsYNcDKwIkXHNwE3AZsXIPvDDjln5C1f6QGOAZam6Px5wCkFyb45sMAh65gi5CwthO/n3BSdPwnYr0D5j3PK278oWUsHsCPwborOHw90LrgOtztl3qdIeUsD0Al4NkXnjyanWX6dOsx2yFzTG2iTBBiWovN/DxTuQwGc4pT7pnp5lnIrmGB101XSXpK20f+dHiVppaRPJH0oabak2fXOzYGzFQw2kzDUzH6eMG1sLnY+/1i9BwrXakkirKP7Suol6Rj5TKpXKyjDZEmTWv1NNrM1wL4KFjVJrHVvMrPBCdJFh+BDOF1So8vOaZL2bcCopBiALYABwEsphub2WAg8BUxPmP72otuoNcAQp/x1h/9CIBhanE262XjWPEJw6yoFwFb4Jn8AhxYt93oA2wGPx+6tyLwIbFV0W7UG+LGzDhOLlnk9CBsw70TvrrjMIHxrSwPhU+n9jF1StNzrAHQHZsburcgsB44suq3aAgxy1mM+8IWi5V4LsDvJJ2J5MrDotmoLsC3+b395Jn+Eg4sXondVfB4vuq1qAfzWWY/lwB5Fy70W4Jb4fRWdeYSNp1JBmDN5zL4B7ixa7rUAPYHVsXsrA75bdFu1hbBUfsZZj5WETS8XmWwFA5tJukvpQ6IsUdjla23R2lnBJSqxU0UrPjCzByPkE5uLJB3nTHO3mf03C2HcAN9L+DauAMYQTLXaNc0mTJAOBE4Hfo3fumcNcEhebdIowF7Ap866LKUs337CunWqswIA/wS6JSzz5gTlXR277mkBPgdMSFCXXxUt+1qAyxJUYAIJd98Ib4zXpHsiJdrqbQG4NkHbzaQs637C2/+RswJTgR1TlHmns7w1wHrx8oqG4JiSZNKcpcu6D+DcBJ3xjRTl7Yvfpn9UzDrHgOCP+ImzHgDjipZ9HfB/v+5LWd4oZ3mLqDO5zBvCZPZNZz0AFlMmez/CjNwT93YlCSd9zeXth9+p4/qYdU4LYad0nLMOLQwoWv51wD8Tvy1leXc4y/uYskyWtHaz5z5nHVoYTwlsFNfSXJkPHBVYQ4rhC9gN/8z/sohVTkVze41wyt/CLEp2ZC3gSGcl/pyyvBuc5U0HtoxV37QAv3PK38Ia4ISi5V8P/Ic+fVOUtTV+z55SfC8Jb/5wp+ytGVJ0HWoCTHFUYhbhrCBpWQOcjTYd2CJmfRPK3YngXJKUsWSweZX6MAjYS5LHPfozSbdDYmtlr6/7rWa2MmlhMSDE53lIUlKH0kmSLjCzpnhSRQK4OIVWZ80cUoRii9Q+XUgXdOojMvTvjzGkHBshj6y4w8yWFVU4cLRCqNekIdoWKtz4Nb3uk0UBvB3hTc2C5RS4XAIGki4AxTJCRPLyAuxEfrdeeflLQW3SGXggpewrKCgAhQvg6NTdlB1N5Lz5A/Qi+BakYQVwel4yp10FlNPxMGCShgMys1TbzvUgzPLHKlz8mIblks42syfTS5UDhOEubYi1rMl0JABOApZEkHMRkFaB8of037s8iK4EQFfiBZT+GPhqTPlyA9gDvxVQEURRAsL5/Q2EWXoM/k2ZzvWTQPD9eyVSg2RJYiUgnN0PxBeitR7jgEKvmYt2pkzYp+7d/Bc1NHoNDlGIJJIEJA1qdGJIOLc4T+FquMQGLDVkuFXSNWa2JlKemxYE6+OkexB1RwLC4U0/fAddjbCQKnZvHMhACQjROH4EvB+nv9fhU2D/Itqqw0IkJQB2BX5JMuvcRplLmUy5OgqkV4Kn8HviJiXWPCIKpfOOSULzhG6Qku1MmqQTJOVlMlaqKCQdQgGk1EqQJ6VySO0wCiBtNEoQ5dq4WHQoBZAKU4JVjmcrBciaHJVgtaRRki53pOlGCYxUW8g9WDQhHs8BknZpLn+VQvDnt81sXqxyzOw2guHpcMWPidwk6RFJQ8zsbXxROTsp1H1mZJnKCSHowfHAXdT3Hnqf4PLVh0jrZdItEduyBngQOKxNGVs58zk8Rt1KDaHjLyJ5TOAphK3Y1J8pkgWtaM0K4G7ggHbKWOzI78S0dSo1QA/gtZSN3sJE4KAIMiWJVzifEKuvbvwdfFvH/dLWJxbRJ4HAhZJekhRrmDtK0qvAeSnzmex4dr6kAZL2MrPBZvZRA2k+ceTvuQ8hU6IqAHCtpHskfT5mvgoh4f4I/CRyvhviUTO708yWOtJ4ns39rsENEU0BgKsk3RArvxqYpGHARRmWkQbPuX5plt9RBAHOkpRXkOKRhGvey4Zn1FudmRROYsywu0sarfzuH9pS0r0UcGVrHTymXYszk8JJKgUgmEvdJ2nbOOI0zGGSynUpguRx4FyYmRRO0o4AAyV9LYYgCRhMwRc5tkDwQdzBkaQ0zp6JFYCwpZsmROkCSVOV/G3YW9KpKcqPiddA9f1MpEhAmhHgGknbJ0j3gKQjzGxHM+um8OYcKelPCfI6O0GaLPDIMc/MPHsG5YMQpcvrGLG63hKOYHfvoZENmpa8PSFlG754AeiGzw38iUbzzoOkI8AVkrzBna8ws9HtPWBmt0u615Hn7hQYJp0Qd3CMJM/x7gsZiZMItwIQPGEvdSZ7zMxGNPisdz+hu/P5KBDmQOMVtqo9jM1AnMQkGQHOl+/bv1BhtdAQZjZF0lxH/ll7Ia0H0EfSa5K+7kw60czezkCkxCRRgP7O54ea2WxnmvmOZ2OfO2wQgg/k/ZKekW/d38LdkUVKjWsdTbhixePGPEeSKzgD4fzfE9sn8yBQBFv+qxXu8km6A/muwkFZqfBupJzvfH6E80RNCkvCzo7nZznzb4hmReyjMN85U+kvwBpsZh7j0VzwKsB3HM+ukJTkHjvvXT7vJCijHicpKFas+cVjkkp5OWXDCgD0kM89eqx3wwPoL+nbjiTvxTQkbUXMwIz/kfR9Myulr4JnEuiNXNXwep4QRPlKSXc4y3ja+XzezJB0mpmV5vSvLZ5PwMmOZ+eqwc4BukgaKek0R/4tPJwgTV5Mk9THzKYWLUh7NDQCADvIt+HxRL0JDyHWzhBJU5Ss8ydL+nuCdHnwvKSeZe98qfERoLd8s+AN7nYBeypsDA2QlPi6OEk3lvC7ulrSLZJ+YWalsfppj0YV4JuOPFdImtD6B8JtHSdJulghZHrac/wXFcKvl4k3JfU3s1eKFsRDox3Ry5HnC2a2pPnM4ERJZyhMIGNFw1om6QcJYudnNVq8pxBA6sFSxvOvQ10FADpL6uHIczvgOQVLodh2e0j6YfN5gZc5kWWZIelGhVu7S7fBEw2gt+OsO2uuSlGPb0Uovwl4GjiTkpijZQ5wZYSGS8saYFDKehjJg1m+DwwFvhSrXTcagDGROjEpC4Aky8RadekOfNhguf8CbqScPgj5AUzOqGMb4W9A18j12RO4n3WjnDcRvJFHAxdSsvuFs6RdZw5CJIulyj+QxFwFo9NRWa31CTF6D1XYDJtsZh4bhE0D4KAc33YIcfd/AySxNq5IQL03O6+hcI7CQdAIM4u9XKtoh3oK8GnG5f9D0h8kPWRmyzMuq6IG9eYAmysYRqTZs2/LFIVTvAfKZiBZUQPgupTf9VXA88BgoFQx8ioagHBTxjhHh68EXgaGAX2BvD2HKxw05NPf/Cn4mYJHUOvPwUqFIX2SpNcVrkl9vcjrWit8uII6NCtCDwXHkJmSpnbog5CKioqKioqKioqKioqKio7F/wDS0UGwYp5sdwAAAABJRU5ErkJggg==",
    "scale": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAALSUlEQVR4nO2dbawdRRnHfw+tIFJaRChSfCkFqbwoBQq+AS1IaDSxolSgBuEYogYTJYio3zQx0Q+GL34RmiIFmiAVBVESW1AoKKa0NOUlLRhariVQS1t5l5b29u+HmVtvbs+eu7M7e/acPfNLNm1P5+XZnWdn5/nv7AwkEolEIpFIJBKJRCKRSCQSieZjdRtQB5ImAxcBp/uf1gL3mNnr9VmV6AqSLpG0Q/uzQ9IldduXqBBJl0ra26bxR9gr6dK67UxUgKTJGXd+u55gct32dosD6jagi3wJODxHusNx44OBYJAc4NSK0vY1g+QAhwSknVSZFT3GIDlAog3JAQac5AADTnKAASc5wICTHCDRbCRNkLQwpwo4wnYvG0+o2/5EQSRNkfQ9SZsCGn4sGyVdM0jScN8j6TRJN0p6o0TDj+V1Sb+SNKvu80u0QdJR/k59PGKjZ7FG0nclTa37vGPQtxNCJL0f+CKwADgP6Pbzehj4K3AX8Acz29rl+qPQNw4gNyA7E5gHfB6YTe9EMXuB1cB9wApgjZkN12tSPnrWASRNwk3Z+hRwNnAOMKVWo/LzKvA3fzwKrDWzt+o1qT21O4CkA4EZwEzgROAUYBbwUbrfrVfFMPAMsA54CtgAPAs8b2bv1GlYYQeQ9F5cYx3UIdnBwLtxDfwB4EBgDzAVOBr4MDCN3unKu81e4EVgM7AFeBmYCOwGXgA2ATuBtzuUsQt42sxeKWJAsANI+iBwA/BlmnOH9jvDwO+B68zshZCMQQ4gaSawEjgqJF+ia2wF5pjZs3kz5HYASQcAj+Oez4neZR1whpntzZM45Nl7If3b+M8BqwLSr8I9f/uRWbi2ykWIA8wJt6VWdgK/wV2MmcATAXmfAD7i897py+oncrfVxIBCDytgSLcZBh4E7gB+Z2avjfyHpKCCfBd6P3C/pMNwiuNC3MXt9cFv7rYKcYChcDu6wts4SfZe4G4z2xa7AjN7FVgMLPbvAC7CydDn48LcXmMob8IQB7gX+Dk9IB4B64EHgOXAQ2b2325VbGYvA4uARZLeg3sPMQ+4ACdk1Y1wbZWL3A5gZhskLQa+UcSqEuzEjWxXAX8HHjGzf3fZhrZ4x7vPH0g6GidZfxr4JG5A1kkoq4LFZrYhb+KQHgDgO8ARuM+sYjOMU8SeBZ72xzpgvZntrqC+6JjZFmCZP5D0LuBk4DT/5ynACcCHqGYccTeujXIT5ABmtkvSxcDFwFW09/CJwKEBxf4R+D4wVLcuHhvvuOv8sQ///uNY4BfAFwKKfAMnpY9mly//ZtzAN2y0GxtJtwROsNghqfJuUtJNATbd1AV7DlLYPEVJuiW2HVFfwsi9wl0QmO1wYH5MO/qE+eT7Wnk0C/w1jkbst3AXU+zDylZkO/qBVoE8k3DXOBqxHaBVMN+FciPogcCfa265dgytiKbEcwBJx1JcLp4IXB7Llj7gcsIjsBHm+GsdhZg9wBWUE4mujGVIH9Aqkddw1zoKURxAkgFfK1nMyZLOjGFPLyPpLOCkksW0/DUvTaweYA5wXIRyBqEXiHGO04FzI5QTzQFiNdzCbmgCdSEnAMVahi7KNS/tAJIOIV5o0nRN4CLgfZHKuiSGJhCjB/gKYdLveDT5MRDz3KLceDEcIHaDzWuiJiDpKIrH/lmUvvalHEDSdCINRkbRVE3gSorH/lnMlTSjTAFle4BWhDKyym0aVTh16fC7cOP5OLSqO/WkJmkCPvb/WEXFl9IEyty95xIW+4euxd8KTN/LhD6rQ67VdEo8hss4QCsw/TW4CQ15uawJmkCB2P8t4NrAagoPBgs5QIHY/z/4qdoBeZqiCYTG/suApcCOgDyFNYGiPUBo7H+Hme0Cbg2spwmaQOg53Oqnxt0ZkKewJlDUAYJPyv+5EtgYkK+vNYECsf8Q8LD/e1dulmAHKBD7rzez1QB+wuLSgLz9rgmExv5LRiZ1mtljuMUk8lJIEyjSA7QC8y1p8+9cX66Oqq9fCXFeAbeP+S3kZimkCQQ5QIHYfw9jTsLMhvh/N5eHvtQECsT+D5nZ2C+Sb2X/aeCdCNYEQnuA0Pf+y/3HEmMZhMFg0XHSPvzScysCyphOoCYQ6gClT8rzW8I0gb6aJ1Aw9s8KkSu9WXI7QMHYv+1Hin7JtCZrAsGxv5m9mfF/91ChJhDSAxSN/bNo8mMgVk9J1ZpAiANEOylPIzWBkrF/FpXdLLkcoEzsn0WDNYHCsX8WVWoCeXuAVkBa2D/275SuaZpA2dg/i0o0gXEbNUbsn0XTNIFIsX8WlWgCee7qWLF/Fk0aDMYeJ+2jKk0gjwNUdlKeRmgCkWP/LKLfLB0dIGbsn0WDNIGYsX8W0TWB8XqA2LF/Fk14DFTdU1aiCYznAJWflKevNYGKYv8sot4smQ5Q4Hv/cWP/LHwcnDccgt7TBK4gcuyfhdcE1gdkmdtpPYFOPUDo9/5LAtK24zZcXJyXXnoMtALSCneuZVgSkLbjegJtHcDHjyEXOHfsn4WZPY97FOSlJ9YTKPC9/0p/rmVYSpgmcEWWJpDVA8zBrWOXlxWBsX8WSwLTtyLUWZbQnmhJ2Qr9tQ7RBGaQoQlkOUDo1KI5cjtrnh6YD3B7AEq6HvhJYNbL5FbjrAVf92WB2X4s6To/cCxS52xJNwJzA7O2fQy07RYkPUfxFT824hZxfhT3AmPT6HjXCybH4LrNs4DP4tbVLbp06hlmtna8RH7xx2/mLHORmX0rR5mn43ZRKcIw8A/gL8BjuJ3EXhy9WqqkQ3E98cdx6w/Pw93NRdhoZseP/TFr5Fpmb4DjgG/7AwBJIztfTQBib8Rc5z4GZeqegNsP8ezRP0oaWQ52ZMe1WLS1NcsBhoi3kgW4E6lqXf1/VVRuXXXHXGxjNEPtfswaA9xTkRGxecrMQgSkqPi6n6yr/kDatmmWA/ySeu+sPAi4vm4jgB8QNqehDjbj2nQ/2jqAmb2O26B5c4VGlWE3cLWZLa/bEG/D1TibepHNwOd8m+5HphJoZutx+wEsq8iwojwDfMbMKl/SPS9mtgg3mPtn3baM4Q7gVN+Wben4MsjMXjGzS3E7hNTdG+wEfgqcVvSdQ5V4jX4W8DPcJg51MgTMN7Ov+g2vyiPpYEk/KrDJQVn2SLpdJRdIVhc3jJA0Q9JSScMVXI9ObJf0Q0nV7WQm6VA5JWtTxSfzhpy6eEIku7u+Y4ikmb7eNyu4PqPZKOlaOeGoO0g6QNIFkn4taVukE9kp6c+Svi4pqmCkGreMkTRZ0lWSVkjaFelabZW0WNL5cvs6F6LwunV+Z80HgAe8AWfgXjh8AiddzgDG0+lfxL3bXo2TjlcWmCbV8/gR+M3AzXJTtObipN3ZOEn8mHGKeAe3l/GTuO3zHgbW5t0guhNRFi70hqz2BwCSJuC2mZ+KkyFH6noVeAV4yczejlF/P+Ed/E/+ANz4CpiGm+84xf+8G7da2FZgq5kNV2FP7JUr9+ENfskfiQ74G2EjYdPiolDFKp+JPiI5wICTHGDASQ4w4CQHGHCSAww4yQEGnOQAA06UzQdj4OXk8/xxZAVVnAOcmDPtBuCRCmzYBjwIPBhDxo1BTziApONxkxdm121Ll1gNLKxzPuMItTuApGnAGqBnvvbtEltw3zTE+KKqML0wBriBwWt8cOd8Q91G1NoD+Hf+2xn/tXFT2Q0caWav1WVA3T3AKQxu44M797wD00qo2wEKLZLQMGq9BukRUC+D/QjwU6VCl0prEnfV2fiQwsA6SWEggJm9hJtMuqZuW7rIauCcuhsfeqAHGKELUvBoDsQtpXqE//d23Nc072Skj0XPScGJRCKRSCQSiUQikUgkEolEYhD4H7hWmX/Y3Ub8AAAAAElFTkSuQmCC",
    "search": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAJ2ElEQVR4nO2da6xdRRXH/1PaWqAUChZBBIvWUj5AVUSo0RbaiA1RAh+8DdI0PmIjWnmYGvBdP5pqY6FRwxe0jY8osU0JiRXoQ4U01igKsaVKa5TovX146TMtvac/P8y5eK2nt3vN2bPnnLPnl5ycL3P2/q+115k9e83sWVImk8lkMplMJpPJZDKZTCaT6X1cagFlAoyX9FZJb5P0BkkTm59zJU2WdETS4eb3oKRDkl6StNM590oKzanp2gAAJkiaJekmSe+SNF3SVEljAw+5R9IOSS9I2iJpk3Nub9tCO5yuCgBgmqQ+SfMkvUfShJinkw+GTZLWSdrinDsZ8XyZVgAXAIuAJ4GTpONlYCUwM7VPagEwA/g+cCzhRT8dzwAfArqqB+0KgGuB1cBQ0ktcjOfxvVPouCMzDPBm4Oek7eZD+TPw/tQ+7EqAccC9wKG017AUHgeuSO3TrgF4H7A98UUrm0PA/cCY1P4tQpJBDH7wdI+k5ZLGpdBQARsl3eWc608tZDQqDwBgiqTVkuZXfe4EDEha6Jx7KrWQ01FpAAA3SFor6dLIpzohaZekfvm073DqdzgtfK6kKZKulDQpspYhSQ86574V+TxBVBYAwHxJj8k7v0yGJP1OPmO3VdJ2Sbudc0MFdV0qaYakmZLmSpot6fySNUrSCklLnXNEOHZnA9wFvFriQOsY8BhwO3BeyVrPAm4AVgD/KlEzwBqgV8c8rQE+Q3nP9juAu4HJFWkfC9yKf7wry4YngNdVoT85wEeARglO+xM+43ZWQluuobwM5dqUtlQCMJ/2u/1+YCEdlHMHZuLnAtplVWpbogHcCBxuwzkNYBVwQWpbWgE44BPA3jaD4CupbSkd4GLgn204ZQD4QGo7itC0dUMbtp4E7khtR2kAY4BftuGQjfjHsq4B3xs8QPjYYBC4MrUdpQB8uY2Lv4ouyZ+3ArgNOBpo+1a6/fEQP7ET+i/4amr9ZdD0wWCgD5an1h8MMB4/Jx7CPan1lwl+Qcu+AD8MAe9MrT8I4AuBF78n/vmnAlxP2PqG39Jtt0HgCsIe+b6bWntMgHmErWf8ZGrtJvDLuKw8SbdFegDAkgDf7AMuTK29EPismDVH3k+XPeq1A/CjgCBYllp3IYCfGQ1rAHNT664S/HsNu4x++jcQe51Ce+DX7Vsneh5OrTsFwBzsPeWDqXWPCvADo0H9dGhuvwqw3woGgHNS624JMBn7CHdhat0pAS4BXukJn+EXZVh4jg6a0k0FsMzotw2pNbcEeNZoSF9qzZ0AcCFw0OC3BvCm1Lr/B2AatgHNX+j11S8GgG8YfAewNKaekGRMn2yriVc45xoB5+lVVkqy+GNBLCFBAE8bovc4cFFqzZ0G8AuDDxtEzAyaegD8atZZhp+sd87tt0mqBWsMbcdImhNLiPUWMEvS2Yb2PzQevy6slX9bqSjRsqfWALAIGZL0tPH4tcA5d1TSrw0/6ZgAuM7Qdptz7pDx+HVik6HtDCJlBa0BcJWhrcXAOrLR0HaMpGkxRBQOAPwmjFMNx95qVlMv/iDpmKH9jBgiLD3ANEmWhM4Oo5Za0cyN/NXwk+kxdFgCwCLghKTdRi115EVDW8vttzCWAHi9oe1LRd/PrzmWALD4vzCWALC8h7/HKqSmWPxU6j4Iw8QKgPz4VwyLn7oqAA5bhdSUg4a2yQPAsrePJc1ZZyx/lIkxBFgC4Lih7XirkJpi2e7e4v/CWAIg+f2qB7H4yXK7KEwOgLQkH1jHCoDaLv82YtntLHkAHDC0jTJx0YNY/BSlqJUlAHYZ2k4CLrGKqSGW9K7F/4WxBIAlbSlFyl33GJb5Fav/C1E4AJp19QYMx+7O3S4qAl9YwrJgNm0ANNlpaHuT8dh1Y56xfUcEwPOGtnPIL4SMxs2GtgPOuX0xRFgDYIuh7fnyFT0zp4B/T9Ky0HNzJCnmANgsX1GzKPmdwNa8V9JlhvaW9YNxwe/cXZQBck29/wN4xOBD8CVzoxDybqBlte/Fkm4JOEfPgi96/WHDT/7hnLOsHTQREgDrjO3vDjhHL3OnbKlyq7/jgt8cebexC8s5Ab1WjuZFo++uj6nJ3AM0ix792PizB6zn6VH6ZMv+7XTObYslJhjgamMUN6h52XV8iVzrfspfTK37tGDfJuYZarxPEPB5o79O0Ml1iPH741v5eGrdKQAux7559KOpdY8KfjD4R6NRe6nRNrHSa35ab/TTEND5s6nAAqNhAJup0RwBcF+Aj6yD7DTgH2tCCkUsS629CvB1A44bfdMArkmtvTDALQEBMATcllp7TIA3An8L8M13Ums3A/w0wNCjwOzU2mMATMKe8AFfK6D7dlUDLsO2A+Ywg8C1qfWXSfPih9ZN/Ghq/cEA9wcavQ+4MbX+MsDXBgi9+L+im/MkhD3uDHMEuDW1De2A3w08pNsHXyBiamob2gaYArwc6IRXgc+mtiEE/Gg/ZMAHft/lD6a2oTSA2fg0ZijrAMsbM0kBFmN/1BvJU6ltKB3gc204BHyNHeuq2UrBp3dDb3kjOQksSW1P6QDfLME5jwOXp7ZlJMBY4F7CnnpOR+8FAX5QuKYE5xwAvkbiWnr4rOedhJfIPRM9GQTjgCdKctBBfLGFSieTgAnAxwgf4VvoySAYD/ykRCc1gN/gB19Rtk1p6r4OWImfwaySk1RYULuSxAO+TOzDkj5d8qGPyO+6vbH5eS60Ogl+8cXcER/Luv2yQdJ9zrmHYp+o0swTvlr4sojnPS7/Dt3O5mdAfiOmQ/Lv108c8blI0lskXS2/Tq/T6vVWEgSVpx7xs4CPqvMcHgPk36O4WWG+jh4ElVfwds6tl/QOSc9Wfe6KOSCpzzk3T9KnZHulbhgn6dtVjgkqA/+EsBy/NqDX2MIpCzrxg1Zr/eBhKh0YVgrwduwrjDuV/fgkUcuelRwErcEnjRYBe0q5DNXTAFYDZ9zRmxwEpwdfWvXr+OnRbqCBz3GY1vCRg2B0gIn4rjR0cUVsGvhlcMFlXMhBcGaAc4CFwAY6Y7C4HfgSJU1QkYOgOPhVtkuBbfh/YFX8HXgIeHcku3IQWAEmA3c0L8wLbTiwFf34+/piIu7McYo9yYKgexchjgA4Wz6de1Xze7r87iTDad/z5PflPTzic0DSoHxxq53yKeQdsXbjOhPAYknfU4dmDDMVkLInyHQIwJI2g2BBahsybdJmEOwHJqW2IdMmtHc7WFT0PJXPBmaK4Zx7ROGziIW348kB0MG0EQQhQZPpVAJuB4VvAZkugeIDwzwI7FUKBEF+DOx1gL7mv7zVP9+8O3tPpILrRrOLv13/Lcvze0nrnHNRiktmMplMJpPJZDKZTCaTyWQymW7nPzV81WLaUqIPAAAAAElFTkSuQmCC",
    "send": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAKv0lEQVR4nO2de6wdVRWHf7sP2vI2yNPyFkJBQASLwRYLlIKYKhZK8QUSY4tGU6MhClGDj+CDxEgUC0UhUuXRAqk2EpNCeEiptIgQAuFZNCIREAqVlvb29n7+sae9zeWcc8+as/fsOffs789zZvb6rVkza/bM3rO2lMlkMplMJpPpOVxqAZlyAOMlfVTSOZJOlLSHpD5JayQtk3Stc+71dAozwQHGAmcBNwJv0pp+4OepNWc6BBgNnAZcB7w2TNAbcX9qHzJGgFHAVOCXwH9KBH0oTTNB7gPUCGCypDmSzpM0MWDT/ZL2cs6tDdhmJgTAscAVwPMBrvRWXNbIfs4ACQAmyV/pcyQdUZHZB51zHx7645iKjPc8wCHyAT9f0jEJJBzU6Md8AkQE2F/SbPnAT04sZ0KjH/MJEBhgbw0G/SRJo9Iq2sbLjX7MJ0AAgD0kzZIP+jRJo5MKaszqRj/mE6AkwG6SPiEf9NMljU2raFiWNvoxPwUYAHaSNFM+6GdKGp9WUdusk7S3c27j0D9yBhgGBgdd5sgHf8e0ikrxh0bBl/IJ0BBgrHxanyPpbEm7plXUMbc2+yPfAgqAUfK99tnyz+p7pVUUjLWS9nHO9TX6s6czwJCgz5a0b0I5GxTn9nJ7s+BLPXoCAEfJB/xzkg5JKOVtSXdLelXS5yPZaJr+pR66BWwX9E9LOiyhlI2S7pK0RNIdhZ5rFCcWr0razznX32yDEZ0BgIPlh1YvlDQpoZQtku6RtEjSUufcOkkC5ipe8CVpcavgj0iAA4D5wAORh1eHo7/QMB/Ys4HOucBAZA1TU8SgcoCJDAY99kFtxRYGg75PC71VBP/f+E5uS7r2FoB///4x+Y7cqUo76PKkpBslLXLOvdRqQ+Kn/a3c4pwbiGyjWoB3ARcAy4C+yFfQcDwBXA6816C/iit/K6mHn8MA7AjMxgd9U0UHrxlbg26exUNnwR8A1hu2XwN07xMeMAGYiZ///lbJgxaKfwBXAVM68KfT4C807nNFu9pq0wcAxkmaIf+sfrakXRLKeVH+GX2JpBXOOco2RGf3fCTNl/Qh434tX/7UBvxHD1OAa4E3Sl4hofhvoWMKgdInnV/5X8XfAv9n2O+pENqjgf/oYQo+rb5c8uCE4nX8bWYmEDQbEiD4RTuzjfteHtKPYADH44P+UsmDEoo3GAx6lFk8BAp+0dZtxv2PjOFTKYCj8L3m2B89DMcG/FPEBUDUCR2EDf4uhfZ2eTSmb23BYNCfKXkQQvE2g0HfuSLfgwW/aO+zxjYurcLPRo4fDHwTeLKk86HoB5bjg17pLB4CB79oc5mxnbZfSoVwuCsGXSo6FjGCvzuw0dDOqiocnQh8HXiopLOh2ALcB3wJSDp1iwjBL9q9yNjWN2I6OQu4vzjwqRgAVgJfA94TzVkDRAp+0fafjW3tH8PB3bHfh0LzCL5/cVBwBzuAuMHfE9hsaO+BGA7uBKwq6WCnPAF8Bzg8uGMBIGLwi/YvNrbZsr2yTl5T0sGyPAv8EDg6uDMBIXLwCxv3GNrsp8Xkk7JOHkk19/t/AlcCJwR1IBJUE/z9sB37u2M4emVJJ9vhJfwr4ZPoojFrKgh+YWe+se25MZxdXdLRZryKv6WcAtTx0+mWUFHwC1sPGtruA94dw+EQo3RrgRuAMwk80lYlVBv8A4227ozldKcnwCbg+8B+UQRWBBUGv7B3idHGhbEcD/WmbzNwBz4L1KVUSltQcfALm38z2NiIL1IRHuAnJR1vxQvAZUDKDzDbgjTBP8xop2HFjyAAR+CfL2PQB9wOzKCGWQGYR8XBL+x+22jr/NC+DxX0q5IHwcLzwKWEfpFREhIFv7D9uMHWeny5mnjgJyOuLHkwrPQBS4DTSZQVSBv89xntVTPrF9gVWFryoJTlOeBb+Hp7VfmZLPiF/R8Ybc4K5Xu7As8GVpQ8QGXZBCwGphPxjSGJg19osEyfexNfuKp68J3DH+Pn0FfJs/gh4aCTP6hH8I832v1tCN87FT0eP199ecmDV5ZgWYEaBL/Q8VOj7bNC2A0GMAmfFcosZ9IJpbMC9Qm+w78faZfXgR1C2A4OXZIVqEnwCy0nGe0vDGU7KqTLCs/gs0LD2cHUKPiFnquMGk4LaT86+GllX6D6WcQbgZuAaRRZgfoFfxS+jEu7vEwXDqVvAzgOWMDw692F5ingZmoU/OJ4TDPquDq0hiQAOwNfJPykkxhECX5xHBYYtZwcQ0dSgA/gZwqtCx25AMQM/hjgFYOWF6nhwFkw8F/BzgUeDh/HUgwAX4no7xlGPT+LpaV24N+MLcRWFSMkUYNf+Hi9UdOJMfXUEnxWmIdtlkynVBH8cfj5ku3S3RW/QgB8EL94csysED34hS8zjbp+FFtT14Afnr4Y+Hs3Br/w4XdGbe+vQlfXAUwGfk2YOoKLK9I8AdsTT70rftUBYDfgy8CjHZwAT1PBfRY416jre7E1jSiAE4HfUO67xlMr0LfEqCnlGgfdCzAV2/f1ALdE1mSt+PVYTD0jHuBkbJlgExFrDAGfMWgBuCyWlp4B30m0cElELX80ajk0lpaeAf/+wMJzROgMYq/49VBoDc0YuQMMkpxzqyU9YtjlUEmnRJByjqRxhu0rq/Y9ok+AguuM24cvuOCXoG0XJN0eQUNvgp9/YHn5somA08+xV/z6Syjb7TDiM4Bz7i1Jlke8HRR2Fc9zZVuYozsWe+gmgBMMVyD4j1aDXBzAvQa7W+jyohq1Bftkk45n4AL7YvvM/q4QvloY8beA7UjRGZwjyTKTN6f/WOA7g5YZyX10WM0E2+f1ffjFMCulZzJAic7gWEkXlLUHHCDJMpVruXPutbL2ytIzJ0DBAuP2cynfGfyUbEvF5fRfBdg7g9NL2nnEYCNexa9h6LUMIEnWDyzNnUH8QM5xhl3udM69abWTKQEVdAaB7xraB7C8Kg5Kz2WAojN4s2GXsZKsFTnPM2y7QdKfjO1nOgE41niFrqHNziBwtLHtqDORhqPnMoAkOecek/SwYZeDJbX7ZtCaznPvPwX4L5Qt3NZmu5aKX+uACbF9zTQAe2dwM8MM1GAfdEpe8asnbwHSts7gTYZdxmj4zmBO/90EATuDjKSKX70E9qXxZjRppysrfvXsLWA7rIGY1+T3nP67EXwls446g9grfr1CTdZR6vkM4JxbL+n3hl3G6J1zBj8iyTKV61bnXL9h+0xMgGMMVy/4BS9Hb7e/dZXVKSn9zTQAezHLM4r9xmBbZe1f1KjiV22E1ICyncHpkizfESx2zg0YbWVig18ix1K8aTN+nd8bDPsATE7t6/bkDFDgnNsg+5vBeZI+btjnBUmrLboyFYK9M7jJuH2u+FV3gL8ag2rh2NT+DSXfAt5JrFe0TxfzEGpFPgHeyc2S1kZoN+nMn2bkE2AIzrm3ZesMtksltQgzAcA+r284apf6t5IzQAOcc49LWhmwydqO/OUToDkhO4O1PQEyTcD+ZrAZq1L70oqcAZpQvBlcFKCpfPV3K/j1DzthADgwtR+ZDqCzVdPvTa0/0yHAJzs4AWam1p8JAPY6vwB3pNadCQR+qRpLvZ8VwC6pdWcCgn8svJrWZd/6gV/QRd/79fZyZCUADpd0kaRpkg6SP4ZrJN0n6Xrn3LPJxGUymUwmk8lkMu3wfwP9SA7pCVGpAAAAAElFTkSuQmCC",
    "server": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAH/0lEQVR4nO3dT6hcRRbH8e95GRf+IYJEgoP6HoIgEZUXIoJrdSL+mQkOJoigLy4UFREXihAYQUENLkQloIYnPOUNihAJBpToVlA0DJOF4kJNoo5mEqMTJCRqfrOou3q5/07f29X90ucDvelbt7qae7q6btWtKgghhBBCCBPHumYgaRq4FbgGuBBYXbxCvwT8B9hXvD4F3jaz/3XJdKAAkLQCuBt4AJjtUoDQyTHgHeAlM/tokAzcASDpZuBZYM0gHxiGQsBrwKNmdthzYusAKH71W4FHfGULGR0CNpnZh21PaBUAklYCbwLrByxYyOc4cJeZvdkmcWMAFL/8XcBfOhYs5HMSmDOzhaaEUy0ye464+MvNFPCKpKubEtbWAJJuAXb2VaqQ3X5gbV3DsLIGKKr+p4dRqpDNxcA/6hLU/QXMAZf3WpwwCvdJurTqYF0APDCEwoT8zgC2VB0sbQNImgG+dn7Qb8A/gd3AD85zQzMDLgHuARobd0scBVab2bFWqSU9JJ99kq5wFioMQNKUpC2STjqv0QbPhyw6Mj4h6cohfudQQtJTzgDYVpZPVRvgz46yLJrZv/1fIXT0BPCFI/1M2Zt9BMBuR9rQEzP7HXjZccp02ZtVAbDKkfFBR9rQr88daUuf0agKAM8wsRxpQ79OOtKWXus2YwHhNBYBMOEiACbcn0bxoUoPmPwNWEsKwn8BO8zsyCjKE5aQdMTRwXCdM+/bJR0uyednSXPD+k6nI0nXO67TT2V5ZK0BJG0kjReU3WWcC8xLWmFm23OWa5JlawMU1f42mm8xn5d0foYiBfI2AjcA57VIdzawcchlCYWcAXDVkNKGDnIGgKfHMHoXM8kZAJ4RwxhdzCRnAOwASm9FlviVNAklZJAtAIpZrPfTXL0/bGb/zVCkQOau4GK60ibKa4JfgM3RB5BX9q5gM3tL0ntEV/BYGMlYQPF3sFC8wgjFaOCEiwCYcBEAE64qADo/axay6PzsZtXF+9GR8WWOtKFfntXYSvtWqgLggCPje5Wmkof8PAt3fF/2ZlUAfOPIeA0Nc9BD/yTN4hs2/86T+W2OR42kNFHxcUnRHshA0qykA85rVDrdv2p6+FmkdsA5zrJ9AswDXxFDusNwAXADqTvd04knYNrMTvlrr2xFSnoduNNbwjCWPjOzdWUH6qrsJ0mLPoTl74WqA5UBYGZfAq8OpTghp73AG1UHm5aJWwXsAS7quVAhn/Vm9n7VwdpWu5kdAv5OWn40LD/P1F18aL9W8GZge9v0YSy8C/zVzGq79Vvdt5vZPOmO4EQPBQvDtwu4o+nig2Mgx8wWgZsA13r0IbutpF/+0TaJXT13ZvYBafBnnujoGTd7gRvN7DEz+6PtSQP/p0u6FniQ9GzfmYPmEzrbA7wILLSp8pfqY9OolaQ7hXWklahmSF2WoX8HSRtHfQt8DOw0s/2jLVIIIYQQQgghhBBCCGH89dEVPA3cClwDXEiareKZsRLaEakbeF/x+hR4u5hqP7CBAqCYCXQ3aWu52S4FCJ0cA94BXjKzjwbJwB0Akm4GniXNCArjQcBrwKN128SWaR0Axa9+K/CIr2who0PAJjP7sO0JbZ8JXElaum39gAUL+RwH7ioW5GrUGADFL38XsYX8cnISmDOzxjWY2jwS9hxx8ZebKeAVSY1bzDZNDLkF2NlXqUJ2+4G1dQ3DyhqgqPqfHkapQjYX07B2Q91fwBxwea/FCaNwn6RLqw7WBUDpggJh2TkD2FJ1sGqBiBnga+cH/UbaD2g38IPz3NDMgEuAe4DGxt0SR4HVZnasVWpJDzmXH9kn6QpnocIAJE1J2lIsy+OxwfMhi46MT0i6cojfOZSQ9JQzALaV5dPH9vGLZhY7fOT3BPCFI/1M2Zt9BMBuR9rQEzP7HXjZccp02ZtVAbDKkfFBR9rQr88daUuf0agKgM5r0IYsOq/pHAs7TrgIgAkXATDhRrJnkNIDJrFp1LiSdMTRwXCdM+/bJR0uyednSXPD+k6nI0nXO65T6aadWWsASRtJ4wVldxnnAvOSVsTegflkawMU1f42mm8xn5d0foYiBfI2AjcA57VIdza+jRBCBzkD4KohpQ0d5AwAT49h9C5mkjMAPCOGMbqYSc4A2EH5ruFL/UqahBIyyBYAxSzW+2mu3h82s9I97kL/snYFF9OVNlFeE/wCbI4+gLyydwWb2VuS3iO6gsfCSMYCir+DheIVRihGAydcBMCEiwCYcFUB0PlZs5BF52c3qy7ej46ML3OkDf3yrMZW2rdSFQCnbDJc416lqeQhP8/CHd+XvVkVAN84Ml5Dwxz00D9Js/iGzb/zZH6b41EjKU1UfFxStAcykDQr6YDzGpVO96+aHn4WqR1wjrNsn5C2lPuKGNIdhguAG0jd6Z5OPAHTZnbKX3tlK1LS66TdQsPy95mZrSs7UFdlP0la9CEsfy9UHagMADP7Enh1KMUJOe0F3qg62LRM3CrSzpQX9VyokM/6ui3ka1vtZnaItCvo8b5LFbJ4pu7iQ/u1gjcD29umD2PhXdIu4rXd+q3u281snnRHcKKHgoXh2wXc0WYz6dYdN2a2CNwEuNajD9ltJf3yj7ZJ7Oq5M7MPSIM/80RHz7jZC9xoZo+Z2R9tTxr4P13StcCDpGf7zhw0n9DZHuBFYKFNlb9UH5tGrSTdKawjrUQ1Q+qyDP07SNo46lvgY2Cnme0fbZFCCCGEEMKy839MQCI0p8+ujQAAAABJRU5ErkJggg==",
    "settings": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAP5klEQVR4nO2de7BfVXXHv4tABCKRFAoJj4QQgrSahqeEJmDpIFJGbEcChoK2OFOmFdoy0I7FV60jNW3VAUVa7VSdZhgCGIs8NGBaRiABUmmAkEjIgwCBhIAG8oI8uJ/+sX6/ernNvfesffb5nd9vOJ+Z/JG5++y19tnnd87ea6+H1NDQ0NDQ0NDQ8LbD6lagUwD7SXq/pGmS3ivpcEnvkt+DVyWtk7RM0iJJD5jZ6zWp2pATYBowB9hCcbYA/w68r279GxIB3gvcG5j0wfgx8Jt1j6ehIMBewGeBnRkmv80O4Bpgr7rH1zAEwCjgrowTP5A7gFF1j7NhDwCjgYcqnPw2C4ED6h5vQz+AfYD/7MDkt7kX2LvucTe0AL7awclv8w91j7tBEvB+oK+GB+BNYEbd4y9LTxuC8FX5Y5Km1KTC45JONLO+muSXpte3NReqvsmXpKmSZtYovzS9/gZ4WNKpCZculXSnpBWt/79b0oflJuIoi8xsesJ1DWUAjk34bm8Afn+IPj8CvJTQ7zGdHHuDJOCq4CQ9CxxVoN+jgXXBvq/swJAroZfXAGcE2vZJmmVma4draGZrJF0kiUD/pwfaNuQAWBP4hf5HQv93BPpfVcUYO0FPvgGAfSSND1xyS4KYyDUTgBEJMmqnJx8AScdJitzwxxNkPBZou7eknlwI9uoDMOhKfg/0SVqbIOM5xdYBf5AgoyEKMAbYGPg+rxi+10FlRdYZ64EDc461E1R+ooWbayfLjS1jJR0oabvcD2+1pCfNbEvBvvaRNEfSrwdUWBhS+P9fO7Fg27GSvgfMNLPdRS4ARsutiRMljZG0r/y+bJAbqVaa2ZthresGeAfwUeA2YNMwv5zdwH8Dn2OIfTowDvhJ4BfZ5vwS47goQd58YOwQfU4EPg88ih8oDcVrwLyWHvumjmMospqCcW+Zv5B0paRDErrok3S3pO/LfwG75Kv9cyRdIinqjfOapHGpHr7AOyWtl/TO4KVb5W+q+XJv433kC9eZks5V2trrZUnXSbrezLYlXF8twO8BzyX8YqrkqxnG9fW6BzGA54Fzc8xZFoARwD9Sz5n8UGwFxmUY3xHA9prHMpA+4CtksD2U2gYC75B0m6S/VvedLH7ZzNaX7cTM1kmanUGfnJikqyXdBows21ES+Or+VknJi6wKWSxphpntytFZ6yYvknRSjv4yM0/ShalOKWXeAF9Sd07+S/IbkmXyJcnMdkq6QNIrufrMyPnyuUgi6Q0AnClpgbrPkviqpLPM7NEqOsdDxX4iaXQV/ZegT9IHzOy/oheGH4DW63CppGOj11bMi5I+ZGZLqhQCnCTpLrnhp5tYKWmKme2IXJTyC75M3Tf590k6perJl6TW2+VkSfdXLSvIZEl/Er0o9AbAgyFWSZoQFdRitaTn5a/Q4yTtn9hPm1ckfV7StzrtmdtaBP+ZpC9K+rWS3W2Xh6Zvk3SkpEmJ/TwraVJl5mPg3IQ96048cOOoAX3ti5s4H0zocxXwV3RBiBYekvYpYgdH4Hv5+4FZ+Ha6f58T8H1+SoDrOVUO9jtBZV4BTivQ70Tgz4G5wJPAtn597MBv7t24Df19QLfZHAQYnovgb/GQ8mdaurfZCiwFbgYup5h/4rTWPYzwr1UOMmLq3Q0k+8rhYd7dtsvoOMCM1r0sytqqFDk0oATAtytR5G0I8K3gvS98XB75hUVX/t8Itm8YnBuD7QvPVcQh5NBA201mtjTQviMAE+TRP8fKx9M+5t0qtyA+LWmpmT1Xj4aD8oTcyFXU46jwUXzkAdgv0HZDoG1l4IvFGZIulnS2Cnr3AGsk3SvpJkkLzSziG5gdMwPYqOIPQGSuioF7+BTlhewKxHQdCXwCWBH8du6Jp4BLcXe0Ose0PqBz/oBV4KzgjTssuxLF9PwdYHlQ1yKsBD5Y05iOCOp6ZtG+I4vANUG9Lwm2LwWwH77zuE/Sb1Qg4hhJ84Hr6Hx6mI8H20fnanhwQ8cvA0/hL+jQWwA4HHg8+Cspwz3AmA6N7QiGd6ztzy+pylAG3B68UYuBsnby4XSaDKwN6pWDFUCqzb7o2A7CvYcj/KBKhT6RcKNWAydXpM/huMm1LlZQ0ZsAOIX4+QLAH1WhT1up0bhNO8rrZF6Z4t/8Tr72B2M+mQNDgQvxexZlM+7KXh34IiiF3cBFGfX4dqIeVfC1jOO6mJjtvz//lEuPoRQciz9pKewAIokdBtPhA4nyq6R0niDgTNJzHL8GhINxUn0C/1IepZLCBkm/ZWYvJ8rO4ZK2QdIS/cpiOU7S8Srn5vWAmSU/3MChcpNvSkSVJF1hZt9MlR8CP6q9p8Sv5eYSslMWogC7gO8xRA0A4NRWm9RX8IdLjOu2RJngvhKd9ZHAtyhlTK2/nSDTcNNslGXA8QE5J+JeR1GeJGEigNMTZLVZTodsEntSfHzijQKYnyAv5Ub9FA/Djso6CM8MHiXlwV6QIAfgaeCIqLys4IvChxOU7wMmB2X9S1DGMhImv5+8g3E7RoTQdxg4Lth/m0X4uqF+8HwA30gYxOeCciKGkV3ACRnGdiqxwNeVwf6/GOi7zfWUjAmsBODLwYE8FOh7QrDv72Yc161B2UcG+l4c7Ds5DKxy8FDxnwUGs4MBLtFD9P2h4I06JeO4pgVlF4rfx62ZuwL9PkJmR9msnbUCEv4+cMlIFd/PR/b9GyT9LNB+OB6RtDHQvqiu71bMK+va3AEwVbhd/1ie2qUoRV+XBwf6XJLTjavVVyTsrKhXbuFPhaQdcje1rGR/AFr5eJ4PXFJ0lR6JAnop0LYoET/HorpGxrTOzN4ItC9EVYEXXRe507Bnsj8AwP6SIsaJzQXbbQ30WcX+OHJOUCjvoYqPXZKOxOsfZ6WKN8C58rRoRSn6uYgswk4go2281deJgUuK6hr5VI6Up8vLStYHAHeW/HTgkp3yYIwiFG0n+a812zZQXnE8kp00MqbIgvnTZHY+yf0GmC0pYn17NJDRIhpp9Mlg+6G4PNi+kK6tBXNkd3GypGuDulQPHut/Y9BYAvDZoJyIXX43EHltDybzFIZP6dqfqCn4C4G+29xAQQNa5QCHEbP+tekjWGwJ+OegjNXAQSXGdjBxx8wbgjKOJS3J5mIyJMIsBZ7Y4ZkE5QF+lCBvRoKchSQ8BPjkpxwHD5sQYw+yUp1r1uABr50HOIR0X4A+IFzvD3cI+XmCvNURefhrP3oMDO6ckeIQMj1BVpsVlHjLJYG7hJWp1j2nhOxLE2X24ad609jDJOEP12nAHGLf/P58rMS4bk6UCV7NPGlBn+oUepWk1EzcL0g63sySsm7iUbpL5QcpqWyUr77Xy+9B2yk01SFTkpZLmlq0WMRA8KweTyjdMfVKM7s+8dri4Iu+LYlP6g4yVNymvorhQ/G7GcZ1Bm9NLBVhM0MUqsgGvgVJYRdwQUY9UgNUqqB0XYJ+45pFuldytgCVwZQ7kLemcCvKNiBrVS3c+eTuxBuVk3vIHxp2Pmk1CrZQwgeyiGKXJSi1kgy+eYPoMxp3/KyLp6ioUhhwEmm7rEur0Ket1F1BZR6q6gb102kSeVLBRHkKOLrisY3B3cAihMvkFlVmL+DVgCKv0CG3ZeAAYrV+yzKfDtUIxKul/SKg2yaqiBLCEzFEuCq7EkPrNwL4GtXvDrLU6gmO7W+COh5VhRIfDCjQR01BC7gx54HgDSvCMgLJlzKP6XBiD3bhLWnEehQxN64zsyr88obFzB4ys9Pl9YWXZehyuTxJ01Qzuy9Df2HM7AXFfBILz1XkAYh4+RR1iaoMM7tD0hRJ0+WpVlcHLl8p6ZuSppvZe8xsTqqFLyORe1p4riI+6dsDbbsiZq3lzr2o9U94xE47VexY/SpV7Ba5J/EKeS3jdZ3XdnDwRV3kniZVSh1OiTOC38wp2ZV4m4KHq0coHKEc+QREfPIk6Ypg+4bBibq3ReeqGHjd2qKULRgxGnhXTv3roDWO5MxdeOrbyNlAKEtoNOXpAkl/XLDtCEm3A+eZ2aKhGuKuYefJM3tPlYdMjWz9Dfmx7XJJD8vDoxZ2ukjUcODn8dPlWcmnydPVjlPrLQvskLuBPybpQUl3m9mqYfo8XdIP5PeyKAvCyheFtKJRu/B49kkD+tof+BhuLo7yPPAZ6kqN8tZxjGnpEnk7tlmEp4Xbb0Cfx+CnrpHI4TZnR/RPKRu3WtL4yHX9eE6+n91X7tBR1rN1k6S/k3RDZaXSBgG3Bl4u6QuSyj6Ib8i/25vlJfkiQaP9WSNpcqVvR7y6V7exEEh9KFPuwXjSHEarJhwLkVo69kl5pcpuYqOk88xscZVC8DRzdykWKdQJnpaXjt0ZuSjsSNgS8KfygsXdxCGSFpDgll0UfH+9QN03+X2SLotOvpQYGtaqUj075dqKOUDSnVSQxh3fqdypWEx/p/iSmf005cLkc2N82zNXUjY/v4wskXRatJL2YLQ+e4/IPYe7jbmSLk5d+CUHh7YEXiJpXmofFXKCpGsy9vcZdefk3yLp47XaRHBHjNl0n5v2djJk0QSOJC13f5X0AdfSTaV1gbOpt3rHnvh6hnGlusFXxWrgrBxzJmXO5YNbtC6XdLXSIlzelHSHpO/Lo392yYs9niM3QUddnrdIOszMIull/g+8PP16SaOCl26W9F1J90haKz+fnyJpptzkneJS9qI8GuvGKpJFZQUv3DgTuIXhy5/vxnMNX8MQGTbxYNSoVzLArBLjuDhB3p0M4Q6HG5Guwb19hzvkeRmYi8cJVJIatvJsXrgzw9GSjpO/FUbJTZ+b5abLZWa2rWBfIyTdKukjARXmmFm07l5b3k2S/jBwyTxJHy1qlgZGyR1UjpZvL0fKnTk2SPq5pGfqLlvbdeDHq5EyqhFXsIGyngnIeZEqo3IqontWkQUxs82SIlGwE1u/tBD49z+SfOG6lm49Rc89AC1+GGhrik1km/GKfSJvT5BRO736ADwtKeKlOzVBRiSecZekIZ07upWefABai6xIksWUnUDkmme7zUOpKD35ALR4ItD2PAL+iXgEUKGc/y0eD7TtKnr5AYicfpmkuRQ4JcRP/W5S7Pv/QKBtQw5wv7koG4ELGDxJ1Czc+BKhj4rDxKukp9O6Awslhcu0SXpK7tWzSn4PJslNtCmJp0pVDG0oAV5lu27Or/s+lKHX3wB7SfofpW3zcrBE0sm9ugOQensR2HZKuUL1+Ce+KemTvTz5Uo8/AJJkZg9K+koNomeb2cM1yG0YCLA35aqZR/kRHU4T0zAMeKKoTgRr3E/C4VJDBwBGAbdXOPnz8KJYDd0KbtD5FPBGxol/HbiaKtKvNVQDXpY9xYVsID8EImVrG7oJPOXqd4gludwE/BsVpbjtJt42rzTcqXK63HT8HnkIdnsxt03SOnnQ60JJi1Li7BoaGhoaGhoaGhp6gv8F1jniA8T1LWIAAAAASUVORK5CYII=",
    "share-2": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAL0UlEQVR4nO2da4xeRRnH/0MLCIVS5VbuaMUKorTVcgvITcWIgB9ARIsoCUQRQS5SNEYl5aIICIJQBaKCxoQa5aIg8kFQEaHaKm2hF6pYKN2WW8u20na7+/PDnIW17u47z5w579ntmV/SNNnMzDNzZt5zZp6Z5z9SJpPJZDKZTKZxuLor0ASAEZImSjpS0rskjZe0s6RtJCFpjaQOSQslzZX0kKTZzrmeGqqbSQUwAbgeWIGd5UXeA+puR8YIcCBwf0Sn90cPcB8wue52ZVoAjAF+AHQn6vy+dAPTgdF1tzPTD8D7gKcr6PiNWQi8s+72ZvoAnASsbUPn9/IScEzd7c5IAk4HNrSx83vpAk6tu/2NBjih6Ii6eA04uEwbsh8gEmAfSX+VVPekrEPSZOfcczGZN0tcmUYAbCFphurvfEkaK+mXwMiYzHkAxHGRpKHkoJks6YyYjPkTYATYRdLTkraOLOIZSfdK+rv86xtJu8i7io+XtFdkucsk7eOcWxOZPxMCcFXkhO0J4DhgwB8d4ICPAnMibXytnc+icQCjgFURHXMtsLnBzhbALRF2VgKxb6ZMK4ApEZ3y1RL2pkXYy76BqgDuNnbGT0rac8AMo81fpWpvpg/ASPwrNpTlwJgEdnfE9tl5BX/+IIi8DAxnf0nbGdJf55xbWdaoc+4FSdcasoyRP3QSRB4A4Vh24JBU6vW/EbcVZYYSXNc8AMIZZ0j7pHPu+VSGCzfvfEOW4LrmARCO5Xv+VAX2LWW+OTRhHgDhjDKkfbUC+5b5xDahCfMACGedIe2WFdh/kyHt2tCEeQCE02lIu3cF9t9qSBv8BsoDIJwlhrSTAMsvdlDw7t1JhizPhibMAyCchYa0W0lK6ZI9UbbPyoKEtjOSBIzGdvxrFVD6BwZsBswy2F0PbBtafn4DhHOIpC5D+tGSbkxg92z5swKhzHTOWeYrmcEA3o59E6iXbuDwErYPB9YZbU5N2f7GAmwDXEn58/6vAidE2D8Wv7FjoRvYo4rn0RjwW7BTgKUlO74vG4DvAC09dPiB9y3iYg3uascz2mTBh3c9kq7f/4+XgGuAI+kT4wdsV/ztGuDFEuXnANIYgJ2AW6kmqHMwVhf/UvDzmLY3+lQw/pzeOZK+Idte/1BjlaT9YnYgo4IJNgWAD0m6XrZ9/qFIj6QpsdvPjfMDAOOAOyU9oOHf+ZI01Tn369jMjfkE4P3pF0uaKtvOWsuiVd9zvN05d3qZAtr6CQD2kj9b9w69IZIkSaslLZf3t89xzlk2XlrZdJJOknS1pD1TlVswU9L58n7/LyQuuxU3FraHLvg19eF4aZN/Gma0i4GbgcMYJJImwP4k4I+JZtl9WQp8um/dgDOxe+1i6ALOSdNDFYGPbDkDWJCgwfOBz2KLrNker7CVWrhhXVFuv5stwBH4wVsVz1DCtdwW8A6NJyto/CLg2Ba2RwLnYTu/H8q9wNsC2r8FcC7wQkLbncBlDOWwL2Ar4IcJGz0Q19FPLDxwNPFBlYOxAPhIxPMYDXwTWFLC9hLgUmCHNL1UEcBuwD9KP+pwHqDwqwN7ALdXYOMVYCpeDKLMs9kM/1a8EngU/2seiM4izRX4z0nly/TSyxe8VMqDio9rj2WRpPslnaW0y7oeST+T9GXn3PKE5b4OsKu8skfvXKJTUkfKWIJQSg0AYDdJf1I1hyDrYKakc51zf6m7Iu0i+hUDbCXpPm0anb9U0hRJBzWp86VyjqDrJb0nVUVqYp184OUVzrnVdVemDqI+AcAHJf0ucV3azd2SLnTOLa67InViHgDFrHiOvDs3lg5Js4v/JS+SNEF+YlQ1T0n6knNuuA/gesB7+GLoAn4MHDhI2QcVaaqQXl0JnI/Bo5jZCLxvf37Ew58HTDDYmUQ65e1uvODSTlU+m0aA39ix8jARGvd4f37Z83mPU1JLN9MH/K6ehXkxnd/H3g7E+fWXAmfRBk9ao8C2pdsFWCJaBrL5cWPnvwwEx8dnAgH2MnbEjxLatm7y5OCIQCyvyHcby77JmH4wzjSmt9a1sVgGgGXd3yGvpZ+KxyStMKQv46NoFJYBYNmTnu2cs8iaDUpR1mxDlh1T2d7UsQyA4Jhz+QOeqeloneR1LHVtNHmZ1HAsA8CyW7aztSIBWPYJskBCIJYBYJmETaTEce6NKcqyiCRZ6tpoLAPAIpI0Vv4em1QcLNvEzlLXRmMZAHOMZZ9tTD8YNxvTW+uaCQFb0MMGwPLaHsjmyUYvYHYFVwU+XMvCYmD7EvZiN4OeJ28GpQcfq2flkZhBUHR+2e3gmcAhVTyLRoI/EPJUREcsBg4y2JlMuhi7buA2oIqlafPAB2rG0APcCRxMP0tE/OA6BLiDarR6VgEXUTLSZ1Mj5lDo5vKz7PEl7K6Q9+0vK+rQeyi0Hce2Fkg63zl3fxtsDXlij4UfIen3sfmHCL+RHwiL6q5InUTNkp1zD0v6XuK6tJvjJM0Fvo1BXHlTo4z6xghJ90gyh04PQZZJ+oq85k6ybezhQNng0NGSHpW0X5rq1M7f5IND/1x3RdpFivDwcfJBou0+hVNVeDiSfiEfNhZ880Zw4X4FtLv8hHfbwl6vSNbSYfkGArYF7qlg6TYQv6W4lpXqBCJW4xU+Sl0ABYwAPoCfazwGrGlh8zG8WPTRDCdPZtHQa/Hr/Sq5mn7uxgWOAp6owN4i4OSI57EzcDnwXAnbz+G1gYZPVBPemVOFNNs84KgWtkfi9wBSCjT18iDQ8k5evNz7VcB/Etpeg3+DDJ9NLuAEYG6Cxs8DTqMfYahBbL+FamTi1hfl9issjb/k4dnENvuyBDgyWSdVDd69eyjwfWzBnguBG4FDS9qfCPwhdS/gNf3Po/gU4YWgLqc9cvPrgc+l6aE2e/LwETu9UrFj9YZUbKf8LHiBpLnFZckp7R4v6QalF7KaJekC+ZXIJxOX3YqbJJ3nnNtQppDh7Mo1QXVi0XXyU+fcaXVXYlgB7E41y8a6uLjM82jMG2BjgGPkha5azuyHOD2SPuacuzcmc2MHgOSXjZI+L+lSSS1v9BrCrJI0PkbYcvh4mirAObfBOXeD/KR0uqTuNldhtdIEsWwn6esJymk2wAS8pE1VrMDfH3g4fZw6eFf6+/FezhWRZa/Hy/ZmygKcQjmV743pwqt+t5R7B0YB04hzYs1ox/NpBMDWRaeVdem+TMQlD8CHsR+Jz1fHpgbYG5gR2fldlLjNk7jLoy9J2f5MAX63cbCt3P64JYHdLxptNuYwS1vB3wDSZeiITgwbWIPYHYHfEAulC8MZx0YvA41MlE1d/YKyfnpJcs51S7rMkGWk/BH7IPIACMdy5O01SXcktH2XpLWG9MF1zQMgHMvsepZzztJhg+Kce002kazguuYBEI4lduCZCuz/y5A2WJ43D4BwLIdD11Vg3/JGCd7uzgMgnDWGtP0eGSvJGEPaYEGvPADCWWlIu28F9i3BN6+EJswDIBzL3UL74q/USwLevWuJxg6uax4A4cw3pHWSTk9o+zOynd2w1DUTQuGRs2zOrKC43rak3e3xp5BDeZl+AmcGIr8BAik8cg8ZsuwoaTolBDOLvNMlWTSWHirqmkkN8CnDL7GXaSXsTYuw94mUbc70AX9GYFVEp9yKQZsI2BIvbGVlJf5K30xV4GP0YliIF70c8LOLjzA6kTglNoArre1p9KngGICxkp6WNCqyiH/L6xPN0v/enDpJXrZmz8hyOyWNc869EJk/EwpwSeQvtEoujGlLfgNEgJfKe1TSe+uuS8Hjkg5zznVZM+YBEAmwq/yDT+bxi6RD0uTYgNrsB4jEOfe8pBPlD3/UxVr5sLCk0dQZA8Cp2M4KpqILOKXu9mfkA02Bl9rY+S8CR9fd7kwfgHGkkcRpxQKgjE5zpirwR8evAdZW0PEb8FI7jZW1HTbgL9pOJXvfA9xNgpvYM20G2B/4LtAR0fHL8LqL+1dZx+wHaAN4//8Bko6SVyQZL3+5Zu/rvFN+Pb9Q0lx5Kf4nnHM97a9tJpPJZDKZTGaT579i3lANQgSqOQAAAABJRU5ErkJggg==",
    "shield": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAI7UlEQVR4nO2de4xdRR3Hv9PX0qKhFCq0vKNg06KkUh+NMS62MVaqYvBRI4mlCha0ikkTiVH/0ZhgogmKxRjQGqUUtOKjQZRHJYhUmpIqrS2S1kip7WIfC33srtvuxz/mLmlvt7t3zsw5c/fMfJJN+ph7f6/vnjPnN+eckTKZTCaTyWQyyWFiO1AlwNmSPiLpPZIul3RO47+6JD0r6TFJvzTG7IvjYaYUgLOBO4AeRqYH+B5wVmy/MwEA5gNdLRS+md3AvNj+ZzwAbgb6CxR/kH7gs7HjyDgCTAB+6FH4ZlYAE2LHlWkB4GJgfcDiD/IUcFHs+DLDACwCukso/iD7gY/FjjPTBDAVWF1i4ZtZRb5KiA9ggE8BL1VY/EG6gOuApHopbQMwG3giQuGbWQdcETsfyYCd5K0EjkUt+4kcBe4GLoydn9qCLfydQF+gog0A32z8DAT6zl5sxzELIRTAO4B78GvoNHMY+OhxNj7e+LdQ/A/4GfDWmLkbtQBTgM8BGwMWZZDtDHHOxs4p/lWCvaeBpcDkGLkcNQAXADcCDxLuMN/MA8CZw/gwBfhtSbZ7gd8BnwbOqzK3w1H5JQwwXnYZ9g2SZkmaI2mupDeWaLZH0nJJdxpjGME/I2mZpNsknVaiT9skPSVpg6R/SNouaY8x5miJNk+iNAFg26aLJHXKFneypPGSXlOWzVPwtKTFxpitLh8CZkr6qaxAq+SgpKOS9kt6TtI6SfcZY3ZW7EcxgOnAz7GXRjE5DCwHxnrEMg64FTgSNRI7AV4JnBuyVsEB3gvsi5srwJ7HLwkY1+uBtZFjAvgvMD9UXEEBPoi9/InJFuB9Jca4ENgWOcY+YGFZMRYCmEXY62hXdgI3AOMqiHU8cBOwK2K8B4EZZcfaEtgFmTLW31thJ7AM6IgQ90TgFuIJ4UnaYSEKe1ismr8D19MGd+kAHcBnsKefqnl/7PgF/KaiYHuxa/7vjh3zUGCPhPOAX1BeI6uZX/n67XUIwTZ1uiVN8nXkFAxIelLSakmrjTH7S7ITFOzzB5+Q7YPMVXn9lkOSzvRpHvkKYKakLT7fMQTdsg9o/F7SWmPMnsDfXynAdEkLJS2QbYqFXhOYYYx5ruiHfQVwlWyxirJPttu1WdIzsq3RzcaYAR+/2hVsU+pNskeF2Y0/XyrJ5/ayTmPM40U/7HvZ5PL5Q5I+LOmApFck7TLGHPG0P6owxhyTtKnx8yrAJEnnSTpD9gjxgFpvmXtNhEu/bj6OfmPMIxXaGzU0fhGeH/w70F+V7TFVGcq0J1kAiZMFkDhZAImTBZA4WQCJkwWQOFkAiZMFkDhZAImTBZA4WQCJkwWQOFkAiZMFkDhZAImTBZA4WQCJkwWQOFkAiZMFkDhZAImTBZA4WQCJkwWQOFkAiZMFkDhZAImTBZA4WQCJkwWQOFkAiZMFkDhZAImTBZA4WQDticvb24bdAGMkfAVwrEJbKeHy8i6vF0r5FqXPYWyZ26/UDZdcudTgJHwFcNhhbAeQRTACwOlyOwK41OAkfAXQ7Tg+b7g8MlMcx7vW4AR8BbDXcfw0T3sp4Lql3D4fY14CMMYckn0FbKtc5GMvEVxy1O37ut0QM/MXHcZeGsBe3bnMYaz3VnIhBLDdYezlAezVnZkOY3f4GgshAJd31c8OYK/uXOkwdpuvsRACeNZh7AyG2bs3dYDXyW6p2yrem3WEEMCmkYecYK8t9/xpEzrl1gZ+xtdgCAFsllszYkEAm3XFZdPLl9UOp4DGhkXrHT7yASCvCzSB3fTSZUfQvzR2IPEiVCEedRg7TfZQlzmReZKmOoz32avpVUIJ4CHH8UsC2a0Tix3Hu+a8PBqbJr7gsOFhb2PGm5EETMNts0mX3suwBDkCGGOQtMbhIx2SloWwXRO+ILfdv1xyXQ3AHAcFA3QDritftQM4C3jZMXdXxPZ7SIC/OQZyW2yfYwN81zFnG2P7fEqApY7B9AIuna9aAczAfaPpG2L7fUqA04F9jgE9DJS1uXLbgp04r3PM1UvAxNi+DwvwDcegAJbG9rtqgM8XyNPXY/s9IthJTbdjYIeBWbF9rwrgzcARxxwdAELvPF4OwNccgwPYCpwR2/eyASYD/yyQn6/E9r1lsHOBFwsE+SC2J15LgHHAHwrk5d+0+7m/GWBRgUAB7qKGk0LspO/HBXNybWz/CwE8VDDg2+skgkbxby+Yi7Wx/S8McCHuE8JBfkANlo2BMcCKgjnYD5wfOwYvgOsKBg9wL9ARO4aiAB3AfR7xL4odQxCAlR5JeIJRuHIInAP82SPuH8WOIRjARGCjRzJeAN4WO45WAd4O7PSIdwN1e44SuADY7ZGUPmA5bTwvwE72voR7f/94djHaz/unArgSOOiRHIDHgItjx9IMVuAPe8b2CvCW2LGUCjAf6PFM1EFgGTC2DeIZA9yM+5p+Mz3AvNjxVAJwNX6HyUE2AHMjxjEHWB8gjl7A5Y7g0Q+wAPcFkaEYAFZR4WkBOB/4CXAsgP9HAJfnAeoD8C5ssyMEfcD3gekl+jsV+A5hhAv23ol3luXvqACYCewIlFCw59IVwCUBfZyOLfyhgH7uAFyeBK4v2N+sxwMmF6Af20ks3D8AZgF3Y8/RIfkT4PIQSP0BxmNvkBwInGyAvwKLaWFJteHHtcCjJfgygI2xtkve3gAfAvYGTvwgB4A7gDlD2L0M+Bbwn5Js7wWuiZHT4WjLJVfsRO4ulfsk8VZJqyTtl/RJSXNVXj7+KGmJMWZXSd9fP7Bt1SWEu0qIQTdwIzW6t6FygHOBe+LWsRD3U+IlaXIAncCmyEVthS3A/Nj5qiXYvvv1+C21lsVu7JNReYZfNsBpwC3Anqglt+wFvox9x2+mSoBJwBeJc0ToAm4FXhs7D8kDTMA2elyfTi7C88BNwKTYcWeGALgKWINtBYdiAPtAR37B1WgBu4DzVfwWmrqAbwP5HcejFWxDqRO7ft/KXTt9wK+Ba4Dxsf0vm6Q6VNjFoKsbP7Nl380/VlKX7CtvH5G0xhjj9Q7+TCaTyWQymUymrfk/YcjcdKSG57QAAAAASUVORK5CYII=",
    "shield-check": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAKwklEQVR4nO2df7CVRRnHv6sCgpWIogL+nNQYoAzFimmaULGRtKCxlMompDREIWuYsWmqf2qabKYmzTAdLZqSoDItyVREYtQkGR0TCdTRSiRAAW/8vATcT3/sOc7lcO65Z9/d993znrOfmTsD9+579nn2+Z7dffenlEgkEolEIpHoOExsA4oEOE7SJyVdIGmcpBMqf9osabWkRyX9zhizNY6FiVwAjgNuBfbQP3uAW4BjY9udCAAwGdjcROBr2QhcGNv+hAfAbGBfhuBX2Qd8KbYfCUeAgcBPPQJfy3xgYGy/Ek0AnAasDBj8Kk8Cp8b2L9EAYDrQlUPwq2wDLo/tZ6IGYDiwKMfA17KQ9JYQH8AAnwdeLzD4VTYDVwIdNZbSMgDjgcciBL6W5cDZscujY8B28hYAB6KG/WD2A3cBp8Qun7YFG/jbgL2BgtYDfKfy0xPoM7uxI45JCKEAPgDcjd+ATi27gE/1yuOKyu9C8T/gl8B5McuutADDgOuApwMGpcrL1GmzsX2Kf+aQ31PALGBojLIsDcDJwDXAA4Sr5mu5FzimgQ3DgD/mlHc3cD/wBWBUkWXbiMJfYYABstOwZ0gaK2mCpImS3pVjtnskzZN0mzGGfuwzkuZIuknSkTnatE7Sk5JWSfqHpJclbTLG7M8xz0PITQDYYdPpkibJBneopAGS3pZXnn3wlKQZxpi1Lg8BYyT9QlagRbJD0n5J2yS9IGm5pMXGmPUF25ENYCTwK+yrUUx2AfOAwz18OQL4GrA7qie2A7wAODFkrIIDfATYGresANuOnx7Qr3cCSyL7BPAGMDmUX0EBPo59/YnJGuDiHH28FFgX2ce9wKV5+ZgJYCxh36NdWQ9cDRxRgK8DgGuBDRH93QGMztvXpsBOyOQx/94M64E5wKAIfg8GbiCeEJ6gFSaisNVi0TwHXEULrNIBBgFfxDY/RfPR2P4L+ENBznZj5/w/HNvnemBrwguB35LfQFYtv/e126sKwQ7qdEka4mtIH/RIekLSIkmLjDHbcsonKNj9B5+WHQeZqPzGW3ZKOsZn8MhXAGMkrfH5jDp0yW7Q+LOkJcaYTYE/v1CAkZIulTRFdlAs9JzAaGPMC1kf9hXA+bLByspW2dGu5yU9Izs0+rwxpsfHrlYFOyj1btlaYXzl32dK8lleNskYsyLrw76vTS7P75T0CUlvStouaYMxZrdn/qXCGHNA0rOVn7cAhkgaJelo2RriXjU/ZO7VEc79vbkX+4wxjxSYX2mofBFeqv4f2FdU3ocVlVGiNUkC6HCSADqcJIAOJwmgw0kC6HCSADqcJIAOJwmgw0kC6HCSADqcJIAOJwmgwylyNrCtAN4hO58vSauNMdtj2pOVVAM4gt35dLekNyQ9Xvl5A7uFfWRc69xJAnAAGCu71/AzOnghxsDK754CzohhW1aSAJoEGCe7/K3R1u5RkhbSCuv1myQJoAkqwV8m6fgmkp8nqTRnDCcB9INj8Ku05N6FeiQBNCBj8KXwS79zIwmgDzyCL0n/CmtNfiQB1MEz+D2S7g9rUX4kAdTgGXxJut0Y82JAk3IlCaAXAYL/sKSvhrMof0o5FFwZhp0m6ZzKr56RdJ/PcGyg4E81xnRntaF0ABc5bGUOsrMXuJz65xBtJeNZ/sA4st0rVOUhINiRcth7CZrlIp+8StUEAFfIbhUfVufPwyQtAq53/MyO/uaXRgDYan++Gu9oNpJuaVYEnR58qUQCkN1ZXO+bX0tVBHMbJcIesrRUfsGfVubgS+USgMulDEbSj/oSQSX4yyVlPXixGvw9GZ9vGcokgIZn/NahrghS8A+mTAJ4LsMzB4kgBf9QyiSAe2UPUHalKoLvSVohv+BPbafgSyUSQGWQZ7bcmwLJiuBGdXBvvy9KIwBJMsYsljRL2USQlbbo7fdFqQQgScaYO1ScCGK1+S5LyrzKwVcABwrM6y0qIpirfEUQs813maPxOlDKNyh7HdIGvX7FGHOr8qsJYlf7LmXlEoND8BXALoe0g0JOmEi5NQdRX/WAo+RWA7jE4BB8BdDlmD74hcuBRdAK7/nNDHf3xjUGB+ErgC2O6Ud45leXQH2CVnnPd71SbqtPZl4CMMbslD0CtllO9cmvH1t8+gSx2/zeuJRRl+9xuyF65q85pD0zQH59krE5aIVqvzdnOaT1vkouhABedkg7LkB+DXEUQasFX5LGOKR9xTezEAJwOat+fID8+qUigutkl2j3xZ/UGm1+Lec6pF2XmxXNAsxwWL92gAZ39+Zg2wUcegH1ZuxlkC23IBY4Hrer7D8X22YB73UwGGBaBBtPqYjh7FYMfBXsglcXxsa2uXq96k4Ho2+PbXOrAvzMoRy78LgWt4p3H6ByYdFKh0c+BpRuEipvsDWTy42gf63cQOJFqEAsc0g7QvbypMTBXChpuEN6n7ua3iKUAB50TD8zUL7txAzH9K5lnh/YSxNfdWi/uoGsq3PaDmAEbpdNuoy9NCRIDWCMQdI9Do8MkjQnRN5twly53f7lUtbFAExwUDDYXqzrzFfbARwL/Nex7Fz2SBQH8HdHR26KbXNsgB86ltnTsW3uE2CWozPdlOxcvZAAo3G/aPrq2Hb3CXAU9bduN2IpJTpXLxTYjvNyx7J6HRgc2/aGAN92dApgVmy7iwa4PkM5fSu23f2C7dR0OTq2i1YY1y4I4D3AbscyehMox/FzwDcdnQNYCxwd2/a8AYYCL2Yon6/Htr1psH2B1zI4+QAtPFvnC3bi7KEM5fJvWr3trwWYnsFRgDtpw04httPnMtvXm8ti258J4MGMDt/cTiKoBP/mjGWxJLb9mcEuxHDtEFb5CW0wbQwcBszPWAbbgJNi++AFcGVG5wF+DQyK7UNWgEHAYg//p8f2IQjAAo9CeIwSzhwCJwCPe/h9R2wfggEM5tDFmS68Crwvth/NArwfWO/h7yoC76OMDnAysNGjUPYC82jhfgG2s/cV3Mf3e7OBsrf7fQGcC+zwKByAR4HTYvtSC1bgSz192w6c039uJQaYDOzxLKgdwBwCrIoN4M9hwGzc5/Rr2QOU5q4hL4BL8Ksmq6wCJkb0YwKwMoAf3YDLiuDyA0zBfUKkHj3AQgpsFoCTgJ9jdzr5shu4uCjbWwrgQ7gdjd6IvcCPyfH2TmA48APCCBfs2okP5mVvKQDGAK8EKlCwbel84PSANo7EBt5lB1R/vAK47ARuX7DfrBUBCxdgH3YkMfP4ATAWuAvbRofkL4DLJpD2BxiAXSDpsju2Wf6G3cXc75RqxY7LgGU52NKD9bFtp7y9AaYCWwIXfJU3gVuBCXXyPQv4LvCfnPLeQoSd0f3RklOu2I7cnZKm5JjNWkkLZQ+g/qykicqvPB6WNNMYsyGnz28/sMOqMwn3lhCDLuAa2mhtQ+EAJwJ3x41jJn5Djq+kHQcwCXg2clCbYQ0wOXZ5tSXYcfer8JtqzYuN2J1RqYefN8CRwA3Apqght2wBbsSe8ZsoEmAI8GXi1AjV08beHrscOh5gIHagx3V3chZeAq4FhsT2O1EH4HzgHuxQcCh6sBs60gFXZQE7gfMN/CaaNgPfB3I94ziRI9gBpUnY+ftmVu3sBe4DpgEDYtufNx01QoWdDLqk8jNe9mz+wyVtlrRa0iOS7jHGeJ3Bn0gkEolEIpFItDT/B1wlvBOrKdY1AAAAAElFTkSuQmCC",
    "shopping-cart": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAIaklEQVR4nO2dXaxdRRXHf4vParmFhKCWGntLQWuAKEEQrQZUSBUJxEQ+qg8+mBgefDXGaHwxEo1vGiL6ZIiJAgoYNcUogj5IoaUXlBdtCyVGL9/Ygtj2tvx92GdDuZxz7157z5zZ59z1e9xn9syaff5n1uy1ZuZAEARBEARBsOKw5QpIWg9sBc4DTnLU/SzwR+BXZnaknXlBUSR9TdIhdeMxSZtK9yVwMvjyUzEv6czSfQoaImm9uv/yF3Nr6X4Fb+a4Ede34vP3TbhW0lsS1xl0ZJQAzsvQ1irgwgz1Bh0YJYCjmdr7YKZ6g5aMEsCuTO1dnKneoCVD4wCSzgD2AjOJ23vSzGYT1xl0YOgIYGbPAl8GlLi99ZLekbjOoAOjXABmdivwOeC5xG2GG+gRTULBq4FPs3Qo+Ebg1IZt3mRmX29YNpgEJN3jCAj9obS9weuMdAFOHnSUvVhSqnaDjpQQwAwQyaGekFIAnjeGmAj2hCQCMLPngccdt0REsCek9MUeNxAC6AkpBfCQo+z5g9fLoDClRoATgAsSth20JKUA5oBDjvLhBnpAMgGY2SHgr45b4k2gB6QOyMREcMJILQDPRHC9pLWJ2w+clBwBAC5K3H7gJLUAdgPPO8qHGyhMUgGYmYCdjltiIliYEzLU+SCwpWHZyyWlXnUUVLwEbAe+b2a/GVUoR1rWMxEM8jEDXAH8WtIPJQ1d/JNDAN7MYJCfG4GvDPtg2SVhbZC0B9iYo+6gNQeAdWb28rEXc63MCTfQP9YAly6+mEsA3nhAMB7eufhCjAArixcWX8g1B1gF7Cf9DuOgGxvMbN+xF7KMAGZ2EHg0R91Ba+YXf/mQzwVAuIG+sX3YxZwCiIlgv3hg2MUcoeAarwC+Ctybw5Ap5QPALY7yYxdAnRk8vWH5NWb2cEZ7pgpJH3cUXwCGPttsLqBFZjBSwz4+5Cg7Z2b/G/ZB7j16HjdwUewZdOH5wQydAEJ+AXjeBE4F3pPLkGlC1emtnnMXh/p/yC+A7fgyg+EGmuEZ/qGUAAZ7Bp9w3BIrhJrhEcC8mT056sNx+NxYKp4ejwBG/vqhfwI4X9Jbs1kyBQxOW32f45biAvBMBE8k9gwux4X4kmwj3wBgPALYRewZTIln+F9gmUM/swtgsGfwb45bQgBLc4mj7CNm9spSBcYVeHEdIpXNiunA8wNZ0v9DPwUwK+nt2SyZYAYBoHWOWyZSABCjwCiSBYBqxiWA3cCLjvIxDxhOsgBQzVgEMMgM7nDcEgIYTrIAUM04s2+RGezAYKGtJwC05Pt/TV8FEJnBN+MNAPVuBPAuEo2J4BtJGgCqGZsABn9CEaeJtscjgGUDQDXj9rOeUSAE8EaSBoBqxi2AyAy2IEcAqKbPAojM4Ot4A0CN3gBg/AKYAw47yocbqPAkgJ4etgVsFGMVwGDPYJwm6sczAvzFU3GJYEssEXMwCAC933FLY/8PZQTgeROIzGCmAFBN30cACDfg8f+NA0A1JQTwD3yZwZUuAI//f7RpAKhm7AKIzKAbzwjgGv6hzAgAkRlshKR34QsANX7/ryn1YD0TwdNYuZnB5CuAFjMJIwDAR7NY0X8+4Sj7lJl5tuGVRdJex/8NPybpxNI2jxNJ6yS95HhGd7Vpp6Rv9biBc4EfS8p5oklvkHQa8EvgFMdtrghgcSR9yaHumockXSOp6V/VTxSS1kr6oqR9LZ5Nq8RZloMim6AqwvdPqqxf0I29wDmDV2wXxVyAmT0N3F2q/SnjR22+fCg4AgAMhq2dlJ2LTDrPARvN7ECbm4s+eDObA35S0oYp4Bttv3woPALAazPeOWC2sCmTyD3AlW2Hf+iBAOA1V/BnfK89K509wCWDc5ha0wvfO3AFVwP/LW3LhLAPuKLrlw89EQCAmd0HfAz4V2lbes4OYLNn3d9S9EYAAGa2g2r50+2lbekhC8BNwEfM7N+pKu3FHGAYki4Fvgl4DkWeRhaofhDfMrO/p668twKokbQJuI7qTxAvAFaXtWgsvECV2/8dcIeZzedqqPcCOBZV/365FngbcHxhc3JwEHhmsI8yCIIgCIIgCIIgCIIgCNLR20igpI3AZ6gyhLPAGcABYJ5qCfQ2M7u/lH2pGEQ3LwM+RbUTaC0wQ7XU6wngPuAuM/OcsDa5SDpL0u2SjjRYCj0n6arSNrdF0tWSHmnQzyOSbpO0obTNWZF0g6T9rtXwFTdLOrm0/U2RdLKkW1r08z+SrittfxYkfV7S0RYPpeZuSb1PEEk6XtUI15ZXJX2hdD+SImmzpMMdHkrNd0r3ZTkkfS9BPw9J8u4c7ieSjpO0I8FDkaQFSeeW7tMoJL1XaYQuSbs0DecmSNqa6IHU3FG6T6OQdGfivnaeD/RBQdcnru8qSTOJ6+yMpFOoXvVS0vnZFRWApJOolnqlZBX9XEd4OZVtKdmijucmlB4B1gE5DoQ+O0OdXXl3hjpXUwWOWlNaAJ5DID373zo9lEzk6uuZXkOOpbQAPHjC1r0NcTdkbPaXFoBnubNnYtfH3UW97GsfBJBjP+CeDHV2ZXeGOl8GnupSQenzAQ5TbX5IySvAvYnrTMHvqWxLyTYzW+hSQekRANLvA/ytmfVul/HgDN9tiav9ReL6xo9WVih4k9KFgh/WNISCASR9ONGD+XbpviyHpO8m6OdBSdN1iLaqtQBd0sE/0wT8IlSNeLd16OdRTfGagGtVLXrw8KqkH6gKK08EqhaE3Dyw3cOLkj5b2v6sSJqV9HM1WxK2U9InS9vcFklXqlrWthwLkn6q6v8Dk9LbiJmks4BrqBI7G6gWhe6nih08QPUK9KdyFqZD0mXAFmAzVRh7DfAM1aLQ+4E7Ux0JEwRBEARBEAT8H6PebJactX/eAAAAAElFTkSuQmCC",
    "sparkles": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAANoklEQVR4nO2da6wdVRXHfwuKFAoUkECLyENailDoEwSMEOUpBuQdNYggRhMFAQFBo5/84PsVeUUUECMYREBUVNAYQIpUkcuzL6BQCA+BqrVgaUv/ftjn0ntPZ86dNbPnzMw555ec5N45e/ZeM7PO3mvWXnttGDBgwIABAwYMGDCgl5G0haSPSbpe0gJJy1ufBZKuk3SapAlVyzmMpAmSLpL0V0kvS1ok6QpJe1YtW6OQtEnrRi7X2Lwi6QJJ4yqWeZqkJ1NkfF3Sx6uUrzFImtz6BXmZJ2lSRTJvKemJMeR7Q9KhVcjXGCTtmOFGdmJxFUog6eKM8t3fbdkag0K3f1+Bhz/MX9Tl4UC+HmuXIm1tFEvoGnIusH+Eet4NnB2hHg+7OcruWpYQjUXBes5i8GXlJUmbdVF+j+yHF2mrV3uAE4BtIta3HXBcxPpqQ68qwFEl1Hl0CXVWTq8qwNwS6tyvhDorp1JnRxkoePKmlFD1VEkTzOzVGJVJmghMB5Jsi00cVc2WpLZj64DHzWxZXvkai6QDIxp/7RwYQb63KbidV5co5zD3S3pfJ3l6cQiY5Si7qvUpo+4NkDQFmA98GN+vPC+zgTsknZFWoN8V4GHgEUf5mU5Z3kSSAdcDO+atIycbAZcrZQKpFxXA85CGWp8y6m7nUMoxTrOwKXBe0hc9ZQQquGynO04ZAsxRfh9J48xsrU8yAA7JcU5MEtvvKQUA9gTGO8p7FWA8MA141CNUi61znBOTiUkHe00BPOP/OuChEX9nHQ5nkU8BnspxTkyeTjrYazaAZ4xeYmYrzWwl8HhJbYzk1wRFq4pbkg72swIMpfwds403MbPFwBV5zo3Ak8ClSV/0mgLs6yjbVQVo8TngFwXOz8NS4Ggz+2/Slz1jA0jamTBrl5W8CvBWSW83s2cc5wBgZq8Dp0g6DjiT4KhJcgVvTXbjdCWwpu3YOsKwdjNwSSz3da2R9EGnm3TSiHN3cJ57bMnXMogHyIGna37ezF4Y/sfMXgRe6FC+SFu1pl8VIKnL75YdUCsGCtD5WIy2ak1PKICkrQFPdGxRBdhVUsyQs8roCQUgeOc8Lt0HMh5Lw4AZjvJeVjrKJr7eZaVXFMDTJa8Enkg4/ji+G1/mMPBgxnJrgQVFGupHBXjQzDZwybaOPZRQPkabXq7KWO5GM/tPkYb6UQE6jfV1MQRvYWyP4fPA+UUbarwCSNoUeKfjlFgKsJckz9RzZsxMwKnAZSRPIP0dONjMnivcVtEKqkbSHMINycpcM0tcVClpLvA3R11zzOwfjvJuJE0FjiVEOr8C/Bn4c9Iw1pdIOtPhNl3T6VcraVP5onUbv0a/8UMAvrF4gZmlRgG3JmsWldR2Lek3BcgyxtfFEOwKjVYAhVDrvDEARcoMM0NSo+9ho4UnGEZbOcpn8fZ5PIJbAe9wlK8dTVcAbxecxcPm6QHyyFArmq4AH3CUXWZmy8cq1CrjWVRZxlL0AWMhabpCurSsJEbFptR9i6PeVZL2LvNay6SRPYCk6cBvgLc4TrvXUXa+o+ymwG+arAS1R5JJmiLpFEnXOH/5w2QeqyXNzVH/65Kuask4ReENpfbUTkhJbyGs75vZ+swivOp5rP12HjOzzL/Q1sNbBEwt0OYKwuziA6xfhPqIma0uUGd0KlUAhUiemW2fvYi/dv7TZna5U7bPAt+PLMca4DHWK8QQMGRm/47cTma6pgAKcfvtD9uTDy8vy4BpnVzASSikhVsM7FSKVKNZyoZK0ZX0LtEVQGGJ9jRC1z2yG982dlsZOcHMbs5zoqQTgRsjy5OV5YwePh4AFuVcmp5KIQVQSMi0L+sf8kxgH3xLtMvkKjM7s0gFkq4GTo8jTmFWEbKaDCvEEPBQkZU/LgVoPfATCM6POQQjqa6vkncBR3q7/nYUpo//SEgZW0fWAUsIMRG/B24ys9eitqCQePlCxU2/WiZ3SSry1tB+/RMVkkY3gZclna9YCa4VUq7nybdfFT9WCaFaksZLurraS3Nxj6TJRS96sorl2+8mT0s6PtLz7nRPTpS0rOJrzcoSjbHfQaoNIGkT4G7gXbFvYmQWAJcQDL5C431WFF4RzwQ+Q8hLVGfmAYekvT10UoDPA18vS6oCLCNYwPcCfzAz7/RtVCTNBo4ADiC8Ce1cpTwpXGBm3076IlEBJG1BuNFVrn9bCyxktIPkgSxTulUiaVtG+0BmEnqJKpNxvALsnPR2kKYAHwOuKVmokawk+M1HOj0e6VaXXjYKRul0RivGvsAWXRTjNDP7afvBNK18f4mCvMDoBz1EyGzds3HuLUX+OyPWLyjEEk5htBNtJlDWJlVHARsoQFoPsIDixs06wiLMUe7MkZk5BmxIy2pvd6PvTnGHW+KMaJoCLMc3/r9OSLo88mE/2MrBN6AgLZtsBqOVYjohGCUry83sre0H6+rGHdAlBkNAzShxCHjUzDZIpJ1mBA5RXAE2IkwWTQVOGT4oqe+MwCQqMAIT/SVpPcBpwE9KEiSJwWtg+ZxqZj9rP5imABMIjqCqgjhg4AiKycvALpkdQQCSLgS+UaZUORl2Bc8HbjczT26AqCgEj84luIL3o76u4PPM7HtJX3RSgHGEoIrCO2WVzBLCZNCVZva/bjQoaXPgk4TJoDK2qIvJ3cB7zewN95mSJilsod4EnlGI4SsVSSdLerbia83KQkk7FL3gSQrBBU3hapUTELKZpGurvTQXd0raPtbFj5N0nsIu2k3gboWdOWNd/0Q150fwT0nnSNo4y7V5g0I3A44nTCzsT72DQu8BDosQFLoZ8CfqawsNB4XOB24DbvFcc6+HhV9jZqm7ZmZB0rXARyPJU5Rqw8KzoPotDDnJzH6Z50RJJwM3RJYnK/VfGOJB1S0Ne4awNMz1iqjwqrcYeFspUo2md5aGeVD3FoeebWaXOGU7F/huZDn6d3FoVlTO8vCFZpY5nayCh28xxZw8jVge3gi0PkHEyZJ+pJCWxUvmXUUl7Zej/lWSrpR0khqUIKKRSNpL0lLnA7rYUf8XnXU/KcmTsHpAUSTtrfKSRP3KUe+qwcOvCIWcPFlJ3Dw5pd6nHfX+sMxrLJu6evGy8jtH2Z0V5uo70irjmdK93VG2djRdATxpXSFbVk/PFvR5ZKgVTVeAJwDPnjlZHq5HAVYQduZuLI1WgNbWKrE3enKln2/J0FgarQAtYuf3j73/QK3pNwXYU523jBlPmMgqo+1aEiVSVdKOhKDItwM7EJaVvUiYiHkGmFdivJ7nIYwjuJnTAkn3wXdP+lcBFPLPnA4cR3j4nVyfr0r6A3ATcIOZrcnbbgKPAKvJnjh6FukK4DEAVxMmdvoLSVtIukjSCoezZCSLFHz60XzlkoYc7V/aoZ7LHPU0+vVvGJcNIOk9hD12vwZsmbPNPQhBFrdL2mC1ak5iGYJ9ZQCCQwEkfYKQMLFYmPF6DgPuU5w8+4U3emod26ekNmtLJgWQdD5wJb4NGrKwO3BPBCXwPIwJJM/zT8W3Vq8/FEDSkZSbLWwiYceN7QrUMQR4HDJJxp7HAPQ6oGpLRwWQtDthvM4UY16AXYHr8hqGrZCqpxynJI31nvF/qZn9y1G+tozVA3yVYqFYHg4nrDnIS1FDsO8MQOigAJL2A07qoiwAX1XIUJoHz0OZnXBsoABtfInuB43uAZyc81zPQ9leI3LoKji1PG83va0AClmpjshR33LgVkKSyTsJSR685F3h63XMzEr5u4y2moVCRmwPayV9WW0TLZJ2k3S7s66VCuvx8sjtWbz6hRHneYJAXy56f+tE2hBwuLOes8zsK+2LEs1sKXA0cIejrgnkX4iZZW/gYWam/D0WPfXrT1MAT0zcPDO7Iu3L1lq2T+EbDvKmWcn7JtCXBiCkK8COjjquH6tAqyf4q6POvOvxPA9niqQtFeyd3Utqo/akKYDHIn4qYzlP7Fze+QZP97wR65e2eybFemoISIsHWEH2hIVbZyznyT3sCfQcyULgNWDzjOVn4nvV/R9hzWDPkKb5zznqGPN1sWXVv8dRp6f9N2llwnrUccrwAtSsPBx7fX7VxFCAD2nsnbkvJntP4W2/Ha8h2LcGIKQrgMdg2wS4VdKMpC8lnUXwKmZlHXCfo3w7noc0HfBMRXteMxtBWqrYXQhZKzzj4xrg54SgkVcIlvVH8O86dq+ZHeQ8500kHUjYKasMDjKze0uquxI6ZQr9B34XaQwuMrPcKWoVEletIH7I+zpgqyIJmepIp5v0g65JsZ6VwLVFKmg9oCVxxBnFkl57+NBZAa4lhFx3k29G2lCijATS80uos3JSFaD1SnVRF2V5DvhOpLp+H6mekdxWQp31R9K3HDNleVkt6eCIMm+usIt2LF5UzhnKxiNpY0m3RbyZSRTK5pki9wUR5TsntnyNQmHS5LcRb+gwaxXy8ZUh8zjFSfB8pzImXu5pFHqCr0W4ocOskHRMyTJPVthCPS8LFSvleq8g6WhJDxe4qeskXSdp1y7JO0nSvBxy3q2imy30Kgq9wRmS7nfc0Nck3SxpTgXyjpN0vrIZhi9JOld90u0XjvpVSAJ9DHAAsBNhGnkbwibRzxI2ebqDsMFTpY4UhQTQJxL2O5gBTG599TxhDuF3wE3d2ntowIABAwYMGDBgwIBK+D9C4a9G5YUNlQAAAABJRU5ErkJggg==",
    "star": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAL80lEQVR4nO2daaxdVRXH/4uxDC2UGQtlKi2lylQZpUAEgggGCIPQCEosEEmUaEI0xgSNkviBD5hoFETTEAYpoOJQAhSxDAFqoUgHbEspQ6HQUua26UB/ftjn0dvXd1/vOufsfe477/ySm5fcd89ea+29zjrnrLP32lJDQ0NDQ0NDQ8Ogw6pWIDXAFyR9U9KpkkZmX78uabqkyWY2pyrdGiICDAMmAxtozwbgj8DQqvVtKBFgL2BOPwPfmxeBPavWu6EEgG2ApxyD38PjwNZV699QEODaHIPfw3eq1r+hAMD2wJICDrAE2L5qO2KyVdUKRGaSpBEFjh8h6dsl6dKQEmA74PUCZ38PrwHbVW1PLOocAa6UtH8J7YzM2qoltUwEAdtKWiDpwJKafFXSaDNbV1J7XUNdI8AVKm/wlbV1RYntdQ21iwDANpLmSzq45KZfkTTGzNaX3G6l1DECXK7yB19Zm9+I0G6l1CoCEDJ38ySNjiRikaTD6hQF6hYBJso3+OuzT6ccIukyl0ZdTm0iQHb2z5U0xnHY7dlfzw3eQkljzexTxzENsQEmOhM864ExwChgnfPYWkWBAQ+wFTDbOYh3tBx/h/PYeUDdLp8DF+AS5wB+Cny+5fix2XceLqnS5oYMwID/OgfvT320c7ezjbk0UaB6gAudA7cBOKKPdg7HHwUurMLmhgzC2T/TOWj39tPevc62ZtNEgeoAzncO2AbgyH7aG4c/Cpyf0uaGFoAZzsH6Swdt/tnZ5vNAbfIpAwbgXOdAAXyxg3aPov+p433xtRQ2N7QAPOscpL852n7A2fZzNFEgHcDZzgECOM7R/jH4o8DZMW1uaAF4wjk4/8wh4x9OGc/EsLWhF8CZzoEBOCGHnONyyDkrhs0NLeA/+x8qIOtBp6wny7S1oRfA5c4BATi5gLwTc8ir3ayhroCQ8l3tHIxHS5D7sFPmKuCCMmwe1BDSvCOBC/A/lvVwagl6nJxT9l8Jmcoy1iZEo/JnV2BvhWlcoyUd2uszpEDT083stMIKSiJEki8XaGKVwkyi1s98SQvNbHlxDfOTxAGA3bRxUHsP9LBIYk83s3+V0RBwmqTHymirDz7QRqdYkH0WKjjHh5FkfkZpDgDsrE0HtvWs3r0sOR3ykJl9pcwGgYclnVlmmx2wTBsdY5MIYmaryhDgcgDCkqt24fpzZShUAqskHWVmC8tsFBgjaZakHcpsNydIWqK+LyuLSp+2Tngcug9YmfOGKCWTSjV+0364umrjOuAT4B4cqe/+DN4O+A3+vHhV3FTCOG+pT26u2sgO2QD8irBULpehW5P/8asKbibBGznC4+mvK7bVw/3kmbUE/LRqzTtkLfDdCGO9pf65LpM9EPhxOzv6PGOA/RRuKoo8h6dghqRrzOyFKoQDR0u6VdIWJ5pUzGpJh5jZ0t7/aBcarlR3D/4MSRdJOqGqwZckM5sl6ThJl2Q6dSs7SPpWX/9oFwGekJT7JUoElik8gk2X9ICZzatYnz4Bxkk6T6EM7VGS9qpWo014zMw2y2a2c4B3lF75SjNiMQB2VfsM6K6J1VliZpu9l2jnAO8rjoIrtXnyYoG6ICeeGkIp2nZJtR0jiFxmZnv3/rLdM+JS5XeANQrlVPpKXy7J2WbtyBx+uaSnev+P8Aaxt1OMVqhSkrdk3WY3gG0hVMzOw+ws7DVEABiOr+h1K7/zCDojpxCApxsnKB9gV+CZAuNyilfg9ALCZhJeATeUALAb/jWQrUzLI/Rg4N0CQl+gqblfGGDPrC/zsgw4IK/w8cCKAsLnAPuU3CeDBmAf8l/zIZzAxxRV4nBgaQEl5hNSyw0OssGfW6Dfl9HPSmivMocBbxZQZjFwUCnKDAIIk2EXFujvtwmbY5Wq1GjgjQJKvQYcUqpSNQQ4AFhUoJ+XAofHUu5A4JUCyr0VTbkaABxK8ZNsVGwlRwIvF1DybVoqdDUECDULi2xv8yoQo0Zyn8ruS7EblHfoo1DTYIVQpu6tAv2Z/kYb2Bt/ccZW3gOOTap0F0KoSLK8QD++BFQzI5uwIaO3Rl8r75Nj6XZdIORZiiTb5gH7Vm3EcPzFmlr5ADipUiMqADgJ+LBAv80C9qjaDkmfvah4uoAxnwBF1t4NKIAJwEcF+msmkHq1Vf8Au5Bvi9YeVgKnV21HbAhvWosssHkK2KVqO/oE2JlibxFXUuN7AkLYLzL4/yaswexegB2BaQWMXEoNXyABIwg5kLw8AsSYKlY+wA7A1ALGuqt6dTv46w21MhXohkWpnUPYuLnI0rLzqrahLPBXNG/lAQbqBtbAtoRVxXmYTQ2qbxLWEubNmt5HWJI/cAG2Ae7K2QFnVK1/UYCzctp+F3lX9nYbhJXGk3N0wuSqdS8KcHseuwm7oNUHwsZOtzk74h0G8GUgs9mb5/89dd2EghAJHnN2SNz32xEhvOXzMC314CcVlm22eL3zsIE8gWSs8/fXm9mGKJq0IXmoMbOZkt5wHDIili4J8LymXZwtN09KVdeaVxy/3SmaFvHxpG09fVIaVTmAJ6U5kHfqXuf47dBoWvRDcgcg3NUf5jhkRSxdEuDRfSwVPPFUEQFGyufti2MpkoBFjt8OlZRvCVcBqnAA74KFriwH0yFznL8fF0WLfqjCATxTwt80s/eiaRIZM3tf0luOQ5JPl6/CATxe7j2DuhGPDYMiAniMnBtNi3R4HKDeEYDwgsPzBFAHB/DYMJbEL4FSR4BR8pVbH2yXgCEKfZSM1A7gCXFIeimWIgmZK8mT3096GehmB1hsZh9H0yQRZrZS0quOQ2rtAIPtCaAHz31A0ieB1FOOKnsCyNKsRyrU8T1C0n7auJfRuwpbsLyoUI/4RTOjRPFzJHW6xXzSCJDMAQgzW0c7DiklAhAqlV2rUAG901TrYsJ0tN+WVMLWY8uhwPZmtqYEud0DcIRzdkyhAkeE9Yo3EtYd5uVj4BcULHyZ2vauBJjo6IB15JwLT1iZdB1hPmFZrCDsoJLrlS1hjcQ6h7yJeeTkIeVNoOf6v9AbAglrEK5WKEx9s8otd7+bpBskLQJ+6HXOzBbPNnbJbgRTOoDn5qbjayZh5u3FCm8Nb1Hc/Qv3lPRLSQsIW8h5snZdmRLuVgfo6AmAsHDkOUlTlDaDNlLB2WYDF9PZRA7PU029HADYSdKBjkP6PVuA04FnJD2isDVLVYxVcL6n2XKBi9mOdg8i0TLwVBFgnFNWnw4AHAv8XdI0SceXoVhJHC/pUeBJYEKb33gigCnRdPhUDuAJaWvUayoVYYHFFEnPSjq3TMVK5kuSHies5e/9KPeywvZtnZLkMpAyAnTKvJ7NjwlFKW9RCJ8XK9F29yVwhqTngSlk5XGzRTH/c7SR5EkgVSbQlQImVAj5iaSrlH+PnP54T9L9CqlfKaSIL1K5G2VtpeC05wG3SrpR4dJ2dIfH16eaKr4SqHMplr3rj48ICZ1hfeg4DPgZIfsXg0/w1Ql4M8XYpNhsebjCGVclayVNlnSDmb3d3w8J5deul/Q9+SavxGD32JNiU9wDVBnK1ku6TdIoM7tmS4MvSWa2wsx+JGmMpD+o2pVJ0fsuhQMkn+mqMJvoHknjzOwqM/MsRpUkmdkbZjZJYRCmZG2mJnrfpXAA7xLpokyVNN7MLjWzBUUbM7P5ZvZ1SeMlPVhYOx+eCbS5SOEAwxPIkKQnJU0ws3NiLLM2s1lm9lVJp2SyUhB9/8UUDvBR5PZfkHSOmU0ws+gDY2ZPmNkESedksmMSu++SOECsuX0LJF0q6RgzmxpJRlsymeMlXZbpEgPP+4PuBNgfWF/i8/TrwCS6qIQaoRTeVRTb76c364GBXB1lI8DdJXTIcuD7wJCq7WkHMAT4AcV2AOnhjqrtKQ1CoeS8nfIhcAM5p2NVATA00znvhhDLqHonkLIBTiTsDtIpq4Gb6JZdMXIA7JHZsNph9wfAiVXrHgXCDqRb2l5mLXALNdpuFtgvs2ntFmyfAaTOm6SFUDz5XOBOwiaUqwizbp8Ffk6Nt5gFDiJMVf9PdqavyfrgzqxPBsrr7oaGhoaGhoaGhoaGhoaGhoaGhoYBxf8B76x/+zsAvwsAAAAASUVORK5CYII=",
    "target": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAARrklEQVR4nO2de7BV1X3HvwsFbEBFrIVGUNBoNFSJTWdKoqAQpWkeGmMZTa3G4KONNmp8ILVJdBKTmoSkPkg6ZiIaSGxj04yKadqCIGp9Jjq1SRS1CkJ5qAiEC5QL3E//+O0Tr3jOveu39jpnn4v3O3PnMJx9fo+91977t35PqR/96Ec/+tGPfvSjH287hKoFaBaA35d0hKTDJb1b0ihJQyQNLT73Kz4labOk9cVnR/G5UtJSSc9JejaEsLqV8rcKu8UCAIZLOl7SZEkTZBd8n8xsNsoWw6OSFklaEkJYn5lHy9EnFwCwp6QTi78pksZLGtBiMbokPSVpsaQFkhaFEHa0WIbS6FMLABgn6SxJ50gaUa00b8Hrkn4saV4I4aGqhYlF2y8A4HclnSu78OMqFicWv5Q0T9KtIYR1VQvTJwG8E/gW0EHfxSZgFmaQ9iMGwMHAjcCWKq9cZmwD5gKHVX1+2xbAMGA2sL3SS9VcdAI3AftWfb5raAsbAJgm6Wa1n2HXLKyRdJXMYKRKQSpdAMDhkr4t2869HfGApAtDCL+qSoBW750lSUAAZsqs5bfrxZekSZKeAq4EKrkZW84UOEDSXEkfahHLtZKekXnxXpa0RW92/Upvdg0PkXSQzIV8pKTfa5GcP5V0TgjhtRbxk9TiBQBMknSHpAObxKJD0oMy79yDMh/+hjIEgWGyhTBR5mqeqDdiCLmxUtIn+5IjKRrFY64ZFv4y4DrgA5iLuNl6DASOA74KLG+CPtuBy5utR8uAve+/kfkkbcT21ScCldgxhW4DsMVwSyFTTtxUpW5ZAAwC7sh4UtYC19JG++gagKHAJcDKjPr+EBhYtW5JAIYA/5rpRCzDTu7vVK1Xb8AW/dnA85l0XwjsXbVeLmB3w8MZlO8AZgKDqtbJC2whXA1sznAeHgaaZXjmBWYk/SyD0vOBMVXrUxbAgZi9UhYLgcFV69MjMKOo7Dt/LfDRqnXJDeAU4JWS5+YHVOQwigIWvi2D+4F3Vq1HswCMABaUPEezc8qUbTUBV0r6euLPuyRdJ+lLIYSduWSqAbOkD5E5dEbLHDnDZAmikjmQNhSfKyQ9K+nFEML2Jsiyp6RrJf2N0l3xl4cQvpVNqLIAJpHu5NkMfCSzPCOAM7D9+a+xMKwXncVvbwFOB7JGKrFXQmrOw3bg2JzyJAMYTrpH7HXguExyHABcDDyRKEtv6AIeBz6LxTNyyPzHwGuJ8qzA0uWqA+bluztRgVXA0RlkOA64i7S7PBWdBc/SdyHwHuxipuCnVGkUAlclCr6Ckls8YCpmNFaNxUCpkDYwlnQP4hVleJcRehxpd906LL07le+h2MpvN9wHHFFCr6OwV6IX24AjU/mmChuARQnCbiHxnQ8MxuIAWxP4tgpbgWtIdNhgNkFKFvQSWvkqAM5KEHInidY+lin8SALPqvAL4F2Jup6MnSsv/jyFX4qA+2AGnBfXJvI7lbRHY9XYCJyRqPNXEvitwZJXmgssdduLRcAeCbxSjcx2wvUJeu+JPda9uNHLyyvYWPwOn7U43buYjfHNhBPQrpiNM7kDGAmsdvLpxLm78qZQzUj4zbkhhFWxBxcnaq6kM518YrBD0kuy5NDX9dak0IMljZVfx95wkaRhwNkhhK6YH4QQ1gB/JekuB5+Bkq6Q9NcJMvaMYkV6LfCfJPBJecU0whbMYXMJcDQReQVYHH988Zu7yVuidlPC+bjHyWMrzQio4Y/0bcb5OAK+6OTRCI8C04HSTSKAfYFzC5o5cLWT/0H4t4bfKKv3rkLsj1W6ejDTyWOak349LAYmZ1X+zTJOobz3sQv4hJPv5508OsgUr6gJ4LXGX8KRxoV59zY4eXTHKuDsbAr3Lu/HKJcSvh4Y6+A3GHjZyWNGToWfdjL/jFO5nzvpd8cdZHjUe4H5Q/6phNyP47tJLnbSz1NvCLzPyXg1juxdzL2bgk7gL7MoWQLAZ0iPRH7BwWcv/A64Y3IoeIOTaXRVC/Au0nz7HcCfllYuE7ACld8k6LEFOMTBZ4aTfrmsISzDd62D4UZgaO+Uf0v/35wKgV38D5RSrAnA8hJSUsDvdfDYG99CW0OZcjngQ05lbnXQnuqkDfaobVVVsRvAR0h7HUxx8JjjpD21jEJed+zxDtop26ks73zM2TMaeH/xN5pMxSfAhQl63eegP8VJ+2tllHnSwWgZkf5u7KR78aMSeuwBnIA5s57rgcdzxTHHkxC86sbvhwn6RaWXYXGSlxx0n0hVYji+uPR1Dtp3OeiCpUu5t3rYyZoGvODkB7bvvoC0KOa++C32aLc5Vpoeix3Afl4dBJzmVCDKMMOyd73vyWkJ8r8PeMrJpx6eJGE7BXzSyaeTyCxfYKKT9se98nuDMpuILGPG79BYnCD7NPIUZdawFfiLBDkecPK5KJLuIHzxAX+eAL78+p81iS44LOSC/kzM354bXTjdq8BJTh6POGh7ttCPeeSuMfB0vLgykuYIfBfnUafMZzjpe9GF83WEuXxjsZP414AnPtOwT1Jdqx2LJ3uMrtimRpPlq0f8buyBwB9JmuOk70WQdDs+m+B7jmMHSDoh8tgHHXT3BUY2YlgP73YQl6yYMgaeUO1WSf8ScyC2/fyupFZ0EHmHpNuIT/H6kUyXWMS+8p5x0JSs7d1b0EgJT4HDWsfkjIkOuv8RQtgYeeyZksoHPuIxXlJUxm+hwyIH7UmRdNdLetVBt+41zfEEiLr7sV2CJ1c+6qRhvu4vOejmwpeJ9xF4FsBhxPvvlzro1r2mjRaAp799rBCHypIWY7Ek8rhjJY1x0O2OzuIvBYcUvGPg2coOKmjHwLMA6l7TRgsgOqIn624ZA89TZYfi33GnOOhKEpJulTQ+hDA4hDBY9ki/rfjOg1jev5LpFIu67+s6WOGgWbfTWKMF4GlLtinyuIMcNF8MIcTemR9z0EXS9BDCeSGEp2v/GUJ4OoQwXdJ58i2CKN6FLsscdA+OPC723EttsAA8NJfFHIRF8A510J0TQri90ZchhDmSGn5fB4cS38hxmYNu7LlqiwXQUUaABvhN5HEj5dv3x7hEPbn7AwoZYhCrk9QGC8BjA8QugGbQ9HQd75TNJ+gNT8tnGMbKUPpiNYNm325G7HtfB8U/LTxPlUpHvpRFowUQewdK8Xe2h2bsHRBdcyjbgv5BxHFHy7ddjZWharuqLs1GC8DzaIldAB6asXGINfLdgZdmOqaGrkKGGHhiK31qATTjfTUm5qBie/U/DrrnANMbfQmcL+lTDnovOJpJRlcCaTddAC87aI4lPklzvoNukPQ9LLP2vRR9dYBjgO/LEX308MZ6Bo1x0F0eeVxbLIBRkcfFRgwlq8+P7Xx1t4OuZIvg07LJ39uAbZKelJRSWxjLe5wkT25hrIt3tIOmawGsdhCOdfG+KMnTe/eEyOMekjV9SMFAme89BS9KejjyWE8YvFPx+njc63WN1UYLwBNkiAodF+/KFxx0o+LiRXPp6Bq7jLja0djaswCeDyHExg08C6DuNW20ADyP6xHEd6fyZLFMddC9Q9IvHLTL4ueS7ow5EJtv5MlrfCCS7n6SPD0AXAvA8wSQ4t/XnrDoXpJOizmwmL97gWwoZLOxWdJ5jpm/p8uXqRSbO+DtDvpcvf+suwBCCKvl813HZvoslm/ffn7sgSGEJ2XGXTM9c0j6VAjhvxy/Oc9xbJek+yOPjcocKrAhhBDrrzAAjzmyTj1p4Z4sWYAPOuWeQfPSwl2NmfEXwMYalQL+3UHXlV1dY3CTg8FmIvftWL99D+5PkP3PyF8Y4mpbh5WlPejkE9VZBX9hyA3ec1hr0epBbHFjSmnY6Qny/yHWs7csngDem8D/TCefbcTXBExy0vZmTZmVia849KsO2t7i0P+lXHFoyhDH5VhxqDtiCgzDXxwalQJf0L/eQXcHqT2E8d1By4kvD5/goFtD1LarAb89sILKWcBS6tsIO4vvZhXHlikP/3GCfrFP0AH4OoY9nqqH8A99PsFBe7GTNji6j/XCeyAwijcaRIwi0+Rx/DYOwEIH/Q86abubVXdn5rVi5zhon+ikDWY7fDhZoSYD+Chp09NOcPC43Uk7fZwN1rZ8jYOZt0lUynDpzWSaNJYTpDeJusfBw9skahUlXmU1pt4ewdF7ZaxDaEqbuOyzBssAu/NTLv4WfB1DZzrpz8qh3DFOpmvwNYq8xkm/hu1ENlRoJrB3furQzM87+KQ0ihyfS0lvq9joC4O1ii0z7PFOLODSUmB9gP65hNyP4WsVe6mTfkwGdDTzK53Ml+OYnIUNhVrn5NEdq4GzadHkLKxZtLd5c3d4m0XvhX+4ZFTTjlgBhuNvherti38a5X34DwAnZVP8rTJOBR4qKWMXcKqT7zVOHh3kHiuLOUg8cBk4BY+/dfJohMeB88kwRQvz6l2AP4jVCN45CgfjNzBTp7j3KMhI/ONTPPNuanxudvLoCVuBe4HLMGM2ZmTMYCyOcBk2oTTnoEp3UAaYn6Bz9MgY1zsT+I4krzfu4yGE6MRNmjs0aqfePDSqVqwyVNJwWVXuGPkSOGMxT9I5sUOjJAmbLBIdIygwO4TwWedv4oDNr9nmXJGvAJ4avloQx/vKaWfcjH9s3CjgVSefTiC2tDwN+PIEalhCgq8da4XWzLZvzUYXCRNTMQ9sirFZbj5ApHD7YOFZL76cyO8U+u7oWHeL20JnT7i3htW0yh+CP9kBLNx6ciK/g4CHE3hWhSdwTALZRddTSXvqJc0pTgZwX4KQWwBPq7ju/AYBXyDvIMfc2IKNeEsqNsHyJFJiCgtS+JUC8B78BiHY4zymTLsR30OwrV27YT5Ov8cueh2FeQi92AZ4G3vmAXBF4slaWeZkFbwnk/YUyo2FOOL5DXQZS5pdBfC5MrxLAduuefP7alhFhmgVltHzE9KeRqnYVvCckEH+cfj9/DXcS4viHz0psB82LiYF60m0CerIsT9wEfAIzdk67sQM0QuB/TPJPAF4LVGel3PJURrAsaTHxLeQuDvoQZ4DsGzgfwB+Sdokr23Fb7+D1RlkDa5g1n6qm7kTeH8OObI9PoDLJH0z8eddkq6XdI2jMjYamBNqrKySebSsscIwvdHepkPSBlkN/QpZcexLTZTlOkkzlH7+Lw0h+KeANBv4x8ztiiU43cZ9CZh711sttCvSx8A1G5hR+IOSCr5CSiVLmwP4BH7f/q6YS9VGX2/Acu5Tsn13Ral9dbuAfH6LBWQabtl0AEPJ47qtedai08vaBVga1xfJ47n8T2BI1Tq5AAwhz5MALNP4KuAdVevVG7CEkgtI39vvigWApxtY+wALa96a6USA2QfXUkEWcG/AnnqXkO7Rq4d5xHcjb09ghuHXM54UsMfqnVh2bjOyd2J1G4BVBN2CP2m2N9xIQmVy2wL4HOnOop6wHPg77EI0/W7BopITga+R7zHfHZ3Axc3Wo4aWbimwEuh/lK/BoQebZX0DF8s6kj3jmGhWF1g3riNlfZCmSDpONjquGVgu6YwQgr+lSyJavqfE/Nffl9Sq2r5XZZ69pTIv32aZx29D8W9JGiLzDO5d/Hu0rAffEfK1YiuDeyR9OoTweov4SapgAUhmF0i6XNJXlN6pc3dBp6SZkm5wtJ7Lhkq9SsBhkmZLmlqlHBViiaQLQwi/rkqASq3MEMLzIYQ/kXSyfCPQ+jpWy9rST67y4kttMjImhDBf0lGyoU6ehtJ9DZ2S/l7SESGEuVU88tseWAbwjbR38qcX/4f5Cpq1+9n9gNUizgI2VXrpymET5gQbUfX5bIT2Di3qt9vG6bKBDsnZxC3Gf8vqG28LIayrWpie0PYLoDuAcZLOkhlQsQMbW4V1skLOeSGEh6oWJhZ9agHUgKVVTZF0kmwYwzFqvUG7UzZ2ZpGkhZIWOQZItA365ALYFZi79njZopggm76dO2K4QdZz/1HZRV8SQtiQmUfLsVssgHoARspcuYfL3LoHypJAh6jnpNAOmYt4peyCL5X0bAhhbSvl70c/+tGPfvSjH/3oRz+ahv8HE370DS2Q4iMAAAAASUVORK5CYII=",
    "thumbs-up": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAJDElEQVR4nO2da6wdVRXHf4vb0oeKoKgRo0ZQECr2qaCECrGtUECFGjGGDwb6QY3xC0RJTGioia+EaIwxAYkmJiQYacFiH2JJpPWBgPZWi49ATAW0PIQqXLTt7b1/P8zUXA+dc2fN2XNm5nb9vp6919r77P9+zH5CEARBEARBcMxhTSdg2EgyYA1wDbAImAPsB/YBjwCjwA5g1MzUVDqDGpA0R9JGleMJSV+XtKDpdAeJkHRLycKfyqSkzZKWNp3+YAAkLcwLsyqTkm6VdFLTeQkqIGnDAIU/lcclnd90fgIHGrz293JI0tqm8xWUROlqfy/XNZ23YBqUvvb38smm8xj0QeU/+6oyLun9TeezKjN6IkjSIuC31J/Pp4B3mtnTNftJznFNJ6BmbmA4In8d8M0h+AnKImmR/H3/swN2Bxc2ne8gR/6+/2FJsyVdLmlPRQH8sul8B1Su/VdOiT9b0pcr2JCk5U3mPQAk3ekstD2SXjIekrS2gghuayLPQY6kxRUK7co+9tY5bb0oaf4w8xxMQYlq/xR7I5J+7bR5yTDzHOQoce2fYne50+bXhpHfoAclrv09tkcddu+tO6+pmDETQZIWAx9yRltvZpMlw97psHu6Mx3BoEi6q67an9u/yGH7sMd2MCCq1vd/1OnjTKf9E+rKb0pmikrX4Zvz3wPc4fTxH2f42c7wjdB5AUhaAnzQGe2Ljr7/CK90hj/oDN8InRcAw6n94BvYHTCzsQo+hk6nBZDX/suc0Twj/6mc6wj7WAX7gRdJP3IOzH6viqNzSQ86/Hg+GRulsy3AMGu/pFeQHSMry4NeH4ETSZuctd/13d/jyzMHIMW5gXqRtET+7/6PDODvSw4/ByTNTZnfoAcNsfbn/n7u8PWzhFkNetHwa/88ZbW6LOtT5jfoQcOv/Rc6/a1Mmd9gCqpW+9cM6HOdw9e4si+GoA4k3e0s/IFqf+5zu8Pf/anyGvSgZmr/bEljDn+xE6gu1Eztf4/T56Wp8htMQQ3U/tzv9Q5/E5JelSK/QQ9qoPbnfjc7fO5KkdegB0nL1EztH5H0T4fPOByamrwQdjgLf7eyuwAH9b3Y6bfyZFOTtHY1UNI84PuAd2HlxkQXPHrO+AnYmcDn0JmVypCyBZDTgNcMYHeEbOfNMmAFcIoz/m5827f74RHA08DZks5O5PtoTADPAHtT7jYaqKlUtib/MWAVsICEgqrIGjPbOKiRvAt5ikzMbWMS+COwHbjdzIY/+SRphXwrZMNgVAn6/jx/ZzWdGQcPSfJujPkfrjGApJMl/RD4KXBeVac18YWElzt36Yz/UmCTpC2S3uCN7DkZswTYBbRxtLvJzDYntNclARzhYmCXnLuRSjWZkt4LbAPauNL1d2Bxyhu6JD1KNqDtIgeAK8xsa5nA07YAks4ENtPOwh8DPlzD9WwnJrY3TOYCd0h6V5nAfQWg7KaLjbTzD/kXcImZ1bEDd28NNofJfGCDSqxNTNcC3Ai8PUmS0vJn4Dwz21GT/btqsjtM3gjcNF2gwjGApLcAfwKOT5ioQRknu5DxBjP7d11OlJ3s/R3w5rp8DIlJYJmZFS5U9WsBrqVdhX8QOMPMrquz8AHM7HlgNfDXOv0MgeOAz/cLcNQWIO/7n6TawG+cbHBWhtnAy0uG3W9mQ11vz1uCz5LdPHIqzd6tPJ/sgSsv48ApZvaP0jEkXVZhRupXklZJKn0uXtJKh/3nKmR+xiBplrLLqu6tUDZXF9kt6gIucKbvB8D5ZnaPmY074wYlMLPD+aB3JXCzM/oFRT8UCcBzEPJx4GozO+xJUVCN/HDrZ4CHHdEKy7NIAJ7R73fqHpQF/09e2b7tiFJYnkUC8FyHEnvhmuE3jrAnqGCPZJEARhzGO3EXTnB0WrslLJgWzzjt+aKLMUIA3eViR9i9RT+EADqIpDmA56Wy0aIfQgDdZDnlZ1AB7iv6IQTQTTzvEQj4SdGPIYBustoRdtTM/lb0YwigY0g6DXibI0rfvZIhgO7h3QK+pd+PIYDu4Wn+nwMe6BcgBNAhJL0M35b1rWY20S9ACKBbrMS3KWTasxIhgG7haf4ngHumCxQC6BYfcIS938yenS5QCKAjSFoIvMkRpe/o/wghgO7gfY201FnJEEB38PT/+8jONUxLCKADSDoJOMcR5cdlj8qHALrBRfhuXynV/0MIoCt4+v9DQOm3i0MALSffzOm5gv4+M3uhbOAQQPs5B3itI3zp5h9CAF3AM/qHkp9/RwgBtB9P//8XM3vEYzwE0GIkvR7f9u+7vT5CAO1mNb4j6a7+H0IAbcfT/78IuK/MCQG0lPyeBc/e/+1mdsDrJwTQXpbjO6Trbv4hBNBmvJ9/26o4CQG0F8/n324ze6yKkxBAC1F2Rd8ZjiiVmn8IAbQV7/NzIYAZhqf/3w9UfjQiBNAylL2V9D5HlG2DXNAVAmgfK4B5jvCVm38IAbQRT/M/SYm9//0IAbQPz9UvDwz6VkIIoEVIege+OxoHav4hBNA2atn7348QQLvwNP/7SHBJZwigJUg6Ed9TfFtTPJMXAmgPq6hp738/igRwyGFjqI84zGA8n3/jZI93DkyRAJ5x2FiVIiHHMsoe3vbM/+/Mn7UZmCIBeHaWXiVpQYrEHMN8Dni1I3yS5h+KBeB5i+94YHP+DRs4kfQpYJ0zWjIBFD0atRR4yGnrEHAbWd9UtgtZCnylZNgx4HJnmtrKCPBW4CrgXGfcR83Mc09gX4oEYGRvBp6eylGQjPVm5m0xCjlqF5B/X34rlZMgGQeBW1Ia7DcPcCvZg1BBe7i5372/Veh76kTSFcCGlA6DyjwJnGVm+1Ma7TsTaGYbge+ldBhUYhL4ROrCh3JTwZ8GdqZ2HLi41swK7/wfhGkFkB83upQQQRMIuN7MvlGXg1KLQfm04yrgu3UlJHgJY8DHzeyrdTopvRpoZgfM7BpgDfF1UDdbgYVmdnvdjio9h54vXqwle8PWc4IlKGaCbIr3JjMrfOQpNZUEMBVJi8i6h3eTzRyeDMytYGqE7Cr0WWQt0wTZsudBsr5wpvEC8ATwB+AXwJZBN3gGQRAEQRAEQSn+CwNn9RhTuOC7AAAAAElFTkSuQmCC",
    "timer": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAALJ0lEQVR4nO2da6xdRRXH/wOVAi2F8o7yaAQKWpFXbIhUkKKSIDWgtlqFDyIa/UAK0SiiJqjRkAiGICYqH4wiDyuQFrS2ChQ0pFBeDY+WFgK0pEALpqUPSm17f36YfeXa9LZnzczeM/t0/5Lm3MA+e73mzJ69ZmaN1NHR0dHR0dHRsdvhciuQCuAESZdLOlvSITWJeUPSfEnXO+eeq0lGhxXgEmAzzbEZuCS33R2SgHOAbQ0Gf5BtwDm57Y+l9Y8A4HFJp2YS/4Rz7rRMspPQ6gYAHC3p5cxqjHPOLc+sQzB75FYgkqNzK6AydAim7Q1gbW4FVIYOwbT9EbCnpNclHZxJhTclHe6c25ZJfjSt7gEqx1+bUYVr2xz8vgDYE5id4TVwdtUDdeQGGAFcCaxqIPCrKlkjctudglaPAbanCsoE1ZsKftY5t7Wm+3d0dHR0dHR0dHR0dHR0dHR0dNRDX2UChwM4UtJ4SUdKGiVptKSx1ackbZC0pvrcKOkVScucc680r22z9F0DACbIrww+Q9Lx8oEfFXi7jZKWSVoq6SFJ9zvnFqfQsxRa3wCAMZI+K+lc+cAfVrPI1+WXhs+TdJdzbn3N8jq2B9gDmAT8BljfwAzgcGwCZgJT6JPZwaIBxuCnYl/NGPThWAl8B9gvt5/6DuAg4Grg31lD3BtvAdcAB+b2W+sBRgI/IG83H8o64HvAXrn92EqAjwPP5o1hEpYB5+b2Z2sADgVuzxy0OrgNqGulUjBFvQYCZ0q6TdJ7c+tSE6skXeScuze3IoMUsSwccMAMSfeqf4Mv+RzFXPyAtgjfZ+8B8ImcP0v6VAPiVkpaIZ/hW1N9Sj5TOLb6PErS+xrQZa6kabkTSVkbAHC4pDmSTqnh9i/KZ+welLRYPrffk7Px7/Lj5VcYnymfYXx/DTo+Iek859yqGu5dNsAxwAuJB1qPApfhdw2n1nccMAN4LLHOzwN1NK5yASYAryVy4DrgWuCDDet/HenyE6/hJ7H6H+AIYEUCp2XPuOFT098lTYZyJTAuly2NABwCPBfpqK3A9fjBYxEABwC/rHSLYTFwUG57agEYDSyMdNAC4OTctgwHcCrwSKSNDwOh6xfKBbglwikD+O6++N24+N3KVxPXG9ye246kAN+McMYqoIkcQVKAs4mbtr40tw1JAE7CL5wI4SmgtZlB/ID3mUDb3wZOzG1DFMAoYGmgA/4FjM1tQyzAgcBDgT5YAuyb24Zg8M/tEP4O7JNb/1QA+1Y2hfDT3PoHAYwH3gkweCEwetcS2kXVCEJ6gs34OsjtArgvwNhlwKG5da8L4GB8t27lH7l1NwFMDzByLbtBThw4Fp/FtDI1t+49gV+yHZLtm55b96YApgb4ZzGFrCHYKcC0AONuzK23BeAC4B58/n5l9fcFxnv8OsBPn6vLpmQAjxuNWk5LUp/A3viNIMMxExjZ4732BV4y+upJIPsinmEBzjcaBDAlt969UAV/bg/2/MpwzwsD/HVenXZGAfzVaMys3Dr3Ar0HH/xBEuMN977H6LO767Q1GPxy7i0GQ7bRgkUQ2II/yLcM9/8AthNPtgDJNsCmHFV+WZJlg+SdzrlnE8pPDrC3pFnyO48tjOv1QufcEkmzDfceIWmaUZ/6AZ4w/kqKPmqFsF/+ID8xyjoFP+XdKwvrsjsIfNrXwv25dd4ZxAUf4BMBMh8wykiSNEv1CLCenvX7RHKTQ3i3P8hjku4L+N4fjNeXc2IZO3833p4NFLqHnvhf/mrguEDZY4CNBlm3prY/CPy2rtUGxW/JrfOOwG9Ft76SDWU1kQs4gD8Z5L1OgqRQikfAh2Srzz8vgcyk4LN3d0g6P/AWb0g6xzn3dKQqcw3XHiYpei9EigZgHc0/kEBmMgoKvmQfO0QfmJmiARxvuPZ559yKBDKTUFjwVfnmRcNXLL7fISkaQM9pT0kPJ5CXhNKCP4QFhmujVwqlaAAWJYo4cr3g4Eu+KGWvWH58OySqAeAXKBxj+MqyGHkpKDz4kq0BHBf7JhDbA+wnqaf57wqLccmpgj9bccGfXGPwJZuP9ta79Y6DSNEALOQuhPALhWf4BoP/TEJ9dsRq4/VRSbWmG8CGSHnB4OfovxH49aaCL0nWkjGtaQDbJG2KlBfDFIXZ22TwJV+3yHIecdQ2+dgGYFnLt9E5R6S8GMYFfKfp4Kvy0duGr2QdA2w2XGsZLNbBWuP1jQd/CBZfvRMjKLYBWJ5XI8lbN/dBw7XZgl+9qVj8FFVmrskGIEV2V5HcJz9Xvyty/vIl+6BuXYywphtAtu3e1bN1uqRXd3LZq8obfMnuo6g3qxQNYMBw/bGR8qJwzr0gaaJ8PeItQ/7Xluq/TcwcfMnmowFF9gDRYCv5dllWZYcA7A+cUf3bP7c+gwCXG/z5Uqy8FOfcLJU/jq0XoqcvU+Gce0v+JLDSsPgoOrWeYjbQMsMXvYBhN8BSN7mIBmCZ4fsIhS4ILYHKN5YVVkU0AMugaYSkSQlk9itnyfZYjh6wpmgAD8uWjZqcQGa/Ylnrv0lSGTuEgPmGkety2lDpomHwlVUsb1RJ6galCsR8w7VHyXd1Hf/PZPX+NiXZfD4sqRqAda/fxYnk9hNWn5SzvxIYgd+p0isbgINz610K+NoKlm1hK0lUODtJD+Cc2yqfSu2VUZJmpJDdJ1whyVIS9hbnnGXRSP3ga+VbWAsckFvv3OBT0muNvvtwbr13CPC00ZAf59Y5N8DPjD5blFvnYQGuMBqzCcg6Q5gTwuopl/voxJeHf8NoUHG7hZsCmGP01ZuUXkgb+KHRKIAv5da7aYCLA/x0VW69dwl+ULPGaNh6oJip4rohrGB0LYPm5CnZap6952qZFaMl3UqPZVbbDL4G0R2yr+e/3jlnXdmcB/wRca8YWzjA7yi5Fm4k+HI6Nwf4pTX1lP8HYSXRAa7JrXtdAD8P9MmFuXUPApgXaHDPpVbbAvDtQF/Mya17MMBx+OPPrAwA38+tfyoIezMCPz/Q7jwJ8LVA4wFupMVrB/DP/Osi7P9KbhuSAPwxwgl3UdCy7V7BHyg9K8LuYqupmsG/FYScljXIcuD03Hb0CnAa8EKEvUvpt8WzwInYZ72Gshm4EnhPbluGA9gLuKrSNZQ1QHQByCIBziL8HOFBlgKfzG3L9lS2hZ4TPMgm4GO5bakV/Dk5MUerg39LmAmcVIA9JwN3RtoD3iemk8daC/B1bAckDMcAcDfw0Qw2nIEvLp3Kjv44Lr5XgIuA/yRw3iBLgKup8fRR/HHwM4BFCfXeAny1Lp13Rda8O/BpSTNlWw+3KwYkLZJfNTtf0j+dc0F76PFz72fKL9meLOlkpfXZRklTnXN/S3hPE9knXvDd9z2SDqxJxICkFfJ7GJdWf6+T9JbeLa4wWtL+8jN0R8uXYB0vv4ehrkTUm5LOd849UtP9eyJ7A5Ak4Cj5VcWNP8sz8aikLzjnovf3x1JEmrUqk36WpB/JVnGkbSDpBkmTSgi+VEgPMBTgM5JuknRobl0Ss0rSpc65v+RWZChF9ABDcc7dLV8l4wbZKmaWyoCkmyVNKC34xYM/UHFBwleupnkcmJjbj60Gv236i8BTeWNpYhEwjRZPZRcHfm59CmX3CE/il8IVN7bqK4BJwG+xLz+vgzV4XbrSN02DP+lzKn7hxfoGg76+kvl5/DLv1tI3XRV+ncBE+ZTt2ZJOl7RPottvkq+FdH/1b2G1Jb719E0D2B58AYXBtO4J1ecR8mnfMZIO0LvFqzfIl5NfV/29Uj5t/Jx8Cvll51w/J6g6Ojo6Ojo6Ojp2J/4LMIdvxA3MayYAAAAASUVORK5CYII=",
    "trending-down": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAEt0lEQVR4nO3dTYtcRRjF8VPJaKIyQQJxkYWzEQ1MstCtoCRBl+LCL2DAvaDgfAjXviAu3TuYhROIOyEQspJAREYUgjoSNRtJBjP8XXQ3zHT65b5U3ap7+/xgNs2k6qnU6erqut13JDMzMzMzMzMzMzMzMzMzMzMzMzMzMzMzMzMzMzMzM+uBUOWXgOclvSDpL0m3QwiPklZlZQAuAbc4ag/YAtZy12cJAe8CB8x3FTiZu05LADgHPFww+RM7DsEAAZ9XmHyHYKiAOzUC4BAMDfB7zQA4BD11bM7jvzZo601J2w5Bv8wLwNcN2xtMCIBTwKvjn1O56+kUsA7sNngZmPi2ryEAzgJfAfuHxrM/fuxs7vo6A2wyOvRZmRCMx3x3wZjuApu56+wMo/OA31qE4BrwVO5xVDEea5XN75/Ahdz1diZCCHZKD0GNyXcIhhaCBpO/siE4z8D2BBHGtAeczz2OzjCglYDmz/xpK7cS9D4EESd/wiGoKVsIiD/5EysXgt7tCSLUvMxK7gnaPJs6Wwki1FrVyq0ExYcgQo11OQQ1JQtBhNqaWrkQFLcniFDTzvinqc72BIyC/hnwI3APuAF8BDzTRf+HiyhiJYhVC3AC+KZFO8lXAuAKR69eHvYz8FLK/qeLyR6C2DVQcAiAyyz+9DaMLu2vp+h/XlHZQpCqbwoNAY9/b2Oerdh9Lyus8xCk7pPCQgBs1Oj7Rqx+6xTYWQi66ouCQgC8XqPfvRh9Niky+cR00cdUf0WEAHijRp9/t+2vTaHJJihl20vGlD0E9CUA42KjT1SKNmuOKWsI6FMAxgVHm7CYbbUcU7YQ0LcASNEm7uUIbUQ7eiZTCOhjACQJuDAedFPLDj4WSfIVNuAkHR8bEyEA874ZlFQI4QdJr0n6o2ETTeu+JuntEMLDhv9+rnGbb0m62rCJ5yR9R8cXkLIEQJJCCHckXVTzENQ1mfwHqToIIexLekfNQ3BG0vUuQ5AtAFKnIUg++RN9DEF2tN8TLJLla+t0sCegr5vAWUjzwY2sn0Am8bsDhhQAKXoIsn/8fDymZCFgaAGQor0cFHW3EhK9HDDEAEitV4IinvnTSLASMNQASI1DUOTkTxA5BAw5AFLtEBQ9+RNEDAFDD4BUeU9Q1Gv+MkTaE7AKAZAk4EXg5oxBHQCfACdy11hXpBC81zYAle4WXgLgmKRLGp0cnpa0K2k7hPBT1sJaYLRqbWt0d7Um7kt6tuLv/hNCON2wH0uF9nuCVitA1msBFuXaQSsOQAFyhsABKESuEDgABckRAgegMF2HwAEoUJch8B9+igQ4J+l9jc4pzkRs+j9JT0Rs7wgHIALgiqRPJT2Zu5a6HICWgMuSvlD5L6czT317cxRcKuCWpFdy11HBgxDC09MPOgAtABuSfsldR0WPQgiP7SVKX7ZKt5G7gBpmzrUD0M793AXUMPPvPTsA7dyWdC93ERXtznrQAWghhHAg6ePcdVT0Ye4CBgk4Dmx3cD2/jS9z/z8NGrAGbJH2TuRN/At8sKh2vw2MCFiTtKm4R8F1HNfoncm6pO9DCN3fGs7MzMzMzMzMzMzMzMzMzMzMzMzMzMzMzMzMzMzMzMxS+B9VDtQWPkbC9wAAAABJRU5ErkJggg==",
    "trending-up": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAE3ElEQVR4nO3dz2sdVRjG8ec00VY0RYqCdGEXFoymXbl0UUxNF7rwb1DoXhC0f4m4celOwZQuTEBEoVAQ3aioSIVuWpP6ayM0/vq6mLkSx3tz7505M+fMzPOBbobmnPfmeefOnJm5N5KZmZmZmZmZmZmZmZmZmZmZmZmZmZmZmZmZmZmZmZn1QEhdwJAAq5I2JD2aqIQVSWckrUm6HkK4kaiOcQFWgSvAHnn5DXgt9e9n0IAVYDtx0PO8nfr3NFjAG6nTXdAL0+r3OUADwIqkHyQ9krqWBXwdQni6uvFYikoGZEP9CF+Snpi20Q3QzMOpC1jC6rSNboBmbqUuYAl/T9voBmgghHBL0uep61jQH9M2ugGae10z9q7MHEzb6AZoKITwoaTLkn5PXcscTNvoZWAkwLqkVyU9p7iXgh+SdF+EcX4JIZyKMI51ATgOXIt4IejnafP4EJAh4Lik9yS92PZcboDMdBm+5AbIStfhS26AbKQIX3IDZCFV+JIbILmU4UszbhDkCDgmaVPFOvuUpJuStkMI3yUtrAHghKRtSZdqDvGr+nVDqh7gLPDplLXtX8CbwP2pa1wWcALYabCu3wMuL/H/p14HyB5wnvnP2n1AsTf1AnHCPwdsLfEz/WsA4CngzoIvcBd4IHXN89D8Ct8+cL4ca7gNwHLhT2TdBEQMvxxvmA1AvfAnsmwCIodfjjm8BqBZ+BNZNQEthF+OO6wGIE74E1k0AS2FX449nAYgbvgTSZuAFsMvxx9GA1As9fYb/KKOskOCJSKRlnpz5mjcAMkvBVM8SbOr9j5QeUnSVTp8J6C4vPuu6l/huyvp+RDCl/Gqmi5pA5ThfyTpsZan2pK03UUT0Pza/l1JF0MIX8SrarZkDdBh+BOtN0HfwpcSNQDFic0nqh9+3cewtyS938Y5QTnmVdUPf1/SZpfhSwkaIMIxf1fSM5Lu1Pz56OcEfTrmJwWs02yp9+8ZfTnW7QZjRVki0vJSb87c/VkGEjH8ypjJmoCE4Zfz96MBaCH8ytidNwGJwy9ryL8BaDH8yhydNQEZhF/WkXcDECf8hYKhoyYgk/DLWi4sMe9ejDmXKa6z8CtzttYEZBR+Wc+ZJebu7ivjSBB+Ze7oTUBm4R+q67MF578Se+5ZBSULv1JDtCYg0/DL2i5SPBx7lJvAWhvzV4tJHn6llsZNQMbhH3qtrwAHM+b/HniyzfknRZyj2TdlRn+yN0JNO7R8Szfia10H3gK+BX4EblB8h+GDXU2exZ4/o7Ym7wR1tb7nZ4GMw6/U2GUTOPwFtR5+pdYummA04Wd3zO+g5nk6O+YnRY/2/Bm1t/FOMJo9v7fhV15DzCZw+AtKHv4E8ZpgNOH37pg/T4TX5GP+grLZ86uo/04wmj1/sOFP1GgChz+U8CeWaILRhL/BwI7581CcExzVBLcZyTF/jeIW4mjCnwBOA+/w37trB+W206nr6wTF38CrK8mHMWMDTgLPlv9Opq6nUxS3Ekcb/uhR78TP4Q8F8I3DH4dZnw38eIkxdiW9FEK4F6EeywHFmvie9/wRA17m6KdNrzn8gQM2+f9z53sUy8TefNG0zbbQXw0DHpd0VtJPkr4KIfzZalVmZmZmZmZmZmZmZmZmZmZmZmZmZmZmZmZmZmZmZmZmdtg/CTHncu0305IAAAAASUVORK5CYII=",
    "trophy": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAALaElEQVR4nO2deaxdVRWHf6vQAi3QUlqgUg0yySgzhDCVoSXEWECKglELMgqaFKmYqJFBMEoskAACWoJoUYQyqGCgTBVRAUtlTo0CLVBlkMHXQuf3+ce+z5Zy73t3nWGfc9/bX/L+ePfss9dv773uvufsYW0pkUgkEolEIjHgsKoFSBJwsKRjJe0naStJ61cqqDyWSnpJ0qOSbjezP1Wsp1oHAA6TNE3S7lXqqJA5ks41s4erEjCoCqPAIOBSSfdr4Da+JO0taTZwCVDJlzG60UZBp0v6cmzbNeenks4wM2IaraIH+IZS4zfjNElTYhuN2gMAn5D0lKT1YtrtIJZK2sXMXohlMHYP8G2lxu+N9SV9J6bBaD0AMFzS60oO0BdLJG1uZotiGIvZA4xTavx22EDSIbGMxXSA3SLa6nSivRqvW3SGjde8LRt/G61xaa+ibfVj9gGOWOP/RZJeNbOFRRsq5BkA2FrSJEkTJO0jaeOcWd4t6Yq8umrGFEmfyplHl6THJN0n6VYzm59XVC6Aw4FZQDfFcl2lBSsB4LqC66gbuAcYl0dXpmcAYFtglsJQ7ngV/zZRi0mqgimjjo6U9BDwe0Iv7MbtAMApCoM547MYbJOxJeZdFWWW6ShJTwOTvTe27QCAAZcrjOMP9RpyMqbk/Kug7DINk/Qz4Ec4JpbacoBGhtcp3lj1kEh2YhKrTOdKuqZdJ2i3B/iewmRFLBZHtBWLmGU6Q9J320nYpwMAn5b0rbyKnMyPbC8G8yPbOx84qq9EvToAMFJhnjr2U/m7ke3FIHaZTNL1hDmYlvTVA1woafPCJCViM0bS+b0laOkAwEclnZ5TQJekdxp/y3LmNZBYptX11pUzr7OAj7S62FsPcLayPbn+UdLnJG1qZsPNbKSZjZR0Y4a8Bio39tSbmQ2XNErSiZKyrCJeT9JZrS42dQBgkKQvOQ0tkTTZzA42s1vM7G3n/YkWmNlbZnazmR2osJxuqTOLyY02/RCteoD95Bu4WC5popn93Cks4cTMbpB0jKQVjtvGStqz2YVWDjDOJ0sXmtn9znsSGTGzeyVd4rztsGYftnIAz9z9a5Iuc4pJ5OdSSW860rt6gG0dGd9qZt7fpEROzGyJpJmOW7Zr9mErB9jMkfFfHGkTxeKp+6Zt2soBPIs333KkTRTLG460GzT7sJUDrMybcSIKwxxpm7ZpKwd4x5Hxro60iWL5pCNt0566lQMscGR8rCNtolg8dT+/2YetHOBpR8Z7tjPtmCgWYKJ8PUDTNm3lAI849VwLpFnDSABjJP3YeVvTNm3lAA/IN978MUn3N2YQEyUCbKWwGntLx23vS3qo2YWmDmBmiyX91qltF4WVqecAeTeGJNYCGA5MlfSkpJ2ct99pZu83u9Db1rArJX3WaWiEwrDw94HHJL2o1ZMWBznzGsgctMbmmMGStlGYoMu6ufbKVhdaOoCZPQI8JOnQDAbXV9jhGm2Xaz9jx8ZfEdxnZo+2utjXkrCpklYVJCQRnxUKy8Rb0qsDmNlchVmn2KR9AcXwAzN7prcE7ewLOF9S7Dh2H49sLwaxy/SgpIv6StSnA5jZCkmfkfR8AaLaZdOItmIRs0zPSppkZn3O6bS1M8jM3pJ0uKReu5MC6Y+7g2NFY3lS0hFm1tZ8TtuizOw1SQdK+l1GYR7yLoWuIzGCPt0h6WAze73dG1xeaWZdko5WWDJeZiO9V2LeVVFmmd5V2MNxXKzoYgK2AK4E3i848gVArJ+aaADPl1BP7wFXAJ4VXB8g928tMFrSFyQdL2lfSevkzVNhHmJDM+sXYxDAYIXdwUW8Cq5SWAo2U9KMxvNZZgp92CJsRNxX0g4KkxWbrHF5lMLbRLvsambPFiivMoA9JM113HKbPriA421JCyXNk/R446e4swA2BFY5urczq9ZcFMDXHOVeCURbZhctUGRjhvGfjluOLEtLBUxwpJ3XWPIdhdjBoltOSjRhPOBZ9FhLCFPjR/SZcDWeOspNbAf4gyPtMEnHlSUkIpPkOwOpsuNjSgcYiy+o5J+r1pwX4HFHebuBLarWXCrAHEeFAHTsQhJgnLOs0XdZVXFkzK+d6S8sRUUc+pyNW4tbSlFRJ4AtCa86Ho6pWrcX4DhnGVfQ37v/HoA7nZXzMh200BQYAbzqLONtVeuOBnCos3IAOibGEPDLDOXr2GedTACPZqikmNFKMwGclaFclR8hGx3gyAwVtYxw3GwtASYAyzOU6/CqtVcCcF+GyuoC9q9a+9oABwCLM5Tnnqq1VwawM9m+MV3U6FtD+OYvylCOZcAOVeuvFMLByVlYRg2eCYCvkM2JAS6oWn/lAOsBT2WsQIBf0EdA5JJ0jyDb034PcwgLRRLATmT7/ezhFSBaoArCIM/CHHoXEc5RTvQAnJijQnuYTYnv08AhwMM5NXYDk8rS2NEAF+es3B5eJJydM7oATZsB04CXCtLWa/j2AQ3hYKobCqpoCN+2p4Czgd2APk9KBQYDewBfBZ6h2DMRf4LjQKcY1EqMJAHrSJoh6YQSsl8u6QVJrygsuuxZQ7+RwqLVsQp78cvYyHmTQjT1Wq10rp0DSP93gumSTqpYSlFMl3Rm3Rq/1hB+Di4uuAuOTTdwATXr9jsK4ATyvSJWxSLg+Krrr18A7Aj8reIG9fAEsH3V9davAIYAF5F92DUGywhdfn+McFIPCKOGs6pt56bcy0Cf2IkJYQaujB23Xp4DPBs/EkUBjK+69YHxVddDHqpYFp6oEckBBjh9jo33I7oknaww3DtK0sZafdrJksb1/0h6VdINjev9noHkAKvM7PZ2EgLTyxZTFzr9J2CkI63n8GrPqZwd3VN0ugN4FoZ6Dln0xN2p7TL1duhYBwC2UwhO1S4vlJT2JGBrR/pa0ZEO0Kjwu+Q7su6vjrRzHGmHSrqLcJJHomyAkwj7Arzs7bCxf4b8u4DJZZZ9QENYqnVthoYB+AeOOXlgEDA/o62raWPpWcIBMAy4J2ODAHw9g81v5rB3NzC0jLoYcAAbAY/kaIx/kSHaWMPu6znszs5iN7EGhF1DD+RoBIDP57A/OaftWaS1AdkgrAm8KWcD3FyAjpk5NXRMYItaAZyXs+Ifp4AumPBTMDenlilF1MmAgbDXfkWOCp8DjCpQz2jyrUlcDuxTlJ5+DTCU8NqWlZnAhiXo2hi4I4eueUQMAt2xAD/MWMGLgVNL1mbA6YSDGrJwcZn6Oh5gO8LKWi8LCTH5Y+ncC/h3Bp1L6eB5g9IBbs5Qqa8A21SgdXuyxQiYEVtrR0CIF+Q5UALgHcB7knaRmncB/uvUvJIUHOLDANc7K3IVcFQNdE/Ev3/xmqp11wpgE/ynj11ete4egKuc2hfTQaFvSwc4w1mBC6jRZAvhPCRvbOBTqtZdG/CP93+xas1rA5zsLMO9VWuuBcBwfBs+/04IIFErgHUJsYnaZRklDFh5qcOSsEMkeeLlXVXHSBuNk7qvdtwyRKHslVIHBzjQkXalpF+VJaQAZiic7NkuB5QlpF0KXbpEWAq1syRPeDbP5srnJe0OuHRFZp5CHbTDBOBBR95vSnqu0dsUQiGxaxoNP1XSOZIyH2ScaIs3JF0maVoRjlDE4dHrSLpd0sS8eSVc/EbhuPhcz0NFPANMVWr8Kjha0rl5M8nVAzS+/a8p7LZNxOdNSWPy9AJ5e4CdlRq/SkZLyjUZltcBRuS8P5GfXG2Q1wEW5Lw/kZ9cbZDLAcxsgaS5efJI5OIJM3s5TwZFvAWcJ6m7gHwSProV6j4XuR3AzB6QdJpCKPZEHJZLOtXMPKOITSksijUhSuYUSYfKNxTcwyCFCZLBDV3dCgXtdMca0vgbJAmF8DPLla3XfEPSbElXmNm8ogQmEolEIpFIJAYY/wOJR2vCtxraNwAAAABJRU5ErkJggg==",
    "truck": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAIpUlEQVR4nO2da4xdVRXH/6v2JX1QUkoLNaZSSYX6oYKAMVExUhOCVDAiJlC1kQRJDI/SSI0aSzSmifWDD0xaEDVRiqiloJEIRitSaLWgIhER/ABFZ5raDqQytJTOzw/7NjTTc2fueex9zpm7fkm/3Dl7r7XO/Xffddbe+2zJcRzHcRzHcRzHcRzHcRzHcRzHcRzHmYBYVR0B0yUtljRP0uSq+k3EfknPmdm+uh1JTSkBAGdL+rikD0paqvZ98aN5XtJvJW2RdL+ZvVazP80EuBB4mInNbuA6YGrd9zsmuUYAYJ6kWyVdHsedRvKUpFVmtrNuR2LQswAIw/29kt4Uz53GcljSjWZ2a92OVM2kXi4C3i1pm/rzy5ekKZK+A3y5bkeqZtwRADhT0iOS5sR3pxV8diKNBGMKADhB0mOS3pbGnVZwWNJ7zWxH3Y5UwXg/AbfIv/zRTJF0x0R5OugqAOAtkq5L6EubOFPSZ+p2ogrGGgFukjQhVB6JNcCUup0oS2YO0PntH5Q0q0CfhyX9r4xTiZkqaUbBtivM7BdVOtMIgEsKVM52AMvb+L8CeDtwGzCSM+Yf1O17FIBv5LwRdwFtnwcQcFVOETxXt89RAH6T4yY8T/jJmBAAG3PEDtDq+ki3JHBRjj5uN7PhCnxpCt/Kef2iGE6kopsA8qj68SocaRB/l/RKjusn5AjQ0xxBh0NVONIUzAxJB3M0aV3Seyx5vmhnAuIC6HNcAH2OC6DPaX3xpgFsAQ4XaHdA0jOSHpJ0l5n9s1q3SgDsz1EIWV63v1WTM/4qGAG2AKenjtV/ArKpbL9EDnuXSfor8LGUhl0A2Uyrye5MSZuBVakMugCyqXMdxCRJG4ELUhlzjqfu+zJF0t3Am2MbqjvQpkLdDijssdwCvDGmERdANkUe62JwjqSNMQ24ALJp0gTXSuD6WJ27ALI5UrcDo9gQKyl0AbSDyYqUFHopuDxrFXZP5eFiSTfkbHM0KXyPmeVZsJKfnKXQfi8F544fMODnOWwcy4+qjNV/Amqgs+roUwrLz/JyJRUmhS6AmjCzA5I+IumlAs03AO+vwg8XQI2Y2dOSPiFpJGfTyZJ+Sti/WQoXQM2Y2X2Svlqg6VyFpLDUngwXQDNYp/Bmsrwsk7SpjGEXQAOoMyl0ATSEupJCF0CDqCMpdAE0jNRJoQugmaxToqTQBdBAUiaFLoCGkiopdAE0mBRJoQug4cROCl0A7eAWSb8s0G6ZpO+OdYELoAWY2YikqyQ9XaD5J4FLu/3RBdASzOwlhe1jBwo0Xw9kftcugBZhZk8pJIV59y0skfSBrD+4AFqGmW2V9JUCTVdkfegCaCdFksJzsz50AbSQgklh5pJyF0BL6SSFn8vRZHrWhy6AdlN6f4ALoM9xAfQ5LoA+xwXQ57gA+hwXQJ/TqO3hwHxJ8yWdJGmPpBfMrE0HUB0HMFPhyN35koYk7TGzPfV69Tq1CgA4Q9KHO//OV8a794EhSQ8oHFx9v5m9mNTJnBCOkLlI0qWSliuIefQ1r0raKek+SVvN7NmkTo5Hgv3xS4G7C+yNPwR8EzjuplZJkfiBk4D1wHCBuB4E3lHAz+U5bOzP6iNpDgDMBn4o6QlJlxfoYqrCaabPAtdU6lwJgGsl/UvSzZKKvNbtQkm7gO8DRc5qrBYijADA6cCTOfrthduIcIZvzvjvrTimfwBLevSzHSMAcK6kP0paWnHXV0t6gHqPrcucZy/BEknbgXdW3G8m0QUALFRI4OZGMvE+hdO8U7/hOyZzJW0FTottKKoACK853Srp1Jh2JF0h6QuRbaRmoaR7gMxp3KqIPQKskZRkKJO0Djgrka1UnCdpdUwD0QQAnKwggFS8QcU2UDSdtcApsTqPOQJ8XtLsiP1ncRlwfmKbsZmlfCt/clFFJfB64KOjPjNJV5bo82VJMwq2vR14pIRtSYrxVFEmpmuB2Tp+OfjCHH30vpQc+HdFz7R52AVcDEzr+HAKsJr0BzhVyT7gBjpDODAN+BDwWA2+7M4jgO2JnbuHLgUd4AxgILE/VfAf4K1dYppGiDklf8jypVsOsKNntZRnQNJKM3s1649m9oykTyf0pypWdZvkMbNDklZKGkzoz6NZH3YTQJHXkxTlx+NN+ZrZrxSmh9vCgJn9eqwLOjHfmcgfKdRjjiNTAGa2XdKuqO68zlCP1zV6GngUvb7VI7M+H4E/mVlmYjzWY+BNyv9miiJMpBJuXlLEPqIxikldBWBmD0laH8OjUfQ6tz8nqhfVcmKP10Vd19Dha2b2cKGWwCTClGtMBhhnDpzw6NQ2LhonplnAYGQfNlF2koxwusWNwCsRHd1K5/k/w/4S4t+oGAwQlrxlxTSd6tcRHMswVZ80BiwmrFiJJYTHgUvozH4BC4A1wIuR7KVgiFDMWnDMF78C+HMke8PA98jxytjcwwNhlesFCi8gOlXZ5eSjpeCiJdUyZdMnJZUtBR/LHEmLFGry8ySdXLCfMjENS+p2VtBhhXrCXyRta8wqamBDJJWPxQhwXsSY3lVDTABfjxVTNAirZFPX8X+WIK7UJdwhINZqqnjTwWY2JGlDrP4zeE3SlxLYWduxlYr1ZrYvob3qICQ9OxP9T1mbMK4vJorpUbo8HbUG4DTghcg36k4SLgolPBr/JHJMu4HYaynTAJwD7I10o35HWHyaOqYTgG2RYtoLnJ06pqgQNob8reIbtYkIG0NyxDQZ+HbFMT1BBecBNhJC+fMO4EjJm/Rf4Oq64zkKcE3HpzIcIRRx6t8aFhvgLIptDn2ZsAGzcRNDwEzgZuBAgbgeBJbVHUNyCOXl1cDvgYNdbs5eYDNwBWFhZKMBTuz4upnuo8LBTsyrgcV1+tuouXjC4skFCtOkg5J2m9lwvV6VA5ih8IKIBQoLQAbNbG+9XjmO4ziO4ziO4ziO4ziO4ziO4ziO4zgTnP8Du1H9OcTQJZkAAAAASUVORK5CYII=",
    "users": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAALhUlEQVR4nO2daaxdVRXH/5uhFKx0YFIKtA9KoQzKUAKIIm2hQQNCjCakEgj4SWPACRCHyGT8ZgIoIjQEjKKAtgURbQULbQOUQWYpWAshwGspQ+f59eeHfarPx73vnbXP3ufe+97+Jf3St8/5r73PumfYe621pUwmk8lkMplMJpPJZDKZTCYz+HGtNqAswEckHS5poqRxkkZJGlH8eZ2kDyS9IelVSUucc+tbYWen0dYOAJwo6RxJUyRNlrRLyUO3SXpS0nxJ9zrnnkhj4f8AJkm6VNJnJY2R9NoOfUmLnXOktmFQAIwErgCWEI+XgcuBPRPZ/FVgcz/6/yr0R6XQHxQAewLXAasiXvi+fABcA3w0ot2nAz0l9dcA15LIETsW4GzgjSSXvDFvAxcAlR+BwLMB+t3AjBhj19Hgf/V3RbywVmYDoyvYf3BF/T9U0e9ogE8Ay6pewQgsBY4O7MPUCPrLQvVjsFMrRIFTJS2Q1NUK/T4cImlBYZOV1RH0uyQtAqZEOFf7g39p2hjhlxObjcA0Y1+GAe9H0t8AnJ5q3JtR6x0AmCxplqThdeqWZLik2cBxZQ9wzm2RdH0k/d0lzQFOiHS+UtQ2EQTsL+kfkvarSzOQbknHOeeWl2kM7CrpfknTI+of75zrjnS+fqnFAYCdJT0kP0tWlU2SlsjPtK2WhPy0cJekSZJ2i6Dxd0nTnXM9ZRoDwyR9X34mMMZkzyOSpjrntkc4V+sBrqz4fFwN3AxMAZo+PoDh+DfzXxXHVOHygH7uBkwr9NdW1P9OtVFvE4DxwPrAQdgAXEXAFCowGj/rtyFQex1wUIV+j8LP+oXqrwfGheq3DcAfAwfgMeDgCPqHAIsDbbgrgv4E4PFA/Tur6rcU4Ghge0DHb8e/XMWyYxjw6wA7eoAjI+nfHqh/VIwxaAnAnQGdvoUIc/QNbHHAzAB7fhNR/9YA/Tti6NcOsDewxdjZB/BfDKls2hmYa7RpM7BXRP15Rv1NsfRrBfiGsaPdwN412LUvsMJo29cj6u8HrDTqfy2Wfm0AC42d/EqNtl1otO3hyPoXG/UfjKmfHGAEttv/S0Bt09LATsALBvs242MSY+nvjI9Saol+b1IN+qclWd7if1bnrFehdaPhkGGSTomo3xOgf1Is/d6kcoDSCyryU7v3JLKjP+6WtMXQ3tKnMtwlH7xaluMj60tK5wCHG9o+5pxbk8iOpjjnVklabDjksMj67xn1K89HNCKVA0wwtE0est0PjxnaHppA39L38Qn0kzmAJc7tlUQ2lOFVQ9sxCfSXGNrum0A/mQNYwq7fS2RDbO1ooeS9eMfQdsTATeykcgDLmrzlRSw2mwxtY8QZ9GWroW3ZrCgTqRzAkpeX5Pu2JJZf1boE+pa+b06gn8wBLIO1fyIbynCAoe3aFuvHiED+EKkcwBLPluTzpiRHGNqWihE0YskHeDuBfjIHsLzZfyaRDWWwxCha3tjLYun70gT6yRzAMlhHABMT2dEUfDq3RTfq5yo+0OMQwyHPx9TfQSoHeNTY/qIkVsTVtPZpIC42tm/lhJmNYrXtPcNq1ypqzJ0HxmCLGl5JxNVKfLDMmlbp9ybJSYvVtnmGQ0ZKujqFLU24VpIlR39e5NXKa2SbWHqg43IEgLMMHg6wjbAETatdUwotC5+LqD+V8gUldvD5WPq1AewKLDd29C3gwIQ2jcMXh7DQDUSZhSv0rWPyZiz9RiSLwnHObZX0c+Nh+0uah88jjAowVtJcSR83Hnqjc86ybt9M/wBJf5M9N/KXMfRbAj47JqTmz2vAJyPacQzweoAdHxChng9wbKD++3R6YSngBwEdB5+v/20q3P7wj6HvEl6P4IqKfd8VXx1sU6D+ZVX02wJ8wuQrgQMAPmB0BoZMoWLgZ2ALvOzLP/FZv6F9voBqpe5eDtVvO4DTsL/59mUlPuv2PGAivRwCf8EnFn+7BXvcfV/MXyT4SN+zIup/Kv6VaCH4LN/YrCv+xeZHxr5NAp6PqH9VosvQOvC/kL9GHKRU/BnDrBtwEPBORP0/WfQ7CmAP4NGIgxWbJwBT6BXh6e9R9DsOYB/gmYiDFotnMOYm4tcUrLOKzXiaFiSB1n6rcc6tlF8Ht6wVpGa+pNOcc+8ajztSUoxs5vmSphW5ArXSkmeNc26dpC9I+kUr9Ptwo6QznXMhIVcxSsDfVOivinCuzgP4MvBupNuohZXAFyvaPhrYGqj/TlX9QUMxkNcT73naH9vx5WL2iWT77436PfhKIclrIXQcwJHFxUnhCD34T6yoSZ7AWMqtMG7DO0vn1vypC3xlr6uBf0e48Evxk1CVK471Y++hwJNN9Jfhy8W1Zcm3dt8zyMmHTk+VdJr8W3eXmr95b5P0uqQXJT0sX/HzxTr26ylsnSbpVEl7SVom6SHn3LOptavQ1g7QCPwCyVj5BNQdYVVr5XcNe6so4JzJZDKZTCaTyWQymUwmk8lkelPnrmEj5Kd1D5N0kPxM3h516Q/AevmZxDfkS8e9UMQsDHqSOUAxN36ypHMlnSF/8ZPtBRCZHknPSXpQ0hxJj6deTwAOk/RN+WipMfJrGvMl3eecs1QUbS3AXsD3aI99gWOxFLgCSFEsUsBF9J899CpwGe2cJobPA/wp1bdLa2fWAj8hQr5gr3E7nfJJM6vxO6GlKFoZDj4NKmZsfLuzHDg/0tiFREh3A+fF0K9q/N7AfTFHtsOYQ4VQbqCrov49gKUu84cIjgoGjpXfC/jsKgZ0OOdIeho4JvD4ror6X5L0FGCpN/h/BDkAMF3SAknJqnl0EOMkLQTOCDg2RvXPgyUtAqaEHGz+DATOlP80SlE8uZPZJOlc59zcsgfgo5tWKM6G0xslne2ce8hykOkOAJwsaZbyxW/EcEmzgBPLHlCEr10fSX93SXOAyZaDSt8B8MWbnlKijQsGESskTXbOvVmmcXEXuF9+siwG3YV+qdrCpRwAX6ZlgfzMXlU2ykftvi4//doOjJZ/ITtK/pdclUXyuYY9ZRoXTvBDSZfI10ysykJJU8rqDwjhdX52sBq4CV+jr23LnuA3eZ4K3IytkmcjrgzQH46fGLqF6hNqleob9TZqAuFFjtYBPybizFldACPxySnrA/u+kQrJKPh0uesIL3C1Aaj6mSkBswINWBjFgBaDz1IKLWpReT9EfO2jxYH6d1cVPx6fTGllJoaqXu0OvgjVbQHjsJ0IuYj4R9Mdgfrh9RaxZ76Cf9Z3XMbRQAAO/25g5XcR9WcG6P82VPBj2DaABrifwVrgSP8tdPUX45hsAazlYfvTn2vU30xIKjzwLaNQNy2ocVM3+AUwa8HnSyPq7wusMOpfEiK0yCgyI1Yn2x3gfOPYLIisf6FR/xGrwGhsRRqeYxA+95uB3xHlBcP4bAViTPD01n/JqN/wU7zZ8/oU2eL3bqgjB79dKHbvuMFwyC6SopV+LfQtawi7SDqp0R+aOYBlr/rNkip/73Ygd8u27a1pkSaB/gmN/rOZA1g2VFzknFtjaD8oKMrKWXYSmxRZf5Wkx6vqN3OA8YYTd07IcnwsfR/fjvrNHMCy5Bt1Q8UOw9L3FMvoFv2GcxHNHMCyq/VKQ9vBhqXvKYpAW/QbhpI3cwDLNi2duaFRHCxbuqdYBt9YVX/QTttmypEdYIiTHWCIkx1giJMdYIiTHWCIkx1giJMdYIiTHWCIkx1giNPMATYZzmFpO9ho9ThV1m/mAM+VPOk2+Ty/ocqLKr8WUnZMa9Vv5gC3lTzpbOdcuyR41k7R99klm5cd09br45MQBkoJW45PGR/SAAcycJj4LBIFzSbTB3bDZ6k2Sg17Bl/YMCNf5JHG1b62F2OYtKBGFf0BvQI4XL4YUpekVfK7cc0rIlMzBfiMqOnyu5uNkvSapHudc0uGgn4mk8lkMplMJpPJZDKZTCaTaWf+A0N40l1PtHeZAAAAAElFTkSuQmCC",
    "wallet": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAGuElEQVR4nO3dW6gdVx3H8e+/8WiMrZegooa2AUFFoVJIaFJTbcSiRRELFZHYkocaBbEIohDUKG0ffPACLXmoqBTR2D4pVfBWGkjVSBup4gWtTUlN1Wq9VmyO1vTnw8zRU9k7Z/57Zs3ss+f3gfOy95o1i8wvs2fWzFoLzMzMzMzMzMzMbPFFqYolbQYuAl4BbAE2ldrXOncaeBj4EXA4Iv7e5847DYCkpwN7gL3ATuCsLusfgVPAbcANEXG8jx12EgBJG4B9wAHgBV3UOXLLwIeAT0eESu6odQAknQ8cAi5u3xz7P58F3hURT5TaQatTtKQdwDF88Eu5Briu5A5mPgPUB//bwDndNccmeALYHRFHSlQ+UwDq0/49wPO6bY5NcSwitpeoOP0TUF/wfQkf/D5tk/SqEhXPcg2wDyjSGDujK0pUmgpAfZ9/oERDbE0Xlag0ewbYg+/zh7K1RKVPSZbfO8M+/gXcBTwAFO3UWIe2AG9sWPYZJRuyJkmbJZ1WztckbRm04XNM0mWJf8s/l2hD5gywg9xPxteBt0TE6VyTrE+ZA3ppouzjwLt98OdfJgA7E2UfjIjfZBtj/csEYHOi7OPZhtgwMgHIXMFvzDbEhpEJwHKi7AslPS3bGOtfJgB/S5TdCLw/2RYbQCYAf0jWfZ2kdya3sZ5l+gEeTda9AfiMpKuBW4FfUT3btv+5MFF2SdIu4EREPNRVA7JdwbPYVf9ZO2dTdakj6STwFeBgRNzXplK/tbs+nQtcC/xc0o2SZn7l3gFY3zYA7wWOzvrMxQFYDBcAd0h6TnZDB2BxvAz4XHYjB2CxXCGp6fsFgAOwiPZnCjsAi+fi+rX9RhyAxRMk3t0oFQC/+zeslzYtWKon8IvAnVTjB1bmB/AQsnY2AU2fsDa+HSwVgFMRcQtwS6H6R0fSzVSDcppofGb3NcDIOQAj5wCMnAMwcg7AyDkAI+cAjJwDMHIOwMg5ACPnAIycAzByDsDIOQAj18fIIKtJWgK2Aa8Ezic358IlJdrkAPRA0qVUz/LfRD8vxjSeAtgBKKieUPtT5KbX6cJVkn5Bg/UGfA1QgKQlSZ8Avkf/Bx+q+Rk+STU6+4zH2GeAjkl6FvBVcrOqlXIN8Hvgw9MK+AzQIUnnAN9iPg7+iv2SXj3tSwegI/Wp9hCFJnVu4Syqn4OpX1o33kd1lT+Ppq434AB0QNJ5wPVDt2MNE9cbcAC68RHmf2HMiT9NDkBLkp4PXDV0OxrYOulD3wa293aaD9la7SfAD6nWU5hV6/UGHID2UhMyAL8G9kbE4bY7lnTZDPt/EgeghfrWb0dikz8Cr4mIE2ValOdrgHbOI/dw56PzdPDBAWjrRYmyopoxda44AO1kbv0ejYgi6/604QC0k1kYY5OkpxZryYwcgHYeSZRdAnaXasisHIB2HgD+nSh/ff1a2NxwAFqIiGWqDp2mtgOHJA27COQq7gdo7zvk5v2/Etgp6fPA3eSW4lntbOD1M277Xw5Ae7cCH0xus4XqAVKfJnY5+yegpYi4l3ohhzn38KQPHYBufGzoBjRw76QPHYAORMSdwG1Dt2MN35z0oQPQnfdQPembR48At0/6wgHoSET8CXgzufUV+3JDRJya9IUD0KGI+DFwOfCXoduyyhHg4LQvHYCORcRRqtFAPxu6LcB9wFsj4vS0Ag5AARHxS6pev4/T7pWvNo4Al0TEGVd8dQAKiYhTEbEfeAlwE/1eG3wf2L3WwQf3BBYXEQ8C10r6APBaqmFjF1K9TfRcmv8nzKwX8NOIaLRMrwPQk4j4J/CN+i8tuV5AY/4JGDkHYOQcgJFzAEbOARg5B2DkHICRcwBGzgEYuUwAMusBO1jrROZAPZYo++JsQ2wYmQD8NlF2Vz1xks25TAAyLzgsATdL8sOmOZcJwFFg6pslE7wBuF3SubkmWZ8a/w+NiL9KuovcNKiXA8clfRc4DjR6Rm0TzcV6AV8gPw/uyrDouRsavcD+0bRg9nbty8DvkttY/040LZgKQD0c+kC2Nda7HzQt2HhpkRX11GiHgalTkNugTgJbm74TmO6xqyt+B7DmG6c2iJuaHnyY4QywQtJ24A7gmbPWYZ27H7hg2jCwSWbus4+Ie4DXUS1JYsNbBvZkDj60fGhTh2Ab1SgUG84y8LaIuDu7YeundhHxENU9/j58iziE+6nmH544/HstM18DTCJpI9X06VdT9Vxt6LJ+e5KTwI3Awexpf7VOA7CapGdTrVLxcqpJkfpYMXPRPUbVyXMUOJa52jczMzMzMzMzM7Mx+w8z0NPH8YcsLQAAAABJRU5ErkJggg==",
    "workflow": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAIEklEQVR4nO2da4xeRRnHf08vlC4gBNsUBQppAREbvFBJKmpCAkkFL1EgVgopqTZqTLTwoUETo2iieIGEGI0piBU/EJNCBBNpE8sHFbBig0aTttyKK8TGC7S2Bbq0+/fDvG3Wurx75uyZd8555/kl+2Uzc85/5/z3nJlnZp4Bx3Ecx3Ecx3Ecx3Ecx3GGH0t1YUlzgUuAM4AzgQXA7AS32guMAs8BO83sqQT3cKogabak6yTdL+mA8vCEpBslnZa7PYpC0ockPZnpoU/GQUm3Sjo+d9sMNZLmSdqS91n3ZYekZbnbaSiRdF6vgdvOq5Kuyd1ebaR2J1DSEuAR4A3NyUnKYeDTZvaj3ELaRC0DSHojsBVY3Kyc5IwDV5jZ5txC2kK0ASQZsAW4tHk5A+GfwFIzG80tpA3MqFHnarr78AHmA3fmFtEWot4AkmYBfwbOTyNnoCz3T0H8G2AFw/HwAb7V+5wVzazI8h+rcY9DhNHCLmCsRv2pOAW4DDg1st7bgaXA440rGkYkzZW0P3L8vUnSWQPQNiLpthrxgW+n1jY0SLossnE3K/QZBqlxXaTGJwepr9NIWh3RsK9JOjuTzl9G6DwkKcUMZWeI6QS+KaLso2b2XKSWprglouxMwnR1saQywLOxQhrk98DLEeXPTCWkC8QYIOZVmaK3XwkzE3AwosqcVFq6QJ1IoDNEuAEKxw1QOG6AwnEDFI4boHDcAIXjBigcN0DhuAEKxw1QOG6AwnEDFI4boHDcAIXjBigcN0DhuAEKxw1QOG6AwokxwN6IsvNjhTSFpBnAiRFV9qXS0gViDPC3iLKXS4rdq9cUFxK3gvmFVEK6QCoDnAh8M1JLU6yNKDsO7E4lZKiQdLqkwxHbriTppgFrXBOpb+cg9XUeSb+NbGBJulfS2xLreqekn9TQ9t2UurpAbIaQzwN31LzXy/TfsTOTsEtndk/XOGGH0UFAfeqdABxXU9P7zew3NesOBbEGOBl4GpiXRs5AeRY4z8wO5xaSk6g4gJntBb6eSMug+VLpDx/qpYk7Dvgj8Nbm5QyMPwAX9zaSFk10JNDMxoCr6G4A5QDwKX/4gVqhYDPbDqymf+esjQhYZWZ/yi2kLdSeCzCzjQQTHGpOTlIErDOz+3ILGSokXS7pPzXG4IPkFUkrc7fV0CLpLWrvmQGPS3pH7jYqAkkrJT2V+YEf4RlJK+TZQPvSeOP0GnwZcD3wAUIWrplN32cSxglBql8ADxAylSUZ5ysMhRcQ/rYTEtxijDD59ryZvZbg+kdJ/t+hkCzyDGAh/x+yvQi4teKl9gMfneT3BwiNtdvMknRIJc0E3gd8GPggcG6K+0zCYcJ09Sbgp8AjQzV8VehAVuXFDPpmKJyCtqvZr1NtnpZ0vRr8rPmSsNdB0kXANsJ/3tl51RxlMXAPsEVSI1nb3QCToHDA1K+Bto4eLgW2SbpyuhdyAxyDpJuBnwEjubVMwQjwc0mrpnMRN8AEJK0AvsEAOscNMQv4saQ65zgAboCjSFoK3E13Hv4RDNggqdbsrBuAo0vJ7wLm5tZSk5OAjaqR+t4NEFhJOEKmy1wAfCa2UvEGUAhUfS23job4sqSok1wHeqRLS7mE+uP8MUIksmnmUG8UMh9YBXyvWTmJUAsigZJuj4zGjUu6S9KFKfRM0HWBwpL6WLqzylntMEDM7OW4pBtS6Oij76aYp6+weef0qtcvug+gMKsXcwD23Wa2IZGcSTGz24H1EVVmAO+KKVwypxE37q+7KWa6fJEwG1qVymc1lm6AN0eUHQP+kkpIP8zsReBXEVUWVi2Y2wAxCzbmqPnVPSdFlD2QeS7+HxFlT65aMLcB9kSUHSEEO5wGyW2Av0aWj9n771QgqwHM7N/AaESVT0q6LpWeEsn9BgB4OKKsAfdIWi9pSSpBJdGGUPD9wA0R5Q1YA6yRNFXOgako+uBoaIcBHiKs6q1zhu8I7V+502qyfwJ6S7m7kKqlawtFKpHdAD1+CGzPLWIKhvKQ6VYYoJdzYDWQdBfMNKmbh6jVtMIAAGb2O+DG3Dr60Jq2apJW/VFm9n3gq7l1vA7DsyVrAq0yAICZ3QJ8jvZ9DtqmpxFaZwAAM/sBYTPmjtxaJjCdeENraaUBAMxsK2Gl7lrg+cxyIG7msjO01gAQRgdmdgewCPgIsIG4uQNnCtoQCZySXpKEB3s/SDqFsOplHtML0MTkJxhKOmGAYzGzPcStJZiUBAtMOkerPwFOetwAheMGKJzSDRAztk+xKDWGmLmIykGr0g0QE18YAd6dSkgFzoko+/eqBUs3wAuE/IJV+UoqIf2QtIiQe7EqlU9CK9oAZnYQ2BVR5QpJ61LpmQyFpA/riUu26cGyqkj6TuTmS0m6TVLypWiSFik+B/M+ScdXvUfxgRBJFwNba1R9ibBd66VmFQGhw7cYeA/xaXbvM7OrqxZ2A4Se/RN0P0XMEa41s3urFi7eAACSlhNWJ3edncCSmJzJRXcCj2Bmm4jbfdtWbo5NmO1vgB6SFhJyA3f1TMTNZrY8tpK/AXqY2Sjwcbq58GMXUGvPpBtgAmb2MPBZumWCPcCVZvavOpXdAMdgZncSzkV8JbeWCjwDvLd3jJ/TJJKWSdpeI0g0KDZLOjV3Ow01kmZL+oKk/Xmf9f+wQ9I18tVMg0PSAklrJW3L9ND3Sdoo6RMKqW0bw10UiaTFhIOzz+r9VE7IFMEYsJswXT0KPGZmrya4j+M4juM4juM4juM4juM4RfBf3MAUnL/R1RkAAAAASUVORK5CYII=",
    "wrench": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAJCElEQVR4nO2da4xeRRnH/+NuoUZBS7VIsZIioPaCAUSiSRWQIghaFTXRRBIUjUItUEFNNJh+0JhoNNpa0aAiRdGISkxMihgpRqOhkciWmlQFLJVbr0ChNNvu/vwwu7AsZ7vnmTnvey4zv+T98mbOPM/Mf845c+byjJTJZDKZTCaTSQ5XtwNtBDhB0mmSFkl6raR5kuZImiVpcELS/ZJ2SnpY0lZJWyQNSdronHuinz5PRW4AJQBmS7pA0rsknSnp6MgsRyVtknSbpHXOuXsj88tUDTAD+ADwO+AAvWMUuAk4qu4yZyQBRwJfAB7qoehF/AuYV3f5kwV4MfA5YGefhZ/IPcDMuusiOYBlwAM1Cj+Ra/tZ9qQ7gcArJK2V9MG6fZnAbknHOuf298PYi/phpIkAZ8t/kjVJfEk6Sv5Loy8k2QCAqyX9XtIxdfsyBaf0y9Dg9Em6AzAo6XuSLq3bl2mY3S9DyTQA4DBJP5f0vrp9KcG+fhlKogEAMyT9UtJ7Ks56m6SNkv4pP9S7Xc8Xb6Wk8wPy3RLvWkaSBDj8SFsVjAIbgOXA8dPYvWwsvZWDQFP7Ju0DWFWB8HuBb+IngcrYDBUf4JZe10kyABdFCAEwDHwbeKXBZoz4+4CTelknyQAcDzwRIf5fgYVGmzHijwIX96o+kgIYAP4SKMQI/rUxYLR5eaT4n+lVfSQHcFWgEHuBCwPsZfGbAjAXeDJAiF3A6QH2svhNAvhxgBB7gFMDbGXxmwSwAP8OtzCMnxiy2uqE+F0bCfyi7BNcVznn/mi5ALhc0mqFTacj6Qrn3OqAayf6MFfSUkknSJohaa+kuyVtcM49E5N3KwGOw75279cBdmq984HXA7cy9ZNuN3Atqa0sAr5qFGMnMMdoI1b85ZFlvBh4pqS9e0hljSH+u/9hoyCXGW3ULX7IINMWYFaM3VYAnBNQMaX7P/jl4W0Tf5wbY2y3AuA6Y6VcYsh7Lv4zMYS6xR/3YUGMD40GP927zVAhjwKHG/L/QUTF1y3+OF+L8aPR4HvFFr5hyHsWfobOSpPEB/h7jC+NBrjUWBlvMuT9kYDKbpr44BeZHFFkqwurgi3j9w9JstwNZxh9QdIK59wa43XPZeC/Ttao2j0bA/K7mV9AFxrAyYa0dzjnMKQ/zpC2qeKPs7jozy40gNcZ0v7NmHfpzqL8wtDrjfk/C354uVfiS1PUU6sbAPAy+aAMZdlkNPG4Ie1CSbcSMARL3NxCWeYX/dnqBiDpVcb0/zGmty7PfqeMjaBP4ktT7IJqewOwBFUYkfSYMf8NxvSSoRFEio+k2w3puxeAAlhq+BTaE5D/ALA18NNr/aEaARXMKgLnG67ZUeRH258AlvUMw9bMnXMjkr5uvW6MKZ8EVLeewFKmGQG2mg1wruEO2BVoYxDYaLAzmec9CahwPQF9KH+jAZYYKmAYCOpoAfOBRwy2JrMemEnFU8rA+w3X/6+aWm8QwGJjJQZ3hIBFwHajvYncS8VTysCnDHlsLipX2/sA1l594bdwGcZi+Z0tqbAzVYKFCn/nTzXCeMgNqpPYXvRn2xvADkmWBZBR8+IVNAIr0w0vW8rzYNGfrW4AY+P69xkuMW/8KLDZr0ZQZm6h9MympH8X/dnqBjCGZXj37VUY7EMjmFZ84A2yhawtDEfbhQZwtyHtYipaKdvDRlB2VtEaeaSbi0KAtxp71Csrth/7dTCR0otJ8NvXy7K1yjI3CnxQZ8tm0KEe+FBFI7CIb10G95Oqy9wogN8YK+TMHviwCL/gNIQR4NMGW6uN+X+46vI2Cnz8HgvX9ciP+diHjXcByww2jgaeNuS/H79uopsQtoDyjh76Mzjm03SziE/j72Tr9jTr3X/I/Y+tDhZN+Bq625xz5/XApWfBh5hZIv/peaL8kTJPSbpf0l2S1jvnnjTmebJ8b94yC7rMOfdbi51WQNzS6VV1+28F/2S5y1jOrRi2wLWGSPH3Aa+puwxWgK8ElPXquv2unEjxR4GmB4p+AcB7sUc+2QG8tG7fK6UC8VfUXQYrwFuw9frHuaZu3yslUfHfTNjO5PvoUoSQRMVfSli4O4AL6va/MlITH7/dfSXhZxXeXHcZKiNB8efiD6oMZRtdOYAyJfGBw4ArgMcjxD8ALKm7LJVQt/jAPODL+LvxduC7QCWLSSbZmQl8Arg/Qvhxrqzav1qgfvE/xtRRQTbh38/BA0n4d/zpwLfw3+pVsDamzI2B+sVfbrA/BHwHH7vvNPwhlJPzGwBeDZwFrAB+Stz+giJ+gTG0fSOhXeJPxUF8pM7dxL3Py3IL/kCsdkM3xO83P6ILEz1k8a2M4mMBt3oqX1IWP4DdQNVnH9YDWXwrG2jhNHYhZPEt7MHXVxf2cGTxDRwAvg9Ydv80G7L4ZRgGbgBOrKreGwFZ/OnYhj+v8Niq6rwx0G7xR/HDtr8ifH5+Kv6Ln2d4B115x0+G9ou/fEJeg/i9iCuBnwH/oHw08Z3An4G1wCV07RFfBB0Sfxo7c4CF+JhF50z4nQGcBLwkphytBD/NWaf4n4y034hz/VpJFj9hsvgJk8VPmCx+wmTxEyaLnzBZ/ITJ4idMFj9hsvgJk8VPmCx+wlRQ+a2Y1csUgA9UcDCLnyD4ee7HsviJAqzL4icK8MZAAbL4XQC4MYufKMCRlF/smMXvGsBFWfyEwUe+sPClSHtZ/CYB3GkQYDMRIUqy+A0EeMAgwucj7GTxmwj+GJOyWI8vG7eRxW8qvW4AWfyGg+0VYDqUIIvfAoA/GUQZouTO1ix+SwDWGMX5bIk8s/htAfiQUaAR4BoKngT4LdWrsvjNpDDWHPBySY9KOtyY32ZJ6/TcSdWnSvqo/LFpIZQ9SDlTNcDNgXdsVeQ7v07wAZDriqmTxW8C+EjUWfxUwR9hYhkUyuJ3DeDd2A8tzOJ3CeDKLH7ijDWCXjwJDgAfr7t8mRIAF+Lj3VXFI8BZdZcrYwA4Bh8kMeYTcQS4noJzdjItATgFuAnb4tGngB8CC+r2P+OJPnYEOELSuZLeJmmxpHmSZksalbRL0oOShiTdKekPzrl9sTYzmUwmk8lkMplMBP8HGLTprEuXspMAAAAASUVORK5CYII=",
    "x-circle": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAMqklEQVR4nO2da6xdRRXH/1OhqLeUVtQCRRATaKEiqX4QoS0pApKggBoBNVDUhCgBK++GEMUoyKNgC0XEBI1thcSoQRDQtBZrKwGNkqA8Ch9sofYB9AG992Jb2p8f5hy41HPvPWv2zJ597tm/pLlNu++e/5q1ztkza2avkWpqampqampqaroOl1tAKoADJU2WdISkSZIOltQjaUzj5/jGT0nqk7Sl8bO38XOtpFWSnpP0rHNufZn6y2JEBADwHkknSJop6Vh5h4+N3Myr8sHwmKRlkpY757ZEbqN0OjIAgL0kndT4c6KkYySNKlnGbklPSHpE0hJJy5xzb5SsoTAdFQDAFEnnSjpf0oS8av6PzZJ+JWmRc25lbjHtUvkAAN4r6Wvyjp+SWU67/EvSIkl3O+c25RbTkQAHAbcCvXQu24C5+AFpTTsAhwLzgf6cnovMdmAhcHju/q0swDhgAbAzq6vSsgO4Ddgvd383qcQYAPiCpNtVvYFdKjZIukp+wEhOIVkDADhC0h3y07lu5M+SLnTOPZVLQNlzZ0kS4IA58qPlbnW+JM2Q9ARwBZDlw1h6o8D7JC2UdGpJTW6U9Ix8Fu8FSf16e+pXentquEfSIfIp5CMlvb8knQ9KOt8590pJ7UkqOQCAGZLukTQxURO9klbIZ+dWyOfwtxa5ITBOPhCmy6eap+utNYTYrJX0xU5KJLVN42suxQh/NfB94Dh8iji1HXsD04DrgTUJ7NkJXJbajtLAP+9vjtxJr+Ln1ScBWcYxDdtG4YPhroammNyW07YoAKOBeyJ2ykbgWio0j24CjAFmA2sj2vsLYO/ctgUB9AAPReqI1fjOfVduu4YDH/TnAc9Hsn0psG9uu0zgPw2PRjC+F5gDjM5tkxV8IFwN9EXoh0eBVAPPuOAHSQ9HMPoB4IO57SkKMBE/XinKUmCf3PYMCX5QVPSZvxH4dG5bYgOcAbxUsG8Wkylh1Bb45dsi/Ak4KLcdqQAmAEsK9tGC3Ha0BD/PD2UX8F3gHbntSA2wFz5vsatAf12a2463AcwgPMnTB5yW24aywT8SQvc87ASOz22DJL8jl/CM2GZgWm4bcgF8HHglsO9exG+Xy2qAA34baMA64CNZDagAwFENZ4bwIDkHhcBVgcJfZARM8WIBHEZ4BvHyXKKn4Lc4WdmE395dMwDgaPwj0cp24MiyxTpgWYDYfrr4mT8c+DFByC7o5ZT5KADODRC5iy4c7VsBTidsivilsgSOxQ/grFxbisARAHBdQP9uwG9eSS5uQYC4ZXRBkicW+GTR8oB+np9a2GHYEz4bGcHp3VQABwDrjX29g5SzK+BOoyAYgQs7ZYHPFlpJs1aAj8jXjWJ+k0RMFwHcb+zz10nxjYt9pa+POtlTGOAQ7FPDm2OL2B//pquFOVFFdDHANca+78W/fxFNgDXl+286cBtXVQH2AV4w+uDKdu7dVvYIeFLS0QbNFzrn7jRcbwIYK+lMSR9t/NM/JN3nnHstVZu52we+KckyzXvaOVc85Q58zBh560m4exc4C7+esCebgLNStZu7feCd2BNwU2M0PM/YaLK3WoCzgd1DtL0buChh+xe10f7ZCdu/0uiLW4s2uDc+kdMurwJjItm7p5axtP7ktXJC9CBgeOc32YR/REQH2Bd4zeCPDRR5XQ441dAYwN0R7d1TyyyDjt34Z2asti+gPec3OS9W2y20/NSgA+CUoe433LtnJxv1LTReb+EYw7VO0rwYQQBcIOnHsr1JbdFqZbHx+k8O9Z/DBcBMQ0Nr5F/JToW1lErhIAh0vmTXauERSasN15841H8OGgD48quWSF7snNttuN7KkwG/ExwEBZwvhWlti0ZNoXsNvzIVGG9uCPi88VlznLkRm552B4GtMA0MaX/A14pkg8AB+qYbNZ0Z0ohl3X8bJbzGzPDTwKFoKwgo5vyk08ABGkdjWx+w7xMA/mZo4OEEdg6myzoiH8iQs4MI9764xH74vUHb4yENWCpeXJHAxqG0RQ+CCPcszfkNvZb1GVudJHydXgufSGTnUBqjBUGEe5Xq/Ibm44w6D7DcfKbx5vZRZgSIEAQR7lG68xu2jzdqnWG5+TcMN96Q0M52tBZ1YMc5f4DtlnoDF7S6x2B5gEkGHc/apcfDOfcTSV9XWPLFKWyej6TZzrnbA343JqsM17b06WABYKlvbxGRhIJBYKUqzpdsfd/Sp4MFgGVFb63h2mSUFARVcr4kvWi4tmWlscECwFKWbJvh2qQkDoKqOV+y9f3IDwApWRBU0flSRQKg13BtKUQOgqo6X0oYAJYxQOUCQIoWBFV2vpQwAGreorp1+SIwWABYPtVJ9gAWhWLr+U2i7SxKROGx2mABYPlqqVwARHJ+kyoHQSUCoFKVrCM7v0lVg6AOgIEkcn6TKgZBJQLgYMO1yUjs/CZVC4IPGK41BcB6w40tC0dJKMn5TaoUBJa+X9fqHwcLAMsiw2TDtdEp6HwUvopYhSCwBEBLnw4WAJYl3gmUUZ2qBRGcP1vFlpKzBQF+E46lBkD7H2rgQOPmiE7cEnZxxHuVHgSk3BLWaMCyKbStYgSxIME2rgj3LDUI8GcptcuWkAYeNzTQSdvCB93GFeHepQUB8AeDtsdCGrjN0EAfJZSEoX4xpKnR+mLIvJBGPms0PukJFtSvhg3UN8Oo6YyQRsZjK1p8fQJbB+qZZTS6SdDuXYo9DpLVB2hou8Gg5Q1CZ2nA3w0NrSHhWbeEnUhWaOs24UFwS0zb99A0ClvFsL8Odb/hHLbMoO0QSe2/fGDHOlcvvJmjwKaSlBtTZ8qWAh7Sh8MFwBJDQ5KU8qvP8s59tJ08gUGQrD6ApHON1y8NbglftnyD4eumLhJVrSJR6yhaph/7szfZIUbUZeIsyR+AuTEanWpsdAN1ocgU7YYUihy2xE+qUrEXOefuMFxvgu4sFfstST80/MpTzrkPx2rcei7wGqp+1HkH0fj0Ww+XjFe0A388rGXwAXB1NAFdDvAdY9/3EvtYWWCuUUQ/cFhUEV0IcCh+rcXCTSmEHID9tOv7ogvpMoAHjH2e5siYhpgfGcVAyEJEjSQJ+FxAf6d7jQ1/fs12o6CXgInJRI1QgIOBl419vQM4NLUwyz6BJsspUra8y8BnYFcG9HOx8wHaFDcW+E+AuO8lFzdCwLbc22Q9sF9ZAr8cIHAXcHopAjsY/EackCXoc8oW+scAkf3A9FKFdhDAsdinfADWVdsoYo/CPiAE2AzESVGOIICjgS0B/bkdyPN2FnB5gGCAtdRJojfBH8odMq4CuCSncAfcFyh8HW2sVo10gCnY8/xNfgfkrWCC3zy6OtCALXTxmAD/zH8lsO9eAPbPbYMkCTge2BloSD9dODvAj/atp7E32UGG1/GGBLg00BjwU8Tr6IJkET7JcwPhW84BZue2oyXALQWMAp8xHLFpY3x6d0XBProxtx2Dgh8ULi5o4EuMwAUk/MKONbe/JwvJPegbDvxxsw8VNBT8MmjHTxWBD+FH60VZQgnvX0YBGAM8GsHofuAaOnB7GX4b17ex76FoxV+Antw2mQB6iPNNAH6n8VXAu3PbNRzAPvhXykLn9nuyBKhMJTYT+BHv3ZE6Avz44FrKWvUygP/Wm014Rq8ViyjhPMak4AeGN0XsFPBfq78EPkPRN1+K2TYKmAbchX3T7HDMJ+HLtqUDXEJ4smgo1gA/wDuijNNLR+OPbr2ReF/zA9lBiZVGSp1S4ItI3Cvb260W+iStlD9he4WkZ5xz9vo4A8BX4zpS0nT5k7inSUo1Flkj6RznnL2kSyClzynx+eufSzqtpCZfli97t0r+jJ0++aqZWxt/l6QeSePkS6/2yAfoJPkaiJZSbEW4X9JXnHObS2pPUqZa+PhkxmWSrpPUGXPbdOyQNEfSvMbR8KWSNasEHC5pgaRTcurIyHJJFzrnns4lIOso0zn3vHPuU5JOl+0ItE5nvaRZkmbmdL5UkSNjnHMPyL99PF/SzsxyUrJD/g3fyc65hTm+8isP/uWT+cRJoVaF/+JzBalmPyMP/LuIc4FtWV1XjG34JNiE3P05GNVeWtSb08avyheg6pTdxP+UtFDSz5xzm3KLGYrKB8BAgCnyVbJmSbJVv07PJkm/lrTIObcyt5h26agAaILfOnaipJPl6+ZNVfkD2l2SnpCvw7dU0jLn3K6SNRSmIwNgT/Dp2hPkg+JYSUdIir1iuFXSc5Iek3f6cufc1shtlM6ICIBW4A9ImCwfDJMkTZQ/47BHPuU7Tm+dedgr7+Btjb/3SVor7/BVkp51zm0sU39NTU1NTU1NTU1NMv4H+MwCLuY4eboAAAAASUVORK5CYII=",
    "zap": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAABmJLR0QA/wD/AP+gvaeTAAAKAklEQVR4nO2dfbBVVRnGn8VVREMQc6QoDQtR8QMba5TMVDBHhkYSIZnECig1NT+YUXKaxCZ1zEGH1Bs02hV0cmiymFIcRiJgZAgaRRttYswYIxQ/ICjk44L3/vpjHabb6X6ctdZee59z9vrN3P/uevf77mefc/b6epaUSCQSiUQikUiUHWA0MA94Cdhe+XsRmAucUnR+iUgARwILgU56pgNYABxRdL6JDAGGABt6Eb6adcCgovNOZAAwCFjvIP5BlhSdeyIQYCCwxkP8g1xadA0JT4AjgJUB4gOsLLqOhAfAAOC5QPEBPgCOLLqeEPoVnUDeAP0lPSXpixmEa5F0UgZxCqNUDwBwqKTFkiZkGPbwDGPlTmkeAKBF0hOSLss49HsZx0tkDdAPWJTBb341OysPVsPS9N8AgJH0U0lfixD+WWNMR4S4iSwADPBwhE/+QS4vusZELwD3RhR/D/ChomtM9ABwV0TxIQ0F1y/AnMjiA8R4p0iEAtySg/j7gaOLrjVRBXBjDuIDLC+61qxomm4gMFPSvJwul37/6wlgBr2v5OmNPcBmh//vAIYVXXOiAjAFOyvnQztwOfBvhzZri645Sw4pOoEQgEmSnpSdlXPlgKSvSOqU5DKlm77+6wFgPLDP85P/ATC1EqfNse3IomsvPcDFwF5P8TuAaZU4LcB7Dm3/VHTtpQcYh31x86ETuLpLrLGO7e8ssPQE8DlgV4D411XFe8gxxhlF1V56gHNwe1uv5taqeAa37t+momovPcCZwD8DxL+9m5hnO8a4r4jaSw9wBrAtQPw7eojrOlU8Ju/aSw9wErA1QPy5vcTe6BBnK9A0w+YNAXAi8GaA+D3OCwCnOcZqzbP20gMcD7wRIP4j2LWAPcX/vmO8i/Ksv9QAxwGbAsRvo4+va9x2Au/AbiZJxAYYittvczWLahB/OG4zhwtzKr/cAMcCfw4Q/ymgz8ktYJZj3Il51F9qgGOAVwLEX4Ld9lXLtZ53iLub5AYSF2Aw8EKA+MuAw2q81lDc1g78Mnb9pQZ/d46DLAcGOFzvGsf4X41Zf6nBGjSsDhB/DY4bM7DfFrWyHxgSq/5SQ7g7xxpgoOM1h2CXgNXKslj1lxrC3TnW4+HQhV3758LVfUdNOAH0B54JEH8Dnl/LwA8drtMBDM26/nokt0WhhLtzvCrpYmPMDs/2H3X4302STgdO97xWLNolbTbG/D2rgLk8AIS7c2yUdJExZltIGg7/O0JS3e7+Ad6S9GtJrcaYjSGxok9xYodmH5N0hWeI12XFfycwlbcD29cTwyTdIOlVoJWAbepRHwD+685xlWeINySNM8a8mUE6L2cQo95okXSdpHXAcT4BepwyDaUi/kOSrvcMsUXS+caYTNbiAUfJfgvUNGrYgPxV0hhjzHaXRjG/AebKX/ytsp/8zBZiGmN2Smrmod0TJS0sOglJEnBPQFfvXWBUpLxG4jYY1IgUO3tJmDvHdmB05PxmZ3Kb65c/xLx/MW/uDuCsHHI0wGNZ3Ok6ZkSt9yOzdwDgZkn3ejbfJWm8MebFrPLpCWMMkmZKuk9uYwONxPm1/mMmDwDwbUkPeDbfLWmCMWZdFrnUgjGm0xgzW9JYSdEfugI4udZ/DB4JBGZIapVfl3KvpInGmOdD8/DBGLNK0meAcyVNknS2pBNUnwbQLnMgg6Nl0RXgSuzEiQ/7gEtySbTBAe52vLcP5pHUFOCAp/j7Scet1ARwh8f9nRU7qYkVEX04AEyOmmCTANzmeY/j7WMEPk+YNcuV0ZJrIoCbPO/xFmJZ2GMPV3TZU9+VDmB6lMSaDOBa/G3vvhszsTmeSXUC10ZLrIkApuP/Yr2JWPsYgEOAtz0TuylKUk0Gtlfl63m4D9udjZbcFzwTuy1aUk0EMJmwXlXcSSD83ki7dedI/C/ApYT1qqbkkWSrY2J3R0+qCQAuwb9X1UFgr6rm4VugTZLLW7zv6t2YtMsuM1snu6hyTWVyqBCAcZKelt/QM5K+aYxpyzarnq4GD3g+pfXMBuDCXG7g/9/P84D3PfPuxE7ABeMyGxi0/LhO+bSkFcCPyNEECjhH0lJJvqt5Zxlj5meYUt9gl1M1M2304iuU4X08C7v4xZfZsXPsLfmQrdyNQNQuKzAau+zNlzkx86ulgC9nc5/rlnYi2cEDo7ALXn25J0ZezgBPZ3Kr65fHI9yzkcBbATndn3VOB3H+zQOOkbRWdh16M9Iuaagx5l9ZBAM+KWm1pI97hmiV9J1Y3VXnN9/KBs2xkl7JPp264DDZ+oIBjpe0Qv7iP6qI4kuei0KNMVskjZH0E0nNeHp28LkAwMdkxR/uGeIJSdfEHqjy7vsaY3YbY66XNErSg5L+kVlWxfORkMZYc4nfyW4z9+EXkqYbYzpD8qiFTPu92PP0hkuqR2+9+apdkAXGGK+Rtso70kpJp/m0lz2V7ApjzAHP9olqsP6ALoss7vS8zhDcfIireYacfYnL4oE/UW61vuR6Aaxx1TLZ4WUflkuabIzZ79k+0RO4+QO2A04bK4CBWOs6X1aSLGnjAByF25bwZx3jHw78PkD8tYDLyaUJF4BpjoJ8yyF2f2BpgPjetneJGgF+5SBIB1BTFxA4FPhtgPgvA0fHrr/UYC1pXRZdrK4xbguwOED8v1ASI8pCAS5zFObmGmK2AD8PEP81wMW0MuEL8LijOCf0Ec9gD6Ty5XXsEHEiNtjfaJfTRl/oI54B5geIv5k+HrBEhmCPmHfhe73EMsDDAeJvAT6VZ/2lB/dPa4/WdLgfM9uVd4BT8qy99AD9cFuB81ovse4KEP9d4NQ8a09IAs51FKrb9XaEeR7mYnuX6AZgrqNYn+0mxi0B4u/sLmYiJ4C/OYi1har9AMCNAeK/D5xXVO2lBzjTUbAfV7Wfib87x27ggoJKT0gS8ANH0S7o0vbrhNnejS+w9IQk4XYE7TYq5w1jbe983TnagS8VXXvpAUY4Cvdopd0k6tmdI1EbuLuVTwDGE2Z7N7XouhMVgHUO4u3CWrPs9RS/A5hWdM2JCsAw3F7gNgN7PMXvxJpkNzy5HRyZA5PktvLX65QtWWuWG3KzZknUBrDC89Psyq1F15qoAvgw/m/xLtxedK2JbgC+kYP4yfOwXgF+E1n8uUXXmOgB7K4c37f5WphXdI2JXsAO4cbiEXJwDksEADyZxC8pWAv7nRHEX0SOxpEJT4BTI4i/mFhHrtQhjf6UH5txvCWSrjLGNKPvUbc0+gOwJ8NYSyVNTdYsDQQwGP9FHF15DhhQdD0JD4BVgeKvJLlzNC7YQyx9WQMMLLqGRCD4DQWvxxo7JRodYBBuq4H+SLJmaS6wbiDz6X1VUAewAPA9qaPpaLqhTuBkSTMkXSjpE7IreDZLWiXpZ8aYZjz6JpFIJBKJRCKRqJ3/AN0DbECSXu4RAAAAAElFTkSuQmCC",
}  # baked Lucide icons (white glyph, transparent)

_ICON_ALIASES = {
    "chat": "message-circle", "message": "message-circle", "comment": "message-square",
    "people": "users", "team": "users", "user": "users",
    "graph": "bar-chart", "chart": "bar-chart", "bar": "bar-chart",
    "analytics": "activity", "metrics": "activity", "data": "database",
    "ai": "cpu", "secure": "shield", "security": "shield-check", "privacy": "lock",
    "ok": "check-circle", "approved": "check-circle", "done": "check-circle",
    "warning": "alert-triangle", "info-circle": "info", "time": "clock",
    "schedule": "calendar", "fast": "zap", "speed": "zap", "quality": "award",
    "premium": "star", "money": "dollar-sign", "price": "dollar-sign",
    "revenue": "dollar-sign", "growth": "trending-up", "shipping": "truck",
    "company": "building", "goal": "target", "idea": "lightbulb", "launch": "rocket",
    "process": "workflow", "gear": "settings", "config": "settings",
    "cloud": "cloud", "global": "globe", "world": "globe", "doc": "file-text",
    "document": "file-text", "book": "book-open", "flag": "flag",
    "shield-check": "shield-check", "puzzle": "puzzle", "layers": "layers",
    "scale": "scale", "balance": "scale", "lock": "lock", "key": "key",
    "eye": "eye", "star": "star", "heart": "heart", "trophy": "trophy",
    "check": "check-circle", "phone": "phone", "mail": "mail",
}

_ICON_FALLBACK_GLYPH = {
    "check-circle": "\u2713", "check": "\u2713", "star": "\u2605", "heart": "\u2665",
    "shield": "\u25C6", "shield-check": "\u25C6", "lock": "\u25AE", "target": "\u25CE",
    "bar-chart": "\u2759", "trending-up": "\u2197", "trending-down": "\u2198",
    "zap": "\u26A1", "rocket": "\u25B2", "lightbulb": "\u25CF", "flag": "\u2691",
    "clock": "\u25F7", "calendar": "\u25A6", "dollar-sign": "$", "users": "\u25CF",
    "cpu": "\u25A3", "cloud": "\u2601", "globe": "\u25CB", "settings": "\u2699",
    "award": "\u2606", "trophy": "\u2605", "layers": "\u2263", "scale": "\u2696",
    "workflow": "\u21C4", "arrow": "\u2192", "info": "i", "alert-triangle": "!",
}


def _icon_name(raw: str) -> str:
    if not raw:
        return ""
    n = str(raw).strip().lower().replace("_", "-").replace(" ", "-")
    return _ICON_ALIASES.get(n, n)


def _icon_png_bytes(name: str, color_hex: str, px: int = 220) -> Optional[bytes]:
    """Return a tinted PNG of the baked icon glyph, or None."""
    if not _HAS_PIL:
        return None
    b64 = _ICON_PNG.get(name)
    if not b64:
        return None
    try:
        raw = base64.b64decode(b64)
        glyph = Image.open(BytesIO(raw)).convert("RGBA")
        glyph = glyph.resize((px, px), Image.LANCZOS)
        r, g, b = _to_rgb(color_hex)
        solid = Image.new("RGBA", glyph.size, (r, g, b, 0))
        alpha = glyph.split()[3]
        solid.putalpha(alpha)
        out = BytesIO()
        solid.save(out, format="PNG")
        return out.getvalue()
    except Exception:
        return None


def _draw_icon(slide, theme, name, x, y, d, *, circle_fill=None, glyph_color=None,
               idx=0):
    """Draw an icon-in-circle at (x,y) with diameter d (inches)."""
    name = _icon_name(name)
    circle_fill = circle_fill or theme["accent_soft"]
    glyph_color = glyph_color or theme["accent"]
    _oval(slide, x, y, d, fill=circle_fill)
    pad = d * 0.28
    png = _icon_png_bytes(name, glyph_color)
    if png:
        try:
            slide.shapes.add_picture(BytesIO(png), Inches(x + pad), Inches(y + pad),
                                     Inches(d - 2 * pad), Inches(d - 2 * pad))
            return
        except Exception:
            pass
    # fallback: glyph char or first letter
    glyph = _ICON_FALLBACK_GLYPH.get(name) or (name[:1].upper() if name else "\u25CF")
    tb, tf = _textbox(slide, x, y - 0.02, d, d, anchor="middle")
    _add_para(tf, glyph, first=True, size=int(d * 26), color=glyph_color,
              bold=True, font=theme["body_font"], align="center", space_after=0)


# ============================================================================
# Image pipeline (base64 / URL / Unsplash / AI-gen → resized bytes)
# ============================================================================

def _cover_crop(raw: bytes, target_ratio: float, max_px: int) -> Optional[bytes]:
    """Center-crop to ``target_ratio`` (w/h) and downscale. Returns PNG bytes."""
    if not _HAS_PIL:
        return raw
    try:
        im = Image.open(BytesIO(raw)).convert("RGB")
        w, h = im.size
        cur = w / h
        if cur > target_ratio:
            nw = int(h * target_ratio)
            x0 = (w - nw) // 2
            im = im.crop((x0, 0, x0 + nw, h))
        else:
            nh = int(w / target_ratio)
            y0 = (h - nh) // 2
            im = im.crop((0, y0, w, y0 + nh))
        if im.width > max_px:
            im = im.resize((max_px, int(max_px / target_ratio)), Image.LANCZOS)
        out = BytesIO()
        im.save(out, format="JPEG", quality=86)
        return out.getvalue()
    except Exception:
        return None


def _decode_data_image(val: str) -> Optional[bytes]:
    if not isinstance(val, str):
        return None
    s = val.strip()
    if s.startswith("data:image"):
        try:
            return base64.b64decode(s.split(",", 1)[1])
        except Exception:
            return None
    # bare base64 (heuristic: long, no spaces)
    if len(s) > 256 and re.fullmatch(r"[A-Za-z0-9+/=\s]+", s):
        try:
            return base64.b64decode(re.sub(r"\s", "", s))
        except Exception:
            return None
    return None


async def _fetch_url(url: str) -> Optional[bytes]:
    if not (_HAS_HTTPX and isinstance(url, str) and url.startswith("http")):
        return None
    try:
        async with httpx.AsyncClient(follow_redirects=True, timeout=25) as c:
            r = await c.get(url)
            if r.status_code == 200 and r.content:
                return r.content
    except Exception:
        return None
    return None


async def _unsplash(query: str, key: str) -> Optional[bytes]:
    if not (_HAS_HTTPX and key and query):
        return None
    try:
        async with httpx.AsyncClient(timeout=25) as c:
            r = await c.get(
                "https://api.unsplash.com/photos/random",
                params={"query": query, "orientation": "landscape"},
                headers={"Authorization": f"Client-ID {key}"},
            )
            if r.status_code == 200:
                data = r.json()
                u = (data.get("urls") or {}).get("regular")
                if u:
                    return await _fetch_url(u)
    except Exception:
        return None
    return None


async def _ai_image(prompt: str, request, user_dict) -> Optional[bytes]:
    if not (_HAS_OWUI_IMAGES and request is not None and user_dict and prompt):
        return None
    try:
        user_model = Users.get_user_by_id(user_dict["id"]) if _HAS_OWUI_FILES else None
        res = await _owui_image_generations(
            request=request, form_data=_OwuiImageForm(prompt=prompt), user=user_model
        )
        if isinstance(res, list) and res:
            u = res[0].get("url") if isinstance(res[0], dict) else None
            if u and u.startswith("http"):
                return await _fetch_url(u)
            if u and u.startswith("/"):
                # local cache path
                path = u.lstrip("/")
                for base in ("/app/backend/data", ""):
                    fp = os.path.join(base, path) if base else path
                    if os.path.isfile(fp):
                        with open(fp, "rb") as fh:
                            return fh.read()
    except Exception:
        return None
    return None


async def _resolve_one_image(item, valves, request, user_dict, ratio=1.4):
    """item may be a dict (image spec) or a string (url / base64 / hint)."""
    if isinstance(item, str):
        item = {"url": item} if item.startswith(("http", "data:")) else {"hint": item}
    if not isinstance(item, dict):
        return None
    raw = None
    for k in ("base64", "b64", "data", "image_base64"):
        if item.get(k):
            raw = _decode_data_image(item[k])
            if raw:
                break
    if not raw:
        url = _first(item, "url", "src", "image_url", "href", default="")
        if url:
            raw = _decode_data_image(url) or await _fetch_url(url)
    if not raw:
        hint = _first(item, "hint", "image_hint", "query", "keywords", default="")
        if hint and valves.unsplash_access_key:
            raw = await _unsplash(str(hint), valves.unsplash_access_key)
    if not raw:
        prompt = _first(item, "generate", "image_generate", "prompt", "ai", default="")
        if prompt and valves.image_generation:
            raw = await _ai_image(str(prompt), request, user_dict)
    if not raw:
        return None
    return _cover_crop(raw, ratio, valves.max_image_px)


# ============================================================================
# Deck builder + renderers
# ============================================================================

class _Deck:
    def __init__(self, prs, theme, footer_label, image_resolver=None):
        self.prs = prs
        self.theme = theme
        self.footer_label = footer_label
        self.image_resolver = image_resolver
        self.total = 0

    def blank(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    # -- chrome helpers --------------------------------------------------
    def _content_head(self, slide, slide_dict, page, *, dark=False):
        eb = _first(slide_dict, "eyebrow", "kicker", "category", "section")
        title = _first(slide_dict, "title", "heading", default="")
        y_title = 0.98 if eb else 0.8
        if eb:
            _eyebrow(slide, self.theme, eb, dark=dark)
        if title:
            _title(slide, self.theme, title, y=y_title, dark=dark)
        _footer(slide, self.theme, page=page, label=self.footer_label, dark=dark)
        return (y_title + (0.95 if len(_strip_md(str(title))) < 46 else 1.35)) if title else CONTENT_TOP


# ---- Cover -----------------------------------------------------------------

def _r_cover(deck, slide_dict, page):
    t = deck.theme
    slide = deck.blank()
    _set_bg(slide, t["bg_dark"])
    # layered decorative circles (right side)
    _oval(slide, SLIDE_W_IN - 3.2, -2.0, 5.6, fill=_lighten(t["bg_dark"], 0.08))
    _oval(slide, SLIDE_W_IN - 1.1, 4.2, 3.4, fill=_lighten(t["bg_dark"], 0.05))
    _oval(slide, SLIDE_W_IN - 4.4, 5.2, 1.8, fill=t["accent"], alpha=22)
    # brand mark (icon circle) top-left
    _draw_icon(slide, t, _first(slide_dict, "icon", default="sparkles") or "sparkles",
               MARGIN, 0.7, 0.95, circle_fill=_lighten(t["bg_dark"], 0.12),
               glyph_color=t["accent"])
    eb = _first(slide_dict, "eyebrow", "kicker", "category", default="")
    if eb:
        _eyebrow(slide, t, eb, x=MARGIN, y=1.95, dark=True)
    title = _first(slide_dict, "title", "heading", default="")
    tb, tf = _textbox(slide, MARGIN, 2.35, 8.7, 2.3, anchor="top")
    _add_para(tf, _strip_md(str(title)), first=True, size=42, color=t["on_dark"],
              bold=True, font=t["head_font"], align="left", line=1.02, space_after=0)
    sub = _first(slide_dict, "subtitle", "lead", "description", "summary", default="")
    if sub:
        tb2, tf2 = _textbox(slide, MARGIN, 4.75, 8.4, 1.2, anchor="top")
        _add_para(tf2, _strip_md(str(sub)), first=True, size=15,
                  color=t["on_dark_soft"], font=t["body_font"], align="left",
                  line=1.3, space_after=0)
    # chips row (author / tags)
    chips = _as_list(_first(slide_dict, "chips", "tags", "authors", default=[]))
    author = _first(slide_dict, "author", default="")
    if author and not chips:
        chips = [c.strip() for c in re.split(r"[,/|]", str(author)) if c.strip()]
    cx = MARGIN
    for chip in chips[:4]:
        txt = _strip_md(str(chip))
        w = min(3.2, 0.62 + len(txt) * 0.085)
        _card(slide, cx, 6.05, w, 0.5, fill=_lighten(t["bg_dark"], 0.12),
              radius=0.5, shadow=False, line=_lighten(t["bg_dark"], 0.22), line_w=1.0)
        tbc, tfc = _textbox(slide, cx, 6.05, w, 0.5, anchor="middle")
        _add_para(tfc, txt, first=True, size=11, color=t["on_dark_soft"],
                  font=t["body_font"], align="center", space_after=0)
        cx += w + 0.18
    date = _first(slide_dict, "date", "footer", default=deck.footer_label)
    if date:
        tbd, tfd = _textbox(slide, MARGIN, 6.95, 6.0, 0.35, anchor="middle")
        _add_para(tfd, _strip_md(str(date)), first=True, size=9.5,
                  color=t["on_dark_faint"], font=t["body_font"], align="left",
                  space_after=0)


# ---- Section divider -------------------------------------------------------

def _r_section(deck, slide_dict, page):
    t = deck.theme
    slide = deck.blank()
    _set_bg(slide, t["bg_dark"])
    _oval(slide, SLIDE_W_IN - 2.6, -1.6, 4.6, fill=_lighten(t["bg_dark"], 0.07))
    num = _first(slide_dict, "number", "chapter", "index", default="")
    if num:
        tbn, tfn = _textbox(slide, MARGIN, 1.7, 3.0, 1.6, anchor="top")
        _add_para(tfn, str(num), first=True, size=76, color=t["accent"],
                  bold=True, font=t["head_font"], align="left", space_after=0)
    eb = _first(slide_dict, "eyebrow", "kicker", "category", default="")
    if eb:
        _eyebrow(slide, t, eb, x=MARGIN, y=(3.4 if num else 2.6), dark=True)
    title = _first(slide_dict, "title", "heading", default="")
    tb, tf = _textbox(slide, MARGIN, (3.75 if num else 3.0), 10.5, 1.6, anchor="top")
    _add_para(tf, _strip_md(str(title)), first=True, size=38, color=t["on_dark"],
              bold=True, font=t["head_font"], align="left", line=1.03, space_after=0)
    lead = _first(slide_dict, "lead", "subtitle", "description", "summary", default="")
    if lead:
        tbl, tfl = _textbox(slide, MARGIN, (5.35 if num else 4.6), 9.5, 1.4)
        _add_para(tfl, _strip_md(str(lead)), first=True, size=14,
                  color=t["on_dark_soft"], font=t["body_font"], line=1.3, space_after=0)
    _footer(slide, t, page=page, label=deck.footer_label, dark=True)


# ---- Title only / title + body --------------------------------------------

def _r_title_body(deck, slide_dict, page):
    t = deck.theme
    slide = deck.blank()
    _set_bg(slide, t["bg_light"])
    top = deck._content_head(slide, slide_dict, page)
    body = _first(slide_dict, "body", "content", "text", "description", "summary",
                  "paragraph", "lead", default="")
    bullets = _harvest_bullets(slide_dict)
    y = top + 0.15
    if body and isinstance(body, str) and not bullets:
        tb, tf = _textbox(slide, MARGIN, y, SLIDE_W_IN - 2 * MARGIN, 4.4)
        for i, para in enumerate([p for p in str(body).split("\n") if p.strip()]):
            _add_para(tf, _md_runs(para), first=(i == 0), size=15, color=t["muted"],
                      font=t["body_font"], line=1.36, space_after=10, align="left")
    elif bullets:
        _bullet_block(slide, t, bullets, MARGIN, y, SLIDE_W_IN - 2 * MARGIN, 4.6)


def _r_title_only(deck, slide_dict, page):
    t = deck.theme
    slide = deck.blank()
    _set_bg(slide, t["bg_light"])
    eb = _first(slide_dict, "eyebrow", "kicker", default="")
    if eb:
        _eyebrow(slide, t, eb, y=2.6)
    title = _first(slide_dict, "title", "heading", default="")
    tb, tf = _textbox(slide, MARGIN, 2.95, SLIDE_W_IN - 2 * MARGIN, 1.6, anchor="top")
    _add_para(tf, _strip_md(str(title)), first=True, size=40, color=t["ink"],
              bold=True, font=t["head_font"], line=1.03, space_after=0)
    sub = _first(slide_dict, "subtitle", "lead", "body", default="")
    if sub:
        tbs, tfs = _textbox(slide, MARGIN, 4.55, 10.5, 1.4)
        _add_para(tfs, _strip_md(str(sub)), first=True, size=15, color=t["muted"],
                  font=t["body_font"], line=1.35, space_after=0)
    _footer(slide, t, page=page, label=deck.footer_label)


def _bullet_block(slide, theme, bullets, x, y, w, h, *, size=14, anchor="top",
                  fill=False):
    tb, tf = _textbox(slide, x, y, w, h, anchor=anchor)
    items = [b for b in bullets if str(b).strip()]
    n = len(items)
    if fill and n:
        # distribute so the block breathes and fills the available height
        line_pt = size * 1.35 * max(1, 1)
        avail_pt = h * 72
        leftover = avail_pt - n * line_pt
        sp = int(max(8, min(30, leftover / max(1, n))))
    else:
        sp = 9 if n <= 5 else (6 if n <= 8 else 3)
    first = True
    for b in items:
        txt = _clean_bullet(b) if not str(b).strip().startswith("**") else _strip_md(b)
        if not txt:
            continue
        _add_bullet(tf, txt, first=first, size=size, color=theme["muted"],
                    accent=theme["accent"], font=theme["body_font"], space_after=sp,
                    line=1.2)
        first = False


def _r_title_bullets(deck, slide_dict, page):
    t = deck.theme
    slide = deck.blank()
    _set_bg(slide, t["bg_light"])
    top = deck._content_head(slide, slide_dict, page)
    bullets = _harvest_bullets(slide_dict)
    if not bullets:
        body = _first(slide_dict, "body", "content", "text", default="")
        bullets = [ln for ln in str(body).split("\n") if ln.strip()] if body else []
    n = len([b for b in bullets if str(b).strip()])
    # sparse lists get larger type and breathe to fill the canvas
    size = 18 if n <= 4 else (16 if n <= 6 else 14)
    y = top + 0.3
    h = FOOTER_Y - y - 0.35
    _bullet_block(slide, t, bullets, MARGIN, y, SLIDE_W_IN - 2 * MARGIN, h,
                  size=size, anchor=("middle" if n <= 6 else "top"), fill=(n <= 8))


# ---- Two column / comparison ----------------------------------------------

def _card_column(slide, theme, x, y, w, h, card, *, highlight=False, badge=None):
    t = theme
    fill = t["primary"] if highlight else t["card_soft"]
    heading_color = t["on_dark"] if highlight else t["ink"]
    body_color = t["on_dark_soft"] if highlight else t["muted"]
    _card(slide, x, y, w, h, fill=fill, radius=0.06,
          line=None if highlight else t["card_border"], line_w=1.0,
          shadow=highlight)
    if badge:
        bw = 1.7
        _card(slide, x + (w - bw) / 2, y - 0.18, bw, 0.42, fill=t["accent"],
              radius=0.5, shadow=False)
        tbb, tfb = _textbox(slide, x + (w - bw) / 2, y - 0.18, bw, 0.42, anchor="middle")
        _add_para(tfb, _strip_md(str(badge)).upper(), first=True, size=9,
                  color=_on(t["accent"]), bold=True, font=t["body_font"],
                  align="center", tracking=1.2, space_after=0)
    pad = 0.36
    cy = y + pad
    heading = _first(card, "heading", "title", "label", "name", default="")
    icon = _first(card, "icon", default="")
    if icon:
        idd = 0.66
        _draw_icon(slide, t, icon, x + pad, cy, idd,
                   circle_fill=(_lighten(t["primary"], 0.14) if highlight else t["accent_soft"]),
                   glyph_color=(t["light_accent"] if highlight else t["accent"]))
        if heading:
            tbh, tfh = _textbox(slide, x + pad + idd + 0.22, cy, w - 2 * pad - idd - 0.22,
                                idd, anchor="middle")
            _add_para(tfh, _strip_md(str(heading)), first=True, size=19,
                      color=heading_color, bold=True, font=t["head_font"], space_after=0)
        cy += idd + 0.22
    elif heading:
        tbh, tfh = _textbox(slide, x + pad, cy, w - 2 * pad, 0.6)
        _add_para(tfh, _strip_md(str(heading)), first=True, size=19,
                  color=heading_color, bold=True, font=t["head_font"], space_after=0)
        cy += 0.62
    sub = _first(card, "subtitle", "value", "price", "tagline", default="")
    if sub:
        tbs, tfs = _textbox(slide, x + pad, cy, w - 2 * pad, 0.5)
        _add_para(tfs, _strip_md(str(sub)), first=True, size=15,
                  color=(t["accent"] if not highlight else t["light_accent"]),
                  bold=True, font=t["body_font"], space_after=0)
        cy += 0.5
    desc = _first(card, "description", "detail", "note", default="")
    if desc:
        tbd, tfd = _textbox(slide, x + pad, cy, w - 2 * pad, 0.6)
        _add_para(tfd, _strip_md(str(desc)), first=True, size=11.5,
                  color=body_color, font=t["body_font"], line=1.25, space_after=0)
        cy += 0.5
    points = card.get("points") or card.get("bullets") or card.get("items") or []
    pts = _harvest_bullets({"points": points}) if points else []
    if pts:
        tbp, tfp = _textbox(slide, x + pad, cy + 0.1, w - 2 * pad, h - (cy - y) - pad)
        first = True
        for p in pts:
            _add_bullet(tfp, _clean_bullet(p), first=first, size=12,
                        color=body_color, accent=(t["light_accent"] if highlight else t["accent"]),
                        font=t["body_font"], space_after=6)
            first = False


def _split_cards(slide_dict):
    """Return a list of card dicts for two_column/comparison."""
    cards = []
    for side in ("left", "right"):
        c = slide_dict.get(side)
        if isinstance(c, dict):
            cards.append(c)
    if cards:
        return cards
    # columns/cards/items array
    for key in ("columns", "cards", "options", "tiers", "plans", "items"):
        arr = slide_dict.get(key)
        if isinstance(arr, list) and arr and isinstance(arr[0], dict):
            return arr
    # flat left_*/right_*
    flat = {}
    for side in ("left", "right"):
        d = {}
        for f in ("heading", "title", "points", "bullets", "subtitle", "description"):
            v = slide_dict.get(f"{side}_{f}")
            if v is not None:
                d["heading" if f in ("heading", "title") else f] = v
        if d:
            flat[side] = d
    if flat:
        return list(flat.values())
    # fallback: split bullets in half
    bl = _harvest_bullets(slide_dict)
    if bl:
        mid = (len(bl) + 1) // 2
        return [{"heading": "", "points": bl[:mid]}, {"heading": "", "points": bl[mid:]}]
    return []


def _r_comparison(deck, slide_dict, page):
    t = deck.theme
    slide = deck.blank()
    _set_bg(slide, t["bg_light"])
    top = deck._content_head(slide, slide_dict, page)
    cards = _split_cards(slide_dict)
    if not cards:
        return
    n = min(len(cards), 3)
    cards = cards[:n]
    gap = 0.4
    total_w = SLIDE_W_IN - 2 * MARGIN
    cw = (total_w - gap * (n - 1)) / n
    y = max(top + 0.25, 1.95)
    # Adaptive height: description-only cards are shorter and vertically
    # centred; cards with bullet lists (pricing tiers etc.) fill the slide.
    has_points = any(
        (c.get("points") or c.get("bullets") or c.get("items"))
        for c in cards if isinstance(c, dict)
    )
    if has_points:
        h = FOOTER_Y - y - 0.25
    else:
        h = min(FOOTER_Y - y - 0.25, 3.0)
        y = top + max(0.25, (FOOTER_Y - top - h) / 2)  # centre vertically
    for i, card in enumerate(cards):
        x = MARGIN + i * (cw + gap)
        hl = bool(card.get("highlight") or card.get("featured") or card.get("recommended"))
        badge = card.get("badge")  # only explicit badges (never auto)
        _card_column(slide, t, x, y, cw, h, card, highlight=hl, badge=badge)


# ---- KPI row ---------------------------------------------------------------

def _r_kpi(deck, slide_dict, page):
    t = deck.theme
    slide = deck.blank()
    _set_bg(slide, t["bg_light"])
    top = deck._content_head(slide, slide_dict, page)
    stats = _first(slide_dict, "stats", "kpis", "metrics", "numbers", "items",
                   "cards", default=[])
    stats = [s for s in _as_list(stats) if isinstance(s, dict)] or _as_list(stats)
    if not stats:
        return
    n = min(len(stats), 4)
    stats = stats[:n]
    gap = 0.4
    total_w = SLIDE_W_IN - 2 * MARGIN
    cw = (total_w - gap * (n - 1)) / n
    y = max(top + 0.35, 2.2)
    h = 2.7
    for i, s in enumerate(stats):
        x = MARGIN + i * (cw + gap)
        _card(slide, x, y, cw, h, fill=t["card_soft"], radius=0.06,
              line=t["card_border"], line_w=1.0, shadow=False)
        if isinstance(s, dict):
            value = _first(s, "value", "number", "stat", "amount", default="")
            label = _first(s, "label", "name", "title", "caption", default="")
            change = _first(s, "change", "delta", "trend", default="")
        else:
            value, label, change = str(s), "", ""
        pad = 0.3
        vstr = _strip_md(str(value))
        # scale the big number down when long so it stays on one line
        vsize = 40 if len(vstr) <= 5 else (33 if len(vstr) <= 7 else 26)
        tbv, tfv = _textbox(slide, x + pad, y + 0.45, cw - 2 * pad, 1.0,
                            anchor="middle", wrap=False)
        _add_para(tfv, vstr, first=True, size=vsize, color=t["accent"],
                  bold=True, font=t["head_font"], align="left", space_after=0)
        tbl, tfl = _textbox(slide, x + pad, y + 1.55, cw - 2 * pad, 0.9)
        _add_para(tfl, _strip_md(str(label)), first=True, size=12.5, color=t["muted"],
                  font=t["body_font"], align="left", line=1.2, space_after=0)
        if change:
            tbc, tfc = _textbox(slide, x + pad, y + h - 0.5, cw - 2 * pad, 0.35)
            up = not str(change).strip().startswith("-")
            _add_para(tfc, ("\u25B2 " if up else "\u25BC ") + _strip_md(str(change)),
                      first=True, size=11, color=(t["accent"] if up else "C0392B"),
                      bold=True, font=t["body_font"], space_after=0)


# ---- Timeline / process ----------------------------------------------------

def _r_timeline(deck, slide_dict, page, *, process=False):
    t = deck.theme
    slide = deck.blank()
    _set_bg(slide, t["bg_light"])
    top = deck._content_head(slide, slide_dict, page)
    steps = _first(slide_dict, "steps", "phases", "milestones", "stages", "items",
                   "events", default=[])
    steps = _as_list(steps)
    if not steps:
        return
    n = min(len(steps), 5)
    steps = steps[:n]
    gap = 0.35
    total_w = SLIDE_W_IN - 2 * MARGIN
    cw = (total_w - gap * (n - 1)) / n
    # estimate needed card height from richest step, then centre vertically
    has_desc = any(isinstance(s, dict) and _first(s, "description", "detail", "text",
                   "body", default="") for s in steps)
    h = 2.5 if has_desc else 1.7
    y = top + max(0.6, (FOOTER_Y - top - h) / 2)
    # connector line
    _rect(slide, MARGIN + cw / 2, y - 0.32, total_w - cw, 0.03, fill=t["card_border"])
    for i, s in enumerate(steps):
        x = MARGIN + i * (cw + gap)
        # node circle with number
        nd = 0.5
        _oval(slide, x + cw / 2 - nd / 2, y - 0.32 - nd / 2, nd, fill=t["accent"])
        tbn, tfn = _textbox(slide, x + cw / 2 - nd / 2, y - 0.32 - nd / 2, nd, nd, anchor="middle")
        _add_para(tfn, str(i + 1), first=True, size=15, color=_on(t["accent"]),
                  bold=True, font=t["body_font"], align="center", space_after=0)
        _card(slide, x, y, cw, h, fill=t["card_soft"], radius=0.06,
              line=t["card_border"], line_w=1.0, shadow=False)
        if isinstance(s, dict):
            label = _first(s, "title", "label", "heading", "name", "phase", default=f"Step {i+1}")
            when = _first(s, "date", "when", "time", "period", default="")
            desc = _first(s, "description", "detail", "text", "body", default="")
        else:
            label, when, desc = str(s), "", ""
        pad = 0.26
        cy = y + pad
        if when:
            tbw, tfw = _textbox(slide, x + pad, cy, cw - 2 * pad, 0.3)
            _add_para(tfw, _strip_md(str(when)).upper(), first=True, size=9,
                      color=t["accent"], bold=True, font=t["body_font"],
                      tracking=1.0, space_after=0)
            cy += 0.34
        tbt, tft = _textbox(slide, x + pad, cy, cw - 2 * pad, 0.6)
        _add_para(tft, _strip_md(str(label)), first=True, size=14, color=t["ink"],
                  bold=True, font=t["head_font"], line=1.05, space_after=0)
        cy += 0.6
        if desc:
            tbd, tfd = _textbox(slide, x + pad, cy, cw - 2 * pad, h - (cy - y) - pad)
            _add_para(tfd, _strip_md(str(desc)), first=True, size=11, color=t["muted"],
                      font=t["body_font"], line=1.22, space_after=0)


# ---- Quote -----------------------------------------------------------------

def _r_quote(deck, slide_dict, page):
    t = deck.theme
    slide = deck.blank()
    _set_bg(slide, t["bg_dark"])
    _oval(slide, -1.4, -1.4, 3.4, fill=_lighten(t["bg_dark"], 0.06))
    # big quotation mark
    tbq, tfq = _textbox(slide, MARGIN, 0.9, 2.0, 1.6)
    _add_para(tfq, "\u201C", first=True, size=110, color=t["accent"], bold=True,
              font=t["head_font"], space_after=0)
    quote = _first(slide_dict, "quote", "text", "content", "message", default="")
    tb, tf = _textbox(slide, MARGIN, 2.4, SLIDE_W_IN - 2 * MARGIN - 1.5, 3.0, anchor="top")
    _add_para(tf, _strip_md(str(quote)), first=True, size=27, color=t["on_dark"],
              italic=True, font=t["head_font"], line=1.22, space_after=0)
    author = _first(slide_dict, "author", "by", "source", default="")
    role = _first(slide_dict, "role", "title", "position", default="")
    if author:
        tba, tfa = _textbox(slide, MARGIN, 5.7, 9.0, 0.8)
        line = f"\u2014 {_strip_md(str(author))}"
        _add_para(tfa, line, first=True, size=15, color=t["accent"], bold=True,
                  font=t["body_font"], space_after=2)
        if role:
            _add_para(tfa, _strip_md(str(role)), size=12, color=t["on_dark_faint"],
                      font=t["body_font"], space_after=0)
    _footer(slide, t, page=page, label=deck.footer_label, dark=True)


# ---- Closing ---------------------------------------------------------------

def _r_closing(deck, slide_dict, page):
    t = deck.theme
    slide = deck.blank()
    _set_bg(slide, t["bg_dark"])
    _oval(slide, SLIDE_W_IN - 3.0, -2.0, 5.2, fill=_lighten(t["bg_dark"], 0.07))
    _oval(slide, SLIDE_W_IN - 1.0, 4.6, 3.0, fill=t["accent"], alpha=18)
    eb = _first(slide_dict, "eyebrow", "kicker", default="")
    if eb:
        _eyebrow(slide, t, eb, y=1.4, dark=True)
    title = _first(slide_dict, "title", "heading", default="Thank you")
    tb, tf = _textbox(slide, MARGIN, 1.85, 9.5, 1.4)
    _add_para(tf, _strip_md(str(title)), first=True, size=40, color=t["on_dark"],
              bold=True, font=t["head_font"], line=1.03, space_after=0)
    takeaways = _first(slide_dict, "takeaways", "key_takeaways", "next_steps",
                       "points", default=[])
    takeaways = _harvest_bullets({"points": takeaways}) if takeaways else _harvest_bullets(slide_dict)
    y = 3.5
    if takeaways:
        tbk, tfk = _textbox(slide, MARGIN, y, 8.5, 2.6)
        first = True
        for b in takeaways[:5]:
            _add_bullet(tfk, _clean_bullet(b), first=first, size=15,
                        color=t["on_dark_soft"], accent=t["accent"],
                        font=t["body_font"], space_after=10)
            first = False
    contact = _first(slide_dict, "contact", "email", "cta", "footer", default="")
    if contact:
        tbc, tfc = _textbox(slide, MARGIN, 6.4, 10.0, 0.5)
        _add_para(tfc, _strip_md(str(contact)), first=True, size=13,
                  color=t["accent"], bold=True, font=t["body_font"], space_after=0)
    _footer(slide, t, page=page, label=deck.footer_label, dark=True)


# ---- Alert / callout -------------------------------------------------------

_ALERT_COLORS = {
    "info": "2563EB", "tip": "059669", "success": "059669", "note": "2563EB",
    "warning": "D97706", "danger": "DC2626", "error": "DC2626",
}


_ALERT_ICONS = {
    "info": "info", "note": "info", "tip": "lightbulb", "success": "check-circle",
    "warning": "alert-triangle", "danger": "alert-octagon", "error": "alert-octagon",
}


def _r_alert(deck, slide_dict, page):
    import math
    t = deck.theme
    slide = deck.blank()
    _set_bg(slide, t["bg_light"])
    top = deck._content_head(slide, slide_dict, page)
    level = str(_first(slide_dict, "level", "severity", "variant", default="info")).lower()
    color = _ALERT_COLORS.get(level, t["accent"])
    body = _first(slide_dict, "body", "content", "text", "message", "description", default="")
    bullets = _harvest_bullets(slide_dict)
    heading = _first(slide_dict, "callout_title", "alert_title", "heading", default="")

    y = max(top + 0.35, 2.2)
    avail = FOOTER_Y - y - 0.4
    x = MARGIN
    card_w = SLIDE_W_IN - 2 * MARGIN
    icon_d = 0.95
    text_x = x + 0.55 + icon_d + 0.5
    text_w = x + card_w - 0.6 - text_x

    # Size the card to its content instead of a fixed height, so short
    # callouts don't leave a large empty coloured area.
    def _est_lines(s: str, size: float) -> int:
        cpl = max(12, int(text_w / (size * 0.0072)))
        return max(1, math.ceil(len(str(s)) / cpl))

    lines = 0
    if heading:
        lines += 1
    if body and not bullets:
        lines += _est_lines(body, 16)
    else:
        for b in bullets:
            lines += _est_lines(_clean_bullet(b), 15)
    content_h = lines * 0.34 + (0.4 if heading else 0.0) + 0.8
    h = max(1.5, min(avail, content_h))

    _card(slide, x, y, card_w, h, fill=_lighten(color, 0.9),
          radius=0.05, shadow=False)
    _rect(slide, x, y, 0.1, h, fill=color)
    _draw_icon(slide, t, _ALERT_ICONS.get(level, "info"),
               x + 0.55, y + h / 2 - icon_d / 2, icon_d,
               circle_fill=color, glyph_color=_on(color))

    tb, tf = _textbox(slide, text_x, y, text_w, h, anchor="middle")
    first = True
    if heading:
        _add_para(tf, _strip_md(str(heading)), first=True, size=18, color=t["ink"],
                  bold=True, font=t["head_font"], space_after=6)
        first = False
    if body and not bullets:
        _add_para(tf, _md_runs(str(body)), first=first, size=16, color=t["ink"],
                  font=t["body_font"], line=1.35, space_after=0)
    else:
        for b in bullets:
            _add_bullet(tf, _clean_bullet(b), first=first, size=15, color=t["ink"],
                        accent=color, font=t["body_font"], space_after=8)
            first = False


# ---- Table -----------------------------------------------------------------

def _r_table(deck, slide_dict, page):
    t = deck.theme
    slide = deck.blank()
    _set_bg(slide, t["bg_light"])
    top = deck._content_head(slide, slide_dict, page)
    rows = _first(slide_dict, "rows", "data", "table", default=[])
    headers = _first(slide_dict, "headers", "columns", "header", default=[])
    rows = _as_list(rows)
    # normalize rows of dicts
    if rows and isinstance(rows[0], dict):
        if not headers:
            headers = list(rows[0].keys())
        rows = [[r.get(h, "") for h in headers] for r in rows]
    if not rows:
        return
    ncols = max(len(headers) if headers else 0, max(len(r) for r in rows))
    nrows = len(rows) + (1 if headers else 0)
    y = max(top + 0.25, 2.0)
    h = min(FOOTER_Y - y - 0.2, 0.5 * nrows + 0.2)
    tbl_shape = slide.shapes.add_table(nrows, ncols, Inches(MARGIN), Inches(y),
                                       Inches(SLIDE_W_IN - 2 * MARGIN), Inches(h))
    table = tbl_shape.table
    r0 = 0
    if headers:
        for c in range(ncols):
            cell = table.cell(0, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(t["primary"])
            cell.text = _strip_md(str(headers[c])) if c < len(headers) else ""
            para = cell.text_frame.paragraphs[0]
            for run in para.runs:
                run.font.size = Pt(12)
                run.font.bold = True
                run.font.color.rgb = _rgb(_on(t["primary"]))
                run.font.name = t["body_font"]
        r0 = 1
    for ri, row in enumerate(rows):
        for c in range(ncols):
            cell = table.cell(r0 + ri, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(t["card_soft"] if ri % 2 else "FFFFFF")
            cell.text = _strip_md(str(row[c])) if c < len(row) else ""
            para = cell.text_frame.paragraphs[0]
            for run in para.runs:
                run.font.size = Pt(11.5)
                run.font.color.rgb = _rgb(t["ink"])
                run.font.name = t["body_font"]


# ---- Icon list / icon grid -------------------------------------------------

def _icon_items(slide_dict):
    items = _first(slide_dict, "items", "cards", "features", "points", "bullets",
                   "pillars", default=[])
    out = []
    for it in _as_list(items):
        if isinstance(it, dict):
            out.append(it)
        elif isinstance(it, str):
            out.append({"title": it})
    return out


def _r_icon_list(deck, slide_dict, page):
    t = deck.theme
    slide = deck.blank()
    _set_bg(slide, t["bg_light"])
    top = deck._content_head(slide, slide_dict, page)
    items = _icon_items(slide_dict)[:5]
    if not items:
        return
    y = max(top + 0.25, 2.05)
    avail = FOOTER_Y - y - 0.1
    rh = min(1.0, avail / max(1, len(items)))
    for i, it in enumerate(items):
        ry = y + i * rh
        d = min(0.72, rh - 0.18)
        _draw_icon(slide, t, _first(it, "icon", default=""), MARGIN, ry, d,
                   circle_fill=t["accent_soft"], glyph_color=t["accent"], idx=i)
        tx = MARGIN + d + 0.35
        title = _first(it, "title", "heading", "label", "name", default="")
        desc = _first(it, "description", "detail", "text", "body", default="")
        tbt, tft = _textbox(slide, tx, ry, SLIDE_W_IN - MARGIN - tx, rh)
        _add_para(tft, _strip_md(str(title)), first=True, size=16, color=t["ink"],
                  bold=True, font=t["head_font"], space_after=2)
        if desc:
            _add_para(tft, _strip_md(str(desc)), size=12, color=t["muted"],
                      font=t["body_font"], line=1.22, space_after=0)


def _r_icon_grid(deck, slide_dict, page, cols=3):
    t = deck.theme
    slide = deck.blank()
    _set_bg(slide, t["bg_light"])
    top = deck._content_head(slide, slide_dict, page)
    items = _icon_items(slide_dict)
    if not items:
        return
    n = len(items)
    cols = 2 if n <= 4 and cols != 3 else cols
    cols = min(cols, 3)
    rows = (n + cols - 1) // cols
    gap = 0.35
    total_w = SLIDE_W_IN - 2 * MARGIN
    cw = (total_w - gap * (cols - 1)) / cols
    y0 = max(top + 0.25, 2.0)
    avail_h = FOOTER_Y - y0 - 0.1
    ch = (avail_h - gap * (rows - 1)) / rows
    for i, it in enumerate(items[:cols * rows]):
        r, c = divmod(i, cols)
        x = MARGIN + c * (cw + gap)
        yy = y0 + r * (ch + gap)
        _card(slide, x, yy, cw, ch, fill=t["card_soft"], radius=0.06,
              line=t["card_border"], line_w=1.0, shadow=False)
        pad = 0.3
        _draw_icon(slide, t, _first(it, "icon", default=""), x + pad, yy + pad, 0.66,
                   circle_fill=t["accent_soft"], glyph_color=t["accent"], idx=i)
        title = _first(it, "title", "heading", "label", "name", default="")
        desc = _first(it, "description", "detail", "text", "body", default="")
        tbt, tft = _textbox(slide, x + pad, yy + pad + 0.82, cw - 2 * pad, ch - 1.2)
        _add_para(tft, _strip_md(str(title)), first=True, size=15, color=t["ink"],
                  bold=True, font=t["head_font"], space_after=3)
        if desc:
            _add_para(tft, _strip_md(str(desc)), size=11.5, color=t["muted"],
                      font=t["body_font"], line=1.24, space_after=0)


# ============================================================================
# The Tools class (OpenWebUI native function calling)
# ============================================================================

class Tools:
    def __init__(self):
        self.valves = self.Valves()

    class Valves(BaseModel):
        default_theme: str = Field(
            default="auto",
            description="Default theme (auto | midnight | forest | ocean | ...).",
        )
        footer_label: str = Field(
            default="",
            description="Default footer label (overridden by spec.footer).",
        )
        unsplash_access_key: str = Field(
            default="", description="Unsplash key for stock images (optional)."
        )
        image_generation: bool = Field(
            default=False, description="Enable AI image generation via Open WebUI."
        )
        max_image_px: int = Field(
            default=1600, description="Maximum image width (px)."
        )
        emit_status: bool = Field(default=True, description="Emit status events.")
        pptx_export_dir: str = Field(
            default="/app/backend/data/cache/files",
            description="Fallback directory for saving.",
        )

    # -- status / link helpers -------------------------------------------
    async def _emit(self, emitter, desc, *, done=False):
        if emitter and self.valves.emit_status:
            try:
                await emitter({"type": "status",
                               "data": {"description": desc, "done": done}})
            except Exception:
                pass

    @staticmethod
    async def _emit_link(emitter, fname, url, *, slides=0, kb=0):
        if not emitter:
            return
        msg = (
            f"\n\n---\n\n\U0001F4CA **Presentation ready** · {slides} slides · {kb} KB\n\n"
            f"\U0001F4E5 [Download {fname}]({url})\n\n---\n"
        )
        try:
            await emitter({"type": "message", "data": {"content": msg}})
        except Exception:
            pass

    def _save(self, data: bytes, *, title, request, user_dict):
        slug = _slugify(title)
        day = datetime.now(timezone.utc).strftime("%Y%m%d")
        short = uuid.uuid4().hex[:6]
        filename = f"presentation-{slug}_{day}_{short}.pptx"
        if _HAS_OWUI_FILES and request is not None and user_dict:
            try:
                user_model = Users.get_user_by_id(user_dict["id"])
                if user_model:
                    upload = UploadFile(
                        file=BytesIO(data),
                        filename=filename,
                        headers=Headers({
                            "content-type":
                                "application/vnd.openxmlformats-officedocument."
                                "presentationml.presentation"
                        }),
                    )
                    item = upload_file_handler(request=request, file=upload,
                                               metadata={}, process=False,
                                               user=user_model)
                    fid = getattr(item, "id", None) if item else None
                    if fid:
                        return filename, f"/api/v1/files/{fid}/content", None
            except Exception as exc:
                print(f"[generate_slides] Files API save failed: {exc}")
        export_dir = (self.valves.pptx_export_dir or "").strip() or \
            "/app/backend/data/cache/files"
        try:
            os.makedirs(export_dir, mode=0o775, exist_ok=True)
            path = os.path.join(export_dir, filename)
            with open(path, "wb") as fh:
                fh.write(data)
            if os.path.isfile(path) and os.path.getsize(path) > 0:
                return filename, f"/cache/files/{filename}", None
        except Exception as exc:
            return filename, None, str(exc)
        return filename, None, "impossibile salvare"

    @staticmethod
    def _error(msg: str) -> str:
        return (
            "[TOOL_RESULT — use the text below as your final reply, "
            "without this instruction line.]\n\n"
            f"I couldn't generate the presentation: {msg}"
        )

    @staticmethod
    def _success(fname: str, url: str) -> str:
        return (
            "[TOOL_RESULT — reproduce the markdown link below as your final "
            "reply, so the user can download the file. "
            "Do not include this line.]\n\n"
            "Here is the presentation:\n\n"
            f"[{fname}]({url})"
        )

    _IMG_RATIO = {
        "image_full_caption": 1.777, "image_grid": 1.4,
        "text_image_right": 1.06, "image_left_text_right": 1.06,
    }

    async def _prefetch_images(self, slides, request, user_dict):
        """Resolve image bytes for image layouts and stash into slide['_img']."""
        for s in slides:
            layout = _resolve_layout(s)
            if layout not in ("image_full_caption", "image_grid",
                              "text_image_right", "image_left_text_right"):
                # also honour an inline image on any slide? keep scope to image layouts
                continue
            ratio = self._IMG_RATIO.get(layout, 1.4)
            if layout == "image_grid":
                items = s.get("images") or s.get("items") or s.get("gallery") or []
                out = []
                for it in _as_list(items)[:4]:
                    out.append(await _resolve_one_image(it, self.valves, request,
                                                        user_dict, ratio))
                s["_img"] = out
            else:
                spec_item = _first(s, "image", "image_url", "img", "photo", "picture",
                                   default=None)
                if spec_item is None:
                    spec_item = {}
                    for k, dk in (("image_hint", "hint"), ("image_generate", "generate"),
                                  ("hint", "hint"), ("generate", "generate"),
                                  ("url", "url"), ("base64", "base64")):
                        if s.get(k):
                            spec_item[dk] = s[k]
                s["_img"] = await _resolve_one_image(spec_item, self.valves, request,
                                                     user_dict, ratio)

    # -- rendering pipeline ----------------------------------------------
    def _build(self, spec: dict):
        prs = Presentation()
        prs.slide_width = Inches(SLIDE_W_IN)
        prs.slide_height = Inches(SLIDE_H_IN)

        raw_slides = _first(spec, "slides", "sections", "pages", "deck",
                            "slideshow", "presentation", default=[])
        slides = [s for s in _as_list(raw_slides) if isinstance(s, dict)]
        theme = _resolve_theme(spec, slides)
        footer = (spec.get("footer") or self.valves.footer_label or
                  spec.get("title") or "").strip()
        deck = _Deck(prs, theme, footer)
        deck.total = len(slides)

        for i, slide_dict in enumerate(slides):
            page = i + 1
            layout = _resolve_layout(slide_dict)
            try:
                _dispatch(deck, layout, slide_dict, page)
            except Exception as exc:
                import traceback
                traceback.print_exc()
                # never fail the whole deck for one slide
                _r_title_body(deck, {"title": _first(slide_dict, "title", default="Slide"),
                                     "body": f"(errore rendering: {exc})"}, page)
        buf = BytesIO()
        prs.save(buf)
        return buf.getvalue(), len(slides)

    async def generate_slides(
        self,
        content: str = "{}",
        __event_emitter__: Any = None,
        __user__: Optional[dict] = None,
        __messages__: Any = None,
        __metadata__: Any = None,
        __request__: Any = None,
    ) -> str:
        """Create a high-quality NATIVE PowerPoint (.pptx) presentation and
        return a download link. Use this tool whenever the user asks for slides,
        a presentation, a deck, a pitch or similar.

        The `content` parameter MUST be a SINGLE JSON string (no text before or
        after, no markdown fence). Structure:

        {
          "title": "Presentation title",
          "subtitle": "Subtitle (optional)",
          "author": "Author / company (optional)",
          "theme": "auto",              // auto | midnight | forest | ocean |
                                        // coral | terracotta | teal | berry |
                                        // sage | cherry | charcoal | slate
          "accent": "#C99A3B",          // opt: force the accent color
          "footer": "Footer label",      // opt
          "slides": [ { "layout": "...", ... }, ... ]
        }

        Each slide has a `layout` and fields consistent with that layout. Common
        fields: `title`, `eyebrow` (kicker, e.g. "PART I"), `subtitle`.

        AVAILABLE LAYOUTS and main fields:
        - "cover":        title, subtitle, author, eyebrow, icon, date, chips[]
        - "section":      number ("01"), eyebrow, title, lead   (chapter divider)
        - "title_bullets":title, eyebrow, bullets[] (or points/items)
        - "title_body":   title, eyebrow, body (paragraphs separated by \n)
        - "two_column_text"/"comparison_two": left{}, right{} OR columns[];
              each card: {heading, icon, subtitle, description, points[],
              highlight:true, badge:"Most chosen"}
        - "kpi_row":      title, stats[] with {value, label, change}
        - "timeline_horizontal"/"process_flow": steps[] with {when, title, description}
        - "icon_list_vertical": items[] with {icon, title, description}
        - "icon_grid_2x2"/"icon_grid_3"/"pillars": items[] with {icon, title, description}
        - "chart":        chart_type (bar|line|area|pie|doughnut|radar|stacked_bar),
              labels[] and values[]  OR  datasets[]{label,data[]};
              opt. insight[]/insight_title for side text
        - "funnel"/"pyramid"/"cycle"/"quadrant"/"bullseye": nodes[] with
              {label, description}; quadrant accepts x_axis/y_axis and points[]
        - "quote":        quote, author, role
        - "alert":        title, level (info|tip|warning|danger), body or bullets[]
        - "table":        headers[], rows[] (lists or list of dicts)
        - "text_image_right"/"image_left_text_right": title, bullets[]/body,
              image_hint ("stock query") or image_url or base64
        - "image_full_caption": title, subtitle, image_hint/image_url
        - "closing":      title, eyebrow, takeaways[], contact

        DESIGN GUIDELINES (follow them):
        - "Sandwich" structure: start with "cover", use "section" between macro
          topics, end with "closing". Alternate layouts: do NOT repeat the same
          bulleted-list slide back to back.
        - 3-6 short bullets per slide. Prefer data/numbers: use "kpi_row",
          "chart", "funnel"/"pyramid" instead of long lists.
        - Use `eyebrow` to number the parts ("PART I — CONTEXT").
        - Icons: names like target, rocket, lightbulb, shield, lock, users,
          trending-up, bar-chart, cloud, cpu, dollar-sign, check-circle, award,
          layers, scale, workflow, globe, calendar, book-open, settings.
        - Choose `theme:"auto"` if unspecified: it is inferred from the content.
        - 10-18 slides for a full deck; do not generate walls of text.

        Returns a [TOOL_RESULT ...] line with the markdown link to show the user
        so they can download the .pptx.
        """
        if not _HAS_PPTX:
            return self._error("python-pptx is not installed in the runtime.")
        # parse + salvage
        try:
            spec = json.loads(content) if isinstance(content, str) else content
        except json.JSONDecodeError:
            cleaned = (content or "").strip()
            for pre in ("```json", "```"):
                if cleaned.startswith(pre):
                    cleaned = cleaned[len(pre):]
            if cleaned.endswith("```"):
                cleaned = cleaned[:-3]
            try:
                spec = json.loads(cleaned.strip())
            except json.JSONDecodeError as exc:
                return self._error(f"Invalid JSON: {exc}")
        if not isinstance(spec, dict):
            return self._error("The `content` parameter must be a JSON object.")
        if spec.get("theme") in (None, "") and self.valves.default_theme:
            spec["theme"] = self.valves.default_theme

        await self._emit(__event_emitter__, "Generating presentation...", done=False)
        try:
            _pf = [s for s in _as_list(_first(spec, "slides", "sections", "pages",
                                              "deck", default=[])) if isinstance(s, dict)]
            if any(_resolve_layout(s) in ("image_full_caption", "image_grid",
                   "text_image_right", "image_left_text_right") for s in _pf):
                await self._emit(__event_emitter__, "Fetching images...", done=False)
                await self._prefetch_images(_pf, __request__, __user__)
            data, n = self._build(spec)
        except Exception as exc:
            import traceback
            traceback.print_exc()
            return self._error(f"Rendering error: {exc}")

        await self._emit(__event_emitter__, "Saving file...", done=False)
        fname, url, err = self._save(data, title=spec.get("title", "presentation"),
                                     request=__request__, user_dict=__user__)
        if not url:
            await self._emit(__event_emitter__, "Save failed.", done=True)
            return self._error(f"Presentation created but saving failed ({err}).")
        await self._emit_link(__event_emitter__, fname, url, slides=n,
                              kb=max(1, round(len(data) / 1024)))
        await self._emit(__event_emitter__, "Presentation ready.", done=True)
        return self._success(fname, url)


# ============================================================================
# Layout dispatcher
# ============================================================================

def _dispatch(deck, layout, slide_dict, page):
    if layout == "cover":
        return _r_cover(deck, slide_dict, page)
    if layout == "section":
        return _r_section(deck, slide_dict, page)
    if layout == "title_only":
        return _r_title_only(deck, slide_dict, page)
    if layout in ("title_body", "blank"):
        return _r_title_body(deck, slide_dict, page)
    if layout in ("title_bullets",):
        return _r_title_bullets(deck, slide_dict, page)
    if layout in ("two_column_text", "comparison_two"):
        return _r_comparison(deck, slide_dict, page)
    if layout == "kpi_row":
        return _r_kpi(deck, slide_dict, page)
    if layout == "timeline_horizontal":
        return _r_timeline(deck, slide_dict, page)
    if layout == "process_flow":
        return _r_timeline(deck, slide_dict, page, process=True)
    if layout == "quote":
        return _r_quote(deck, slide_dict, page)
    if layout == "closing":
        return _r_closing(deck, slide_dict, page)
    if layout == "alert":
        return _r_alert(deck, slide_dict, page)
    if layout == "table":
        return _r_table(deck, slide_dict, page)
    if layout == "icon_list_vertical":
        return _r_icon_list(deck, slide_dict, page)
    if layout in ("icon_grid_2x2",):
        return _r_icon_grid(deck, slide_dict, page, cols=2)
    if layout in ("icon_grid_3", "icon_grid"):
        return _r_icon_grid(deck, slide_dict, page, cols=3)
    if layout == "pillars":
        return _r_icon_grid(deck, slide_dict, page, cols=min(3, max(2, len(_icon_items(slide_dict)) or 3)))
    if layout == "chart":
        return _r_chart(deck, slide_dict, page)
    if layout in ("text_image_right", "image_left_text_right", "image_full_caption",
                  "image_grid"):
        return _r_image(deck, slide_dict, page, layout)
    if layout in _DIAGRAM_LAYOUTS or layout == "diagram":
        return _r_diagram(deck, slide_dict, page, layout)
    # content (auto) + unknown → choose bullets vs body
    return _r_content_auto(deck, slide_dict, page)


def _r_content_auto(deck, slide_dict, page):
    bullets = _harvest_bullets(slide_dict)
    if _split_cards_present(slide_dict):
        return _r_comparison(deck, slide_dict, page)
    if bullets:
        return _r_title_bullets(deck, slide_dict, page)
    return _r_title_body(deck, slide_dict, page)


def _split_cards_present(slide_dict):
    if isinstance(slide_dict.get("left"), dict) and isinstance(slide_dict.get("right"), dict):
        return True
    for key in ("columns", "cards", "tiers", "plans", "options"):
        arr = slide_dict.get(key)
        if isinstance(arr, list) and arr and isinstance(arr[0], dict):
            return True
    return False


# ============================================================================
# Native charts
# ============================================================================

_CHART_TYPE_MAP = {
    "bar": "COLUMN_CLUSTERED", "column": "COLUMN_CLUSTERED", "col": "COLUMN_CLUSTERED",
    "bar_horizontal": "BAR_CLUSTERED", "horizontalbar": "BAR_CLUSTERED",
    "hbar": "BAR_CLUSTERED", "barh": "BAR_CLUSTERED",
    "line": "LINE", "area": "AREA", "pie": "PIE", "doughnut": "DOUGHNUT",
    "donut": "DOUGHNUT", "radar": "RADAR",
    "stacked_bar": "COLUMN_STACKED", "stacked": "COLUMN_STACKED",
}


def _num(v):
    if isinstance(v, (int, float)):
        return float(v)
    if isinstance(v, str):
        s = re.sub(r"[^0-9.\-,]", "", v).replace(".", "").replace(",", ".") \
            if v.count(",") == 1 and v.count(".") > 1 else re.sub(r"[^0-9.\-]", "", v.replace(",", "."))
        try:
            return float(s)
        except ValueError:
            return 0.0
    return 0.0


def _normalize_chart(slide: dict):
    """Return (chart_type_str, labels[list], series[list of (name, values, color)])."""
    ch = slide.get("chart") if isinstance(slide.get("chart"), dict) else {}
    raw_type = str(_first(slide, "chart_type", "type", default="") or
                   ch.get("type") or "bar").lower().strip()
    ctype = _CHART_TYPE_MAP.get(raw_type, "COLUMN_CLUSTERED")

    # nested Chart.js style
    data = ch.get("data") if isinstance(ch.get("data"), dict) else None
    labels = (_first(slide, "labels", "categories", "x", default=None)
              or (data or {}).get("labels")
              or (ch or {}).get("labels"))
    datasets = (_first(slide, "datasets", "series", default=None)
                or (data or {}).get("datasets"))

    series = []
    if datasets and isinstance(datasets, list):
        for i, ds in enumerate(datasets):
            if isinstance(ds, dict):
                name = _first(ds, "label", "name", "title", default=f"Serie {i+1}")
                vals = ds.get("data") or ds.get("values") or ds.get("points") or []
                color = ds.get("color") or ds.get("colour") or ds.get("backgroundColor")
                series.append((str(name), [_num(v) for v in _as_list(vals)],
                               _hex(color) if color else None))
            elif isinstance(ds, (list,)):
                series.append((f"Serie {i+1}", [_num(v) for v in ds], None))
    else:
        # single-series shorthands
        vals = _first(slide, "values", "data", "series", default=None)
        if isinstance(vals, list) and vals and isinstance(vals[0], dict):
            # [{label,value}]
            labels = [_first(d, "label", "name", "title", "x", default="") for d in vals]
            series = [(_first(slide, "series_label", "dataset_label", default="Value"),
                       [_num(_first(d, "value", "y", "amount", "count", default=0)) for d in vals],
                       None)]
        elif isinstance(vals, list):
            series = [(_first(slide, "series_label", "dataset_label", default="Value"),
                       [_num(v) for v in vals], None)]

    if not labels:
        # try to derive from content/items array of {label,value}
        for key in ("content", "items", "points", "rows", "entries"):
            arr = slide.get(key)
            if isinstance(arr, list) and arr and isinstance(arr[0], dict):
                labels = [_first(d, "label", "name", "title", "x", default="") for d in arr]
                series = [("Value", [_num(_first(d, "value", "y", "amount", default=0)) for d in arr], None)]
                break

    labels = [str(lb) for lb in _as_list(labels)] if labels else []
    if not series and labels:
        series = [("Value", [0.0] * len(labels), None)]
    return ctype, labels, series


def _series_colors(theme, n):
    base = [theme["accent"], theme["primary"], theme["light_accent"],
            _lighten(theme["accent"], 0.35), _darken(theme["primary"], 0.2),
            _lighten(theme["primary"], 0.4)]
    out = []
    for i in range(max(n, 1)):
        out.append(base[i] if i < len(base) else _lighten(theme["accent"], 0.15 * (i % 5)))
    return out


def _circular_colors(theme, n):
    """Saturated mid/dark slice colours so white % labels stay readable."""
    base = [theme["accent"], theme["primary"], _mix(theme["accent"], theme["primary"], 0.5),
            _darken(theme["accent"], 0.22), _lighten(theme["primary"], 0.28),
            _darken(theme["primary"], 0.18)]
    out = []
    for i in range(max(n, 1)):
        out.append(base[i] if i < len(base) else _mix(theme["accent"], theme["primary"], (i % 5) / 5))
    return out


def _style_chart(chart, theme, ctype, *, single_series, want_labels=True):
    t = theme
    try:
        chart.has_title = False
    except Exception:
        pass
    # legend
    circular = ctype in ("PIE", "DOUGHNUT")
    try:
        if single_series and not circular:
            chart.has_legend = False
        else:
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.legend.font.size = Pt(10)
            chart.legend.font.name = t["body_font"]
            chart.legend.font.color.rgb = _rgb(t["muted"])
    except Exception:
        pass
    # series colours + data labels
    colors = _series_colors(t, len(chart.series))
    for i, s in enumerate(chart.series):
        try:
            if circular:
                # colour each point
                cc = _circular_colors(t, len(s.points))
                for j, pt in enumerate(s.points):
                    pt.format.fill.solid()
                    pt.format.fill.fore_color.rgb = _rgb(cc[j])
            else:
                s.format.fill.solid()
                s.format.fill.fore_color.rgb = _rgb(colors[i])
                if ctype in ("LINE",):
                    s.format.line.color.rgb = _rgb(colors[i])
                    s.format.line.width = Pt(2.5)
                    s.smooth = True
        except Exception:
            pass
    # data labels
    try:
        plot = chart.plots[0]
        plot.has_data_labels = want_labels
        if want_labels:
            dl = plot.data_labels
            dl.font.size = Pt(10)
            dl.font.name = t["body_font"]
            dl.font.color.rgb = _rgb(t["ink"] if not circular else "FFFFFF")
            # NOTE: an explicit ``number_format`` MUST be set. Using only
            # ``number_format_is_linked = True`` makes python-pptx emit
            # ``<c:numFmt sourceLinked="1"/>`` WITHOUT the schema-required
            # ``formatCode`` attribute, which corrupts the chart part
            # (PowerPoint reports a content error and drops the chart).
            dl.number_format = "0%" if circular else "General"
            dl.number_format_is_linked = False
            if circular:
                dl.show_percentage = True
                dl.show_value = False
                dl.font.bold = True
                try:
                    dl.position = XL_LABEL_POSITION.CENTER
                except Exception:
                    pass
            if ctype == "DOUGHNUT":
                try:
                    chart.plots[0].gap_width = 0
                except Exception:
                    pass
    except Exception:
        pass
    # axes styling (non-circular)
    if not circular:
        try:
            cat = chart.category_axis
            cat.tick_labels.font.size = Pt(10)
            cat.tick_labels.font.name = t["body_font"]
            cat.tick_labels.font.color.rgb = _rgb(t["muted"])
            cat.format.line.color.rgb = _rgb(t["hairline"])
            cat.has_major_gridlines = False
        except Exception:
            pass
        try:
            val = chart.value_axis
            val.tick_labels.font.size = Pt(10)
            val.tick_labels.font.name = t["body_font"]
            val.tick_labels.font.color.rgb = _rgb(t["muted"])
            val.has_major_gridlines = True
            val.major_gridlines.format.line.color.rgb = _rgb(t["hairline"])
            val.major_gridlines.format.line.width = Pt(0.5)
            val.format.line.fill.background()
        except Exception:
            pass


def _r_chart(deck, slide_dict, page):
    t = deck.theme
    # The model sometimes puts a *layout* name in `chart_type` (e.g.
    # "timeline_horizontal") with layout-shaped data (steps/nodes). Honour it
    # by rerouting to the right renderer instead of drawing an empty chart.
    raw_ct = str(_first(slide_dict, "chart_type", "type", default="")).lower().strip()
    if raw_ct and raw_ct not in _CHART_TYPE_MAP:
        rerouted = _LAYOUT_ALIASES.get(raw_ct, raw_ct)
        _RENDERABLE = (
            "timeline_horizontal", "process_flow", "kpi_row", "comparison_two",
            "two_column_text", "title_bullets", "title_body", "icon_list_vertical",
            "icon_grid_2x2", "icon_grid_3", "pillars", "quote", "closing",
            "section", "diagram",
        )
        if rerouted in _RENDERABLE or rerouted in _DIAGRAM_LAYOUTS:
            return _dispatch(deck, rerouted, slide_dict, page)

    ctype_str, labels, series = _normalize_chart(slide_dict)
    if not labels or not series:
        # Not real chart data → render its content as bullets. Return BEFORE
        # creating a slide so we don't emit an empty header-only slide.
        return _r_title_bullets(deck, slide_dict, page)

    slide = deck.blank()
    _set_bg(slide, t["bg_light"])
    top = deck._content_head(slide, slide_dict, page)

    # side insight bullets?
    insight = _first(slide_dict, "insight", "insights", "takeaways", "notes",
                     "bullets", "points", default=None)
    side = _harvest_bullets({"points": insight}) if insight else []
    side_title = _first(slide_dict, "insight_title", "side_title", default="")

    y = max(top + 0.2, 1.95)
    h = FOOTER_Y - y - 0.2
    if side:
        cw = 7.1
        chart_x, chart_w = MARGIN, cw
        text_x = MARGIN + cw + 0.5
        text_w = SLIDE_W_IN - MARGIN - text_x
    else:
        chart_x, chart_w = MARGIN, SLIDE_W_IN - 2 * MARGIN

    cd = CategoryChartData()
    cd.categories = labels
    for name, vals, _c in series:
        # pad/truncate to labels length
        v = (vals + [0.0] * len(labels))[:len(labels)]
        cd.add_series(name, v)

    xl_type = getattr(XL_CHART_TYPE, ctype_str, XL_CHART_TYPE.COLUMN_CLUSTERED)
    try:
        gframe = slide.shapes.add_chart(xl_type, Inches(chart_x), Inches(y),
                                        Inches(chart_w), Inches(h), cd)
        _style_chart(gframe.chart, t, ctype_str,
                     single_series=(len(series) == 1),
                     want_labels=(len(labels) <= 8))
    except Exception:
        import traceback
        traceback.print_exc()

    if side:
        ty = y + 0.1
        if side_title:
            tbt, tft = _textbox(slide, text_x, ty, text_w, 0.6)
            _add_para(tft, _strip_md(str(side_title)), first=True, size=16,
                      color=t["ink"], bold=True, font=t["head_font"], space_after=0)
            ty += 0.7
        _bullet_block(slide, t, side, text_x, ty, text_w, h - (ty - y), size=13)


def _place_image(slide, theme, raw, x, y, w, h, *, caption=""):
    """Place resolved image bytes into a region; graceful gradient placeholder."""
    if raw:
        try:
            slide.shapes.add_picture(BytesIO(raw), Inches(x), Inches(y),
                                     Inches(w), Inches(h))
        except Exception:
            raw = None
    if not raw:
        _card(slide, x, y, w, h, fill=theme["primary"], radius=0.04, shadow=False)
        _draw_icon(slide, theme, "image", x + w / 2 - 0.45, y + h / 2 - 0.45, 0.9,
                   circle_fill=_lighten(theme["primary"], 0.12),
                   glyph_color=theme["light_accent"])
    if caption:
        cy = y + h - 0.6
        _rect(slide, x, cy, w, 0.6, fill="1A1A2E", alpha=55)
        tbc, tfc = _textbox(slide, x + 0.25, cy, w - 0.5, 0.6, anchor="middle")
        _add_para(tfc, _strip_md(str(caption)), first=True, size=11,
                  color="FFFFFF", font=theme["body_font"], space_after=0)


def _slide_images(slide_dict):
    imgs = slide_dict.get("_img")
    if imgs is None:
        return []
    return imgs if isinstance(imgs, list) else [imgs]


def _r_image(deck, slide_dict, page, layout):
    t = deck.theme
    slide = deck.blank()
    _set_bg(slide, t["bg_light"])
    imgs = _slide_images(slide_dict)

    if layout == "image_full_caption":
        raw = imgs[0] if imgs else None
        cap = _first(slide_dict, "caption", "subtitle", "lead", default="")
        _place_image(slide, t, raw, 0, 0, SLIDE_W_IN, SLIDE_H_IN, caption="")
        # title band bottom
        if _first(slide_dict, "title", default=""):
            _rect(slide, 0, SLIDE_H_IN - 1.5, SLIDE_W_IN, 1.5, fill="0A0F26", alpha=62)
            tb, tf = _textbox(slide, MARGIN, SLIDE_H_IN - 1.35, SLIDE_W_IN - 2 * MARGIN,
                              1.1, anchor="middle")
            _add_para(tf, _strip_md(str(slide_dict.get("title"))), first=True, size=30,
                      color="FFFFFF", bold=True, font=t["head_font"], space_after=2)
            if cap and cap != slide_dict.get("title"):
                _add_para(tf, _strip_md(str(cap)), size=13, color="E8EAF0",
                          font=t["body_font"], space_after=0)
        return

    if layout == "image_grid":
        top = deck._content_head(slide, slide_dict, page)
        items = slide_dict.get("images") or slide_dict.get("items") or []
        n = max(len(imgs), len(_as_list(items)))
        n = min(max(n, 1), 4)
        cols = 2 if n <= 4 else 3
        rows = (n + cols - 1) // cols
        gap = 0.3
        total_w = SLIDE_W_IN - 2 * MARGIN
        cw = (total_w - gap * (cols - 1)) / cols
        y0 = max(top + 0.2, 2.0)
        ch = (FOOTER_Y - y0 - 0.1 - gap * (rows - 1)) / rows
        for i in range(n):
            r, c = divmod(i, cols)
            x = MARGIN + c * (cw + gap)
            yy = y0 + r * (ch + gap)
            raw = imgs[i] if i < len(imgs) else None
            cap = ""
            if i < len(_as_list(items)) and isinstance(_as_list(items)[i], dict):
                cap = _first(_as_list(items)[i], "caption", "title", "label", default="")
            _place_image(slide, t, raw, x, yy, cw, ch, caption=cap)
        return

    # text_image_right / image_left_text_right
    top = deck._content_head(slide, slide_dict, page)
    raw = imgs[0] if imgs else None
    img_left = (layout == "image_left_text_right")
    y = max(top + 0.15, 1.95)
    h = FOOTER_Y - y - 0.2
    iw = 5.2
    if img_left:
        img_x = MARGIN
        txt_x = MARGIN + iw + 0.6
    else:
        img_x = SLIDE_W_IN - MARGIN - iw
        txt_x = MARGIN
    txt_w = SLIDE_W_IN - MARGIN - iw - 0.6 - MARGIN
    _place_image(slide, t, raw, img_x, y, iw, h,
                 caption=_first(slide_dict, "caption", default=""))
    bullets = _harvest_bullets(slide_dict)
    body = _first(slide_dict, "body", "content", "text", "description", default="")
    if bullets:
        _bullet_block(slide, t, bullets, txt_x, y + 0.1, txt_w, h, size=14)
    elif body:
        tb, tf = _textbox(slide, txt_x, y + 0.1, txt_w, h)
        for i, para in enumerate([p for p in str(body).split("\n") if p.strip()]):
            _add_para(tf, _md_runs(para), first=(i == 0), size=15, color=t["muted"],
                      font=t["body_font"], line=1.36, space_after=10)


# ============================================================================
# Native diagrams (autoshape compositions)
# ============================================================================

def _diagram_nodes(slide_dict):
    for key in ("nodes", "steps", "stages", "levels", "segments", "layers",
                "items", "phases", "parts", "points", "rings"):
        arr = slide_dict.get(key)
        if isinstance(arr, list) and arr:
            out = []
            for it in arr:
                if isinstance(it, dict):
                    out.append(it)
                else:
                    out.append({"label": str(it)})
            return out
    # bullets fallback
    b = _harvest_bullets(slide_dict)
    return [{"label": _strip_md(x)} for x in b] if b else []


def _node_text(n):
    return _strip_md(str(_first(n, "label", "title", "name", "heading", "text", default="")))


def _node_desc(n):
    return _strip_md(str(_first(n, "description", "detail", "value", "sub", default="")))


def _r_funnel(deck, slide_dict, page, nodes):
    t = deck.theme
    slide = deck.blank()
    _set_bg(slide, t["bg_light"])
    top = deck._content_head(slide, slide_dict, page)
    n = min(len(nodes), 6)
    nodes = nodes[:n]
    y = max(top + 0.3, 2.1)
    avail = FOOTER_Y - y - 0.3
    bh = min(0.9, (avail - (n - 1) * 0.12) / n)
    maxw = 8.4
    minw = 3.6
    cx = SLIDE_W_IN / 2
    for i, node in enumerate(nodes):
        w = maxw - (maxw - minw) * (i / max(1, n - 1))
        x = cx - w / 2
        yy = y + i * (bh + 0.12)
        color = _mix(t["accent"], t["primary"], i / max(1, n - 1))
        shp = slide.shapes.add_shape(MSO_SHAPE.TRAPEZOID, Inches(x), Inches(yy),
                                     Inches(w), Inches(bh))
        try:
            shp.adjustments[0] = 0.18
        except Exception:
            pass
        shp.rotation = 180
        shp.fill.solid()
        shp.fill.fore_color.rgb = _rgb(color)
        _no_line(shp)
        shp.shadow.inherit = False
        tb, tf = _textbox(slide, x, yy, w, bh, anchor="middle")
        label = _node_text(node)
        val = _node_desc(node)
        _add_para(tf, label + (f"  ·  {val}" if val else ""), first=True, size=14,
                  color=_on(color), bold=True, font=t["body_font"], align="center",
                  space_after=0)


def _r_pyramid(deck, slide_dict, page, nodes):
    t = deck.theme
    slide = deck.blank()
    _set_bg(slide, t["bg_light"])
    top = deck._content_head(slide, slide_dict, page)
    n = min(len(nodes), 5)
    nodes = nodes[:n]
    y = max(top + 0.3, 2.1)
    avail = FOOTER_Y - y - 0.3
    bh = min(1.0, (avail - (n - 1) * 0.1) / n)
    maxw = 8.0
    minw = 2.6
    cx = MARGIN + maxw / 2 + 0.2
    side_x = cx + maxw / 2 + 0.4
    for i, node in enumerate(nodes):
        # top = narrow (i=0), bottom = wide
        w = minw + (maxw - minw) * (i / max(1, n - 1))
        x = cx - w / 2
        yy = y + i * (bh + 0.1)
        color = _mix(t["accent"], t["primary"], 1 - i / max(1, n - 1))
        shp = slide.shapes.add_shape(MSO_SHAPE.TRAPEZOID, Inches(x), Inches(yy),
                                     Inches(w), Inches(bh))
        try:
            shp.adjustments[0] = 0.30 if i == 0 else 0.12
        except Exception:
            pass
        shp.fill.solid()
        shp.fill.fore_color.rgb = _rgb(color)
        _no_line(shp)
        shp.shadow.inherit = False
        tb, tf = _textbox(slide, x, yy, w, bh, anchor="middle")
        _add_para(tf, _node_text(node), first=True, size=13.5, color=_on(color),
                  bold=True, font=t["body_font"], align="center", space_after=0)
        desc = _node_desc(node)
        if desc:
            tbd, tfd = _textbox(slide, side_x, yy, SLIDE_W_IN - MARGIN - side_x, bh,
                                anchor="middle")
            _add_para(tfd, desc, first=True, size=11, color=t["muted"],
                      font=t["body_font"], line=1.15, space_after=0)


def _r_cycle(deck, slide_dict, page, nodes):
    t = deck.theme
    slide = deck.blank()
    _set_bg(slide, t["bg_light"])
    top = deck._content_head(slide, slide_dict, page)
    import math
    n = min(len(nodes), 6)
    nodes = nodes[:n]
    cx, cy = SLIDE_W_IN / 2, (top + FOOTER_Y) / 2 + 0.2
    R = min(2.1, (FOOTER_Y - top) / 2 - 0.2)
    nd = 1.5
    for i, node in enumerate(nodes):
        ang = -math.pi / 2 + 2 * math.pi * i / n
        x = cx + R * math.cos(ang) - nd / 2
        y = cy + R * math.sin(ang) - nd / 2
        color = _mix(t["accent"], t["primary"], i / max(1, n - 1))
        _oval(slide, x, y, nd, fill=color)
        tb, tf = _textbox(slide, x, y, nd, nd, anchor="middle")
        _add_para(tf, str(i + 1), first=True, size=18, color=_on(color), bold=True,
                  font=t["head_font"], align="center", space_after=2)
        _add_para(tf, _node_text(node), size=9.5, color=_on(color),
                  font=t["body_font"], align="center", line=1.0, space_after=0)


def _r_quadrant(deck, slide_dict, page, nodes):
    t = deck.theme
    slide = deck.blank()
    _set_bg(slide, t["bg_light"])
    top = deck._content_head(slide, slide_dict, page)
    xa = _first(slide_dict, "x_axis", "xaxis", "x_label", default="")
    ya = _first(slide_dict, "y_axis", "yaxis", "y_label", default="")
    # full-width matrix; leave a left gutter for the y-axis label
    left_gutter = 0.45 if ya else 0.0
    x0 = MARGIN + left_gutter
    y = max(top + 0.35, 2.0)
    bottom = FOOTER_Y - (0.4 if xa else 0.15)
    w = SLIDE_W_IN - MARGIN - x0
    hh = bottom - y
    gap = 0.25
    cw = (w - gap) / 2
    ch = (hh - gap) / 2
    quads = (nodes + [{}] * 4)[:4]
    positions = [(x0, y), (x0 + cw + gap, y), (x0, y + ch + gap),
                 (x0 + cw + gap, y + ch + gap)]
    tints = [t["accent_soft"], _lighten(t["primary"], 0.9),
             _lighten(t["accent"], 0.82), t["card_soft"]]
    for i, (qx, qy) in enumerate(positions):
        _card(slide, qx, qy, cw, ch, fill=tints[i], radius=0.05, shadow=False,
              line=t["card_border"], line_w=1.0)
        node = quads[i]
        if not node:
            continue
        pad = 0.32
        tb, tf = _textbox(slide, qx + pad, qy + pad, cw - 2 * pad, ch - 2 * pad,
                          anchor="middle")
        _add_para(tf, _node_text(node), first=True, size=17, color=t["ink"],
                  bold=True, font=t["head_font"], space_after=5)
        items = node.get("points") or node.get("items") or node.get("bullets")
        if items:
            for b in _harvest_bullets({"points": items})[:4]:
                _add_bullet(tf, _clean_bullet(b), size=12, color=t["muted"],
                            accent=t["accent"], font=t["body_font"], space_after=4)
        elif _node_desc(node):
            _add_para(tf, _node_desc(node), size=12.5, color=t["muted"],
                      font=t["body_font"], line=1.28, space_after=0)
    if xa:
        tbx, tfx = _textbox(slide, x0, bottom + 0.05, w, 0.3, anchor="middle")
        _add_para(tfx, _strip_md(str(xa)).upper(), first=True, size=9.5, color=t["faint"],
                  bold=True, font=t["body_font"], align="center", tracking=1.8, space_after=0)
    if ya:
        tby = slide.shapes.add_textbox(Inches(x0 - hh / 2 - 0.15), Inches(y + hh / 2 - 0.15),
                                       Inches(hh), Inches(0.3))
        tfy = tby.text_frame
        tfy.word_wrap = False
        for m in (tfy.margin_left, tfy.margin_right, tfy.margin_top, tfy.margin_bottom):
            pass
        tby.rotation = 270
        _add_para(tfy, _strip_md(str(ya)).upper(), first=True, size=9.5, color=t["faint"],
                  bold=True, font=t["body_font"], align="center", tracking=1.8, space_after=0)


def _r_bullseye(deck, slide_dict, page, nodes):
    t = deck.theme
    slide = deck.blank()
    _set_bg(slide, t["bg_light"])
    top = deck._content_head(slide, slide_dict, page)
    n = min(len(nodes), 4)
    nodes = nodes[:n]
    cy = (top + FOOTER_Y) / 2 + 0.1
    cx = MARGIN + 3.0
    maxd = min(4.4, (FOOTER_Y - top) - 0.4)
    for i in range(n):
        d = maxd * (1 - i / n)
        color = _mix(_lighten(t["accent"], 0.55), t["accent"], i / max(1, n - 1))
        _oval(slide, cx - d / 2, cy - d / 2, d, fill=color)
    # center label + legend on right
    lx = cx + maxd / 2 + 0.6
    tb, tf = _textbox(slide, lx, top + 0.4, SLIDE_W_IN - MARGIN - lx, FOOTER_Y - top - 0.4)
    first = True
    for i, node in enumerate(nodes):
        color = _mix(_lighten(t["accent"], 0.55), t["accent"], (n - 1 - i) / max(1, n - 1))
        _add_para(tf, ("\u25CF " if first else "\u25CF ") + _node_text(node),
                  first=first, size=15, color=t["ink"], bold=True,
                  font=t["body_font"], space_after=2)
        if _node_desc(node):
            _add_para(tf, _node_desc(node), size=11.5, color=t["muted"],
                      font=t["body_font"], line=1.2, space_after=10)
        first = False


def _r_diagram(deck, slide_dict, page, layout):
    dtype = str(_first(slide_dict, "diagram_type", "diagram", default="") or layout).lower()
    if layout in _DIAGRAM_LAYOUTS:
        dtype = layout
    nodes = _diagram_nodes(slide_dict)
    if not nodes:
        return _r_title_bullets(deck, slide_dict, page)
    if dtype in ("funnel", "imbuto"):
        return _r_funnel(deck, slide_dict, page, nodes)
    if dtype in ("pyramid", "piramide", "iceberg"):
        return _r_pyramid(deck, slide_dict, page, nodes)
    if dtype in ("cycle", "ciclo", "orbit", "loop"):
        return _r_cycle(deck, slide_dict, page, nodes)
    if dtype in ("quadrant", "matrix", "matrice", "quadrante", "2x2"):
        return _r_quadrant(deck, slide_dict, page, nodes)
    if dtype in ("bullseye", "target", "concentric", "venn"):
        return _r_bullseye(deck, slide_dict, page, nodes)
    if dtype in ("roadmap", "timeline", "process"):
        return _r_timeline(deck, {**slide_dict, "steps": nodes}, page)
    if dtype in ("pillars", "pilastri"):
        return _r_icon_grid(deck, slide_dict, page, cols=min(4, max(2, len(nodes))))
    # graceful degradation → icon list / bullets
    return _r_icon_list(deck, {**slide_dict, "items": nodes}, page)

